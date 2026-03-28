import streamlit as st
import streamlit.components.v1 as components
import pandas as pd
import numpy as np
import json
import os
import base64
import sys
from io import BytesIO
from textwrap import shorten
from datetime import datetime
from tc_core.utils.portabilidade import get_base_path, get_data_root
from versionamento import obter_versao_atual

# Diretório raiz do projeto
_ROOT = str(get_base_path())
_DATA_ROOT = str(get_data_root())

# Função para obter mês atual em português
def obter_mes_atual():
    """Retorna o mês atual em português"""
    meses = {
        1: "Janeiro", 2: "Fevereiro", 3: "Março", 4: "Abril",
        5: "Maio", 6: "Junho", 7: "Julho", 8: "Agosto",
        9: "Setembro", 10: "Outubro", 11: "Novembro", 12: "Dezembro"
    }
    agora = datetime.now()
    return meses[agora.month]

# Função para obter data e hora de atualização dos dados
def obter_data_atualizacao_dados():
    """Retorna a data e hora da última atualização dos arquivos de dados"""
    try:
        arquivos_dados = [
            os.path.join(_DATA_ROOT, "TC_Ext", "historico_consolidado", "df_final_historico.parquet"),
            os.path.join(_DATA_ROOT, "TC_Ext", "historico_consolidado", "df_vol_historico.parquet"),
            os.path.join(_DATA_ROOT, "TC_Ext", "historico_consolidado", "BUD", "df_final_historico_BUD.parquet"),
        ]
        
        data_atualizacao = None
        for arquivo in arquivos_dados:
            if os.path.exists(arquivo):
                try:
                    data_modificacao = os.path.getmtime(arquivo)
                    if data_modificacao and data_modificacao > 0:
                        if data_atualizacao is None or data_modificacao > data_atualizacao:
                            data_atualizacao = data_modificacao
                except (OSError, ValueError):
                    continue
        
        if data_atualizacao and data_atualizacao > 0:
            try:
                dt = datetime.fromtimestamp(data_atualizacao)
                meses = {
                    1: "Janeiro", 2: "Fevereiro", 3: "Março", 4: "Abril",
                    5: "Maio", 6: "Junho", 7: "Julho", 8: "Agosto",
                    9: "Setembro", 10: "Outubro", 11: "Novembro", 12: "Dezembro"
                }
                return f"{dt.day:02d} de {meses[dt.month]} de {dt.year} às {dt.hour:02d}:{dt.minute:02d}"
            except (ValueError, OSError):
                return "Não disponível"
        return None
    except Exception:
        return None

# Cabeçalho compacto com data de atualização
mes_atual = obter_mes_atual()
ano_atual = datetime.now().year
versao_atual = obter_versao_atual()
data_atualizacao = obter_data_atualizacao_dados()

# Montar textos do cabeçalho
texto_esquerda = f"📚 Stellantis Cost Intelligence (SCI) | Versão {versao_atual} | {mes_atual} {ano_atual} | Desenvolvido por Hudson Cardin, Lauro Paiva e Frederico Cesar de Jesus"
texto_direita = f"📅 Dados atualizados em: {data_atualizacao}" if data_atualizacao else ""

st.markdown(f"""
<div style='display: flex; justify-content: space-between; align-items: center; color: #fff; padding: 8px 10px; font-size: 0.85rem; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); border-bottom: 1px solid #5a4fcf; margin-bottom: 10px;'>
    <div style='flex: 1;'>{texto_esquerda}</div>
    <div style='flex: 0 0 auto; margin-left: 20px;'>{texto_direita}</div>
</div>
""", unsafe_allow_html=True)


# CSS para melhorar visualização
st.markdown("""
    <style>
        .block-container {
            padding-top: 4rem !important;
            padding-bottom: 0.25rem !important;
        }
        div[data-testid="stVerticalBlock"] {
            gap: 0.38rem !important;
        }
        hr {
            margin: 0.18rem 0 !important;
            opacity: 0.16 !important;
        }
        h1 {
            white-space: nowrap !important;
            overflow: hidden !important;
            text-overflow: ellipsis !important;
        }
        .stTabs [data-baseweb="tab-list"] {
            gap: 8px;
        }
        .stTabs [data-baseweb="tab"] {
            padding: 10px 20px;
            font-weight: 600;
        }
    </style>
""", unsafe_allow_html=True)

st.title("📚 Documentação — Stellantis Cost Intelligence (SCI)")


# Tradutor (PT/EN/FR/ES) — traduz a página inteira via widget client-side
with st.container():
        cols = st.columns([1, 6])
        with cols[0]:
                st.markdown("**🌐 Tradutor**")
        with cols[1]:
                st.markdown(
                        "<div id='sci-translate-container' style='min-height: 28px;'></div>",
                        unsafe_allow_html=True,
                )

        components.html(
                """
                <script>
                (function() {
                    try {
                        const parentWindow = window.parent;
                        const doc = parentWindow.document;

                        // Container visível (criado via st.markdown)
                        const container = doc.getElementById('sci-translate-container');
                        if (!container) return;

                        // Elemento do tradutor
                        let el = doc.getElementById('google_translate_element');
                        if (!el) {
                            el = doc.createElement('div');
                            el.id = 'google_translate_element';
                            container.appendChild(el);
                        }

                        // Callback esperado pelo script do Google Translate
                        if (!parentWindow.googleTranslateElementInit) {
                            parentWindow.googleTranslateElementInit = function() {
                                if (!parentWindow.google || !parentWindow.google.translate) return;
                                new parentWindow.google.translate.TranslateElement(
                                    {
                                        pageLanguage: 'pt',
                                        includedLanguages: 'pt,en,fr,es',
                                        autoDisplay: false
                                    },
                                    'google_translate_element'
                                );
                            };
                        }

                        // Carrega o script uma única vez
                        if (!doc.getElementById('google-translate-script')) {
                            const s = doc.createElement('script');
                            s.id = 'google-translate-script';
                            s.src = 'https://translate.google.com/translate_a/element.js?cb=googleTranslateElementInit';
                            doc.body.appendChild(s);
                        }
                    } catch (e) {
                        // Se a política do navegador/Streamlit bloquear acesso ao parent,
                        // simplesmente não mostra o tradutor.
                    }
                })();
                </script>
                """,
                height=0,
        )

        st.caption(
                "Tradução automática (PT/EN/FR/ES) via Google Translate. "
                "Se o acesso externo estiver bloqueado na rede corporativa, o tradutor pode não aparecer."
        )


# Função para detectar caminho base correto
def get_base_path():
    """Retorna o caminho base correto para LEITURA de dados"""
    import sys
    if hasattr(sys, '_MEIPASS'):
        # Rodando no executável PyInstaller - apontar para _internal
        return sys._MEIPASS
    else:
        # Rodando em desenvolvimento
        return os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def _formatar_mtime(ts: float) -> str:
    try:
        return datetime.fromtimestamp(ts).strftime("%d/%m/%Y %H:%M:%S")
    except Exception:
        return "Não disponível"


@st.cache_data(show_spinner=False)
def _ler_arquivo_texto_cacheado(caminho: str, mtime: float) -> str:
    with open(caminho, "r", encoding="utf-8") as f:
        return _reparar_mojibake(f.read())


def _reparar_mojibake(texto: str) -> str:
    if not texto:
        return ""

    marcadores = ("Ã", "Â", "â", "ðŸ", "ï¸", "â€”", "â†’", "â‚¬")
    original_score = sum(texto.count(m) for m in marcadores)
    if original_score == 0:
        return texto

    try:
        reparado = texto.encode("latin-1", errors="ignore").decode("utf-8", errors="ignore")
    except Exception:
        return texto

    reparado_score = sum(reparado.count(m) for m in marcadores)
    return reparado if reparado_score < original_score else texto


def _carregar_markdown(caminho: str) -> tuple[str | None, str | None, float | None]:
    if not os.path.exists(caminho):
        return None, f"Arquivo não encontrado: {caminho}", None
    try:
        mtime = os.path.getmtime(caminho)
        return _ler_arquivo_texto_cacheado(caminho, mtime), None, mtime
    except Exception as e:
        return None, f"Erro ao carregar arquivo: {caminho} ({e})", None


def _extrair_secao_por_heading(md: str, headings: list[str]) -> str:
    """Extrai o conteúdo de uma seção markdown (sem o heading).

    Procura o primeiro heading encontrado em `headings` e retorna até o próximo
    heading de nível 2 ("## ").
    """
    if not md:
        return ""
    start = -1
    heading_encontrado = None
    for h in headings:
        start = md.find(h)
        if start != -1:
            heading_encontrado = h
            break
    if start == -1 or heading_encontrado is None:
        return ""

    start_line_end = md.find("\n", start)
    if start_line_end == -1:
        return ""
    start_content = start_line_end + 1
    end = md.find("\n## ", start_content)
    if end == -1:
        end = len(md)
    return md[start_content:end].strip()


def _quebrar_markdown_nivel2(md: str) -> list[tuple[str, str]]:
    if not md:
        return []

    secoes: list[tuple[str, str]] = []
    atual_titulo: str | None = None
    atual_linhas: list[str] = []

    for linha in md.splitlines():
        if linha.startswith("## "):
            if atual_titulo and any(l.strip() for l in atual_linhas):
                secoes.append((atual_titulo, "\n".join(atual_linhas).strip()))
            atual_titulo = linha[3:].strip()
            atual_linhas = []
            continue
        if atual_titulo is not None:
            atual_linhas.append(linha)

    if atual_titulo and any(l.strip() for l in atual_linhas):
        secoes.append((atual_titulo, "\n".join(atual_linhas).strip()))

    return secoes


def _renderizar_markdown_em_expanders(md: str, expanded_first: bool = True) -> None:
    secoes = _quebrar_markdown_nivel2(md)
    if not secoes:
        st.info("Conteúdo de documentação indisponível para esta seção.")
        return

    for indice, (titulo, conteudo) in enumerate(secoes):
        with st.expander(titulo, expanded=expanded_first and indice == 0):
            st.markdown(conteudo)


def _render_doc_regras_tc_ext() -> None:
    with st.expander("🔢 **Cálculos Principais e Métricas Fundamentais**", expanded=True):
        st.markdown(r"""
### 📊 CPU (Custo por Unidade)

O CPU é sempre calculado como razão entre o custo total agregado e o volume total agregado do mesmo perímetro.

**Ordem exata do cálculo**
1. Filtrar o conjunto de custo conforme oficina, veículo, período, Type 05, Type 06, Account e demais filtros ativos.
2. Filtrar o conjunto de volume usando exatamente o mesmo perímetro lógico do custo.
3. Agregar custos no nível exibido na tela: linha, grupo, período ou total.
4. Agregar volumes no mesmo nível exibido na tela.
5. Calcular o CPU como custo agregado dividido pelo volume agregado.
6. Só depois formatar, converter moeda ou exibir em tabela/gráfico.

**Por que essa ordem é obrigatória**
- CPU é uma razão. Somar ou tirar média de CPUs já calculados gera distorção.
- Se o volume for agregado em um perímetro diferente do custo, o denominador deixa de representar o numerador.
- Se inverter a ordem e calcular CPU por linha antes do agrupamento, o resultado final deixa de ser auditável.

**Fórmula explícita**

$$
CPU = \frac{\sum Custo}{\sum Volume}
$$

Onde:
- $\sum Custo$ é a soma dos custos após todos os filtros e agrupamentos.
- $\sum Volume$ é a soma dos volumes do mesmo perímetro.

**O que entra no cálculo**
- Valores monetários já filtrados.
- Volumes das mesmas combinações exibidas.

**O que não entra**
- Volume de oficinas, veículos ou períodos que ficaram fora do filtro.
- Média de CPU por linha.

**Nível em que o cálculo acontece**
- Linha detalhada: custo da linha dividido pelo volume da linha, se a linha já tiver granularidade compatível.
- Grupo/período/total: sempre por agregação prévia de custo e volume.

**Exemplo numérico completo**

Considere duas linhas no mesmo período:
- Linha A: custo de R$ 120.000 e volume de 3.000 unidades.
- Linha B: custo de R$ 80.000 e volume de 1.000 unidades.

Passo a passo:
1. Soma do custo: R$ 120.000 + R$ 80.000 = R$ 200.000.
2. Soma do volume: 3.000 + 1.000 = 4.000 unidades.
3. CPU correto: R$ 200.000 / 4.000 = R$ 50,00 por unidade.

Se alguém invertesse a ordem:
1. CPU da linha A = R$ 120.000 / 3.000 = R$ 40,00.
2. CPU da linha B = R$ 80.000 / 1.000 = R$ 80,00.
3. Média simples incorreta = (40 + 80) / 2 = R$ 60,00.

O valor correto é R$ 50,00, não R$ 60,00.

### 💰 Custo Total, deltas e ratios

O custo total exibido em qualquer visual é sempre a soma simples dos valores monetários válidos do perímetro filtrado.

**Ordem exata do cálculo**
1. Aplicar todos os filtros da tela.
2. Agrupar no nível de visualização.
3. Somar o campo monetário.
4. Calcular deltas e ratios sobre os totais já agregados.

**Fórmulas explícitas**

$$
Custo\ Total = \sum Custo_i
$$

$$
Delta_{Flex-BUD} = Flex\ Bud - BUD
$$

$$
Delta_{Real-Flex} = Real - Flex\ Bud
$$

$$
Ratio = \frac{Real}{Flex\ Bud}
$$

**Exemplo numérico completo**
- Budget original: R$ 600.000.
- Flex Bud: R$ 645.000.
- Real: R$ 670.000.

Resultados:
- Delta Flex vs BUD = 645.000 - 600.000 = R$ 45.000.
- Delta Real vs Flex = 670.000 - 645.000 = R$ 25.000.
- Ratio = 670.000 / 645.000 = 1,0388 = 103,88%.

Interpretação:
- O orçamento flexível ficou R$ 45.000 acima do budget original por efeito de volume.
- O realizado ainda ficou R$ 25.000 acima do esperado flexível.
        """)

    with st.expander("🔄 **Cálculo de Flex Bud (Budget Flexível)**", expanded=False):
        st.markdown(r"""
### 📋 Conceito operacional

O Flex Bud do TC Ext ajusta apenas a parcela de Custo Variável do budget pela proporção entre volume realizado e volume budget. A lógica implementada em produção segue a mesma ideia da função de cálculo do módulo compartilhado:

$$
Flex\ Bud = Custo\ Fixo + (Custo\ Vari\'{a}vel \times \frac{Volume\ Actual}{Volume\ Budget})
$$

### 1. Real x Budget

**Ordem exata do cálculo**
1. Filtrar o dataframe principal de budget no perímetro selecionado.
2. Classificar cada linha como fixa ou Custo Variável pela coluna Custo.
3. Somar o custo fixo por período.
4. Somar o custo total por período.
5. Calcular Custo Variável como total menos fixo.
6. Somar volume budget por período.
7. Somar volume actual por período.
8. Calcular a proporção $Volume\ Actual / Volume\ Budget$.
9. Aplicar a proporção somente ao Custo Variável.
10. Somar custo fixo e Custo Variável flexibilizado.

**Por que essa ordem é necessária**
- A separação entre fixo e Custo Variável precisa ocorrer antes da flexibilização; se a proporção for aplicada ao total, o fixo é inflado ou reduzido indevidamente.
- O volume é agregado por período antes da divisão para evitar distorção de granularidade.
- O Custo Variável é derivado do total menos o fixo para garantir fechamento do número final.

**Fórmulas explícitas**

$$
Custo\ Vari\'{a}vel = Custo\ Total\ Budget - Custo\ Fixo
$$

$$
Propor\c{c}\tilde{a}o = \frac{Volume\ Actual}{Volume\ Budget}
$$

$$
Flex\ Bud = Custo\ Fixo + (Custo\ Vari\'{a}vel \times Propor\c{c}\tilde{a}o)
$$

**O que entra no cálculo**
- Custo budget do período.
- Classificação entre custo fixo e Custo Variável.
- Volume budget agregado do período.
- Volume actual agregado do mesmo período.

**O que não entra**
- Volume de outro período.
- Linhas fora do filtro.
- Reclassificação manual de custo fora da coluna Custo.

**Nível do cálculo**
- Primeiro por período.
- Depois pode ser convertido para CPU dividindo o Flex Bud pelo volume actual do período.

**Exemplo numérico completo**
- Período: Março.
- Custo fixo budget: R$ 180.000.
- Custo total budget: R$ 500.000.
- Portanto Custo Variável budget: R$ 320.000.
- Volume budget: 40.000 unidades.
- Volume actual: 46.000 unidades.

Passo a passo:
1. Proporção = 46.000 / 40.000 = 1,15.
2. Parcela fixa permanece R$ 180.000.
3. Parcela de Custo Variável flexibilizada = 320.000 × 1,15 = R$ 368.000.
4. Flex Bud = 180.000 + 368.000 = R$ 548.000.

Se o realizado do período foi R$ 560.000:
- Delta Real vs Flex = 560.000 - 548.000 = R$ 12.000.

Em CPU:
- CPU Flex = 548.000 / 46.000 = R$ 11,91 por unidade.
- CPU Real = 560.000 / 46.000 = R$ 12,17 por unidade.

### 2. Real x Real no Waterfall

No waterfall entre dois meses reais, a lógica é análoga, mas o custo de referência é o mês inicial e a proporção de volume compara mês final contra mês inicial.

**Ordem exata do cálculo**
1. Filtrar o mês inicial.
2. Separar custo fixo e Custo Variável do mês inicial.
3. Agregar volume do mês inicial.
4. Agregar volume do mês final.
5. Calcular a proporção $V_2 / V_1$.
6. Manter o fixo do mês inicial.
7. Recalcular apenas a parcela de Custo Variável com a nova proporção.
8. Somar as duas parcelas e comparar contra o mês final realizado.

**Exemplo numérico completo**
- Mês 1: custo fixo R$ 90.000, Custo Variável R$ 210.000, volume 30.000.
- Mês 2: volume 36.000.

Passo a passo:
1. Proporção = 36.000 / 30.000 = 1,20.
2. Fixo flexado = R$ 90.000.
3. Custo Variável flexado = 210.000 × 1,20 = R$ 252.000.
4. Flex do mês 1 no volume do mês 2 = 90.000 + 252.000 = R$ 342.000.

Se o mês 2 real foi R$ 350.000, o delta operacional puro é de R$ 8.000.
        """)

    with st.expander("📊 **Cálculo de Volumes**", expanded=False):
        st.markdown(r"""
### 📁 Papel do volume no sistema

Volume não é dado auxiliar decorativo. Ele sustenta CPU, Flex Bud, gráficos comparativos e análises de diluição de custo.

**Ordem exata do cálculo de volume para uso analítico**
1. Carregar a base de volume correta do contexto: histórico, ano corrente, budget ou actual.
2. Aplicar o mesmo perímetro lógico dos dados de custo.
3. Agrupar volume no nível necessário: período, veículo, oficina ou combinação equivalente.
4. Somar os volumes válidos.
5. Entregar esse volume agregado para CPU, Flex Bud ou forecast.

**Por que essa ordem é necessária**
- O volume é o denominador do CPU e a base de proporção do Flex Bud.
- Se o volume for filtrado com um universo maior ou menor que o custo, o resultado monetário por unidade fica contaminado.
- Se o agrupamento ocorrer em um nível incompatível, o sistema compara numerador e denominador de naturezas diferentes.

**Fórmula explícita**

$$
Volume\ Total = \sum Volume_i
$$

**O que entra**
- Apenas linhas de volume alinhadas ao perímetro de custo.
- Volume do período efetivamente selecionado.

**O que não entra**
- Volume de oficinas ou veículos que não aparecem na seleção atual.
- Volume de períodos futuros quando a tela é histórica.

**Exemplo numérico completo**

Seleção na tela:
- Oficina = Pintura.
- Período = Abril.
- Veículos = A e B.

Volumes da base filtrada:
- Veículo A: 12.000 unidades.
- Veículo B: 8.000 unidades.
- Veículo C: 5.000 unidades, mas fora do filtro.

Resultado correto:
1. Volume total = 12.000 + 8.000 = 20.000 unidades.
2. Se o custo total filtrado for R$ 900.000, o CPU correto = 900.000 / 20.000 = R$ 45,00.

Erro clássico:
1. Usar também o veículo C no denominador.
2. Volume incorreto = 25.000 unidades.
3. CPU incorreto = 900.000 / 25.000 = R$ 36,00.

O desvio de R$ 9,00 por unidade nasce apenas de um perímetro incorreto.
        """)

    with st.expander("💱 **Moeda e Taxas de Câmbio**", expanded=False):
        st.markdown(r"""
### 💱 Ordem entre fator, moeda e cálculo

No SCI, fator visual e moeda são transformações de apresentação. O cálculo econômico continua ancorado no valor base em BRL e na sequência correta de agregação.

**Ordem exata**
1. Carregar o valor monetário base em BRL.
2. Agregar no nível exibido.
3. Converter moeda, se a visualização pedir USD ou EUR.
4. Aplicar fator visual K ou M somente para exibição de custo total.
5. Se a visualização for CPU, dividir custo agregado pelo volume agregado e manter fator visual como Nenhum.

**Por que essa ordem é necessária**
- Converter moeda antes ou depois da soma preserva o total desde que a taxa seja única, mas aplicar fator visual antes do CPU quebra a escala do resultado.
- CPU já é uma razão monetária por unidade; usar K ou M nele altera a interpretação da unidade.

**Fórmula explícita de conversão**

$$
Valor_{moeda} = Valor_{BRL} \times Taxa_{moeda}
$$

**Exemplo numérico completo**
- Custo agregado em BRL: R$ 2.400.000.
- Taxa USD: 0,20.

Conversão:
1. Valor em USD = 2.400.000 × 0,20 = US$ 480.000.
2. Se a visualização pedir fator K, o número exibido vira 480 K.

Se o volume do mesmo período for 60.000 unidades:
1. CPU em BRL = 2.400.000 / 60.000 = R$ 40,00.
2. CPU em USD = 480.000 / 60.000 = US$ 8,00.

O que não pode ser feito:
1. Aplicar fator K no custo antes do CPU.
2. Fazer 2.400 K / 60.000 = 0,04 e exibir como se fosse US$ 8,00 ou R$ 40,00.
        """)

    with st.expander("🔍 **Filtros e Perímetros de Análise**", expanded=False):
        st.markdown(r"""
### 🎛️ Regras de perímetro

Toda métrica do TC Ext depende do mesmo perímetro lógico entre custo, volume e budget.

**Ordem exata de aplicação**
1. Definir ano e período.
2. Aplicar filtros estruturais: oficina, veículo, USI, centro de custo, conta e classificações.
3. Aplicar filtros textuais ou avançados, quando existirem.
4. Gerar o dataframe principal filtrado.
5. Extrair dele o universo de dimensões necessário para sincronizar budget e volume.
6. Recalcular totais, CPU, Flex e gráficos apenas com esse perímetro sincronizado.

**Por que essa ordem é necessária**
- O dataframe principal define o universo econômico da análise.
- Budget e volume são subordinados a esse universo; se forem filtrados antes por critérios diferentes, a comparação deixa de ser justa.

**Exemplo numérico completo**

Filtro desejado:
- Ano 2025.
- Oficina = Montagem.
- Type 06 = Energia.

Após filtrar, o sistema encontra:
- Real: R$ 300.000.
- Budget: R$ 280.000.
- Volume actual: 15.000.
- Volume budget: 14.000.

Proporção correta de Flex = 15.000 / 14.000 = 1,0714.

Se o budget tivesse sido filtrado sem respeitar a oficina, poderia entrar volume budget de 20.000 e o Flex cair para uma proporção de 0,75, completamente incompatível com o real analisado.
        """)


def _render_doc_regras_tc_veiculos() -> None:
    with st.expander("💰 Cadeia de Custos", expanded=True):
        st.markdown(r"""
### 🔗 Sequência lógica da cadeia de custos do TC Veículos

O TC Veículos não começa no CPU. Ele começa na despesa primária, passa pelo rateio FA, isola o Fluxo Principal, separa a D&A dedicada e só depois rateia por veículo.

**Ordem exata do cálculo**
1. Carregar a despesa primária da oficina e período.
2. Calcular o Rateio FA de cada oficina e período.
3. Calcular o Custo FA como despesa primária multiplicada pelo Rateio FA.
4. Calcular o Custo FP como despesa primária menos custo FA.
5. Carregar ou distribuir D&A dedicada.
6. Calcular FP sem Dedicada como Custo FP menos D&A dedicada.
7. Calcular o percentual de rateio por veículo via tempo de produção.
8. Ratear o FP sem Dedicada por veículo.
9. Somar a D&A dedicada do veículo ao valor rateado.
10. Só então calcular o CPU por veículo.

**Por que essa ordem é necessária**
- Se o rateio por veículo for feito antes de separar FA e FP, o veículo absorve custo em uma base errada.
- Se a D&A dedicada não for destacada antes do rateio, ela é distribuída duas vezes ou no universo errado.
- O CPU depende do custo final do veículo, portanto só pode vir no fim.

**Fórmulas explícitas**

$$
Custo\ FA = Despesa\ Primaria \times Rateio\ FA
$$

$$
Custo\ FP = Despesa\ Primaria - Custo\ FA
$$

$$
FP\ sem\ Dedicada = Custo\ FP - D\&A\ dedicada
$$

$$
Custo\ Rateado_{veic} = FP\ sem\ Dedicada \times Percentual_{veic}
$$

$$
Custo\ FP\ Veiculo = Custo\ Rateado_{veic} + D\&A\ dedicada_{veic}
$$

**Exemplo numérico completo**
- Despesa primária da oficina no mês: R$ 1.000.000.
- Rateio FA: 25%.
- D&A dedicada do veículo X: R$ 30.000.
- Percentual do veículo X no rateio por tempo: 40%.

Passo a passo:
1. Custo FA = 1.000.000 × 25% = R$ 250.000.
2. Custo FP = 1.000.000 - 250.000 = R$ 750.000.
3. FP sem Dedicada = 750.000 - 30.000 = R$ 720.000.
4. Custo rateado do veículo X = 720.000 × 40% = R$ 288.000.
5. Custo FP Veículo X = 288.000 + 30.000 = R$ 318.000.

Esse é o custo que seguirá para o CPU do veículo X.
        """)

    with st.expander("🚗 Rateio por Veículo", expanded=False):
        st.markdown(r"""
### ⏱️ Regra de rateio por tempo

O rateio por veículo distribui o custo da oficina usando o tempo consumido por cada veículo dentro da mesma oficina e do mesmo período.

**Ordem exata do cálculo**
1. Calcular Tempo Veic de cada veículo como EST × Volume.
2. Somar o tempo total da oficina no período.
3. Calcular o percentual de cada veículo dividindo seu tempo pelo total.
4. Aplicar esse percentual ao FP sem Dedicada.
5. Validar se a soma dos percentuais da oficina no período fecha em 100%.

**Fórmulas explícitas**

$$
Tempo\ Veic = EST \times Volume
$$

$$
Percentual_{veic} = \frac{Tempo\ Veic_{veic}}{\sum Tempo\ Veic_{oficina, periodo}}
$$

$$
Custo\ Rateado_{veic} = FP\ sem\ Dedicada \times Percentual_{veic}
$$

**O que entra**
- Tempo do veículo na mesma oficina e período.
- FP sem Dedicada do mesmo grupo.

**O que não entra**
- Tempo de outra oficina.
- Tempo de outro período.
- D&A dedicada, que já é adicionada depois.

**Exemplo numérico completo**

Oficina Solda em Maio:
- Veículo A: EST 2,0 horas, volume 5.000 → tempo = 10.000.
- Veículo B: EST 1,5 horas, volume 4.000 → tempo = 6.000.
- Veículo C: EST 1,0 hora, volume 2.000 → tempo = 2.000.

Passo a passo:
1. Tempo total = 10.000 + 6.000 + 2.000 = 18.000.
2. Percentual A = 10.000 / 18.000 = 55,56%.
3. Percentual B = 6.000 / 18.000 = 33,33%.
4. Percentual C = 2.000 / 18.000 = 11,11%.

Se o FP sem Dedicada da oficina for R$ 900.000:
- A recebe R$ 500.040.
- B recebe R$ 299.970.
- C recebe R$ 99.990.

O fechamento fica compatível com o total original, salvo pequenas diferenças de arredondamento visual.
        """)

    with st.expander("📊 Flex Budget", expanded=False):
        st.markdown(r"""
### 📈 Flex Budget no TC Veículos

No TC Veículos, o budget flexível segue a mesma lógica estrutural do TC Ext: custo fixo fica constante e o Custo Variável é ajustado pela relação entre volume actual e volume budget.

**Ordem exata do cálculo**
1. Agregar volume budget por período.
2. Agregar volume actual por período.
3. Agregar custo budget por período.
4. Identificar a parcela fixa pela coluna Custo.
5. Calcular a parcela não fixa como total menos fixo.
6. Calcular a proporção de volume.
7. Aplicar a proporção somente à parcela não fixa.
8. Somar fixo e Custo Variável flexibilizado.

**Fórmula explícita**

$$
Flex\ Bud = Custo\ Fixo + (Custo\ N\tilde{a}o\ Fixo \times \frac{Vol\ Actual}{Vol\ Budget})
$$

**Exemplo numérico completo**
- Custo budget fixo: R$ 250.000.
- Custo budget total: R$ 700.000.
- Custo Variável: R$ 450.000.
- Volume budget: 50.000 unidades.
- Volume actual: 55.000 unidades.

Passo a passo:
1. Proporção = 55.000 / 50.000 = 1,10.
2. Custo Variável flexibilizado = 450.000 × 1,10 = R$ 495.000.
3. Flex Bud total = 250.000 + 495.000 = R$ 745.000.

Se o realizado for R$ 760.000:
- Delta Real vs Flex = 760.000 - 745.000 = R$ 15.000.
        """)

    with st.expander("📈 CPU (Custo por Unidade)", expanded=False):
        st.markdown(r"""
### 🧮 CPU do TC Veículos

O CPU do TC Veículos é calculado depois de concluído o custo final do veículo. Em dados consolidados, ele usa custo agregado e volume agregado. Em dados rateados, ele usa o custo final do veículo dividido pelo volume do próprio veículo.

**Ordem exata do cálculo**
1. Concluir cadeia de custos até Custo FP Veículo.
2. Carregar o volume do mesmo veículo e período.
3. Agregar custo e volume no nível exibido.
4. Dividir custo agregado por volume agregado.

**Fórmula explícita**

$$
CPU_{veiculo} = \frac{Custo\ FP\ Veiculo}{Volume}
$$

**Exemplo numérico completo**
- Custo FP Veículo X: R$ 318.000.
- Volume do veículo X: 6.000 unidades.

Resultado:
- CPU = 318.000 / 6.000 = R$ 53,00 por unidade.

Se dois veículos forem agrupados:
- Veículo X: R$ 318.000 e 6.000 unidades.
- Veículo Y: R$ 182.000 e 4.000 unidades.

CPU correto do grupo:
1. Custo agregado = 318.000 + 182.000 = R$ 500.000.
2. Volume agregado = 6.000 + 4.000 = 10.000.
3. CPU do grupo = 500.000 / 10.000 = R$ 50,00.

Não é correto fazer média simples entre R$ 53,00 e R$ 45,50.
        """)

    with st.expander("🎯 KPIs (Topo e Resumo)", expanded=False):
        st.markdown(r"""
### 📌 O que cada KPI mede

Os KPIs do topo e do resumo não são números independentes. Eles seguem a mesma cadeia econômica do pipeline.

**Ordem exata dos KPIs de topo**
1. Somar despesa primária.
2. Somar custo FA.
3. Somar custo FP.
4. Somar D&A dedicada.
5. Somar FP sem Dedicada.
6. Separar, quando necessário, linhas de Redis como subconjunto específico.

**Exemplo numérico completo**
- Despesa primária: R$ 3.000.000.
- Custo FA: R$ 800.000.
- Custo FP: R$ 2.200.000.
- D&A dedicada: R$ 150.000.
- FP sem Dedicada: R$ 2.050.000.

Fechamento esperado:
- Custo FP = 3.000.000 - 800.000 = 2.200.000.
- FP sem Dedicada = 2.200.000 - 150.000 = 2.050.000.

Se qualquer KPI não fechar nessa sequência, o erro está na cadeia anterior, não no card em si.
        """)

    with st.expander("🎛️ Filtros", expanded=False):
        st.markdown(r"""
### 🔍 Regras dos filtros do TC Veículos

**Ordem exata**
1. Escolher se a análise é consolidada ou por veículo.
2. Se veículo = Todos, usar bases consolidadas.
3. Se veículo específico, usar bases rateadas por veículo.
4. Aplicar oficina, período e demais filtros sobre a base já correta.
5. Sincronizar volume e budget com o mesmo universo.

**Por que essa ordem é necessária**
- O filtro de veículo muda a fonte de dados. Não é apenas um filtro visual.
- Se o usuário escolher um veículo específico, o sistema precisa sair do consolidado e entrar no parquet rateado; filtrar depois não corrige essa diferença estrutural.

**Exemplo prático**
- Veículo = Todos → usar df_principal e df_principal_BUD.
- Veículo = CC21 → usar df_veiculos_custo_fp e df_veiculos_custo_fp_BUD.

Se essa troca de base não acontecer, o usuário enxerga custo consolidado com rótulo de veículo específico, o que é incorreto.
        """)

    with st.expander("🔮 Best Estimate — Premissas", expanded=False):
        st.markdown(r"""
### 🔮 Lógica de forecast do Best Estimate

O Best Estimate do TC Veículos parte apenas da base histórica real válida, projeta os meses futuros com regra econômica explícita e depois consolida tudo em arquivos reutilizados pela análise.

Nos meses históricos, o comportamento esperado é simples: o Best Estimate deve reproduzir o Real. Nos meses futuros, o sistema calcula o forecast a partir da média histórica, do volume projetado, da sensibilidade e do bloco monetário de inflação e produtividade.

**Ordem exata do cálculo**
1. Carregar a base histórica real e remover qualquer linha antiga de BE, BE Manual ou Forecast da base usada para média.
2. Aplicar os filtros do simulador sem eliminar os períodos históricos necessários para compor a referência.
3. Calcular a média histórica por combinação de chaves do custo.
4. Calcular o volume médio histórico da mesma combinação lógica.
5. Ler o volume futuro configurado para o mês projetado.
6. Calcular a proporção de volume futuro versus volume médio histórico.
7. Transformar essa proporção em variação percentual.
8. Aplicar a sensibilidade correspondente ao tipo de custo ou ao Type 06.
9. Aplicar inflação e produtividade no bloco monetário final.
10. Gravar o forecast linha a linha.
11. Persistir três camadas no consolidado: Histórico, BE e BE Manual.
12. Gerar o forecast por veículo com a mesma lógica do Real para FP sem Dedicada e D&A dedicado.
13. Validar convergência: meses históricos do arquivo rateado devem fechar com o Real.

**Fórmulas explícitas**

$$
Propor\c{c}\tilde{a}o\ de\ Volume = \frac{Volume\ Futuro}{Volume\ M\acute{e}dio\ Hist\acute{o}rico}
$$

$$
Varia\c{c}\tilde{a}o\ Percentual = Propor\c{c}\tilde{a}o - 1
$$

$$
Varia\c{c}\tilde{a}o\ Ajustada = Varia\c{c}\tilde{a}o\ Percentual \times Sensibilidade
$$

$$
Fator\ de\ Varia\c{c}\tilde{a}o = 1 + Varia\c{c}\tilde{a}o\ Ajustada
$$

$$
Fator\ Monet\acute{a}rio = 
(1 + \frac{Infla\c{c}\tilde{a}o}{100}) \times (1 - \frac{Produtividade}{100})
$$

$$
Forecast = M\acute{e}dia\ Hist\acute{o}rica \times Fator\ de\ Varia\c{c}\tilde{a}o \times Fator\ Monet\acute{a}rio
$$

**Exemplo numérico completo**
- Média histórica: R$ 100.000.
- Volume médio histórico: 10.000 unidades.
- Volume futuro: 11.500 unidades.
- Sensibilidade: 80%.
- Inflação: 4%.
- Produtividade: 3%.

Passo a passo:
1. Proporção = 11.500 / 10.000 = 1,15.
2. Variação percentual = 1,15 - 1 = 0,15.
3. Variação ajustada = 0,15 × 0,80 = 0,12.
4. Fator de variação = 1,12.
5. Fator monetário = (1 + 0,04) × (1 - 0,03) = 1,0088.
6. Forecast = 100.000 × 1,12 × 1,0088 = R$ 112.985,60.

Na prática operacional do sistema, a produtividade reduz o custo projetado após o ajuste de volume e é aplicada como redutor multiplicativo, não como subtração manual no final.

### 🚗 Rateio do forecast por veículo

Quando o fluxo exige granularidade veicular, o forecast consolidado é convertido em `forecast_veiculos_custo_fp.parquet` usando a mesma regra econômica do processamento Real.

**Regra correta do rateio por veículo**

$$
Custo\ Rateado = FP\ sem\ Dedicada \times Percentual
$$

$$
Custo\ FP\ Veiculo = Custo\ Rateado + D\&A\ dedicado
$$

Isto é propositalmente diferente de `Custo FP × Percentual`. O D&A dedicado não pode ser espalhado proporcionalmente junto com o restante do custo, porque ele já pertence ao veículo correto.

**Fontes usadas no rateio do BE**
- Percentual de rateio: `df_veiculos_percentual_rateio.parquet` carregado por `load_percentual_rateio_veiculos_real(ano)`.
- D&A dedicado: `df_dea_dedicado.parquet` carregado por `load_dea_dedicado_real(ano)`.
- Base de custo: `forecast_completo.parquet`.

**Fallback correto para linhas sem correspondência direta**
- Primeiro fallback: distribuição média por período entre veículos conhecidos, normalizada para fechar 100% no período.
- Último fallback: distribuição uniforme somente para períodos totalmente órfãos, sem nenhum veículo conhecido naquele período.

### ✅ Garantia dos meses históricos na análise

Na análise de Best Estimate, os meses marcados como `Histórico` precisam ser numericamente idênticos ao Real.

Por isso o sistema faz duas proteções complementares:
- Na geração do arquivo por veículo, valida a convergência entre meses históricos do BE e do Real.
- Na camada de análise, quando a série de BE contém linhas `Histórico`, esses meses são sobrepostos pelos dados reais antes da exibição.

Essa segunda proteção elimina divergência visual mesmo quando um parquet de forecast antigo ainda existir no ambiente. Isso é especialmente importante no Databricks, onde o deploy sincroniza o código do app, mas os arquivos em `dados/` são gerados no próprio ambiente.

### 💾 Persistência e arquivos gerados

O simulador persiste a configuração aplicada em `config_forecast.json`, não em `premissas.json`.

Arquivos principais gerados em `dados/TC_Principal/Forecast/`:
- `forecast_historico.parquet`: histórico isolado.
- `forecast_previsao.parquet`: apenas meses projetados.
- `forecast_completo.parquet`: consolidado com Histórico, BE e BE Manual.
- `forecast_veiculos_custo_fp.parquet`: consolidado rateado por veículo.
- `custos_especificos.parquet`: camada manual do BE.
- `config_forecast.json`: premissas persistidas da última execução.

### 💰 BE Manual

O BE Manual entra como linha separada no consolidado final. Ele não substitui o forecast calculado; ele é somado como camada adicional e, quando necessário, também participa do fluxo de rateio por veículo no consolidado final.
        """)


def _render_doc_tabelas_graficos_tc_ext() -> None:
    with st.expander("📌 CPU e regra de agregação", expanded=True):
        st.markdown(r"""
### 📊 Por que tabela e gráfico divergem quando o CPU é calculado errado

No TC Ext, tabela e gráfico só fecham entre si quando ambos usam o mesmo dataframe agregado e aplicam a mesma regra de CPU: custo agregado dividido por volume agregado.

**Ordem exata do cálculo para qualquer visualização em CPU**
1. Aplicar os filtros da tela.
2. Agregar o custo no nível da visualização.
3. Agregar o volume no mesmo nível.
4. Calcular o CPU sobre os agregados.
5. Renderizar tabela, gráfico ou KPI.

**Por que essa ordem é necessária**
- Tabela e gráfico podem ter apresentações diferentes, mas não podem ter bases diferentes.
- Se um visual usar CPU por linha e outro usar CPU agregado, ambos mostrarão números distintos para o mesmo período.

**Fórmula explícita**

$$
CPU = \frac{\sum Custo}{\sum Volume}
$$

**Exemplo numérico completo**

Abril filtrado:
- Linha 1: R$ 300.000 e 6.000 unidades.
- Linha 2: R$ 150.000 e 2.000 unidades.

Tabela e gráfico corretos:
1. Custo total = 300.000 + 150.000 = R$ 450.000.
2. Volume total = 6.000 + 2.000 = 8.000.
3. CPU do período = 450.000 / 8.000 = R$ 56,25.

Erro clássico:
1. CPU linha 1 = 50,00.
2. CPU linha 2 = 75,00.
3. Média simples = 62,50.

Se um visual usar R$ 56,25 e o outro R$ 62,50, o problema não está no gráfico: está na ordem de cálculo.

### 📋 O que entra e o que não entra

**Entra**
- Custo e volume do mesmo filtro.
- Agregação no nível exibido.

**Não entra**
- Média de CPUs já calculados.
- Volume fora do universo filtrado.
- Fator K ou M aplicado antes do CPU.
    """)

    with st.expander("📌 Governança do ano completo (12 meses)", expanded=False):
        st.markdown(r"""
### 📅 Regra de ano completo

Quando a visualização trabalha com ano completo, o sistema precisa comparar os 12 meses usando uma convenção única de agregação e ordenação.

**Ordem exata da governança**
1. Identificar os 12 meses válidos do ano selecionado.
2. Garantir a ordenação calendário de Janeiro a Dezembro.
3. Agregar custo e volume mês a mês.
4. Calcular métricas mensais.
5. Calcular o total anual somente após concluir os 12 meses.

**Por que essa ordem é necessária**
- O total anual não é um décimo terceiro mês; ele é a soma coerente dos 12 meses.
- Se faltar mês, o ano completo fica parcial.
- Se a ordenação ficar alfabética, a leitura temporal vira inconsistente.

**Exemplo numérico completo**

Suponha quatro meses já calculados:
- Janeiro: R$ 100.000.
- Fevereiro: R$ 120.000.
- Março: R$ 110.000.
- Abril: R$ 130.000.

O acumulado parcial é R$ 460.000. O total anual correto só existe depois da soma de Janeiro a Dezembro. O sistema não deve misturar acumulado parcial com ano fechado.

### 🔒 Regra de fechamento entre tabela e gráfico

Tabela e gráfico devem mostrar os mesmos 12 pontos e o mesmo total anual. Se um deles excluir meses vazios e o outro mantiver zero explícito, a leitura comparativa do ano muda.
    """)


def _render_doc_tabelas_graficos_tc_veiculos() -> None:
    with st.expander("📊 Análise Flex por Categoria", expanded=True):
        st.markdown(r"""
### 🧾 Como ler a tabela Flex por Account

Na análise por categoria, a tabela não mostra apenas colunas decorativas. Cada coluna nasce de uma ordem lógica fixa.

**Ordem exata do cálculo da tabela**
1. Carregar o budget no perímetro da análise.
2. Carregar o realizado no mesmo perímetro.
3. Calcular o Flex Bud do período.
4. Agregar tudo por Type 05, Type 06 e Account.
5. Exibir as colunas derivadas: BUD, Flex Bud, Real, deltas e ratio.

**Fórmulas explícitas**

$$
Flex\ Bud - BUD = Flex\ Bud - BUD
$$

$$
Total - Flex\ Bud = Real - Flex\ Bud
$$

$$
Total / Flex\ Bud = \frac{Real}{Flex\ Bud}
$$

**Exemplo numérico completo**
- Account 450001.
- BUD = R$ 200.000.
- Flex Bud = R$ 218.000.
- Real = R$ 226.000.

Resultados:
- Flex Bud - BUD = 218.000 - 200.000 = R$ 18.000.
- Total - Flex Bud = 226.000 - 218.000 = R$ 8.000.
- Total / Flex Bud = 226.000 / 218.000 = 103,67%.

**O que entra**
- Valor agregado da conta no período.
- Classificação já consolidada por Type 05 e Type 06.

**O que não entra**
- Mistura entre contas diferentes antes da agregação da linha exibida.
    """)

    with st.expander("📈 Gráficos do TC Veículos", expanded=False):
        st.markdown(r"""
### 📉 Regra de fechamento entre barras, linha e delta

No gráfico principal do TC Veículos, a barra mostra o real, a linha mostra o Flex Bud e o delta é a diferença entre os dois. Todos precisam sair do mesmo agregado por período.

**Ordem exata do cálculo do gráfico**
1. Agregar Real por período.
2. Agregar Flex Bud por período.
3. Calcular delta por período como Real menos Flex Bud.
4. Enviar esses três vetores para o gráfico.

**Fórmulas explícitas**

$$
Delta_{periodo} = Real_{periodo} - Flex\ Bud_{periodo}
$$

**Exemplo numérico completo**

Junho:
- Real = R$ 1.120.000.
- Flex Bud = R$ 1.050.000.
- Delta = R$ 70.000.

Leitura correta do visual:
- Barra em 1.120.000.
- Linha em 1.050.000.
- Faixa de delta em +70.000.

Se o gráfico estiver em CPU, os mesmos valores precisam antes ser convertidos para custo agregado dividido por volume agregado do período. Não se deve desenhar a barra com CPU médio de linha.

### 🎨 Cores e semântica

- Cores ajudam leitura, mas não alteram a regra econômica.
- Histórico e BE podem ter cores diferentes; a base numérica continua sendo a mesma lógica de agregação.
    """)

    with st.expander("📋 Tabs Disponíveis", expanded=False):
        st.markdown(r"""
### 🗂️ O que cada tab consome

Cada tab do TC Veículos existe para uma granularidade analítica específica. A escolha da tab altera o recorte do dado, não a regra central de cálculo.

**Resumo funcional**
- TC Veículos: KPIs e comparação Real versus Flex por período.
- Análise Flex: hierarquia Type 05 → Type 06 → Account.
- Volume: confronto entre budget e realizado, inclusive por veículo.
- Custos por Oficina: distribuição do custo no nível da oficina.
- Tempo de Produção: base do rateio percentual por veículo.
- Dados Detalhados: rastreabilidade linha a linha.

**Por que isso importa para auditoria**
- Um número que aparece em uma tab precisa reconciliar com outra, mas nem sempre no mesmo nível de agrupamento.
- O usuário deve primeiro verificar a granularidade da tab antes de comparar dois valores aparentemente iguais.
    """)


def _render_doc_arquitetura_tc_ext() -> None:
    with st.expander("🏗️ Camadas da Arquitetura do TC Ext", expanded=True):
        st.markdown("""
### Visão estrutural

O TC Ext é organizado em quatro camadas que precisam permanecer coerentes entre si:

1. Entrada e processamento: arquivos Excel, notebooks de extração e geração de parquets.
2. Persistência: pastas anuais, histórico consolidado, Budget e Forecast.
3. Camada compartilhada: utilitários em tc_core para paths, filtros, moeda, fator, cache e componentes de UI.
4. Consumo analítico: Home, Waterfall, Best Estimate, alertas e documentação.

### Estrutura lógica do projeto

```text
TC/
├── app.py
├── pages/
├── tc_ext/
├── tc_principal/
├── tc_core/
├── tc_copilot/
├── alertas/
└── dados/
    ├── TC_Ext/
    │   ├── {ANO}/
    │   ├── Forecast/
    │   └── historico_consolidado/
    └── TC_Principal/
```

### Papel de cada bloco
- app.py: portal, menu, controle de navegação e bootstrap do ambiente.
- pages/: páginas legadas, incluindo Waterfall e esta documentação.
- tc_ext/: regras e páginas do TC Ext.
- tc_core/: contratos comuns de paths, schemas, cache, componentes e helpers de cálculo.
- dados/TC_Ext/: fonte persistida de Real, Budget, histórico e Forecast.
        """)

    with st.expander("🧱 Fluxo arquitetural do dado ao número exibido", expanded=False):
        st.markdown("""
### Sequência obrigatória

1. O usuário ou job processa Excel bruto.
2. O processamento grava parquets anuais e históricos.
3. As páginas carregam parquets via loaders compartilhados.
4. Os filtros recortam o mesmo perímetro para custo e volume.
5. Os agrupamentos são feitos antes de qualquer cálculo derivado de CPU ou Flex.
6. Só depois a interface monta KPI, tabela, gráfico e exportação.

### Por que isso importa
- Se a UI calcular antes da agregação, o número fica visualmente plausível, mas matematicamente errado.
- Se loaders diferentes usarem pastas diferentes, Real, Budget e histórico deixam de conversar.
- Se a camada compartilhada for ignorada, surge divergência entre página, executável e cloud.
        """)

    with st.expander("📁 Contratos de pastas e artefatos", expanded=False):
        st.markdown("""
### Contrato de persistência do TC Ext

Ano corrente:
- dados/TC_Ext/{ANO}/df_final.parquet
- dados/TC_Ext/{ANO}/df_vol.parquet
- dados/TC_Ext/{ANO}/df_ke5z_group.parquet

Budget:
- dados/TC_Ext/{ANO}/BUD/df_final_BUD.parquet
- dados/TC_Ext/{ANO}/BUD/df_vol_BUD.parquet
- dados/TC_Ext/{ANO}/BUD/df_ke5z_group_BUD.parquet

Histórico consolidado:
- dados/TC_Ext/historico_consolidado/df_final_historico.parquet
- dados/TC_Ext/historico_consolidado/df_vol_historico.parquet
- dados/TC_Ext/historico_consolidado/BUD/df_final_historico_BUD.parquet

Forecast:
- dados/TC_Ext/Forecast/forecast_completo.parquet
- dados/TC_Ext/Forecast/custos_especificos.parquet

Esses caminhos formam o contrato usado por desenvolvimento local, executável e cloud.
        """)


def _render_doc_arquitetura_tc_veiculos() -> None:
    with st.expander("🚗 Arquitetura funcional do TC Veículos", expanded=True):
        st.markdown("""
### Visão de alto nível

O TC Veículos tem uma arquitetura mais analítica porque precisa preservar a cadeia de composição de custo até o veículo.

### Camadas
1. Entrada: Reporting veículos.xlsx e insumos auxiliares por ano.
2. Processamento Budget: processamento_dados_veiculos_BUD.py.
3. Processamento Real: processamento_dados_veiculos.py.
4. Persistência: parquets principais, por veículo, CPU, tempos, D&A e comparativos.
5. Consumo: páginas em tc_principal/pages, alertas, Home e Best Estimate.
        """)

    with st.expander("📂 Parquets e contratos críticos", expanded=False):
        st.markdown("""
### Budget
- df_principal_BUD.parquet
- df_vol_veiculos_BUD.parquet
- df_tempo_veiculos_BUD.parquet
- df_dea_dedicado_BUD.parquet
- df_veiculos_custo_fp_BUD.parquet
- df_veiculos_cpu_BUD.parquet

### Real
- df_principal.parquet
- df_tc_sapiens.parquet
- df_vol_veiculos_actual.parquet
- df_tempo_veiculos.parquet
- df_dea_dedicado.parquet
- df_veiculos_custo_fp.parquet
- df_veiculos_cpu.parquet

Quando o filtro de veículo é aplicado, o app troca o dataset consumido e passa da visão consolidada para a visão rateada.
        """)

    with st.expander("⚙️ Ordem arquitetural do pipeline", expanded=False):
        st.markdown("""
1. Ler abas obrigatórias do Excel.
2. Normalizar nomes, períodos e tipos.
3. Calcular custos intermediários: despesa primária, FA, FP, D&A dedicado e FP sem dedicada.
4. Calcular tempo por veículo.
5. Transformar tempo em percentual de rateio.
6. Ratear custo ao veículo.
7. Calcular custo FP do veículo.
8. Calcular CPU por veículo a partir de custo agregado e volume agregado.
9. Persistir parquets do ano e consolidar histórico.
        """)


def _render_doc_especificacao_tc_ext() -> None:
    with st.expander("🧾 Especificação técnica do TC Ext", expanded=True):
        st.markdown(r"""
### Objetivo técnico

O TC Ext deve responder três perguntas sem exigir leitura do código:
1. Qual é o custo total do perímetro filtrado?
2. Como esse custo se compara a Budget, Flex Bud ou período anterior?
3. Qual volume sustenta o cálculo de CPU exibido?

### Entidades principais
- Custo: valor monetário base para agrupamentos.
- Volume: denominador dos cálculos de CPU.
- Período: mês/ano ou agregações maiores.
- Custo fixo/variável: classificação usada no Flex Bud.
- Dimensões analíticas: Oficina, Type 05, Type 06, Account e demais filtros.

### Regras técnicas imutáveis
1. CPU sempre nasce de $\sum custo / \sum volume$.
2. Flex Bud sempre é calculado em custo total antes de qualquer conversão para CPU.
3. O filtro aplicado ao custo precisa ter equivalente lógico no volume.
4. Conversão de moeda e fator são etapas finais de apresentação, não de negócio.
        """)

    with st.expander("📐 Sequência técnica de cálculo e renderização", expanded=False):
        st.markdown(r"""
### Pipeline da consulta
1. Carregar Real, Budget e volumes.
2. Aplicar filtros de negócio.
3. Agregar no nível da visualização.
4. Calcular Flex Bud em custo total.
5. Calcular deltas monetários.
6. Se o modo for CPU, dividir cada bloco agregado por seu volume agregado.
7. Só então formatar e plotar.

### Exemplo de auditoria

Suponha:
- Real = R$ 900.000
- Budget = R$ 840.000
- Volume Real = 45.000
- Volume Budget = 42.000
- Budget fixo = R$ 240.000
- Budget variável = R$ 600.000

Passos:
1. $\rho = 45.000 / 42.000 = 1,0714286$
2. Flex variável = 600.000 × 1,0714286 = R$ 642.857,16
3. Flex total = 240.000 + 642.857,16 = R$ 882.857,16
4. CPU Real = 900.000 / 45.000 = R$ 20,00
5. CPU Flex = 882.857,16 / 45.000 = R$ 19,62
6. Delta Real vs Flex = R$ 17.142,84 em total ou R$ 0,38 em CPU
        """)


def _render_doc_especificacao_tc_veiculos() -> None:
    with st.expander("🚗 Especificação técnica do TC Veículos", expanded=True):
        st.markdown(r"""
### Objetivo técnico

No TC Veículos, a especificação precisa preservar a rastreabilidade entre oficina, conta, veículo e volume produzido.

### Unidades lógicas do sistema
- Linha de custo original: nasce da extração Real ou Budget.
- Linha rateada por veículo: valor redistribuído pela lógica de tempo/proporção.
- Linha de D&A dedicado: permanece vinculada diretamente ao veículo/oficina correspondente.
- Linha de CPU: nasce sempre da relação entre custo FP do veículo e volume do veículo.

### Contrato funcional
1. Sem filtro de veículo o usuário enxerga visão consolidada.
2. Com filtro de veículo, a análise passa a operar sobre bases rateadas.
3. Flex Bud respeita a mesma classificação Fixo/Variável na visão por veículo.
4. Alertas e Best Estimate reutilizam o mesmo contrato de granularidade.
        """)

    with st.expander("🔎 Exemplo auditável por veículo", expanded=False):
        st.markdown(r"""
### Cenário

Uma oficina possui:
- FP sem dedicada = R$ 500.000
- D&A dedicado do veículo A = R$ 80.000
- Tempo total da oficina = 10.000 horas
- Tempo do veículo A = 2.500 horas
- Volume do veículo A = 5.000 unidades

### Ordem exata
1. Percentual de rateio do veículo A = 2.500 / 10.000 = 25%
2. Custo rateado ao veículo A = 500.000 × 25% = R$ 125.000
3. Custo FP do veículo A = 125.000 + 80.000 = R$ 205.000
4. CPU do veículo A = 205.000 / 5.000 = R$ 41,00 por unidade

Se a D&A for rateada junto com FP sem dedicada ou se o CPU for calculado antes da consolidação do custo por veículo, o número final deixa de fechar com a lógica implementada.
        """)


def _render_doc_visao_geral_tecnica() -> None:
    metricas_core = {
        "Arquivos Core": 62,
        "Linhas Core": 84674,
        "Telas": 15,
        "Blocos": 7,
    }
    metricas_workspace = {
        "Arquivos Workspace": 296,
        "Linhas Workspace": 306426,
    }

    df_arquivos = pd.DataFrame(
        {
            "Camada": ["Núcleo", "Workspace ampliado"],
            "Arquivos Python": [
                metricas_core["Arquivos Core"],
                metricas_workspace["Arquivos Workspace"],
            ],
        }
    ).set_index("Camada")
    df_loc = pd.DataFrame(
        {
            "Camada": ["Núcleo", "Workspace ampliado"],
            "Linhas de código": [
                metricas_core["Linhas Core"],
                metricas_workspace["Linhas Workspace"],
            ],
        }
    ).set_index("Camada")

    st.markdown(
        """
    <style>
        .sci-tech-hero {
            background: radial-gradient(circle at top left, #284b63 0%, #13293d 42%, #0b1f2a 100%);
            border: 1px solid rgba(168, 218, 220, 0.18);
            border-radius: 24px;
            padding: 26px 28px;
            color: #f8fafc;
            box-shadow: 0 18px 45px rgba(11, 31, 42, 0.28);
            margin-bottom: 1rem;
        }
        .sci-tech-chip {
            display: inline-block;
            margin: 0 10px 10px 0;
            padding: 8px 14px;
            border-radius: 999px;
            font-size: 0.80rem;
            font-weight: 700;
            letter-spacing: 0.03em;
            color: #dbeafe;
            background: rgba(125, 211, 252, 0.10);
            border: 1px solid rgba(125, 211, 252, 0.22);
        }
        .sci-kpi-card {
            min-height: 152px;
            border-radius: 22px;
            padding: 18px 18px 16px 18px;
            background: linear-gradient(160deg, #102a43 0%, #1f4e5f 52%, #2c7a7b 100%);
            border: 1px solid rgba(148, 210, 189, 0.18);
            box-shadow: 0 14px 35px rgba(16, 42, 67, 0.18);
            color: white;
            margin-bottom: 0.8rem;
        }
        .sci-kpi-value {
            font-size: 2.15rem;
            line-height: 1;
            font-weight: 800;
            margin: 8px 0 6px 0;
            letter-spacing: -0.03em;
        }
        .sci-kpi-label {
            font-size: 0.86rem;
            color: #d9f3f4;
            font-weight: 700;
            text-transform: uppercase;
            letter-spacing: 0.08em;
        }
        .sci-kpi-note {
            font-size: 0.84rem;
            color: #e6fffb;
            opacity: 0.92;
        }
        .sci-stack-card {
            min-height: 128px;
            border-radius: 20px;
            padding: 16px 16px 14px 16px;
            background: linear-gradient(180deg, #f8fbff 0%, #eef6fb 100%);
            border: 1px solid rgba(31, 78, 95, 0.12);
            box-shadow: 0 10px 24px rgba(15, 23, 42, 0.06);
            margin-bottom: 0.8rem;
        }
        .sci-stack-title {
            font-size: 1rem;
            font-weight: 800;
            color: #12344d;
            margin-bottom: 0.2rem;
        }
        .sci-stack-note {
            font-size: 0.84rem;
            color: #486581;
            line-height: 1.45;
        }
        .sci-flow-wrap {
            border-radius: 24px;
            padding: 18px;
            background: linear-gradient(180deg, #fbfdff 0%, #f2f7fb 100%);
            border: 1px solid rgba(18, 52, 77, 0.10);
            margin: 0.8rem 0 1rem 0;
        }
        .sci-flow-row {
            display: flex;
            flex-wrap: wrap;
            gap: 10px;
            align-items: stretch;
            justify-content: space-between;
        }
        .sci-flow-node {
            flex: 1 1 140px;
            min-height: 112px;
            border-radius: 18px;
            padding: 16px 14px;
            background: linear-gradient(160deg, #0f4c5c 0%, #2c7a7b 100%);
            color: white;
            box-shadow: 0 12px 24px rgba(15, 76, 92, 0.14);
        }
        .sci-flow-title {
            font-size: 0.95rem;
            font-weight: 800;
            margin-bottom: 0.3rem;
        }
        .sci-flow-note {
            font-size: 0.82rem;
            opacity: 0.92;
            line-height: 1.35;
        }
        .sci-flow-arrow {
            display: flex;
            align-items: center;
            justify-content: center;
            color: #2c7a7b;
            font-size: 1.6rem;
            font-weight: 800;
            min-width: 26px;
        }
        .sci-stage-card {
            min-height: 120px;
            border-radius: 18px;
            padding: 16px 14px;
            background: linear-gradient(180deg, #fffdf7 0%, #fff4d6 100%);
            border: 1px solid rgba(194, 120, 3, 0.16);
            margin-bottom: 0.8rem;
        }
        .sci-stage-title {
            font-size: 0.92rem;
            font-weight: 800;
            color: #8d5b00;
        }
        .sci-stage-note {
            font-size: 0.83rem;
            color: #7c5e10;
            margin-top: 0.35rem;
            line-height: 1.4;
        }
        .sci-value-card {
            min-height: 148px;
            border-radius: 20px;
            padding: 18px 16px;
            background: linear-gradient(160deg, #111827 0%, #1f2937 100%);
            color: #f9fafb;
            border: 1px solid rgba(96, 165, 250, 0.18);
            box-shadow: 0 12px 26px rgba(17, 24, 39, 0.18);
        }
        .sci-value-title {
            font-size: 1rem;
            font-weight: 800;
            margin-bottom: 0.35rem;
        }
        .sci-value-note {
            font-size: 0.84rem;
            color: #d1d5db;
            line-height: 1.45;
        }
        .sci-subtitle {
            font-size: 0.92rem;
            color: #486581;
            margin-bottom: 0.4rem;
        }
    </style>
        """,
        unsafe_allow_html=True,
    )

    st.markdown(
        """
    <div class="sci-tech-hero">
        <div style="font-size: 0.82rem; font-weight: 800; letter-spacing: 0.12em; text-transform: uppercase; color: #8bd3dd; margin-bottom: 8px;">
            Plataforma técnica integrada
        </div>
        <div style="font-size: 2.55rem; line-height: 1.02; font-weight: 900; max-width: 860px; margin-bottom: 12px;">
            O SCI já opera como sistema de engenharia, dados, cloud e inteligência aplicada.
        </div>
        <div style="font-size: 1rem; line-height: 1.55; color: #dbeafe; max-width: 920px; margin-bottom: 16px;">
            Mais do que um painel, o SCI conecta processamento, persistência, análise executiva,
            simulação, automação e IA em uma única base técnica coerente.
        </div>
        <span class="sci-tech-chip">Dados estruturados</span>
        <span class="sci-tech-chip">Cloud ready</span>
        <span class="sci-tech-chip">Forecast + alertas</span>
        <span class="sci-tech-chip">TC Copilot</span>
        <span class="sci-tech-chip">Governança e rastreabilidade</span>
    </div>
        """,
        unsafe_allow_html=True,
    )

    st.markdown("### Escala do sistema")
    st.markdown(
        "<div class='sci-subtitle'>Números conservadores do projeto atual para mostrar dimensão real, não marketing vazio.</div>",
        unsafe_allow_html=True,
    )

    col1, col2, col3, col4 = st.columns(4)
    cards_kpi = [
        (
            col1,
            "Arquivos core",
            f"{metricas_core['Arquivos Core']}",
            "núcleo principal da aplicação",
        ),
        (
            col2,
            "Linhas core",
            "84k+",
            "código Python no núcleo do produto",
        ),
        (
            col3,
            "Telas",
            f"{metricas_core['Telas']}",
            "jornadas Streamlit em operação",
        ),
        (
            col4,
            "Blocos",
            f"{metricas_core['Blocos']}",
            "camadas centrais de responsabilidade",
        ),
    ]

    for coluna, rotulo, valor, nota in cards_kpi:
        with coluna:
            st.markdown(
                f"""
            <div class="sci-kpi-card">
                <div class="sci-kpi-label">{rotulo}</div>
                <div class="sci-kpi-value">{valor}</div>
                <div class="sci-kpi-note">{nota}</div>
            </div>
                """,
                unsafe_allow_html=True,
            )

    col_chart1, col_chart2 = st.columns(2)
    with col_chart1:
        st.markdown("**Arquivos Python**")
        st.bar_chart(df_arquivos, color="#2c7a7b")
        st.caption("62 no núcleo principal e 296 no workspace ampliado atual.")

    with col_chart2:
        st.markdown("**Linhas de código Python**")
        st.bar_chart(df_loc, color="#1d4ed8")
        st.caption("84 mil+ no núcleo do produto e 306 mil+ no ecossistema ampliado.")

    st.markdown("### Stack tecnológico")
    st.markdown(
        "<div class='sci-subtitle'>Ferramentas escolhidas para sustentar interface, cloud, dados, automação, exportação e IA sem inflar a complexidade para o usuário final.</div>",
        unsafe_allow_html=True,
    )

    tecnologia_cards = [
        ("🐍 Python", "Linguagem-base do SCI, unificando regras, processamento, integração e app."),
        ("📺 Streamlit", "Camada de interface analítica com velocidade alta de evolução de produto."),
        ("☁️ Databricks Apps", "Execução cloud do app com publicação e operação corporativa."),
        ("❄️ Snowflake", "Camada de dados e integração analítica no ecossistema ampliado."),
        ("🧬 Git / GitHub Enterprise", "Versionamento, governança de mudança e rastreabilidade do código."),
        ("🤖 Serving Endpoints", "Base de APIs e recursos de IA usados pelo TC Copilot."),
        ("🐼 pandas + numpy", "Transformação, modelagem e manipulação intensiva de dados."),
        ("📈 Altair + Plotly", "Visualização interativa, comparativos e leitura executiva dos números."),
        ("📦 openpyxl + PyArrow", "Exportação Excel e persistência rápida em Parquet."),
    ]

    for linha_inicio in range(0, len(tecnologia_cards), 3):
        cols = st.columns(3)
        for coluna, (titulo, descricao) in zip(cols, tecnologia_cards[linha_inicio:linha_inicio + 3]):
            with coluna:
                st.markdown(
                    f"""
                <div class="sci-stack-card">
                    <div class="sci-stack-title">{titulo}</div>
                    <div class="sci-stack-note">{descricao}</div>
                </div>
                    """,
                    unsafe_allow_html=True,
                )

    st.markdown("### Arquitetura em uma leitura")
    st.markdown(
        "<div class='sci-subtitle'>Sem entrar em regra funcional detalhada, esta é a espinha dorsal que sustenta o SCI como plataforma.</div>",
        unsafe_allow_html=True,
    )

    st.markdown(
        """
    <div class="sci-flow-wrap">
        <div class="sci-flow-row">
            <div class="sci-flow-node">
                <div class="sci-flow-title">📥 Fonte</div>
                <div class="sci-flow-note">Excel, bases corporativas e insumos operacionais entram no pipeline com estrutura controlada.</div>
            </div>
            <div class="sci-flow-arrow">→</div>
            <div class="sci-flow-node">
                <div class="sci-flow-title">⚙️ Processamento</div>
                <div class="sci-flow-note">Normalização, transformação, persistência e publicação em fluxo repetível.</div>
            </div>
            <div class="sci-flow-arrow">→</div>
            <div class="sci-flow-node">
                <div class="sci-flow-title">💾 Armazenamento</div>
                <div class="sci-flow-note">Parquets, históricos, budget e forecast sustentam leitura rápida e auditável.</div>
            </div>
            <div class="sci-flow-arrow">→</div>
            <div class="sci-flow-node">
                <div class="sci-flow-title">📊 Análise</div>
                <div class="sci-flow-note">Camadas analíticas transformam o dado bruto em sinal executivo e operacional.</div>
            </div>
            <div class="sci-flow-arrow">→</div>
            <div class="sci-flow-node">
                <div class="sci-flow-title">🖥️ Interface</div>
                <div class="sci-flow-note">O usuário recebe uma experiência contínua entre leitura, comparação, forecast, alertas e IA.</div>
            </div>
        </div>
    </div>
        """,
        unsafe_allow_html=True,
    )

    col_data, col_rules, col_view, col_sim, col_ai = st.columns(5)
    pilares = [
        (col_data, "Dados", "bases, histórico e persistência"),
        (col_rules, "Regras", "contratos e lógica compartilhada"),
        (col_view, "Visualização", "dashboards e leitura executiva"),
        (col_sim, "Simulação", "forecast e cenários"),
        (col_ai, "IA", "copilot e resposta inteligente"),
    ]
    for coluna, titulo, nota in pilares:
        with coluna:
            st.markdown(
                f"""
            <div style="border-radius: 16px; padding: 14px 12px; background: #f8fbff; border: 1px solid rgba(18, 52, 77, 0.10); text-align: center; margin-bottom: 0.6rem;">
                <div style="font-size: 0.95rem; font-weight: 800; color: #12344d;">{titulo}</div>
                <div style="font-size: 0.78rem; color: #627d98; margin-top: 4px;">{nota}</div>
            </div>
                """,
                unsafe_allow_html=True,
            )

    st.markdown("### Maturidade do SCI")
    st.markdown(
        "<div class='sci-subtitle'>O projeto já combina operação real, cloud, automação, forecast, alertas e IA. A leitura correta é sistema em produção e evolução contínua.</div>",
        unsafe_allow_html=True,
    )

    etapas = [
        ("📚 Base estruturada", "Dados consolidados, histórico, budget e governança de persistência."),
        ("☁️ Cloud ativa", "Execução no Databricks Apps com publicação e sincronização controladas."),
        ("🔮 Forecast", "Simulação e projeção integradas ao produto."),
        ("🚨 Alertas", "Monitoramento ativo e priorização do que mais importa."),
        ("🤖 TC Copilot", "IA aplicada para ampliar interpretação, produtividade e explicação."),
    ]
    cols = st.columns(5)
    for coluna, (titulo, nota) in zip(cols, etapas):
        with coluna:
            st.markdown(
                f"""
            <div class="sci-stage-card">
                <div class="sci-stage-title">{titulo}</div>
                <div class="sci-stage-note">{nota}</div>
            </div>
                """,
                unsafe_allow_html=True,
            )

    st.markdown("### Por que isso tem valor técnico")
    col_v1, col_v2, col_v3 = st.columns(3)
    valores = [
        (
            col_v1,
            "📄 Sai da planilha",
            "Substitui controles dispersos por uma base reproduzível, organizada e muito menos frágil.",
        ),
        (
            col_v2,
            "🔒 Reduz risco operacional",
            "Centraliza regra, histórico, visualização e publicação em um fluxo governado.",
        ),
        (
            col_v3,
            "📈 Escala com consistência",
            "A arquitetura permite crescer para novas plantas, novos módulos e novos ambientes sem recomeçar do zero.",
        ),
    ]
    for coluna, titulo, nota in valores:
        with coluna:
            st.markdown(
                f"""
            <div class="sci-value-card">
                <div class="sci-value-title">{titulo}</div>
                <div class="sci-value-note">{nota}</div>
            </div>
                """,
                unsafe_allow_html=True,
            )


def _render_doc_tc_cloud() -> None:
    with st.expander("🏗️ Arquitetura cloud validada", expanded=True):
        st.markdown("""
### Separação operacional

O ambiente Databricks estável foi estruturado em dois blocos:

```text
/Workspace/.../Drafts/sci
├── dados/
├── notebooks/
├── jobs/
├── src/
└── workspace_publish/

/Workspace/.../Drafts/sci_app/sci_app
├── app.py
├── pages/
├── tc_core/
├── tc_principal/
├── tc_ext/
├── tc_copilot/
└── alertas/
```

Essa separação impede mistura entre app, pipeline e dados pesados.
        """)

    with st.expander("🛡️ Contratos dos parquets otimizados", expanded=False):
        st.markdown("""
Os parquets `agg` e `thin` são contratos operacionais, não apenas otimizações.

Regras que não podem regredir:
1. `df_veiculos_agg_home` e `df_veiculos_agg_home_BUD` precisam manter `Type 05`, `Type 06`, `Account` e `Custo`.
2. `forecast_agg` precisa manter `Type 05` e `Type 06` para o tooltip do Best Estimate.
3. `df_final_agg` e `df_final_agg_BUD` precisam manter `Type 05`, `Type 06`, `Account` e `Custo`.
4. toda mudança de groupby precisa nascer em `tc_core/parquet_schemas.py`.

Sintomas clássicos de regressão:
- waterfall com aviso `Type 06 não encontrada`;
- tooltip do BE concentrando tudo em `Outros`;
- diferença entre local e Databricks com a mesma funcionalidade.
    """)

    with st.expander("🚗 Fluxo do TC Veículos no Databricks", expanded=False):
        st.markdown("""
1. Excel é colocado em dados/TC_Principal/{ANO}/.
2. Scripts Real e Budget processam as abas necessárias.
3. Parquets são gravados no workspace.
4. O app sobe definindo SCI_SHARED_DATA_ROOT antes de importar páginas.
5. A Home do TC Veículos lê os mesmos parquets usados localmente.
        """)

    with st.expander("🧰 Notebooks, jobs e checklist anti-regressão", expanded=False):
        st.markdown("""
Notebooks validados:
1. 00_validar_ambiente_databricks.py
2. 01_criar_tabelas_delta.py
3. 03_processar_e_publicar_delta.py
4. 05_validacao_pos_job.py
5. 06_ui_consulta_workspace.py

Regras para não regredir:
- app.py precisa configurar o ambiente antes de importar páginas;
- validar cenários críticos com `SCI_USE_OPTIMIZED_PARQUETS=true`;
- manter sincronizados `tc_core/parquet_schemas.py` e os AGG regenerados;
- preservar o fallback do `data_router.py` para FULL quando o AGG estiver desatualizado;
- uploads no Workspace devem continuar via SDK com remoção prévia quando necessário;
- o fluxo local não deve sobrescrever silenciosamente o que está estável no cloud.
        """)

    with st.expander("🔁 Sincronização local e remota", expanded=False):
        st.markdown("""
Quando o Databricks App estiver mais atualizado ou mais estável que a cópia local, o remoto passa a ser a referência operacional.

Fluxo seguro:
1. exportar o workspace remoto com `databricks workspace export-dir`;
2. salvar o espelho em `Databricks/pulled_from_workspace`;
3. propagar para a raiz do repositório e espelhos locais sem apagar artefatos extras úteis;
4. validar que os arquivos puxados e os arquivos locais estão idênticos.
        """)


def _render_doc_apresentacao_styles() -> None:
    st.markdown(
        """
    <style>
        .sci-presentation-toolbar {
            padding: 1rem 1.1rem;
            border-radius: 18px;
            background: linear-gradient(180deg, #f8fbff 0%, #eef6fb 100%);
            border: 1px solid rgba(15, 23, 42, 0.08);
            margin-bottom: 1rem;
        }
        .sci-presentation-toolbar-note {
            font-size: 0.86rem;
            color: #52637a;
            line-height: 1.55;
            margin-top: 0.2rem;
        }
        .sci-slide-shell {
            max-width: 1240px;
            margin: 0.2rem auto 0 auto;
            padding: 0;
            border-radius: 0;
            background: transparent;
            border: none;
            box-shadow: none;
        }
        .sci-slide-shell-stage-1 {
            padding: 0;
            border-radius: 0;
            background: transparent;
            border: none;
            box-shadow: none;
        }
        .sci-slide-shell-stage-2 {
            padding: 0;
            border-radius: 0;
            background: transparent;
            border: none;
            box-shadow: none;
        }
        .sci-slide-note-pill {
            display: inline-block;
            padding: 0.42rem 0.72rem;
            border-radius: 999px;
            background: rgba(20, 184, 166, 0.08);
            border: 1px solid rgba(20, 184, 166, 0.18);
            color: #0f766e;
            font-size: 0.74rem;
            font-weight: 800;
            letter-spacing: 0.04em;
            text-transform: uppercase;
            margin-bottom: 0.85rem;
        }
        .sci-slide-hero {
            padding: 1.15rem 1.2rem;
            border-radius: 18px;
            background: linear-gradient(135deg, #0b1f2a 0%, #164e63 52%, #7dd3c8 100%);
            color: white;
            margin-bottom: 0.9rem;
        }
        .sci-slide-hero-stage-1 {
            background: linear-gradient(135deg, #0b1f2a 0%, #164e63 48%, #2c7a7b 100%);
        }
        .sci-slide-hero-stage-2 {
            background: linear-gradient(135deg, #5f370e 0%, #b45309 50%, #f59e0b 100%);
        }
        .sci-slide-kicker {
            font-size: 0.77rem;
            text-transform: uppercase;
            letter-spacing: 0.11em;
            font-weight: 800;
            color: #ccfbf1;
            margin-bottom: 0.45rem;
        }
        .sci-slide-title {
            font-size: 1.6rem;
            line-height: 1.1;
            font-weight: 900;
            margin-bottom: 0.35rem;
        }
        .sci-slide-headline {
            font-size: 0.98rem;
            line-height: 1.5;
            font-weight: 700;
            color: #f0fdfa;
            margin-bottom: 0.25rem;
        }
        .sci-slide-subtitle {
            font-size: 0.9rem;
            line-height: 1.55;
            color: #ecfeff;
        }
        .sci-slide-section-title {
            font-size: 0.79rem;
            text-transform: uppercase;
            letter-spacing: 0.08em;
            font-weight: 800;
            color: #0f4c5c;
            margin: 0.9rem 0 0.5rem 0;
        }
        .sci-slide-bullets {
            margin: 0;
            padding-left: 1.1rem;
        }
        .sci-slide-bullets li {
            margin-bottom: 0.42rem;
            color: #31475f;
            line-height: 1.5;
        }
        .sci-slide-card {
            min-height: 116px;
            border-radius: 16px;
            padding: 14px 14px 13px 14px;
            background: linear-gradient(180deg, #fbfdff 0%, #eef6fb 100%);
            border: 1px solid rgba(22, 78, 99, 0.12);
            box-shadow: 0 8px 22px rgba(15, 23, 42, 0.06);
            margin-bottom: 0.75rem;
        }
        .sci-slide-card-strong {
            background: linear-gradient(180deg, #fff8ee 0%, #fff0d8 100%);
            border-color: rgba(217, 119, 6, 0.22);
        }
        .sci-slide-card-title {
            font-size: 0.92rem;
            font-weight: 800;
            color: #12344d;
            margin-bottom: 0.38rem;
        }
        .sci-slide-card-body {
            font-size: 0.83rem;
            color: #486581;
            line-height: 1.48;
        }
        .sci-slide-flow-step {
            min-height: 90px;
            border-radius: 16px;
            padding: 12px 12px 11px 12px;
            background: linear-gradient(160deg, #12344d 0%, #1f6f8b 100%);
            color: white;
            box-shadow: 0 10px 24px rgba(18, 52, 77, 0.12);
            margin-bottom: 0.75rem;
        }
        .sci-slide-flow-badge {
            display: inline-flex;
            width: 1.45rem;
            height: 1.45rem;
            align-items: center;
            justify-content: center;
            border-radius: 999px;
            background: rgba(255, 255, 255, 0.18);
            font-size: 0.72rem;
            font-weight: 900;
            margin-bottom: 0.35rem;
        }
        .sci-slide-flow-title {
            font-size: 0.9rem;
            font-weight: 800;
            margin-bottom: 0.22rem;
        }
        .sci-slide-flow-body {
            font-size: 0.8rem;
            line-height: 1.42;
            color: #e6fffb;
        }
        .sci-slide-formula-wrap {
            border-radius: 16px;
            padding: 0.75rem 0.9rem 0.4rem 0.9rem;
            background: linear-gradient(180deg, #f9fbfd 0%, #edf4f8 100%);
            border: 1px solid rgba(18, 52, 77, 0.10);
            overflow-x: auto;
        }
        .sci-slide-example {
            border-radius: 16px;
            padding: 14px 15px;
            background: linear-gradient(180deg, #fffdf7 0%, #fff5d6 100%);
            border: 1px solid rgba(194, 120, 3, 0.18);
        }
        .sci-slide-example-title {
            font-size: 0.88rem;
            font-weight: 800;
            color: #8d5b00;
            margin-bottom: 0.4rem;
        }
        .sci-slide-example-body {
            font-size: 0.83rem;
            color: #7c5e10;
            line-height: 1.5;
        }
        .sci-slide-metric {
            min-height: 110px;
            border-radius: 18px;
            padding: 14px;
            background: linear-gradient(160deg, #102a43 0%, #1f4e5f 55%, #2c7a7b 100%);
            color: white;
            box-shadow: 0 12px 28px rgba(16, 42, 67, 0.14);
            margin-bottom: 0.75rem;
        }
        .sci-slide-metric-label {
            font-size: 0.72rem;
            text-transform: uppercase;
            letter-spacing: 0.08em;
            color: #d5f3f4;
            font-weight: 800;
            margin-bottom: 0.25rem;
        }
        .sci-slide-metric-value {
            font-size: 1.8rem;
            line-height: 1;
            font-weight: 900;
            margin-bottom: 0.25rem;
        }
        .sci-slide-metric-note {
            font-size: 0.8rem;
            color: #e6fffb;
            line-height: 1.42;
        }
        .sci-slide-chip {
            display: inline-block;
            margin: 0 0.45rem 0.45rem 0;
            padding: 0.42rem 0.7rem;
            border-radius: 999px;
            background: rgba(20, 184, 166, 0.10);
            border: 1px solid rgba(20, 184, 166, 0.16);
            color: #115e59;
            font-size: 0.78rem;
            font-weight: 700;
        }
        .sci-slide-compare-card {
            min-height: 184px;
            border-radius: 18px;
            padding: 14px;
            background: linear-gradient(180deg, #fbfdff 0%, #f2f7fb 100%);
            border: 1px solid rgba(18, 52, 77, 0.10);
            box-shadow: 0 10px 24px rgba(15, 23, 42, 0.06);
            margin-bottom: 0.75rem;
        }
        .sci-slide-compare-title {
            font-size: 0.92rem;
            font-weight: 800;
            color: #12344d;
            margin-bottom: 0.45rem;
        }
        .sci-slide-compare-row {
            display: grid;
            grid-template-columns: 96px 1fr 70px;
            gap: 8px;
            align-items: center;
            margin-bottom: 0.42rem;
        }
        .sci-slide-compare-label {
            font-size: 0.78rem;
            color: #486581;
            font-weight: 700;
        }
        .sci-slide-compare-track {
            height: 10px;
            background: #e7eef5;
            border-radius: 999px;
            overflow: hidden;
        }
        .sci-slide-compare-fill {
            height: 100%;
            border-radius: 999px;
        }
        .sci-slide-compare-value {
            font-size: 0.78rem;
            color: #12344d;
            font-weight: 800;
            text-align: right;
        }
        .sci-slide-compare-caption {
            font-size: 0.78rem;
            color: #627d98;
            line-height: 1.45;
            margin-top: 0.45rem;
        }
        .sci-slide-pillar-card {
            min-height: 84px;
            border-radius: 16px;
            padding: 14px 12px;
            background: #f8fbff;
            border: 1px solid rgba(18, 52, 77, 0.10);
            text-align: center;
            margin-bottom: 0.6rem;
        }
        .sci-slide-pillar-title {
            font-size: 0.95rem;
            font-weight: 800;
            color: #12344d;
        }
        .sci-slide-pillar-note {
            font-size: 0.78rem;
            color: #627d98;
            margin-top: 4px;
            line-height: 1.42;
        }
        .sci-slide-stage-card {
            min-height: 118px;
            border-radius: 18px;
            padding: 16px 14px;
            background: linear-gradient(180deg, #fffdf7 0%, #fff4d6 100%);
            border: 1px solid rgba(194, 120, 3, 0.16);
            margin-bottom: 0.75rem;
        }
        .sci-slide-stage-title {
            font-size: 0.9rem;
            font-weight: 800;
            color: #8d5b00;
            margin-bottom: 0.3rem;
        }
        .sci-slide-stage-note {
            font-size: 0.8rem;
            color: #7c5e10;
            line-height: 1.42;
        }
        .sci-slide-team-card {
            min-height: 298px;
            border-radius: 18px;
            padding: 14px;
            background: linear-gradient(180deg, #fbfdff 0%, #edf5fb 100%);
            border: 1px solid rgba(18, 52, 77, 0.10);
            box-shadow: 0 10px 24px rgba(15, 23, 42, 0.06);
            margin-bottom: 0.75rem;
        }
        .sci-slide-team-photo {
            width: 118px;
            height: 132px;
            margin: 0 auto 0.75rem auto;
            border-radius: 14px;
            overflow: hidden;
            background: linear-gradient(180deg, #dce9f5 0%, #cbdceb 100%);
            display: flex;
            align-items: center;
            justify-content: center;
        }
        .sci-slide-team-photo img {
            width: 100%;
            height: 100%;
            object-fit: cover;
        }
        .sci-slide-team-name {
            font-size: 0.96rem;
            font-weight: 900;
            color: #12344d;
            text-align: center;
            margin-bottom: 0.35rem;
        }
        .sci-slide-team-role {
            font-size: 0.74rem;
            font-weight: 800;
            text-transform: uppercase;
            letter-spacing: 0.08em;
            color: #0f766e;
            text-align: center;
            margin-bottom: 0.45rem;
        }
        .sci-slide-team-desc {
            font-size: 0.82rem;
            color: #4b5f75;
            line-height: 1.48;
            text-align: center;
            margin-bottom: 0.7rem;
        }
        .sci-slide-team-focus {
            border-radius: 14px;
            background: rgba(15, 118, 110, 0.07);
            border: 1px solid rgba(15, 118, 110, 0.12);
            padding: 10px 11px;
            font-size: 0.8rem;
            color: #215260;
            line-height: 1.45;
        }
        .sci-slide-footer {
            margin-top: 0.9rem;
            padding: 0.7rem 0.85rem;
            border-radius: 16px;
            background: linear-gradient(180deg, #f8fbfd 0%, #edf4f8 100%);
            border: 1px solid rgba(18, 52, 77, 0.10);
            color: #41576f;
            font-size: 0.82rem;
            line-height: 1.5;
        }
        .sci-slide-notes {
            border-radius: 14px;
            padding: 11px 13px;
            background: #f8fafc;
            border: 1px dashed rgba(71, 85, 105, 0.28);
            font-size: 0.82rem;
            color: #475569;
            margin-top: 0.75rem;
        }
        @media (max-width: 1100px) {
            .sci-slide-shell {
                padding: 1rem;
            }
            .sci-slide-title {
                font-size: 1.42rem;
            }
        }
        @media (max-width: 768px) {
            .sci-slide-shell {
                padding: 0.95rem;
            }
            .sci-slide-title {
                font-size: 1.26rem;
            }
            .sci-slide-headline,
            .sci-slide-subtitle,
            .sci-slide-card-body,
            .sci-slide-flow-body,
            .sci-slide-team-desc,
            .sci-slide-team-focus,
            .sci-slide-example-body,
            .sci-slide-notes {
                font-size: 0.78rem;
            }
        }
    </style>
        """,
        unsafe_allow_html=True,
    )


def _render_doc_apresentacao_cards(cards, columns_per_row=3) -> None:
    if not cards:
        return
    chunk = max(1, min(columns_per_row, len(cards)))
    for start in range(0, len(cards), chunk):
        row = cards[start:start + chunk]
        cols = st.columns(len(row))
        for col, card in zip(cols, row):
            with col:
                extra_class = ' sci-slide-card-strong' if card.get('emphasis') else ''
                st.markdown(
                    f"""
                <div class="sci-slide-card{extra_class}">
                    <div class="sci-slide-card-title">{card['title']}</div>
                    <div class="sci-slide-card-body">{card['body'].replace(chr(10), '<br>')}</div>
                </div>
                    """,
                    unsafe_allow_html=True,
                )


def _render_doc_apresentacao_flow(steps) -> None:
    if not steps:
        return
    per_row = 3 if len(steps) > 4 else len(steps)
    for start in range(0, len(steps), per_row):
        row = steps[start:start + per_row]
        cols = st.columns(len(row))
        for offset, (col, step) in enumerate(zip(cols, row), start=start + 1):
            with col:
                st.markdown(
                    f"""
                <div class="sci-slide-flow-step">
                    <div class="sci-slide-flow-badge">{offset:02d}</div>
                    <div class="sci-slide-flow-title">{step['title']}</div>
                    <div class="sci-slide-flow-body">{step['body']}</div>
                </div>
                    """,
                    unsafe_allow_html=True,
                )


def _render_doc_apresentacao_example(title, body) -> None:
    if not title and not body:
        return
    st.markdown(
        f"""
    <div class="sci-slide-example">
        <div class="sci-slide-example-title">{title}</div>
        <div class="sci-slide-example-body">{body}</div>
    </div>
        """,
        unsafe_allow_html=True,
    )


def _render_doc_apresentacao_formula(formula: str) -> None:
    if not formula:
        return
    st.latex(formula)


def _render_doc_apresentacao_metric_cards(metrics) -> None:
    if not metrics:
        return
    cols = st.columns(len(metrics))
    for col, metric in zip(cols, metrics):
        with col:
            st.markdown(
                f"""
            <div class="sci-slide-metric">
                <div class="sci-slide-metric-label">{metric['label']}</div>
                <div class="sci-slide-metric-value">{metric['value']}</div>
                <div class="sci-slide-metric-note">{metric['note']}</div>
            </div>
                """,
                unsafe_allow_html=True,
            )


def _render_doc_apresentacao_bar_compare(charts) -> None:
    if not charts:
        return
    cols = st.columns(len(charts))
    for col, chart in zip(cols, charts):
        bars = chart.get('bars', [])
        max_value = max((item.get('value', 0) for item in bars), default=1) or 1
        bars_html = []
        for item in bars:
            width_pct = max(10, int((item.get('value', 0) / max_value) * 100))
            bars_html.append(
                (
                    '<div class="sci-slide-compare-row">'
                    f'<div class="sci-slide-compare-label">{item["label"]}</div>'
                    '<div class="sci-slide-compare-track">'
                    f'<div class="sci-slide-compare-fill" style="width:{width_pct}%; background:{item.get("color", "#2c7a7b")};"></div>'
                    '</div>'
                    f'<div class="sci-slide-compare-value">{item["display"]}</div>'
                    '</div>'
                )
            )
        with col:
            st.markdown(
                (
                    '<div class="sci-slide-compare-card">'
                    f'<div class="sci-slide-compare-title">{chart["title"]}</div>'
                    f'<div class="sci-slide-compare-bars">{"".join(bars_html)}</div>'
                    f'<div class="sci-slide-compare-caption">{chart.get("caption", "")}</div>'
                    '</div>'
                ),
                unsafe_allow_html=True,
            )


def _render_doc_apresentacao_pillar_cards(items) -> None:
    if not items:
        return
    cols = st.columns(len(items))
    for col, item in zip(cols, items):
        with col:
            st.markdown(
                f"""
            <div class="sci-slide-pillar-card">
                <div class="sci-slide-pillar-title">{item['title']}</div>
                <div class="sci-slide-pillar-note">{item['body']}</div>
            </div>
                """,
                unsafe_allow_html=True,
            )


def _render_doc_apresentacao_stage_cards(items) -> None:
    if not items:
        return
    cols = st.columns(len(items))
    for col, item in zip(cols, items):
        with col:
            st.markdown(
                f"""
            <div class="sci-slide-stage-card">
                <div class="sci-slide-stage-title">{item['title']}</div>
                <div class="sci-slide-stage-note">{item['body']}</div>
            </div>
                """,
                unsafe_allow_html=True,
            )


def _get_doc_apresentacao_team_members():
    dados_equipe = carregar_dados_equipe()
    membros = [
        {
            'key': 'hudson',
            'nome': 'Hudson Cardin',
            'foco': 'Arquitetura, lógica de negócio, interface e integração ponta a ponta do SCI.',
        },
        {
            'key': 'lauro',
            'nome': 'Lauro Paiva Junior',
            'foco': 'Experiência analítica, visualização, consistência dos indicadores e evolução funcional.',
        },
        {
            'key': 'frederico',
            'nome': 'Frederico Cesar de Jesus',
            'foco': 'Direção funcional, aderência à controladoria e validação executiva do produto.',
        },
    ]

    resultado = []
    for membro in membros:
        dados_m = dados_equipe.get(membro['key'], {})
        papel = dados_m.get('papel_projeto') or dados_m.get('cargo') or 'Equipe SCI'
        descricao = dados_m.get('descricao_papel') or membro['foco']
        foto_src = None
        if dados_m.get('foto'):
            foto_src = 'data:image/jpeg;base64,' + dados_m['foto']

        resultado.append(
            {
                'nome': membro['nome'],
                'papel': papel,
                'descricao': descricao,
                'foco': membro['foco'],
                'foto': foto_src,
            }
        )
    return resultado


def _build_doc_apresentacao_team_slide():
    membros = _get_doc_apresentacao_team_members()
    return {
        'numero': 2,
        'titulo': '👥 Quem Sustenta a Plataforma',
        'kicker': 'EQUIPE DO PROJETO · SCI',
        'headline': 'O SCI foi construído por um núcleo enxuto que cobre produto, lógica econômica, experiência analítica e direção funcional.',
        'subtitle': 'ETAPA 1 — Plataforma e Fundamentos',
        'sections': [
            {
                'kind': 'bullets',
                'title': 'Mensagem do slide',
                'items': [
                    'Time pequeno, mas com cobertura direta de produto, cálculo, interface e direção funcional.',
                    'Cada integrante fecha uma lacuna crítica para transformar dado industrial em narrativa executiva confiável.',
                    'A combinação do time explica velocidade de execução com aderência ao negócio.',
                ],
            },
            {
                'kind': 'team',
                'title': 'Núcleo que sustenta o SCI',
                'items': membros,
                'columns': 3,
            },
        ],
        'notes': 'Apresentar a equipe como prova de execução: produto, motor analítico e direção funcional convivendo no mesmo núcleo.',
        'layout': 'team-grid',
    }


def _build_doc_apresentacao_tech_slides():
    return [
        {
            'numero': 3,
            'titulo': '🧱 Arquitetura da Plataforma',
            'kicker': 'PLATAFORMA TÉCNICA INTEGRADA',
            'headline': 'O SCI já opera como sistema de engenharia, dados, cloud e inteligência aplicada.',
            'subtitle': 'ETAPA 1 — Plataforma e Fundamentos',
            'sections': [
                {
                    'kind': 'metrics',
                    'title': 'Escala do sistema',
                    'items': [
                        {'label': 'Arquivos core', 'value': '62', 'note': 'núcleo principal da aplicação'},
                        {'label': 'Linhas core', 'value': '84k+', 'note': 'código Python no núcleo do produto'},
                        {'label': 'Telas', 'value': '15', 'note': 'jornadas Streamlit em operação'},
                        {'label': 'Blocos', 'value': '7', 'note': 'camadas centrais de responsabilidade'},
                    ],
                },
                {
                    'kind': 'bar_compare',
                    'title': 'Comparativo visual',
                    'items': [
                        {
                            'title': 'Arquivos Python',
                            'caption': '62 no núcleo principal e 296 no workspace ampliado atual.',
                            'bars': [
                                {'label': 'Núcleo', 'value': 62, 'display': '62', 'color': '#2c7a7b'},
                                {'label': 'Workspace', 'value': 296, 'display': '296', 'color': '#164e63'},
                            ],
                        },
                        {
                            'title': 'Linhas de código Python',
                            'caption': '84 mil+ no núcleo do produto e 306 mil+ no ecossistema ampliado.',
                            'bars': [
                                {'label': 'Núcleo', 'value': 84674, 'display': '84k+', 'color': '#1d4ed8'},
                                {'label': 'Workspace', 'value': 306426, 'display': '306k+', 'color': '#2563eb'},
                            ],
                        },
                    ],
                },
                {
                    'kind': 'cards',
                    'title': 'Stack tecnológico',
                    'items': [
                        {'title': '🐍 Python', 'body': 'Linguagem-base do SCI, unificando regras, processamento, integração e app.'},
                        {'title': '📺 Streamlit', 'body': 'Camada de interface analítica com velocidade alta de evolução de produto.'},
                        {'title': '☁️ Databricks Apps', 'body': 'Execução cloud do app com publicação e operação corporativa.'},
                        {'title': '❄️ Snowflake', 'body': 'Camada de dados e integração analítica no ecossistema ampliado.'},
                        {'title': '🧬 Git / GitHub', 'body': 'Versionamento, governança de mudança e rastreabilidade do código.'},
                        {'title': '🤖 Serving Endpoints', 'body': 'Base de APIs e recursos de IA usados pelo TC Copilot.', 'emphasis': True},
                        {'title': '🐼 pandas + numpy', 'body': 'Transformação, modelagem e manipulação intensiva de dados.'},
                        {'title': '📈 Altair + Plotly', 'body': 'Visualização interativa, comparativos e leitura executiva dos números.'},
                        {'title': '📦 openpyxl + PyArrow', 'body': 'Exportação Excel e persistência rápida em Parquet.'},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'chips',
                    'title': 'Atributos da base',
                    'items': [
                        'Dados estruturados',
                        'Cloud ready',
                        'Forecast + alertas',
                        'TC Copilot',
                        'Governança e rastreabilidade',
                    ],
                    'intro': 'A base técnica já foi desenhada para suportar leitura executiva rápida sem perder profundidade operacional nem histórico de decisão.',
                },
            ],
            'notes': 'Abrir a dimensão técnica com a tese central da Visão Geral Técnica: o SCI já é uma plataforma integrada, não apenas uma interface analítica.',
            'layout': 'tech-overview',
        },
        {
            'numero': 4,
            'titulo': '🧱 Plataforma em Operação e Escala',
            'kicker': 'ESPINHA DORSAL E MATURIDADE',
            'headline': 'A espinha dorsal do SCI conecta pipeline, persistência, análise, interface e governança em uma mesma trilha de produto.',
            'subtitle': 'ETAPA 1 — Plataforma e Fundamentos',
            'sections': [
                {
                    'kind': 'flow',
                    'title': 'Arquitetura em uma leitura',
                    'items': [
                        {'title': 'Fonte', 'body': 'Excel, bases corporativas e insumos operacionais entram com estrutura controlada.'},
                        {'title': 'Processamento', 'body': 'Normalização, transformação, persistência e publicação em fluxo repetível.'},
                        {'title': 'Armazenamento', 'body': 'Parquets, históricos, budget e forecast sustentam leitura rápida e auditável.'},
                        {'title': 'Análise', 'body': 'Camadas analíticas transformam o dado bruto em sinal executivo e operacional.'},
                        {'title': 'Interface', 'body': 'A experiência do usuário integra leitura, comparação, forecast, alertas e IA.'},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'pillar_cards',
                    'title': 'Pilares do produto',
                    'items': [
                        {'title': 'Dados', 'body': 'bases, histórico e persistência'},
                        {'title': 'Regras', 'body': 'contratos e lógica compartilhada'},
                        {'title': 'Visualização', 'body': 'dashboards e leitura executiva'},
                        {'title': 'Simulação', 'body': 'forecast e cenários'},
                        {'title': 'IA', 'body': 'copilot e resposta inteligente'},
                    ],
                },
            ],
            'notes': 'Fechar a dupla técnica mostrando a espinha dorsal do produto e os pilares visíveis que organizam o SCI como plataforma.',
            'layout': 'tech-overview',
        },
    ]


def _hex_to_rgb_tuple(hex_value: str):
    value = hex_value.lstrip('#')
    return tuple(int(value[index:index + 2], 16) for index in (0, 2, 4))


def _format_doc_formula_for_ppt(formula: str) -> str:
    if not formula:
        return ''
    import re

    formatted = formula
    while '\\frac{' in formatted:
        formatted = re.sub(r'\\frac\{([^{}]+)\}\{([^{}]+)\}', r'(\1 / \2)', formatted)
    replacements = {
        '\\times': ' × ',
        '\\Delta': 'Δ',
        '\\bar{H}': 'H médio',
        '\\': ' ',
        '{': '',
        '}': '',
        '_': ' ',
    }
    for old, new in replacements.items():
        formatted = formatted.replace(old, new)
    return ' '.join(formatted.split())


def _ppt_add_text(shape, text, font_size=12, bold=False, color=(23, 32, 51), level=0):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    text_frame = shape.text_frame
    if not text_frame.paragraphs:
        paragraph = text_frame.add_paragraph()
    else:
        paragraph = text_frame.paragraphs[0]
    paragraph.level = level
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return paragraph


def _ppt_create_box(slide, left, top, width, height, fill_rgb, line_rgb, radius=True):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches

    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(*line_rgb)
    shape.text_frame.word_wrap = True
    shape.text_frame.margin_left = Inches(0.06)
    shape.text_frame.margin_right = Inches(0.06)
    shape.text_frame.margin_top = Inches(0.04)
    shape.text_frame.margin_bottom = Inches(0.04)
    return shape


def _ppt_add_picture(slide, image_src: str, left, top, width, height):
    from pptx.util import Inches

    if not image_src or ',' not in image_src:
        return None
    try:
        raw = base64.b64decode(image_src.split(',', 1)[1])
        return slide.shapes.add_picture(BytesIO(raw), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    except Exception:
        return None


def _ppt_add_textbox(slide, left, top, width, height):
    from pptx.util import Inches

    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    textbox.text_frame.word_wrap = True
    textbox.text_frame.margin_left = Inches(0.02)
    textbox.text_frame.margin_right = Inches(0.02)
    textbox.text_frame.margin_top = Inches(0.01)
    textbox.text_frame.margin_bottom = Inches(0.01)
    return textbox


def _ppt_add_avatar_fallback(slide, left, top, width, height, initials, fill_rgb, text_rgb):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    shape = _ppt_create_box(slide, left, top, width, height, fill_rgb, fill_rgb)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = initials
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*text_rgb)


def _ppt_render_section_label(slide, top, title, color):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    textbox = _ppt_add_textbox(slide, 0.28, top, 12.68, 0.15)
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title.upper()
    run.font.size = Pt(7.4)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*color)


def _ppt_render_bullets(slide, top, items, line_rgb, body_rgb, variant=None):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    stage_variant = variant in {'stage-1', 'stage-2'}
    if stage_variant:
        height = 0.56 + 0.38 * len(items)
        box_fill = _hex_to_rgb_tuple('#173554')
        font_size = 16
        font_color = (255, 255, 255)
        line_color = _hex_to_rgb_tuple('#4c6b8b')
    else:
        height = 0.32 + 0.22 * len(items)
        box_fill = (255, 255, 255)
        font_size = 8.8
        font_color = body_rgb
        line_color = line_rgb

    box = _ppt_create_box(slide, 0.28, top, 12.68, height, box_fill, line_color)
    tf = box.text_frame
    tf.clear()
    for index, item in enumerate(items):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = RGBColor(*font_color)
    return height


def _ppt_render_cards(slide, top, items, columns, soft_rgb, accent_rgb, line_rgb, ink_rgb, body_rgb, compact=False):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    columns = max(1, min(columns, 4, len(items)))
    rows = (len(items) + columns - 1) // columns
    gap = 0.1 if compact else 0.14
    width = (12.68 - (columns - 1) * gap) / columns
    height = 0.68 if compact else 0.92
    title_font = 7.8 if compact else 8.8
    body_font = 6.6 if compact else 7.3
    row_gap = 0.1 if compact else 0.12
    for index, item in enumerate(items):
        row = index // columns
        col = index % columns
        left = 0.28 + col * (width + gap)
        current_top = top + row * (height + row_gap)
        fill = accent_rgb if item.get('emphasis') else soft_rgb
        box = _ppt_create_box(slide, left, current_top, width, height, fill, line_rgb)
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = item['title']
        r1.font.size = Pt(title_font)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(*ink_rgb)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = item['body'].replace('\n', ' ')
        r2.font.size = Pt(body_font)
        r2.font.color.rgb = RGBColor(*body_rgb)
    return rows * height + max(0, rows - 1) * row_gap


def _ppt_render_metrics(slide, top, items, teal_rgb):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    gap = 0.12
    width = (12.68 - (len(items) - 1) * gap) / len(items)
    for index, item in enumerate(items):
        left = 0.28 + index * (width + gap)
        box = _ppt_create_box(slide, left, top, width, 0.86, teal_rgb, teal_rgb)
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = item['label'].upper()
        r1.font.size = Pt(7.2)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(213, 243, 244)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = item['value']
        r2.font.size = Pt(15.5)
        r2.font.bold = True
        r2.font.color.rgb = RGBColor(255, 255, 255)
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.LEFT
        r3 = p3.add_run()
        r3.text = item['note']
        r3.font.size = Pt(7.0)
        r3.font.color.rgb = RGBColor(230, 255, 251)
    return 0.86


def _ppt_render_bar_compare(slide, top, charts, line_rgb, ink_rgb, body_rgb, compact=False):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    if not charts:
        return 0

    gap = 0.18
    width = (12.68 - (len(charts) - 1) * gap) / len(charts)
    height = 1.08 if compact else 1.34
    for index, chart in enumerate(charts):
        left = 0.28 + index * (width + gap)
        box = _ppt_create_box(slide, left, top, width, height, (251, 253, 255), line_rgb)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = chart['title']
        r.font.size = Pt(8.6)
        r.font.bold = True
        r.font.color.rgb = RGBColor(*ink_rgb)

        bars = chart.get('bars', [])
        max_value = max((item.get('value', 0) for item in bars), default=1) or 1
        current_top = top + (0.28 if compact else 0.34)
        for item in bars:
            textbox = _ppt_add_textbox(slide, left + 0.12, current_top, width - 0.24, 0.14)
            tft = textbox.text_frame
            tft.clear()
            p = tft.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = item['label']
            r.font.size = Pt(6.6 if compact else 7.0)
            r.font.bold = True
            r.font.color.rgb = RGBColor(*body_rgb)

            track_left = left + 0.98
            track_top = current_top + 0.025
            track_width = width - 1.9
            track = slide.shapes.add_shape(1, Inches(track_left), Inches(track_top), Inches(track_width), Inches(0.08))
            track.fill.solid()
            track.fill.fore_color.rgb = RGBColor(231, 238, 245)
            track.line.color.rgb = RGBColor(231, 238, 245)

            fill_width = max(0.35, track_width * (item.get('value', 0) / max_value))
            fill = slide.shapes.add_shape(1, Inches(track_left), Inches(track_top), Inches(fill_width), Inches(0.08))
            fill.fill.solid()
            fill.fill.fore_color.rgb = RGBColor(*_hex_to_rgb_tuple(item.get('color', '#2c7a7b')))
            fill.line.color.rgb = RGBColor(*_hex_to_rgb_tuple(item.get('color', '#2c7a7b')))

            value_box = _ppt_add_textbox(slide, left + width - 0.72, current_top, 0.6, 0.16)
            tfv = value_box.text_frame
            tfv.clear()
            p = tfv.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            r = p.add_run()
            r.text = item['display']
            r.font.size = Pt(6.6 if compact else 7.0)
            r.font.bold = True
            r.font.color.rgb = RGBColor(*ink_rgb)
            current_top += 0.16 if compact else 0.2

        caption_box = _ppt_add_textbox(slide, left + 0.12, top + height - 0.2, width - 0.24, 0.14)
        tfc = caption_box.text_frame
        tfc.clear()
        p = tfc.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = chart.get('caption', '')
        r.font.size = Pt(6.2 if compact else 6.8)
        r.font.color.rgb = RGBColor(98, 125, 152)
    return height


def _ppt_render_pillar_cards(slide, top, items, line_rgb, ink_rgb, body_rgb):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    if not items:
        return 0
    gap = 0.12
    width = (12.68 - (len(items) - 1) * gap) / len(items)
    height = 0.64
    for index, item in enumerate(items):
        left = 0.28 + index * (width + gap)
        box = _ppt_create_box(slide, left, top, width, height, (248, 251, 255), line_rgb)
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = item['title']
        r1.font.size = Pt(8.4)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(*ink_rgb)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = item['body']
        r2.font.size = Pt(6.9)
        r2.font.color.rgb = RGBColor(*body_rgb)
    return height


def _ppt_render_stage_cards(slide, top, items, line_rgb):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    if not items:
        return 0
    gap = 0.1
    width = (12.68 - (len(items) - 1) * gap) / len(items)
    height = 0.86
    for index, item in enumerate(items):
        left = 0.28 + index * (width + gap)
        box = _ppt_create_box(slide, left, top, width, height, (255, 244, 214), line_rgb)
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = item['title']
        r1.font.size = Pt(7.5)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(141, 91, 0)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = item['body']
        r2.font.size = Pt(6.6)
        r2.font.color.rgb = RGBColor(124, 94, 16)
    return height


def _ppt_render_team(slide, top, items, soft_rgb, line_rgb, ink_rgb, body_rgb):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    gap = 0.16
    width = (12.68 - (len(items) - 1) * gap) / len(items)
    height = 2.45
    for index, item in enumerate(items):
        left = 0.28 + index * (width + gap)
        box = _ppt_create_box(slide, left, top, width, height, soft_rgb, line_rgb)
        picture = _ppt_add_picture(slide, item.get('foto'), left + width / 2 - 0.45, top + 0.12, 0.9, 1.0)
        if picture is None:
            initials = ''.join(part[0] for part in item['nome'].split()[:2]).upper()
            _ppt_add_avatar_fallback(slide, left + width / 2 - 0.38, top + 0.14, 0.76, 0.82, initials, (220, 233, 245), ink_rgb)
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = '\n\n\n\n' + item['nome']
        r1.font.size = Pt(9.4)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(*ink_rgb)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = item['papel']
        r2.font.size = Pt(7.0)
        r2.font.bold = True
        r2.font.color.rgb = RGBColor(15, 118, 110)
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = item['descricao']
        r3.font.size = Pt(7.2)
        r3.font.color.rgb = RGBColor(*body_rgb)
        p4 = tf.add_paragraph()
        r4 = p4.add_run()
        r4.text = 'Contribuição central: ' + item['foco']
        r4.font.size = Pt(6.9)
        r4.font.color.rgb = RGBColor(33, 82, 96)
    return height


def _ppt_render_flow(slide, top, items, columns, teal_rgb, compact=False):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    columns = max(1, min(columns, 3, len(items)))
    rows = (len(items) + columns - 1) // columns
    gap = 0.14
    width = (12.68 - (columns - 1) * gap) / columns
    height = 0.72 if compact else 0.8
    row_gap = 0.08 if compact else 0.1
    title_font = 7.8 if compact else 8.0
    body_font = 6.8 if compact else 7.0
    for index, item in enumerate(items):
        row = index // columns
        col = index % columns
        left = 0.28 + col * (width + gap)
        current_top = top + row * (height + row_gap)
        box = _ppt_create_box(slide, left, current_top, width, height, teal_rgb, teal_rgb)
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = f"{index + 1:02d}  {item['title']}"
        r1.font.size = Pt(title_font)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(255, 255, 255)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = item['body']
        r2.font.size = Pt(body_font)
        r2.font.color.rgb = RGBColor(236, 254, 255)
    return rows * height + max(0, rows - 1) * row_gap


def _ppt_render_formula(slide, top, formula, soft_rgb, line_rgb, ink_rgb, height=0.64):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    box = _ppt_create_box(slide, 0.28, top, 12.68, height, soft_rgb, line_rgb)
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = _format_doc_formula_for_ppt(formula)
    r.font.size = Pt(8.8)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*ink_rgb)
    return height


def _ppt_render_example(slide, top, title, body, accent_rgb, line_rgb, height=0.68):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    box = _ppt_create_box(slide, 0.28, top, 12.68, height, accent_rgb, line_rgb)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = title + ': '
    r1.font.size = Pt(8.4)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(141, 91, 0)
    r2 = p.add_run()
    r2.text = body
    r2.font.size = Pt(7.8)
    r2.font.color.rgb = RGBColor(124, 94, 16)
    return height


def _ppt_render_chips(slide, top, intro, items, line_rgb, ink_rgb, body_rgb):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    current_top = top
    if intro:
        textbox = _ppt_add_textbox(slide, 0.28, current_top, 12.68, 0.26)
        tf = textbox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = intro
        r.font.size = Pt(7.8)
        r.font.color.rgb = RGBColor(*body_rgb)
        current_top += 0.24
    left = 0.28
    top_chip = current_top
    for item in items:
        width = max(1.15, min(2.0, 0.09 * len(item) + 0.5))
        if left + width > 12.75:
            left = 0.28
            current_top += 0.28
        chip = _ppt_create_box(slide, left, current_top, width, 0.22, (241, 250, 249), line_rgb)
        tf = chip.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = item
        r.font.size = Pt(7.0)
        r.font.bold = True
        r.font.color.rgb = RGBColor(*ink_rgb)
        left += width + 0.08
    return current_top - top_chip + 0.22 + (0.24 if intro else 0)


def _render_doc_apresentacao_team_slide(membros, columns_per_row=3) -> None:
    if not membros:
        return
    cols = st.columns(min(columns_per_row, len(membros)))
    for col, membro in zip(cols, membros):
        foto_html = (
            f'<img src="{membro["foto"]}" alt="{membro["nome"]}" />'
            if membro.get('foto')
            else '<div style="font-size:2.3rem;color:#5f7287;">👤</div>'
        )
        with col:
            st.markdown(
                f"""
            <div class="sci-slide-team-card">
                <div class="sci-slide-team-photo">{foto_html}</div>
                <div class="sci-slide-team-name">{membro['nome']}</div>
                <div class="sci-slide-team-role">{membro['papel']}</div>
                <div class="sci-slide-team-desc">{membro['descricao']}</div>
                <div class="sci-slide-team-focus"><strong>Contribuição central:</strong> {membro['foco']}</div>
            </div>
                """,
                unsafe_allow_html=True,
            )


def _render_doc_apresentacao_chips(intro, chips) -> None:
    if intro:
        st.markdown(f'<div class="sci-slide-card-body" style="margin-bottom:0.55rem;">{intro}</div>', unsafe_allow_html=True)
    if not chips:
        return
    chips_html = ''.join([f'<span class="sci-slide-chip">{chip}</span>' for chip in chips])
    st.markdown(chips_html, unsafe_allow_html=True)


def _render_doc_apresentacao_section(section) -> None:
    title = section.get('title')
    if title:
        st.markdown(f'<div class="sci-slide-section-title">{title}</div>', unsafe_allow_html=True)

    kind = section.get('kind')
    if kind == 'bullets':
        bullets_html = ''.join([f'<li>{item}</li>' for item in section.get('items', [])])
        st.markdown(f'<ul class="sci-slide-bullets">{bullets_html}</ul>', unsafe_allow_html=True)
    elif kind == 'cards':
        _render_doc_apresentacao_cards(section.get('items', []), columns_per_row=section.get('columns', 3))
    elif kind == 'metrics':
        _render_doc_apresentacao_metric_cards(section.get('items', []))
    elif kind == 'bar_compare':
        _render_doc_apresentacao_bar_compare(section.get('items', []))
    elif kind == 'pillar_cards':
        _render_doc_apresentacao_pillar_cards(section.get('items', []))
    elif kind == 'stage_cards':
        _render_doc_apresentacao_stage_cards(section.get('items', []))
    elif kind == 'team':
        _render_doc_apresentacao_team_slide(section.get('items', []), columns_per_row=section.get('columns', 3))
    elif kind == 'flow':
        _render_doc_apresentacao_flow(section.get('items', []))
    elif kind == 'formula':
        _render_doc_apresentacao_formula(section.get('value', ''))
    elif kind == 'example':
        _render_doc_apresentacao_example(section.get('example_title') or section.get('title') or 'Exemplo executivo', section.get('body', ''))
    elif kind == 'chips':
        _render_doc_apresentacao_chips(section.get('intro'), section.get('items', []))


def _get_doc_apresentacao_ppt_bytes() -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    primary = _hex_to_rgb_tuple('#0b1f2a')
    teal = _hex_to_rgb_tuple('#164e63')
    mint = _hex_to_rgb_tuple('#7dd3c8')
    ink = _hex_to_rgb_tuple('#12344d')
    body = _hex_to_rgb_tuple('#486581')
    line = _hex_to_rgb_tuple('#d7e2eb')
    soft = _hex_to_rgb_tuple('#eef6fb')
    accent = _hex_to_rgb_tuple('#fff1da')
    note_fill = _hex_to_rgb_tuple('#f8fafc')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = _get_doc_apresentacao_slides()
    total_slides = len(slides)
    for slide_spec in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        shell_variant = slide_spec.get('shell_variant')
        hero_fill = primary
        hero_line = mint
        section_label_color = (15, 76, 92)
        if shell_variant == 'stage-1':
            stage_blue = _hex_to_rgb_tuple('#122A46')
            background = _ppt_create_box(slide, 0.08, 0.08, 13.17, 7.34, stage_blue, stage_blue, radius=False)
            background.line.color.rgb = RGBColor(*stage_blue)
            hero_fill = _hex_to_rgb_tuple('#173554')
            hero_line = _hex_to_rgb_tuple('#8fb3d9')
            section_label_color = (223, 235, 255)
        elif shell_variant == 'stage-2':
            stage_blue = _hex_to_rgb_tuple('#122A46')
            background = _ppt_create_box(slide, 0.08, 0.08, 13.17, 7.34, stage_blue, stage_blue, radius=False)
            background.line.color.rgb = RGBColor(*stage_blue)
            hero_fill = _hex_to_rgb_tuple('#1c3f63')
            hero_line = _hex_to_rgb_tuple('#8fb3d9')
            section_label_color = (223, 235, 255)

        hero = _ppt_create_box(slide, 0.28, 0.22, 12.68, 1.34, hero_fill, hero_fill)
        hero.fill.fore_color.rgb = RGBColor(*hero_fill)
        hero.line.color.rgb = RGBColor(*hero_line)
        tf = hero.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = slide_spec.get('kicker', '')
        r1.font.size = Pt(8.4)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(223, 235, 255) if shell_variant in {'stage-1', 'stage-2'} else RGBColor(204, 251, 241)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = slide_spec.get('titulo', '')
        r2.font.size = Pt(17)
        r2.font.bold = True
        r2.font.color.rgb = RGBColor(255, 255, 255)
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.LEFT
        r3 = p3.add_run()
        r3.text = slide_spec.get('headline', '')
        r3.font.size = Pt(9.2)
        r3.font.bold = True
        r3.font.color.rgb = RGBColor(245, 249, 255) if shell_variant in {'stage-1', 'stage-2'} else RGBColor(240, 253, 250)
        p4 = tf.add_paragraph()
        p4.alignment = PP_ALIGN.LEFT
        r4 = p4.add_run()
        r4.text = slide_spec.get('subtitle', '')
        r4.font.size = Pt(8.0)
        r4.font.color.rgb = RGBColor(214, 230, 255) if shell_variant in {'stage-1', 'stage-2'} else RGBColor(236, 254, 255)

        current_top = 1.72
        sections = slide_spec.get('sections', [])
        compact_slide_titles = {
            '🌊 Como o Desvio é Explicado',
            '🧮 TC Veículos: Rateio e Real',
            '🔮 Best Estimate',
        }
        compact_slide = slide_spec.get('titulo', '') in compact_slide_titles
        tech_overview_slide = slide_spec.get('titulo', '') == '🧱 Arquitetura da Plataforma'
        for index, section in enumerate(sections):
            kind = section.get('kind')
            show_section_title = section.get('title') and not (compact_slide and kind in {'formula', 'example'})
            if tech_overview_slide and kind in {'metrics', 'bar_compare'}:
                show_section_title = False
            if show_section_title:
                _ppt_render_section_label(slide, current_top, section['title'], section_label_color)
                current_top += 0.15

            if kind == 'bullets':
                used_height = _ppt_render_bullets(slide, current_top, section.get('items', []), line, body, variant=shell_variant)
            elif kind == 'cards':
                used_height = _ppt_render_cards(
                    slide,
                    current_top,
                    section.get('items', []),
                    section.get('columns', 3),
                    soft,
                    accent,
                    line,
                    ink,
                    body,
                    compact=(tech_overview_slide and section.get('title') == 'Stack tecnológico') or compact_slide,
                )
            elif kind == 'metrics':
                used_height = _ppt_render_metrics(slide, current_top, section.get('items', []), teal)
            elif kind == 'bar_compare':
                used_height = _ppt_render_bar_compare(
                    slide,
                    current_top,
                    section.get('items', []),
                    line,
                    ink,
                    body,
                    compact=tech_overview_slide,
                )
            elif kind == 'pillar_cards':
                used_height = _ppt_render_pillar_cards(slide, current_top, section.get('items', []), line, ink, body)
            elif kind == 'stage_cards':
                used_height = _ppt_render_stage_cards(slide, current_top, section.get('items', []), line)
            elif kind == 'team':
                used_height = _ppt_render_team(slide, current_top, section.get('items', []), soft, line, ink, body)
            elif kind == 'flow':
                used_height = _ppt_render_flow(
                    slide,
                    current_top,
                    section.get('items', []),
                    section.get('columns', 3),
                    teal,
                    compact=compact_slide,
                )
            elif kind == 'formula':
                used_height = _ppt_render_formula(
                    slide,
                    current_top,
                    section.get('value', ''),
                    soft,
                    line,
                    ink,
                    height=0.56 if compact_slide else 0.64,
                )
            elif kind == 'example':
                used_height = _ppt_render_example(
                    slide,
                    current_top,
                    section.get('example_title') or section.get('title') or 'Exemplo executivo',
                    section.get('body', ''),
                    accent,
                    line,
                    height=0.58 if compact_slide else 0.68,
                )
            elif kind == 'chips':
                used_height = _ppt_render_chips(slide, current_top, section.get('intro'), section.get('items', []), line, ink, body)
            else:
                used_height = 0
            current_top += used_height
            if index < len(sections) - 1:
                current_top += 0.08 if compact_slide else 0.12

        footer_box = _ppt_create_box(slide, 0.28, 6.28, 12.68, 0.24, (248, 251, 253), line)
        footer_tf = footer_box.text_frame
        footer_tf.clear()
        footer_p = footer_tf.paragraphs[0]
        footer_p.alignment = PP_ALIGN.LEFT
        footer_r = footer_p.add_run()
        footer_r.text = f"Slide {slide_spec['numero']} / {total_slides}  ·  Layout: {slide_spec['layout']}  ·  Fonte única"
        footer_r.font.size = Pt(7.2)
        footer_r.font.bold = True
        footer_r.font.color.rgb = RGBColor(65, 87, 111)

        note_box = _ppt_create_box(slide, 0.28, 6.58, 12.68, 0.44, note_fill, line)
        tf = note_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = 'Speaker note: '
        r.font.size = Pt(8.2)
        r.font.bold = True
        r.font.color.rgb = RGBColor(71, 85, 105)
        r2 = p.add_run()
        r2.text = slide_spec.get('notes', '')
        r2.font.size = Pt(8.2)
        r2.font.color.rgb = RGBColor(71, 85, 105)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _validate_doc_apresentacao_slides(slides):
    layouts_validos = {
        'opening',
        'team-grid',
        'tech-overview',
        'system-map',
        'flow-formula',
        'stage-break',
        'forecast',
        'ai-layer',
        'monitoring',
        'closing',
    }
    for index, slide in enumerate(slides, start=1):
        obrigatorios = {'numero', 'titulo', 'kicker', 'headline', 'subtitle', 'sections', 'notes', 'layout'}
        faltantes = obrigatorios - set(slide.keys())
        if faltantes:
            raise ValueError(f"Slide {index} sem campos obrigatórios: {sorted(faltantes)}")
        if slide['numero'] != index:
            raise ValueError(f"Numeração inválida no slide {index}: {slide['numero']}")
        if slide['layout'] not in layouts_validos:
            raise ValueError(f"Layout inválido no slide {index}: {slide['layout']}")
    return slides


def _renumber_doc_apresentacao_slides(slides):
    for index, slide in enumerate(slides, start=1):
        slide['numero'] = index
    return slides


def _get_doc_apresentacao_slides():
    slides = [
        {
            'numero': 1,
            'titulo': '🎬 ETAPA 1 — O SCI COMO PLATAFORMA DE DECISÃO',
            'kicker': 'ETAPA 1 · PLATAFORMA E FUNDAMENTOS',
            'headline': 'Como o SCI organiza dados, regras e narrativa para transformar custo industrial em decisão executiva',
            'subtitle': '',
            'shell_variant': 'stage-1',
            'hero_variant': 'stage-1',
            'sections': [
                {
                    'kind': 'bullets',
                    'title': 'O que será abordado nesta etapa',
                    'items': [
                        '1. 🎬 O SCI como Plataforma de Decisão',
                        '2. 👥 Quem Sustenta a Plataforma',
                        '3. 🧱 Arquitetura da Plataforma',
                        '4. 🧱 Plataforma em Operação e Escala',
                        '5. 🧭 Como a Plataforma se Organiza',
                        '6. 📥 Como o Dado Entra e se Torna Confiável',
                        '7. 🌊 Como o Desvio é Explicado',
                    ],
                },
            ],
            'notes': 'Este slide serve como contrato narrativo da Etapa 1. Ele prepara o público para entender a base antes dos diferenciais.',
            'layout': 'opening',
        },
        {
            'numero': 2,
            'titulo': '🎬 O SCI como Plataforma de Decisão',
            'kicker': 'VISÃO GERAL · STELLANTIS COST INTELLIGENCE',
            'headline': 'O SCI transforma dado industrial bruto em plataforma de decisão executiva sobre custos.',
            'subtitle': 'ETAPA 1 — Plataforma e Fundamentos',
            'sections': [
                {
                    'kind': 'bullets',
                    'title': 'Mensagem central',
                    'items': [
                        'Um ponto único para custo, volume, flex, waterfall, forecast, alertas e IA aplicada.',
                        'A mesma base técnica atende leitura operacional, conversa gerencial e escala cloud sem retrabalho.',
                        'O SCI reduz o tempo entre detectar o desvio, explicar a causa e coordenar a reação.',
                    ],
                },
                {
                    'kind': 'cards',
                    'title': 'Quatro papéis do SCI',
                    'items': [
                        {'title': '🔍 Entender', 'body': 'Consolida custo, volume e contexto na mesma conversa executiva.'},
                        {'title': '🎯 Explicar', 'body': 'Waterfall, flex e rastreabilidade mostram causa, não só o valor final.'},
                        {'title': '📈 Projetar', 'body': 'Best Estimate transforma histórico em cenário futuro com premissas explícitas.'},
                        {'title': '🔔 Agir', 'body': 'Alertas, relatórios e Copilot encurtam o ciclo entre análise e decisão.'},
                    ],
                    'columns': 4,
                },
            ],
            'notes': 'Abrir posicionando o SCI como sistema de decisão, não como dashboard isolado. Mostrar que o projeto elimina o ciclo manual de consolidação, interpretação e comunicação de custo.',
            'layout': 'opening',
        },
        _build_doc_apresentacao_team_slide(),
        *_build_doc_apresentacao_tech_slides(),
        {
            'numero': 4,
            'titulo': '🧭 Como a Plataforma se Organiza',
            'kicker': 'NAVEGAÇÃO EXECUTIVA',
            'headline': 'A estrutura do SCI leva o usuário do contexto geral ao diagnóstico e à ação em uma trilha única.',
            'subtitle': 'ETAPA 1 — Plataforma e Fundamentos',
            'sections': [
                {
                    'kind': 'bullets',
                    'title': 'Leitura da navegação',
                    'items': [
                        'Home e módulos centrais concentram a leitura rápida do mês e o ponto de partida da análise.',
                        'Waterfall e Best Estimate explicam presente e futuro pela mesma base econômica.',
                        'Alertas, Copilot e documentação fecham governança, comunicação e autonomia do usuário.',
                    ],
                },
                {
                    'kind': 'cards',
                    'title': 'Mapa dos módulos',
                    'items': [
                        {'title': '🏠 Home', 'body': 'KPIs, contexto do mês e desvio rápido.'},
                        {'title': '🏭 TC Ext', 'body': 'Análise por oficina, account, real, budget e flex.'},
                        {'title': '🚗 TC Veículos', 'body': 'Cadeia de custo até o veículo, rateio e CPU.'},
                        {'title': '🌊 Waterfall', 'body': 'Leitura Budget → Flex → Real pela causa econômica.'},
                        {'title': '🔮 Best Estimate', 'body': 'Forecast com premissas configuráveis e ajuste manual.'},
                        {'title': '🚨 Alertas', 'body': 'Prioridade, ranking e comunicação acionável.'},
                        {'title': '🤖 TC Copilot', 'body': 'Resposta rápida e síntese executiva do SCI.'},
                        {'title': '📚 Documentação', 'body': 'Regras, arquitetura, onboarding e memória técnica.'},
                    ],
                    'columns': 4,
                },
            ],
            'notes': 'Conduzir a leitura como arquitetura de produto, não como lista de telas. Cada bloco tem uma função especializada no ciclo de análise.',
            'layout': 'system-map',
        },
        {
            'numero': 5,
            'titulo': '📥 Como o Dado Entra e se Torna Confiável',
            'kicker': 'BASE CONFIÁVEL ANTES DA ANÁLISE',
            'headline': 'O número só ganha legitimidade quando entra validado, tratado e persistido no formato certo.',
            'subtitle': 'ETAPA 1 — Plataforma e Fundamentos',
            'sections': [
                {
                    'kind': 'bullets',
                    'title': 'Por que este slide importa',
                    'items': [
                        'Entrada de Excel e SAP é normalizada antes de qualquer visualização ou fechamento.',
                        'Parquet vira contrato de velocidade, consistência e reuso entre local, executável e cloud.',
                        'Flex Budget já nasce apoiado em custo fixo, custo variável e volume coerente.',
                    ],
                },
                {
                    'kind': 'flow',
                    'title': 'Pipeline do dado',
                    'items': [
                        {'title': 'Excel / SAP', 'body': 'Fontes operacionais e corporativas'},
                        {'title': 'Validação', 'body': 'Abas, colunas e período corretos'},
                        {'title': 'Tratamento', 'body': 'Normalização, merge, correções e rateios'},
                        {'title': 'Parquet', 'body': 'Persistência auditável e rápida'},
                        {'title': 'App', 'body': 'Consumo unificado nos módulos do SCI'},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'formula',
                    'title': 'Fórmula de leitura',
                    'value': r'Flex\ Bud = C_{Fixo} + (C_{Variavel} \times \frac{Volume\ Real}{Volume\ Budget})',
                },
                {
                    'kind': 'example',
                    'title': 'Exemplo executivo',
                    'example_title': 'Exemplo executivo',
                    'body': 'Com R$ 180 mil fixos, R$ 320 mil variáveis, volume budget de 40 mil e volume real de 46 mil, o flex ajustado chega perto de R$ 548 mil.',
                },
            ],
            'notes': 'Enfatizar que a confiança do SCI nasce do pipeline, não do gráfico.',
            'layout': 'flow-formula',
        },
        {
            'numero': 6,
            'titulo': '🌊 Como o Desvio é Explicado',
            'kicker': 'EXPLICAÇÃO ANTES DA REAÇÃO',
            'headline': 'Waterfall traduz o desvio em causa econômica: quanto veio de volume, quanto veio da operação.',
            'subtitle': 'ETAPA 1 — Plataforma e Fundamentos',
            'sections': [
                {
                    'kind': 'bullets',
                    'title': 'Leitura econômica',
                    'items': [
                        'Budget é a referência original do plano e do compromisso financeiro assumido.',
                        'Flex corrige o esperado para o volume efetivamente realizado, isolando o efeito de mix e produção.',
                        'Real mostra o desvio residual que ainda precisa de explicação operacional.',
                    ],
                },
                {
                    'kind': 'cards',
                    'title': 'Três blocos da ponte',
                    'items': [
                        {'title': 'Budget', 'body': 'Ponto de partida do plano original.'},
                        {'title': 'Flex', 'body': 'Valor esperado quando o volume real entra na conta.', 'emphasis': True},
                        {'title': 'Real', 'body': 'Resultado efetivo que precisa ser explicado.'},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'formula',
                    'title': 'Fórmula de leitura',
                    'value': r'\Delta_{Op} = Real - Flex\ Bud',
                },
                {
                    'kind': 'example',
                    'title': 'Exemplo executivo',
                    'example_title': 'Exemplo executivo',
                    'body': 'Budget de R$ 600 mil, Flex de R$ 645 mil e Real de R$ 670 mil mostram dois sinais: R$ 45 mil explicados por volume e R$ 25 mil ainda ligados à operação.',
                },
            ],
            'notes': 'Fechar o slide dizendo que o waterfall muda a conversa de cobrança para diagnóstico.',
            'layout': 'flow-formula',
        },
        {
            'numero': 7,
            'titulo': '🚀 ETAPA 2 — Funcionalidades Ouro',
            'kicker': 'ETAPA 2 · FUNCIONALIDADES OURO',
            'shell_variant': 'stage-2',
            'hero_variant': 'stage-2',
            'headline': 'Onde a plataforma se torna difícil de substituir: granularidade, antecipação, governança e IA',
            'subtitle': '',
            'sections': [
                {
                    'kind': 'bullets',
                    'title': 'O que será abordado nesta etapa',
                    'items': [
                        '1. 🚀 Funcionalidades Avançadas',
                        '2. 🧮 TC Veículos: Rateio e Real',
                        '3. 🔮 Best Estimate',
                        '4. 🤖 TC Copilot',
                        '5. 🤖 TC Copilot — Entregas Executivas',
                        '6. 🚨 Relatórios e Alertas',
                        '7. 🚨 Alertas — Exemplo de Comunicação',
                        '8. ✅ Conclusão',
                    ],
                },
            ],
            'notes': 'Este slide serve como contrato narrativo da Etapa 2. Ele marca a virada entre a base da plataforma e as funcionalidades ouro do SCI.',
            'layout': 'opening',
        },
        {
            'numero': 8,
            'titulo': '🚀 Funcionalidades Avançadas',
            'kicker': 'ETAPA 2 · FUNCIONALIDADES OURO',
            'headline': 'Após a plataforma, o SCI avança para o que o torna difícil de substituir.',
            'subtitle': 'ETAPA 2 — Funcionalidades Ouro',
            'sections': [
                {
                    'kind': 'stage_cards',
                    'title': 'Maturidade do SCI',
                    'items': [
                        {'title': '📚 Base estruturada', 'body': 'Dados consolidados, histórico, budget e governança de persistência.'},
                        {'title': '☁️ Cloud ativa', 'body': 'Execução no Databricks Apps com publicação e sincronização controladas.'},
                        {'title': '🔮 Forecast', 'body': 'Simulação e projeção integradas ao produto.'},
                        {'title': '🚨 Alertas', 'body': 'Monitoramento ativo e priorização do que mais importa.'},
                        {'title': '🤖 TC Copilot', 'body': 'IA aplicada para ampliar interpretação, produtividade e explicação.'},
                    ],
                },
                {
                    'kind': 'bullets',
                    'title': 'O que muda a partir daqui',
                    'items': [
                        'TC Veículos reconstrói a cadeia de custo até o nível do veículo e do CPU auditável.',
                        'Best Estimate projeta cenários futuros com premissas configuráveis e regra econômica explícita.',
                        'Alertas e relatórios transformam análise em monitoramento acionável de rotina.',
                        'TC Copilot reduz o tempo para responder, sintetizar e comunicar o resultado.',
                    ],
                },
                {
                    'kind': 'cards',
                    'title': 'Quatro diferenciais estruturais',
                    'items': [
                        {'title': '🧮 Granularidade', 'body': 'Rateio e custo real até o veículo e CPU.', 'emphasis': True},
                        {'title': '🔮 Antecipação', 'body': 'Forecast com premissas econômicas auditáveis.', 'emphasis': True},
                        {'title': '🚨 Comunicação', 'body': 'Desvio vira alerta priorizado e mensagem pronta.', 'emphasis': True},
                        {'title': '🤖 IA aplicada', 'body': 'Resposta rápida sobre regras, contexto e narrativas.', 'emphasis': True},
                    ],
                    'columns': 4,
                },
            ],
            'notes': 'Usar este slide como transição clara entre visão macro e os diferenciais do SCI.',
            'layout': 'stage-break',
        },
        {
            'numero': 9,
            'titulo': '🧮 TC Veículos: Rateio e Real',
            'kicker': 'GRANULARIDADE PRODUTIVA',
            'headline': 'No TC Veículos, o SCI reconstrói a cadeia econômica até o veículo para explicar onde o custo realmente nasce.',
            'subtitle': 'ETAPA 2 — Funcionalidades Ouro',
            'sections': [
                {
                    'kind': 'bullets',
                    'title': 'Por que isso é diferencial',
                    'items': [
                        'A cadeia separa o que é Redis, o que é FA, o que é FP e o que permanece dedicado ao veículo.',
                        'O rateio usa tempo de produção para distribuir custo onde ele foi consumido.',
                        'O CPU final passa a ser auditável no nível do veículo e não apenas do consolidado da oficina.',
                    ],
                },
                {
                    'kind': 'flow',
                    'title': 'Cadeia de formação do custo',
                    'items': [
                        {'title': 'Despesa primária', 'body': 'Base econômica da oficina'},
                        {'title': 'Redis', 'body': 'Abatimento do componente Redis'},
                        {'title': 'FA', 'body': 'Separação do fluxo anexo'},
                        {'title': 'FP real', 'body': 'Fluxo principal líquido'},
                        {'title': 'Rateio', 'body': 'Distribuição por tempo de produção'},
                        {'title': 'CPU', 'body': 'Custo final por unidade produzida'},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'formula',
                    'title': 'Fórmula base',
                    'value': r'Custo_{FP} = Despesa\ Primaria - Redis - Custo_{FA}',
                },
                {
                    'kind': 'example',
                    'title': 'Exemplo executivo',
                    'example_title': 'Exemplo executivo',
                    'body': 'Com despesa primária de R$ 1,0 mi, Redis de R$ 70 mil, FA de 25%, D&A dedicada de R$ 30 mil e participação do veículo de 40%, o custo final continua rastreável até o CPU.',
                },
            ],
            'notes': 'Posicionar este bloco como diferencial estrutural do SCI frente a análises apenas consolidadas.',
            'layout': 'flow-formula',
        },
        {
            'numero': 9,
            'titulo': '🔮 Best Estimate',
            'kicker': 'ANTECIPAÇÃO E CENÁRIO',
            'headline': 'Best Estimate transforma histórico em previsão operacional com lógica econômica explícita.',
            'subtitle': 'ETAPA 2 — Funcionalidades Ouro',
            'sections': [
                {
                    'kind': 'flow',
                    'title': 'Regras do forecast',
                    'items': [
                        {'title': 'Base histórica', 'body': 'Parte de média válida, volume futuro e sensibilidade aderente ao comportamento real do custo.'},
                        {'title': 'Regra econômica', 'body': 'Aplica efeito de volume primeiro, depois inflação e produtividade, com BE Manual como camada adicional controlada.'},
                        {'title': 'Saída executiva', 'body': 'Entrega forecast consolidado no mesmo formato de leitura usado para analisar o realizado.'},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'cards',
                    'title': 'Motor do Best Estimate',
                    'items': [
                        {'title': '📊 Entradas', 'body': 'Média histórica válida, volume futuro, sensibilidade, inflação e produtividade.'},
                        {'title': '⚙️ Regra', 'body': 'Efeito de volume primeiro, bloco monetário depois. BE Manual entra como complemento.', 'emphasis': True},
                        {'title': '📤 Saídas', 'body': 'Forecast consolidado, histórico separado e forecast rateado por veículo.'},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'formula',
                    'title': 'Lógica resumida',
                    'value': r'BE = \bar{H} \times F_{vol} \times (1 + Infl) \times (1 - Prod)',
                },
                {
                    'kind': 'example',
                    'title': 'Exemplo executivo',
                    'example_title': 'Exemplo executivo',
                    'body': 'Média de R$ 100 mil, volume +15%, sensibilidade 80%, inflação 4% e produtividade 3% levam o BE para algo próximo de R$ 112,9 mil.',
                },
            ],
            'notes': 'Valorizar base limpa + premissas configuráveis + regra econômica transparente. O BE não é um chute: é uma lógica auditável.',
            'layout': 'forecast',
        },
        {
            'numero': 10,
            'titulo': '🤖 TC Copilot',
            'kicker': 'IA APLICADA AO CONTEXTO SCI',
            'headline': 'TC Copilot adiciona uma camada de interpretação e comunicação rápida sobre a mesma base de regras do SCI.',
            'subtitle': 'ETAPA 2 — Funcionalidades Ouro',
            'sections': [
                {
                    'kind': 'bullets',
                    'title': 'Função da camada de IA',
                    'items': [
                        'Chatbot técnico responde com base na documentação e nos contratos do sistema.',
                        'Relatórios diários e executivos encurtam a preparação da comunicação gerencial.',
                        'Abre caminho para operação multi-planta com a mesma linguagem analítica do SCI.',
                    ],
                },
                {
                    'kind': 'cards',
                    'title': 'Casos de uso imediatos',
                    'items': [
                        {'title': '💬 Chatbot', 'body': 'Responde dúvidas sobre regras, arquitetura e operação do produto.'},
                        {'title': '📅 Relatório diário', 'body': 'Resume variações e pontos de atenção com rotina mais curta.'},
                        {'title': '📋 Relatório executivo', 'body': 'Transforma detalhe técnico em síntese gerencial pronta para uso.'},
                        {'title': '🌐 Escala', 'body': 'Facilita replicação entre áreas e futuras plantas usando a mesma base.'},
                    ],
                    'columns': 4,
                },
            ],
            'notes': 'Posicionar o copilot como acelerador de produtividade e governança do conhecimento.',
            'layout': 'ai-layer',
        },
        {
            'numero': 11,
            'titulo': '🤖 TC Copilot — Entregas Executivas',
            'kicker': 'IA APLICADA AO CONTEXTO SCI',
            'headline': 'O valor do TC Copilot aparece quando ele transforma regra técnica em texto executivo utilizável em segundos.',
            'subtitle': 'ETAPA 2 — Funcionalidades Ouro',
            'sections': [
                {
                    'kind': 'flow',
                    'title': 'Como ele entrega valor',
                    'items': [
                        {'title': 'Pergunta', 'body': 'Recebe uma dúvida técnica, um pedido executivo ou um gatilho de relatório.'},
                        {'title': 'Contexto', 'body': 'Lê documentação, contratos do SCI, histórico e lógica do sistema.'},
                        {'title': 'Resposta', 'body': 'Entrega síntese clara, narrativa pronta e orientação acionável para a gestão.'},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'cards',
                    'title': 'Entregas que valorizam o produto',
                    'items': [
                        {'title': '📅 Relatório diário', 'body': 'Resume variações, highlights e pontos de atenção logo após o fechamento.'},
                        {'title': '📋 Relatório executivo', 'body': 'Converte detalhe analítico em narrativa gerencial pronta para circular.'},
                        {'title': '💬 Resposta orientada', 'body': 'Explica regras, indicadores e leitura do SCI com linguagem adaptada ao público.', 'emphasis': True},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'example',
                    'title': 'Exemplo executivo',
                    'example_title': 'Trecho do resumo executivo gerado',
                    'body': 'No mês de março, o volume real ficou em linha com o budget, enquanto o Custo FP Real fechou abaixo do Budget Flex por ganho operacional concentrado em energy. A principal atenção permaneceu em maintenance, com pressão acima do esperado em materiais de reposição. A leitura executiva recomenda validar as unidades com maior exposição, confirmar a causa dominante e fechar a ação corretiva ainda no fechamento.',
                },
            ],
            'notes': 'Mostrar o Copilot como multiplicador de produtividade: menos tempo para consolidar, explicar e comunicar o resultado com qualidade executiva.',
            'layout': 'ai-layer',
        },
        {
            'numero': 12,
            'titulo': '🚨 Relatórios e Alertas',
            'kicker': 'MONITORAMENTO ACIONÁVEL',
            'headline': 'SCI fecha o ciclo quando transforma desvio em alerta priorizado e comunicação pronta para a organização.',
            'subtitle': 'ETAPA 2 — Funcionalidades Ouro',
            'sections': [
                {
                    'kind': 'bullets',
                    'title': 'Leitura gerencial',
                    'items': [
                        'Alertas priorizam o que mais pesa no resultado em vez de gerar ruído disperso.',
                        'Há validação técnica do cálculo antes do envio para e-mail ou Teams.',
                        'Relatórios garantem repetibilidade e escala da comunicação entre áreas.',
                    ],
                },
                {
                    'kind': 'cards',
                    'title': 'Componentes da central',
                    'items': [
                        {'title': '📊 Dados', 'body': 'Real, budget, flex e ranking consolidado.'},
                        {'title': '⚙️ Regras', 'body': 'Thresholds, filtros, hierarquia e severidade.', 'emphasis': True},
                        {'title': '📣 Comunicação', 'body': 'Teams, e-mail e uso interno no app.'},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'flow',
                    'title': 'Pipeline de monitoramento',
                    'items': [
                        {'title': 'Type 05', 'body': 'Origem do desvio e consolidação inicial'},
                        {'title': 'Type 06', 'body': 'Refino do ranking por agrupamento econômico'},
                        {'title': 'Account', 'body': 'Detalhe acionável para a linha de gestão'},
                    ],
                    'columns': 3,
                },
            ],
            'notes': 'Enquadrar alertas como instrumento de foco gerencial e não só automação de envio.',
            'layout': 'monitoring',
        },
        {
            'numero': 13,
            'titulo': '🚨 Alertas — Exemplo de Comunicação',
            'kicker': 'COMUNICAÇÃO ACIONÁVEL',
            'headline': 'O ganho real do SCI aparece quando o desvio sai do painel e vira mensagem pronta, priorizada e enviada no tempo certo.',
            'subtitle': 'ETAPA 2 — Funcionalidades Ouro',
            'sections': [
                {
                    'kind': 'flow',
                    'title': 'Do cálculo ao alerta',
                    'items': [
                        {'title': 'Leitura', 'body': 'O SCI identifica o desvio mais relevante no fechamento.'},
                        {'title': 'Priorização', 'body': 'Aplica severidade, ranking e contexto econômico antes do envio.'},
                        {'title': 'Comunicação', 'body': 'Entrega texto pronto para e-mail, Teams e acompanhamento gerencial.'},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'cards',
                    'title': 'Por que isso é ouro',
                    'items': [
                        {'title': '⚡ Velocidade', 'body': 'Reduz o tempo entre fechamento, leitura e acionamento da gestão.'},
                        {'title': '🎯 Foco', 'body': 'Evita ruído e leva a atenção para o que mais pesa no resultado.', 'emphasis': True},
                        {'title': '📣 Escala', 'body': 'Mantém padrão de comunicação mesmo com aumento de escopo e usuários.'},
                    ],
                    'columns': 3,
                },
                {
                    'kind': 'example',
                    'title': 'Exemplo executivo',
                    'example_title': 'Trecho do alerta consolidado enviado',
                    'body': 'Relatório de Alertas - Março | Budget Flex x Real | Severidade: critica. Burden concentrou a principal pressão do fechamento, com maintenance acima do esperado, enquanto energy compensou parte relevante do desvio total. O ranking foi aberto por Type 06, depois Account e unidades de maior impacto, permitindo validar causa, prioridade e plano de ação no mesmo fechamento.',
                },
            ],
            'notes': 'Valorizar o alerta como ouro do projeto: ele fecha o ciclo entre cálculo, priorização e comunicação executiva pronta para uso.',
            'layout': 'monitoring',
        },
        {
            'numero': 14,
            'titulo': '✅ Conclusão',
            'kicker': 'MENSAGEM FINAL · COO',
            'headline': 'O SCI se comporta como base de governança de custo industrial com capacidade real de escala.',
            'subtitle': 'ETAPA 2 — Funcionalidades Ouro',
            'sections': [
                {
                    'kind': 'bullets',
                    'title': 'Mensagem final',
                    'items': [
                        'O SCI já entrega leitura diária de custo, fechamento analítico e comunicação pronta em ritmo executivo.',
                        'A mesma base sustenta cálculo de TC, rateios por veículo, forecast, alertas e uso de IA em segundos.',
                        'O sistema mostra capacidade real de ampliar escopo sem perder consistência técnica nem velocidade de resposta.',
                    ],
                },
                {
                    'kind': 'cards',
                    'title': 'Quatro ganhos permanentes',
                    'items': [
                        {'title': '💡 Impacto', 'body': 'Transforma leitura de custo em capacidade real de decisão executiva.', 'emphasis': True},
                        {'title': '🏛 Governança', 'body': 'Mantém regra, histórico, explicação e narrativa no mesmo lugar.', 'emphasis': True},
                        {'title': '📐 Escala', 'body': 'Pronto para crescer em módulos, plantas, canais e rotinas.', 'emphasis': True},
                        {'title': '🔓 Autonomia', 'body': 'Reduz dependência de controles manuais e interpretação dispersa.', 'emphasis': True},
                    ],
                    'columns': 4,
                },
                {
                    'kind': 'flow',
                    'title': 'Capacidade operacional já demonstrada',
                    'items': [
                        {'title': 'Fechamento', 'body': 'Consolida e explica o resultado logo após o fechamento mensal.'},
                        {'title': 'IA em segundos', 'body': 'Gera síntese executiva, resposta contextual e apoio à comunicação.'},
                        {'title': 'Veículos e rateios', 'body': 'Leva o custo até o veículo com rastreabilidade do CPU e do rateio.'},
                    ],
                    'columns': 3,
                },
            ],
            'notes': 'Fechar mostrando capacidade concreta: alertas diários, relatórios detalhados, IA em segundos e cálculo rastreável até veículos e rateios.',
            'layout': 'closing',
        },
    ]
    slides = _renumber_doc_apresentacao_slides(slides)
    return _validate_doc_apresentacao_slides(slides)


def _render_doc_apresentacao_slide(slide, total_slides) -> None:
    shell_variant = slide.get('shell_variant', '')
    shell_class = 'sci-slide-shell'
    if shell_variant:
        shell_class += f' sci-slide-shell-{shell_variant}'
    hero_variant = slide.get('hero_variant', '')
    hero_class = 'sci-slide-hero'
    if hero_variant:
        hero_class += f' sci-slide-hero-{hero_variant}'

    st.markdown(f'<div class="{shell_class}">', unsafe_allow_html=True)
    st.markdown(
        f"""
    <div class="{hero_class}">
        <div class="sci-slide-kicker">{slide.get('kicker', '')}</div>
        <div class="sci-slide-title">{slide.get('titulo', '')}</div>
        <div class="sci-slide-headline">{slide.get('headline', '')}</div>
        <div class="sci-slide-subtitle">{slide.get('subtitle', '')}</div>
    </div>
        """,
        unsafe_allow_html=True,
    )
    for section in slide.get('sections', []):
        _render_doc_apresentacao_section(section)

    st.markdown(
        f"""
    <div class="sci-slide-footer">
        <span style="font-size:0.72rem;font-weight:900;text-transform:uppercase;letter-spacing:0.07em;color:#0f766e;">🎯 Slide {slide.get('numero', '')} / {total_slides}</span>
        &nbsp;·&nbsp;
        <span style="color:#7a96b2;font-size:0.78rem;">layout: {slide.get('layout', 'executive')}</span>
        &nbsp;·&nbsp;
        <span style="color:#42566d;font-size:0.78rem;">Fonte única · UI e PPT</span>
    </div>
        """,
        unsafe_allow_html=True,
    )

    st.markdown(
        f"""
    <div class="sci-slide-notes">
        <strong>Speaker note:</strong> {slide.get('notes', '')}
    </div>
        """,
        unsafe_allow_html=True,
    )
    st.markdown('</div>', unsafe_allow_html=True)


def _render_doc_apresentacao_visual() -> None:
    _render_doc_apresentacao_styles()
    slides = _get_doc_apresentacao_slides()
    total_slides = len(slides)
    ppt_error = None
    ppt_bytes = None
    try:
        ppt_bytes = _get_doc_apresentacao_ppt_bytes()
    except Exception as exc:
        ppt_error = str(exc)

    st.markdown("""
<div style="padding: 1.6rem; background: linear-gradient(135deg, #0f172a 0%, #164e63 55%, #99f6e4 100%); border-radius: 18px; margin-bottom: 1.4rem; color: white; border: 1px solid rgba(153, 246, 228, 0.16);">
    <h2 style="color: white; margin: 0;">🎤 Apresentação executiva do SCI</h2>
    <p style="color: #ecfeff; margin: 0.55rem 0 0 0; max-width: 920px;">
        {total_slides} slides executivos com contrato único entre a visualização desta página e o PPT exportado.
    </p>
</div>
    """, unsafe_allow_html=True)

    st.markdown(
        """
    <div class="sci-presentation-toolbar">
        <div style="font-size:0.86rem; font-weight:800; text-transform:uppercase; letter-spacing:0.08em; color:#0f4c5c; margin-bottom:0.35rem;">Deck reconstruído</div>
        <div class="sci-presentation-toolbar-note">
            O deck foi refeito a partir de um contrato único de slides. Use o exportador para baixar o PPT gerado pela mesma estrutura semântica mostrada aqui.
        </div>
    </div>
        """,
        unsafe_allow_html=True,
    )

    col_dl, col_info = st.columns([1.1, 2.3])
    with col_dl:
        if ppt_bytes is not None:
            st.download_button(
                '📤 Exportar apresentação PPTX',
                data=ppt_bytes,
                file_name='SCI_Apresentacao_Executiva.pptx',
                mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                use_container_width=True,
                key='download_doc_apresentacao_ppt',
            )
        else:
            st.button('📤 Exportar apresentação PPTX', disabled=True, use_container_width=True)
    with col_info:
        if ppt_error:
            st.warning(f'Não foi possível gerar o PPT agora: {ppt_error}')
        else:
            st.info('✅ PPT pronto — gerado a partir da mesma estrutura dos slides exibidos nesta página.')

    tab_roteiro, tab_slides = st.tabs(["🎤 Roteiro", "🧩 Slides"])

    with tab_roteiro:
        stage_2_start = next((index for index, slide in enumerate(slides) if slide['titulo'].startswith('🚀 ETAPA 2')), len(slides))
        stage_1_slides = slides[:stage_2_start]
        stage_2_slides = slides[stage_2_start:]

        def _render_roteiro_stage(header, stage_slides):
            st.markdown(f'### {header}')
            lines = []
            for slide in stage_slides:
                summary = slide.get('headline', '') or slide.get('notes', '')
                summary = summary.rstrip('.')
                lines.append(f"{slide['numero']}\\. **{slide['titulo']}** — {summary}.  ")
            st.markdown('\n'.join(lines))

        _render_roteiro_stage('🎬 Etapa 1 — Plataforma', stage_1_slides)
        _render_roteiro_stage('🚀 Etapa 2 — Funcionalidades Ouro', stage_2_slides)

    with tab_slides:
        st.caption(
            "Cada slide abaixo nasce do mesmo contrato usado pelo exportador PPT. A ideia aqui é estabilidade e clareza, não efeitos visuais frágeis."
        )
        for slide in slides:
            with st.expander(slide['titulo'], expanded=slide['numero'] == 1):
                _render_doc_apresentacao_slide(slide, total_slides)


def _render_doc_build_exe_completo() -> None:
    st.markdown("""
### Guia completo embutido

O build oficial do SCI é feito com PyInstaller usando SCI.spec. O arquivo centraliza dependências de import dinâmico, recursos estáticos e metadados que o build simplificado não garante sozinho.

### Ordem operacional recomendada
1. Ativar o ambiente virtual.
2. Garantir que app.py esteja sem BOM.
3. Executar o build oficial.
4. Validar que dist/Stellantis-Cost-Intelligence/_internal/ contém módulos, dados e recursos.
5. Testar o executável navegando pelas páginas críticas.

### Comando principal

```powershell
C:/User/U235107/GitHub/TC/.venv/Scripts/python.exe -m PyInstaller --clean --noconfirm SCI.spec
```

### Recursos obrigatórios no _internal
- pages/
- tc_core/
- tc_principal/
- tc_ext/
- tc_copilot/
- alertas/
- dados/
- .streamlit/
- arquivos JSON, imagens e scripts de processamento
- pacote st_aggrid e streamlit_aggrid-*.dist-info quando necessário

### Validação final
- abrir o executável;
- confirmar carregamento de páginas, principalmente TC Veículos, Documentação e Extração;
- validar leitura de parquets e ausência de erro de import dinâmico.
    """)

# Funções para persistir dados da equipe
def salvar_dados_equipe(dados):
    """Salva os dados da equipe em arquivo JSON"""
    try:
        base_path = get_base_path()
        dados_path = os.path.join(base_path, 'dados_equipe.json')
        with open(dados_path, 'w', encoding='utf-8') as f:
            json.dump(dados, f, ensure_ascii=False, indent=2)
        return True
    except Exception as e:
        st.error(f"Erro ao salvar dados: {e}")
        return False

def carregar_dados_equipe():
    """Carrega os dados da equipe do arquivo JSON"""
    _estrutura_vazia = {
        'hudson': {
            'cargo': '', 'empresa': '', 'experiencia': '',
            'linkedin': '', 'foto': None,
            'papel_projeto': 'Full-Stack Developer',
            'descricao_papel': (
                'Desenvolvendo tanto a interface quanto a '
                'lógica e os cálculos do sistema'
            ),
        },
        'lauro': {
            'cargo': '', 'empresa': '', 'experiencia': '',
            'linkedin': '', 'foto': None,
            'papel_projeto': 'Full-Stack Developer',
            'descricao_papel': (
                'Desenvolvendo tanto a interface quanto a '
                'lógica e os cálculos do sistema'
            ),
        },
        'frederico': {
            'cargo': 'Manufacturing Finance Controller',
            'empresa': 'Stellantis',
            'experiencia': '',
            'linkedin': '', 'foto': None,
            'papel_projeto': 'Tech Advisor',
            'descricao_papel': (
                'Orientação técnica estratégica, validações '
                'e suporte de alto nível ao projeto'
            ),
        },
    }
    try:
        base_path = get_base_path()
        dados_path = os.path.join(base_path, 'dados_equipe.json')
        if os.path.exists(dados_path):
            with open(dados_path, 'r', encoding='utf-8') as f:
                dados = json.load(f)
            # Garantir que todos os membros e campos existam
            for chave, padrao in _estrutura_vazia.items():
                if chave not in dados:
                    dados[chave] = padrao
                else:
                    for campo, valor in padrao.items():
                        if campo not in dados[chave]:
                            dados[chave][campo] = valor
            return dados
    except Exception as e:
        st.warning(f"Aviso ao carregar dados: {e}")

    return _estrutura_vazia

def salvar_foto_base64(foto_bytes, nome_arquivo):
    """Converte foto para base64 para salvar no JSON"""
    try:
        return base64.b64encode(foto_bytes).decode('utf-8')
    except:
        return None

def carregar_foto_base64(foto_base64):
    """Converte base64 de volta para bytes"""
    try:
        if foto_base64:
            return base64.b64decode(foto_base64)
    except:
        pass
    return None


def _render_doc_equipe_sci() -> None:
    st.header("👥 Equipe do SCI")

    st.markdown("""
    Esta seção apresenta os membros da equipe responsáveis pelo desenvolvimento
    e manutenção do **Stellantis Cost Intelligence (SCI)** — suas funções no projeto e perfis profissionais.
    """)

    st.markdown("""
    <style>
        .team-badge-fullstack {
            display: inline-block;
            background: linear-gradient(135deg, #7C3AED, #6D28D9);
            color: white;
            padding: 4px 14px;
            border-radius: 20px;
            font-size: 0.78rem;
            font-weight: 600;
            letter-spacing: 0.3px;
            margin-bottom: 4px;
        }
        .team-badge-advisor {
            display: inline-block;
            background: linear-gradient(135deg, #2563EB, #1D4ED8);
            color: white;
            padding: 4px 14px;
            border-radius: 20px;
            font-size: 0.78rem;
            font-weight: 600;
            letter-spacing: 0.3px;
            margin-bottom: 4px;
        }
        .team-role-desc {
            font-size: 0.82rem;
            color: #9CA3AF;
            font-style: italic;
            margin-top: 2px;
            margin-bottom: 8px;
        }
        .team-photo-box {
            width: 180px;
            height: 200px;
            border-radius: 12px;
            display: flex;
            align-items: center;
            justify-content: center;
            overflow: hidden;
            margin: 0 auto 8px auto;
            background: transparent;
        }
        .team-photo-box img {
            max-width: 100%;
            max-height: 100%;
            object-fit: cover;
            border-radius: 10px;
        }
        .team-photo-placeholder {
            color: #6B7280;
            font-size: 3rem;
            line-height: 1;
        }
    </style>
    """, unsafe_allow_html=True)

    st.markdown("---")

    dados_equipe = carregar_dados_equipe()
    membros = [
        {
            'key': 'hudson',
            'nome': 'Hudson Cardin',
            'icone': '🔧',
            'badge_class': 'team-badge-fullstack',
        },
        {
            'key': 'lauro',
            'nome': 'Lauro Paiva Junior',
            'icone': '📊',
            'badge_class': 'team-badge-fullstack',
        },
        {
            'key': 'frederico',
            'nome': 'Frederico Cesar de Jesus',
            'icone': '🧭',
            'badge_class': 'team-badge-advisor',
        },
    ]

    cols = st.columns(3)

    for col, membro in zip(cols, membros):
        k = membro['key']
        dados_m = dados_equipe.get(k, {})
        papel = dados_m.get('papel_projeto', '')
        desc_papel = dados_m.get('descricao_papel', '')

        with col:
            st.subheader(f"{membro['icone']} {membro['nome']}")

            if papel:
                st.markdown(
                    f'<span class="{membro["badge_class"]}">'
                    f'{papel}</span>',
                    unsafe_allow_html=True,
                )
            if desc_papel:
                st.markdown(
                    f'<p class="team-role-desc">{desc_papel}</p>',
                    unsafe_allow_html=True,
                )

            with st.expander("📸 Upload da foto", expanded=False):
                foto_up = st.file_uploader(
                    f"📸 Foto de {membro['nome']}",
                    type=['png', 'jpg', 'jpeg'],
                    key=f"foto_{k}",
                    help="Upload da foto de perfil (PNG, JPG, JPEG)",
                )
            _foto_b64_src = None
            if foto_up is not None:
                _raw = foto_up.read()
                dados_equipe[k]['foto'] = salvar_foto_base64(
                    _raw, f"{k}.jpg"
                )
                _foto_b64_src = (
                    'data:image/jpeg;base64,'
                    + base64.b64encode(_raw).decode()
                )
            elif dados_m.get('foto'):
                _foto_b64_src = (
                    'data:image/jpeg;base64,' + dados_m['foto']
                )

            if _foto_b64_src:
                st.markdown(
                    f'<div class="team-photo-box">'
                    f'<img src="{_foto_b64_src}" '
                    f'alt="{membro["nome"]}"/>'
                    f'</div>',
                    unsafe_allow_html=True,
                )
            else:
                st.markdown(
                    '<div class="team-photo-box">'
                    '<span class="team-photo-placeholder">'
                    '👤</span></div>',
                    unsafe_allow_html=True,
                )

            with st.expander(
                f"✏️ Editar informações", expanded=False
            ):
                with st.form(f"form_{k}"):
                    _papel = st.text_input(
                        "🎯 Papel no Projeto:",
                        value=dados_m.get('papel_projeto', ''),
                        key=f"papel_{k}",
                    )
                    _desc_papel = st.text_input(
                        "📝 Descrição do Papel:",
                        value=dados_m.get('descricao_papel', ''),
                        key=f"desc_papel_{k}",
                    )
                    _cargo = st.text_input(
                        "💼 Cargo:",
                        value=dados_m.get('cargo', ''),
                        key=f"cargo_{k}",
                    )
                    _empresa = st.text_input(
                        "🏢 Empresa:",
                        value=dados_m.get('empresa', ''),
                        key=f"empresa_{k}",
                    )
                    _exp = st.text_area(
                        "🎯 Experiência:",
                        value=dados_m.get('experiencia', ''),
                        key=f"exp_{k}",
                    )
                    _linkedin = st.text_input(
                        "🔗 LinkedIn:",
                        value=dados_m.get('linkedin', ''),
                        key=f"linkedin_{k}",
                    )
                    if st.form_submit_button(
                        "💾 Salvar", use_container_width=True
                    ):
                        dados_equipe[k]['papel_projeto'] = _papel
                        dados_equipe[k]['descricao_papel'] = _desc_papel
                        dados_equipe[k]['cargo'] = _cargo
                        dados_equipe[k]['empresa'] = _empresa
                        dados_equipe[k]['experiencia'] = _exp
                        dados_equipe[k]['linkedin'] = _linkedin
                        if salvar_dados_equipe(dados_equipe):
                            st.success("✅ Salvo com sucesso!")
                            st.rerun()

            with st.expander("👨‍💻 Perfil Profissional", expanded=False):
                if dados_m.get('cargo') and dados_m.get('empresa'):
                    st.write(
                        f"💼 **{dados_m['cargo']}** "
                        f"na **{dados_m['empresa']}**"
                    )
                elif dados_m.get('cargo'):
                    st.write(f"💼 **{dados_m['cargo']}**")
                elif dados_m.get('empresa'):
                    st.write(f"🏢 **{dados_m['empresa']}**")
                else:
                    st.write("💼 *Cargo não informado*")

                if dados_m.get('experiencia'):
                    st.write(f"🎯 {dados_m['experiencia']}")
                else:
                    st.write("🎯 *Experiência não informada*")

                if dados_m.get('linkedin'):
                    st.markdown(
                        f"🔗 [Perfil no LinkedIn]"
                        f"({dados_m['linkedin']})"
                    )
                else:
                    st.write("🔗 *LinkedIn não informado*")

    st.markdown("---")

    st.markdown("""
    ### 🎯 Objetivos do Projeto

    O **Stellantis Cost Intelligence (SCI)** é uma plataforma de análise de custos industriais composta por dois módulos
    complementares, cada um atendendo um nível de granularidade diferente:

    **📊 TC Estendido (TC Ext)**
    - Análise de custos por oficina, conta e período
    - Visualização Normal (Custo Total) e CPU (Custo por Unidade)
    - Dashboard interativo com filtros (Ano, Período, Oficina, USI, Veículo)
    - Flex Budget: ajuste do orçamento pela proporção de volume realizado
    - Waterfall Analysis: decomposição de variações entre períodos
    - Exportação Excel completa com formatação profissional

    **🚗 TC Veículos (TC Principal)**
    - Cadeia completa: Despesa Primária → Custo FA → Custo FP → D&A → FP sem Dedicada
    - Rateio proporcional por veículo (tempo de produção)
    - 6 tabs especializadas: TC Veículos, Análise Flex, Volume, Custos por Oficina, Tempo de Produção, Dados Detalhados
    - Best Estimate: simulador de premissas (sensibilidade, inflação, produtividade, volume) com geração de Forecast
    - Análise de Best Estimate: layout da Home alimentado por dados de Forecast

    **🔧 Capacidades Transversais**
    - 🚀 Cache inteligente com TTL e otimização de tipos de dados
    - 📦 Dados em formato Parquet comprimido
    - 💱 Conversão multi-moeda (BRL, USD, EUR) com taxas do banco de dados
    - 📊 Fator de escala configurável (Nenhum / K / M)
    - 🎨 Interface moderna com tabs, gráficos Altair e gradientes
    - ⚡ Performance otimizada para grandes volumes (70%+ redução de memória)
    """)

    st.markdown("<div style='height: 1.0rem;'></div>", unsafe_allow_html=True)

    _c1, c_logo, _c3 = st.columns([1, 3, 1])
    with c_logo:
        try:
            st.image("SCI_faixa.png", width="stretch")
        except TypeError:
            st.image("SCI_faixa.png", use_container_width=True)

# Sidebar com índices
st.sidebar.markdown("## 📑 Índice")
st.sidebar.markdown("---")

modulo_doc = "📌 Ambos (TC Ext + Veículos)"

# Criar índices no sidebar
indice_selecionado = st.sidebar.radio(
    "Selecione a seção:",
    [
        "👥 Equipe do SCI",
        "📘 Visão Geral Técnica",
        "📐 Regras e Cálculo",
        "🧮 Cálculo por Tabelas/Gráficos (Normal vs CPU)",
        "🏗️ Arquitetura e Estrutura",
        "🧾 Especificação Técnica",
        "📥 Guia de Extração de Dados",
        "🔮 Guia de Best Estimate",
        "☁️ TC Cloud",
        "📊 Apresentação Visual",
        "💬 Chatbot de Documentação",
        "🔔 Sistema de Alertas",
        "📦 Guia de Build (EXE)",
        "🚀 Próximos Passos",
    ],
    key="indice_documentacao"
)

st.markdown("---")

# ==========================================
# SEÇÃO 1: EQUIPE DO PROJETO
# ==========================================
if indice_selecionado == "👥 Equipe do SCI":
    _render_doc_equipe_sci()

# ==========================================
# SEÇÃO 1.5: VISÃO GERAL TÉCNICA
# ==========================================
elif indice_selecionado == "📘 Visão Geral Técnica":
    st.header("📘 Visão Geral Técnica")
    _render_doc_visao_geral_tecnica()

# ==========================================
# TC VEÍCULOS: REGRAS E CÁLCULO
# ==========================================
elif indice_selecionado == "📐 Regras e Cálculo" and modulo_doc.startswith("📌 Ambos"):
    st.header("📐 Regras e Cálculo — TC Ext + TC Veículos")

    st.subheader("📊 TC Estendido")
    _render_doc_regras_tc_ext()

    st.markdown("---")

    st.subheader("🚗 TC Veículos")
    _render_doc_regras_tc_veiculos()

    st.stop()

elif indice_selecionado == "📐 Regras e Cálculo" and modulo_doc == "🚗 TC Veículos":
    st.header("📐 Regras e Cálculo — TC Veículos")

    _render_doc_regras_tc_veiculos()

    st.stop()

    st.info(
        "📌 **Módulo TC Veículos** — Regras de cálculo específicas para "
        "análise de custo de produção de veículos."
    )

    with st.expander("💰 **Composição de Custos**", expanded=True):
        st.markdown("""
        ### 🔗 Cadeia de Custos TC Veículos

        ```
        Despesa Primária
                    × Rateio FA
                    = Custo FA (Fluxo Anexo)

                Custo FP (Fluxo Principal)
                    = Despesa Primária − Custo FA

        D&A Dedicado = parcela de D&A atribuída diretamente ao veículo
        FP sem Dedicada = Custo FP − D&A Dedicado
        ```

        **Colunas Monetárias** (recebem conversão de moeda e fator):
        - `Despesa Primaria`, `Custo FA`, `Custo FP`, `D&A dedicado`, `FP sem Dedicada`

        **Redis** — Não é uma coluna nem um `Account` fixo.
        Redis entra como **linhas adicionais** vindas da aba **massa - REDIS**, marcadas com `_fonte_redis=True`.

        **KPI Redis:**
        > Redis = Σ `Despesa Primaria` nas linhas com `_fonte_redis=True` (valores tipicamente negativos por serem receita)
        """)

    with st.expander("🚗 **Rateio por Veículo**", expanded=False):
        st.markdown("""
        ### 📊 Processo de Rateio

        O custo da oficina é **rateado** aos veículos proporcionalmente ao **tempo de produção**:

        - **Percentual(v,o)** = TempoVeic(v,o) / Σ TempoVeic(v,o)
        - **CustoRateado(v,o)** = FPsemDedicada(o) × Percentual(v,o)
        - **CustoFPVeiculo(v,o)** = CustoRateado(v,o) + D&A Dedicado(v,o)

        **Dados Consolidados vs Rateados:**

        | Seleção | Fonte BUD | Fonte Real |
        |---------|-----------|------------|
        | Todos | `df_principal_BUD.parquet` | `df_principal.parquet` |
        | Veículo específico | `df_veiculos_custo_fp_BUD.parquet` | `df_veiculos_custo_fp.parquet` |

        > Quando **Veículo = "Todos"**: dados consolidados.
        > Quando **Veículo = modelo específico**: dados rateados com `Custo FP Veiculo`.
        """)

    with st.expander("📊 **Flex Budget**", expanded=False):
        st.markdown("""
        ### 🔄 Conceito

        O Budget Flex ajusta o orçamento pela proporção de volume realizado:
        - Custos **fixos** permanecem iguais ao Budget
        - Custos **variáveis** são ajustados pela proporção de volume

        ### 📐 Fórmulas

        - **Proporção** = Volume Realizado / Volume Budget
        - **Flex fixo** = BUD fixo (sem alteração)
        - **Flex variável** = BUD variável × Proporção
        - **Flex total** = Flex fixo + Flex variável

        ### 🏷️ Classificação Fixo/Variável

        A coluna `Custo` determina a classificação:
        - Valores que começam com `"Fix"` (case-insensitive) → **Fixo**
        - Todos os demais → **Variável**

        ```python
        mask_fixo = df['Custo'].str.lower().str.startswith('fix')
        ```
        """)

    with st.expander("📈 **CPU (Custo por Unidade)**", expanded=False):
        st.markdown("""
        ### 💲 Fórmula

        **CPU = Custo Total / Volume Total**

        Com proteção contra divisão por zero:
        ```python
        CPU = np.where(volume != 0, custo / volume, 0.0)
        ```

        **Quando o tipo de visualização é CPU:**
        - Cada métrica é dividida pelo volume total
        - O sistema recalcula CPU **após** agregações (nunca soma/média de CPU)
        - O fator K/M é aplicado nas colunas monetárias antes do cálculo; para CPU sem escala, usar `Fator = Nenhum`
        - Volumes de BUD e Actual são usados conforme o contexto
        """)

    with st.expander("🎯 **KPIs do TC Veículos**", expanded=False):
        st.markdown("""
        ### 📊 KPIs do Topo (fora das tabs)

        | KPI | Fórmula |
        |-----|---------|
        | Desp. Primária | Σ Despesa Primaria |
        | Custo FA | Σ Custo FA |
        | Redis | Σ Despesa Primaria (linhas `_fonte_redis=True`, origem: massa - REDIS) |
        | Custo FP | Σ Custo FP |
        | D&A Dedicada | Σ D&A dedicado |
        | FP sem Dedicada | Σ FP sem Dedicada |

        ### 📊 KPIs do Resumo TC Veículos

        | KPI | Fórmula |
        |-----|---------|
        | BUD | BUD fixo + BUD variável |
        | Flex Bud − BUD | Flex total − BUD total |
        | Flex BUD | BUD fixo + BUD variável × Proporção |
        | Real − Flex Bud | Real total − Flex total |
        | Real | Σ Custo FP Real |
        | Real / Flex Bud | Real / Flex BUD (%) |
        """)

    with st.expander("🎯 **Arquitetura de Filtros**", expanded=False):
        st.markdown("""
        ### 🔍 Filtros do TC Veículos

        | Filtro | Tipo | Comportamento |
        |--------|------|---------------|
        | Oficina | multiselect | "Todos" ou seleção múltipla |
        | Tipo Custo | multiselect | Fixo/Variável ou todos |
        | Veículo | **selectbox** | "Todos" (consolidado) ou **1 veículo** (rateado) |
        | Período | multiselect | "Todos" ou seleção de meses |

        **Cascading:** A seleção de Oficina filtra os Veículos disponíveis:
        ```python
        _df_filt_ofi = df[df['Oficina'].isin(oficinas_selecionadas)]
        veiculos = sorted(_df_filt_ofi['Veículo'].dropna().unique())
        ```

        **Filtros globais:** Afetam KPIs, gráficos e Análise Flex simultaneamente.
        """)

    with st.expander("📈 **Sensibilidade e Volume (Best Estimate)**", expanded=False):
        st.markdown("""
        ### 🔮 Premissas do Simulador BE

        O Simulador de Best Estimate permite configurar premissas de **sensibilidade**, **inflação**,
        **produtividade** e **volume** para projetar cenários futuros, mantendo os meses históricos alinhados ao Real:

        **Fórmula Geral:**
        ```
        BE = Média_Histórica × Fator_Variação × Fator_Monetário
        ```

        Onde:
        - `Fator_Variação` = 1 + (Variação_Volume × Sensibilidade)
        - `Fator_Monetário` = (1 + Inflação / 100) × (1 - Produtividade / 100)
        - `Variação_Volume` = (Volume_Futuro / Volume_Médio_Histórico) − 1

        **Sequência correta do motor:**
        - Base histórica real válida → média histórica → volume histórico → volume futuro
        - Ajuste por sensibilidade
        - Aplicação de inflação e produtividade
        - Consolidação de Histórico + BE + BE Manual
        - Rateio por veículo com a mesma regra do Real quando o fluxo exigir granularidade veicular

        **Sensibilidade (impacto do volume no custo):**
        - Controla o quanto a variação de volume afeta o custo
        - Pode ser configurada globalmente por tipo de custo ou de forma detalhada por Type 06
        - Custo Fixo: sensibilidade = 0% → custo não varia com o volume
        - Custo Variável: sensibilidade = 100% → custo varia proporcionalmente ao volume

        **Volume:**
        - Define o volume de produção projetado por veículo
        - Usado para calcular a variação de volume, Flex Budget e CPU do Forecast
        - Quando o custo não tem dimensão Veículo, o volume médio é usado diretamente (`.mean()`)
        - Quando há Veículo, o volume é somado por grupo (`.sum()`)

        **Inflação:**
        - Aplica % de reajuste sobre **todos** os custos (fixos e variáveis)
        - É aplicada **após** o ajuste por sensibilidade
        - Fórmula: `Custo_Final = Custo_Ajustado_Sensibilidade × (1 + Inflação/100)`

        **Produtividade:**
        - Representa ganho de eficiência que **reduz** o custo projetado
        - Pode ser configurada globalmente ou por Type 06
        - É aplicada no mesmo bloco monetário da inflação, como fator redutor multiplicativo
        - Exemplo: produtividade de 5% reduz o custo projetado em 5% após o ajuste de volume

                **Rateio do forecast por veículo:**
                - O arquivo por veículo não usa a simplificação `Custo FP × Percentual`
                - A regra correta é:
                    - `Custo Rateado = FP sem Dedicada × Percentual`
                    - `Custo FP Veiculo = Custo Rateado + D&A dedicado`
                - Os percentuais vêm do Real e o D&A dedicado também vem do Real
                - Se faltar match direto, o fallback prioritário usa distribuição média por período, não distribuição uniforme imediata

        **Resultado por tipo de custo:**
        - **Custo Fixo BE** = Média Histórica × (1 + Inflação%) × (1 - Produtividade%) — sem ajuste de volume
        - **Custo Variável BE** = Média Histórica × (Vol_Futuro / Vol_Histórico) × (1 + Inflação%) × (1 - Produtividade%)

        **Persistência das configurações:**
        - O simulador salva o último conjunto aplicado em `config_forecast.json`
        - O arquivo persiste modo global/detalhado, sensibilidades, inflação, produtividade e configuração de períodos
                - O arquivo canônico de configuração não é `premissas.json`

        ### 📊 Geração de Forecast

        O simulador gera arquivos em `dados/TC_Principal/Forecast/`:
        - `forecast_completo.parquet` — Consolidado final com histórico + BE + BE Manual
        - `forecast_historico.parquet` — Histórico sem os meses previstos, evitando duplicidade com a previsão
        - `forecast_previsao.parquet` — Apenas períodos futuros de BE e BE Manual
        - `forecast_veiculos_custo_fp.parquet` — Forecast rateado por veículo para os fluxos que exigem granularidade veicular
        - `config_forecast.json` — Configurações aplicadas (modo, sensibilidade, inflação, produtividade, períodos)

        **Garantia de consistência na análise:**
        - O gerador valida a convergência dos meses históricos contra o Real no arquivo por veículo
        - A análise de BE sobrepõe os meses `Histórico` com o dado Real antes da exibição
        - Isso garante que meses históricos do gráfico de BE fechem exatamente com o Real, inclusive no Databricks

        Estes dados alimentam a página **Best Estimate (Análise)**, que usa o mesmo
        layout da Home (com gráficos e KPIs) mas com dados de Forecast.
        """)

# ==========================================
# SEÇÃO 2: REGRAS E CÁLCULO — TC ESTENDIDO
# ==========================================
elif indice_selecionado == "📐 Regras e Cálculo":
    st.header("📐 Regras e Cálculo — TC Estendido")

    _render_doc_regras_tc_ext()
    st.stop()

    # Conteúdo antigo removido: esta seção agora é renderizada diretamente do Markdown oficial.
    
    st.markdown("""
    Esta seção documenta todas as regras de cálculo, filtros e metodologias utilizadas no projeto.
    **IMPORTANTE:** Esta documentação serve como referência para garantir que todos os cálculos sejam
    reproduzidos de forma idêntica, permitindo que a IA consulte e refaça qualquer cálculo do sistema.
    
    A documentação está organizada em expanders para facilitar a navegação. Cada seção contém explicações
    detalhadas das regras, fórmulas matemáticas completas e exemplos práticos para facilitar o entendimento.
    """)
    
    st.markdown("---")
    
    # EXPANDER 1: Cálculos Principais
    with st.expander("🔢 **Cálculos Principais e Métricas Fundamentais**", expanded=False):
        with st.expander("📊 **CPU (Custo por Unidade)**", expanded=False):
            st.markdown("""
            ### 📊 CPU (Custo por Unidade)
            
            O **CPU (Custo por Unidade)** é uma métrica fundamental que representa o custo médio por unidade de produção.
            É calculado dividindo o custo total pelo volume de produção.
            
            **Fórmula Matemática:**
            ```
            CPU = Custo_Total / Volume_Total
            ```
            
            Onde:
            - `Custo_Total` = Soma de todos os custos individuais após agrupamento
            - `Volume_Total` = Soma de todos os volumes após agrupamento
            
            **⚠️ REGRA CRÍTICA:** O CPU deve ser calculado **APÓS** o agrupamento dos dados, nunca antes.
            Esta é uma das regras mais importantes do sistema, pois calcular CPU antes de agrupar resulta em valores incorretos.
            
            **Por que calcular após agrupamento?**
            
            A média aritmética de CPUs individuais não é igual ao CPU do total agregado. Isso ocorre porque o CPU é uma
            razão (divisão), e a média de razões não é igual à razão das médias.
            
            **Exemplo Ilustrativo:**
            
            Considere duas linhas de dados:
            - **Linha 1:** Custo Total = R$ 100, Volume = 10 unidades -> CPU = R$ 10,00/unidade
            - **Linha 2:** Custo Total = R$ 200, Volume = 40 unidades -> CPU = R$ 5,00/unidade
            
            **Método Incorreto (calcular CPU antes de agrupar):**
            - CPU médio = (R$ 10,00 + R$ 5,00) / 2 = **R$ 7,50/unidade** [INCORRETO]
            
            **Método Correto (calcular CPU após agrupar):**
            - Custo Total Agregado = R$ 100 + R$ 200 = R$ 300
            - Volume Total Agregado = 10 + 40 = 50 unidades
            - CPU Agregado = R$ 300 / 50 = **R$ 6,00/unidade** [CORRETO]
            
            A diferença entre R$ 7,50 e R$ 6,00 pode parecer pequena, mas em grandes volumes de dados essa discrepância
            se acumula e resulta em análises completamente incorretas.
            
            **Interpretação do CPU:**
            - **CPU baixo:** Indica eficiência operacional, menor custo por unidade produzida
            - **CPU alto:** Indica ineficiência ou custos elevados por unidade produzida
            - **Variação de CPU:** Mudanças no CPU entre períodos indicam variações na eficiência operacional
            """)
        
        with st.expander("💰 **Custo Total**", expanded=False):
            st.markdown("""
            ### 💰 Custo Total
        
        O **Custo Total** representa a soma de todos os custos individuais após a aplicação de filtros e agrupamentos.
        
        **Fórmula Matemática:**
        ```
        Custo_Total = Σ(Custo_Individual)
        ```
        
        Onde `Σ` representa a soma de todos os custos individuais que atendem aos critérios de filtragem.
        
        **Regras de Cálculo:**
        - Sempre somar valores individuais, nunca calcular média
        - Aplicar todos os filtros antes de realizar o agrupamento
        - Considerar apenas valores que atendem aos critérios de seleção
        - Não incluir valores nulos ou zerados no cálculo
        
        **Agrupamento por Dimensões:**
        
        O custo total pode ser calculado para diferentes níveis de agregação:
        - Por período (mês, trimestre, semestre, ano)
        - Por oficina
        - Por veículo
        - Por categoria de custo (Type 05, Type 06, Account)
        - Por combinação de dimensões
        
        **Interpretação:**
        - **Custo Total crescente:** Indica aumento nos gastos operacionais
        - **Custo Total decrescente:** Indica redução nos gastos operacionais
        - **Comparação entre períodos:** Permite identificar tendências e variações
        """)
        
        st.markdown("---")
        
        st.markdown("""
        ### 🔄 Fator de Conversão (K/M)
        
        Os **Fatores de Conversão** são utilizados para facilitar a visualização de valores muito grandes,
        convertendo-os para unidades mais legíveis (milhares ou milhões).
        
        **Fatores Disponíveis:**
        - **K (milhares):** Divide o valor por 1.000
        - **M (Milhões):** Divide o valor por 1.000.000
        - **Nenhum:** Mantém o valor original
        
        **Fórmulas Matemáticas:**
        ```
        Valor_K = Valor_Original / 1.000
        Valor_M = Valor_Original / 1.000.000
        ```
        
        **⚠️ REGRA CRÍTICA:** O fator de conversão **NÃO** deve ser aplicado no modo **CPU (Custo por Unidade)**.
        
        **Por que não aplicar em CPU?**
        
        O CPU já é uma razão (divisão entre Custo Total e Volume). Se aplicarmos o fator de conversão ao Custo Total
        antes de calcular o CPU, estaríamos dividindo duas vezes, o que resultaria em valores completamente incorretos.
        
        **Exemplo:**
        - Custo Total Original: R$ 1.000.000
        - Volume: 10.000 unidades
        - CPU Correto: R$ 1.000.000 / 10.000 = **R$ 100,00/unidade** [CORRETO]
        
        Se aplicássemos o fator K antes:
        - Custo Total com K: R$ 1.000 K
        - CPU Incorreto: R$ 1.000 K / 10.000 = **R$ 0,10/unidade** [INCORRETO] (1000 vezes menor!)
        
        **Ordem de Aplicação das Transformações:**
        
        1. **Primeiro:** Aplicar fator de conversão (K/M) - apenas em modo Custo Total
        2. **Segundo:** Converter moeda (se necessário)
        3. **Terceiro:** Realizar cálculos (CPU, Flex Bud, diferenças, etc.)
        
        Esta ordem garante que todas as transformações sejam aplicadas corretamente e que os resultados finais
        sejam consistentes e precisos.
        """)
        
        st.markdown("---")
        
        st.markdown("""
        ### 📅 Agrupamento por Período
        
        O **Agrupamento por Período** permite consolidar dados em diferentes intervalos de tempo, facilitando
        análises comparativas e identificação de tendências.
        
        **Estrutura de Períodos:**
        
        Quando os dados contêm informação de **Ano**, o sistema cria uma coluna combinada `Período_Ano` que
        agrupa tanto o período quanto o ano:
        ```
        Período_Ano = Período + " " + Ano
        ```
        
        Exemplo: "Janeiro 2024", "Fevereiro 2024", etc.
        
        **Agrupamento com Ano:**
        - Dimensões de agrupamento: `['Ano', 'Período']`
        - Permite comparações ano a ano
        - Facilita análises de tendências de longo prazo
        
        **Agrupamento sem Ano:**
        - Dimensões de agrupamento: `['Período']`
        - Útil quando todos os dados são do mesmo ano
        - Simplifica análises mensais ou trimestrais
        
        **Fórmula de Agregação:**
        ```
        Custo_Total_Agrupado = Σ(Custo_Individual) agrupado por Período
        Volume_Total_Agrupado = Σ(Volume_Individual) agrupado por Período
        ```
        
        **Interpretação:**
        - Permite identificar sazonalidades e padrões temporais
        - Facilita comparações entre períodos equivalentes
        - Suporta análises de tendências e projeções
        """)
        
        st.markdown("---")
        
        st.markdown("""
        ### 📈 Cálculo de Diferenças e Ratios
        
        As **Diferenças e Ratios** são métricas essenciais para análise de desempenho, permitindo comparar
        valores reais com valores planejados ou ajustados.
        
        **1. Diferença Flex Bud - BUD:**
        
        Esta métrica compara o Budget Flexível (ajustado pelo volume real) com o Budget original planejado.
        
        **Fórmula:**
        ```
        Delta_Flex_Bud = Flex_BUD - BUD
        ```
        
        **Interpretação:**
        - **Valor positivo:** Flex Bud > BUD (custo ajustado maior que o planejado)
        - **Valor negativo:** Flex Bud < BUD (custo ajustado menor que o planejado)
        - **Zero:** Flex Bud = BUD (custo ajustado igual ao planejado)
        
        **2. Diferença Total - Flex Bud:**
        
        Esta métrica compara o custo real com o Budget Flexível, indicando a eficiência operacional.
        
        **Fórmula:**
        ```
        Delta_Total_Flex = Total - Flex_BUD
        ```
        
        **Interpretação:**
        - **Valor positivo:** Total > Flex Bud (ineficiência operacional)
        - **Valor negativo:** Total < Flex Bud (eficiência operacional)
        - **Zero:** Total = Flex Bud (desempenho exatamente como esperado)
        
        **3. Ratio Total / Flex Bud:**
        
        Esta métrica expressa o desempenho real como percentual do Budget Flexível.
        
        **Fórmula:**
        ```
        Ratio = Total / Flex_BUD
        Percentual = Ratio * 100%
        ```
        
        **Interpretação:**
        - **< 100%:** Total < Flex Bud (melhor que esperado, eficiência operacional)
        - **= 100%:** Total = Flex Bud (exatamente como esperado)
        - **> 100%:** Total > Flex Bud (pior que esperado, ineficiência operacional)
        
        **Exemplo Prático:**
        - Flex Bud = R$ 500.000
        - Total Real = R$ 520.000
        - Ratio = 520.000 / 500.000 = 1,04 = **104%**
        - Interpretação: O custo real está 4% acima do Budget Flexível, indicando ineficiência operacional
        """)
    
    # EXPANDER 2: Flex Bud
    with st.expander("🔄 **Cálculo de Flex Bud (Budget Flexível)**", expanded=False):
        with st.expander("📋 **Conceito e Regras Fundamentais**", expanded=False):
            st.markdown("""
            ### Conceito
            
            **Flex Bud** (Budget Flexível) é um valor ajustado que considera a variação de volume,
            aplicando regras diferentes para custos fixos e **não‑fixos**.
            
            **IMPORTANTE:** Existem dois contextos diferentes de cálculo:
            1. **Real x Real** (Waterfall): Compara dois períodos reais (Mês 1 vs Mês 2)
            2. **Real x Budget** (TC Ext): Compara período real vs budget planejado
            """)
            
            st.markdown("---")
            
            st.markdown("## 📋 Regras Fundamentais: Fixo vs Não‑Fixo")
            
            st.markdown("""
            ### Regra Geral para Custos Fixos
            
            **Princípio:** Custos fixos NÃO variam com o volume de produção.
            
            **Fórmula Geral:**
            ```
            Flex_Fixo = Valor_Original_Fixo
            ```
            
            **Explicação:**
            - Independente da variação de volume, o custo fixo permanece constante
            - Exemplos: Aluguel, salários fixos, depreciação
            - Sensibilidade ao volume: **0%** (zero por cento)
            """)
            
            st.markdown("---")
            
            st.markdown("""
            ### Regra Geral para Custos Não‑Fixos
            
            **Princípio:** Custos **não‑fixos** variam PROPORCIONALMENTE ao volume de produção.
            
            **Fórmula Geral:**
            ```
            Flex_NãoFixo = Valor_Original_NãoFixo * (Volume_Novo / Volume_Original)
            ```
            
            **Explicação:**
            - Se o volume dobra, o custo **não‑fixo** escala proporcionalmente
            - Se o volume reduz pela metade, o custo **não‑fixo** escala proporcionalmente
            - Exemplos: componentes variáveis e demais classificações que não sejam Fixo
            - Sensibilidade ao volume: **100%** (cem por cento)
            """)
        
        # Ler o conteúdo do Flex Bud que está mais abaixo no arquivo
        # Por enquanto, vou adicionar um placeholder e depois mover o conteúdo correto
        st.info("📚 Conteúdo detalhado do Flex Bud será movido para cá...")
    
    # EXPANDER 3: Volumes
    with st.expander("📊 **Cálculo de Volumes**", expanded=False):
        with st.expander("📁 **Fonte de Dados de Volume**", expanded=False):
            st.markdown("""
            ### 📁 Fonte de Dados de Volume
            
            Os dados de volume são armazenados em arquivos Parquet otimizados para garantir performance e eficiência
            no processamento. O sistema utiliza diferentes arquivos de volume dependendo do contexto de análise.
            
            **Arquivos de Volume:**
            
            - **`df_vol_historico.parquet`:** Dados históricos consolidados de volume
            - **`df_vol.parquet`:** Dados de volume por ano específico
            - **`df_vol_historico_BUD.parquet`:** Dados de volume do budget histórico
            
            **Estrutura dos Dados de Volume:**
            
            Os arquivos de volume contêm as seguintes colunas obrigatórias:
            - **`Volume`:** Quantidade de unidades produzidas
            - **`Período`:** Período de referência (mês, trimestre, etc.)
            - **`Oficina`:** Identificação da oficina
            - **`Veículo`:** Identificação do veículo
            
            E a seguinte coluna opcional:
            - **`Ano`:** Ano de referência (quando disponível)
            
            **Sincronização com Dados de Custo:**
            
            Os dados de volume são estruturados de forma a permitir sincronização perfeita com os dados de custo,
            garantindo que o cálculo de CPU seja preciso e consistente.
            """)
        
        with st.expander("⚠️ **REGRA CRÍTICA: Filtragem de Volumes**", expanded=False):
            st.markdown("""
            ### ⚠️ REGRA CRÍTICA: Filtragem de Volumes
            
            **Princípio Fundamental:** Os volumes devem usar **EXATAMENTE** os mesmos filtros aplicados aos dados
            principais de custo. Esta regra é absolutamente crítica para garantir a precisão do cálculo de CPU.
            
            **Processo de Filtragem:**
            
            O processo de filtragem de volumes segue os seguintes passos:
            
            1. **Criar conjunto filtrado:** Iniciar com uma cópia dos dados de volume originais
            2. **Aplicar filtros sincronizados:** Usar os valores únicos extraídos dos dados de custo filtrados
            3. **Filtrar por período:** Aplicar filtro específico para o período de análise
            4. **Agrupar e somar:** Agrupar por período e somar os volumes
            
            **Fórmula de Agregação:**
            
            ```
            Volume_Total = Σ(Volume_Individual) agrupado por Período
            ```
            
            Onde `Σ` representa a soma de todos os volumes individuais que atendem aos critérios de filtragem.
            
            **Importância da Consistência:**
            
            A consistência entre os filtros aplicados aos dados de custo e volume é essencial porque:
            - O CPU é calculado como Custo Total / Volume Total
            - Se os filtros forem diferentes, o CPU resultante será completamente incorreto
            - Análises baseadas em CPU incorreto podem levar a decisões de negócio equivocadas
            
            **Exemplo de Impacto:**
            
            Se você filtrar os custos para uma oficina específica mas não filtrar os volumes da mesma forma:
            - Custo Total (filtrado): R$ 50.000 (apenas Oficina A)
            - Volume Total (não filtrado): 100.000 unidades (todas as oficinas)
            - CPU Incorreto: R$ 50.000 / 100.000 = R$ 0,50/unidade [INCORRETO]
            
            O correto seria:
            - Custo Total (filtrado): R$ 50.000 (Oficina A)
            - Volume Total (filtrado): 10.000 unidades (Oficina A)
            - CPU Correto: R$ 50.000 / 10.000 = R$ 5,00/unidade [CORRETO]
            """)
    
    # EXPANDER 4: Moeda e Taxas
    with st.expander("💱 **Moeda e Taxas de Câmbio**", expanded=False):
        with st.expander("💱 **Moedas Suportadas**", expanded=False):
            st.markdown("""
            ### 💱 Moedas Suportadas
            
            O sistema suporta conversão entre diferentes moedas para facilitar análises internacionais e comparações
            com dados de outras unidades de negócio. As moedas disponíveis são:
            
            - **BRL (R$):** Real Brasileiro - moeda base do sistema
            - **USD ($):** Dólar Americano
            - **EUR:** Euro
            
            **Moeda Base:**
            
            O Real Brasileiro (BRL) é a moeda base do sistema. Todos os valores são originalmente armazenados em BRL,
            e as conversões para outras moedas são realizadas multiplicando os valores pela taxa de câmbio correspondente.
            """)
        
        with st.expander("📊 **Taxas de Câmbio**", expanded=False):
            st.markdown("""
            ### 📊 Taxas de Câmbio
            
            As **Taxas de Câmbio** definem a relação de conversão entre a moeda base (BRL) e as outras moedas suportadas.
            
            **Definição Matemática:**
            
            As taxas são definidas como a quantidade de moeda estrangeira equivalente a 1 Real Brasileiro:
            ```
            1 BRL = Taxa_USD USD
            1 BRL = Taxa_EUR EUR
            ```
            
            **Exemplo Prático:**
            
            Se a taxa de câmbio USD for 0,20, isso significa que:
            - 1 Real Brasileiro = 0,20 Dólares Americanos
            - Para converter R$ 100,00 para USD: R$ 100,00 * 0,20 = $ 20,00
            
            **Fórmula de Conversão:**
            
            Para converter um valor de BRL para outra moeda:
            ```
            Valor_Convertido = Valor_BRL * Taxa_Cambio
            ```
            
            Onde:
            - `Valor_BRL` = Valor original em Real Brasileiro
            - `Taxa_Câmbio` = Taxa de câmbio da moeda de destino
            - `Valor_Convertido` = Valor convertido para a moeda de destino
            
            **Ordem de Aplicação das Transformações:**
            
            Quando múltiplas transformações são aplicadas (fator de conversão K/M e conversão de moeda), a ordem
            é crítica para garantir resultados corretos:
            
            1. **Primeiro:** Aplicar fator de conversão (K/M) - apenas em modo Custo Total
            2. **Segundo:** Converter moeda (se necessário)
            3. **Terceiro:** Realizar cálculos (CPU, Flex Bud, diferenças, etc.)
            
            **Exemplo Completo de Transformação:**
            
            Considere um valor original de R$ 1.000.000,00:
            
            - **Passo 1 (Fator K):** R$ 1.000.000,00 / 1.000 = R$ 1.000 K
            - **Passo 2 (Conversão USD, taxa 0,20):** R$ 1.000 K * 0,20 = $ 200 K
            - **Resultado Final:** $ 200 K (duzentos mil dólares)
            """)
        
        with st.expander("💾 **Persistência e Atualização de Taxas**", expanded=False):
            st.markdown("""
            ### 💾 Persistência e Atualização de Taxas
            
            As taxas de câmbio são armazenadas de forma persistente para garantir que as conversões sejam
            consistentes entre diferentes sessões de análise.
            
            **Armazenamento:**
            
            - As taxas são salvas em banco de dados ou arquivo de configuração
            - Valores padrão são utilizados caso não existam taxas salvas
            - As taxas podem ser atualizadas a qualquer momento através da interface do sistema
            
            **Atualização de Taxas:**
            
            As taxas de câmbio podem ser atualizadas para refletir as condições de mercado atuais. Quando uma
            nova taxa é definida, ela é aplicada a todos os cálculos subsequentes, garantindo que as análises
            estejam sempre baseadas nas taxas mais recentes.
            
            **Importância da Atualização:**
            
            Manter as taxas de câmbio atualizadas é essencial para garantir a precisão das análises, especialmente
            em períodos de alta volatilidade cambial. Taxas desatualizadas podem resultar em comparações e
            análises completamente incorretas.
            """)
    
    # EXPANDER 5: Filtros e Perímetros
    with st.expander("🔍 **Filtros e Perímetros de Análise**", expanded=False):
        with st.expander("🎯 **Sistema de Filtros da Interface**", expanded=False):
            st.markdown("""
            ### 🎯 Sistema de Filtros da Interface
            
            O sistema possui um conjunto abrangente de filtros que permitem refinar a análise de dados de forma
            precisa e flexível. Os filtros são aplicados sequencialmente, criando um perímetro de análise cada vez
            mais específico conforme o usuário seleciona diferentes critérios.
            
            **Ordem de Aplicação dos Filtros:**
            
            Os filtros são aplicados na seguinte ordem hierárquica, garantindo que cada filtro refine o resultado
            do filtro anterior:
            
            1. **Ano** - Seleção do ano de análise (Radio button)
            2. **Oficina** - Seleção de uma ou mais oficinas (Multiselect)
            3. **Veículo** - Seleção de um ou mais veículos (Multiselect)
            4. **USI** - Seleção de unidades de serviço (Multiselect)
            5. **Período** - Seleção de período específico (Selectbox)
            6. **Centro cst** - Seleção de centro de custo (Selectbox)
            7. **Conta contábil** - Seleção de contas contábeis (Multiselect)
            8. **Type 5** - Seleção de categorias Type 05 (Multiselect)
            9. **Type 6** - Seleção de categorias Type 06 (Multiselect)
            10. **Fornecedor** - Seleção de fornecedores (Multiselect)
            11. **Fornec.** - Seleção adicional de fornecedores (Multiselect)
            12. **Tipo** - Seleção de tipos de custo (Multiselect)
            13. **Filtros Avançados:**
                - **Usuário** - Filtro por usuário responsável
                - **Material** - Filtro por material utilizado
                - **Dt.lçto.** - Filtro por data de lançamento
                - **Texto breve** - Filtro por texto descritivo
                - **Account** - Filtro por conta contábil específica
            
            **Princípio de Funcionamento:**
            
            Cada filtro atua como um "funil" que reduz progressivamente o conjunto de dados analisados. Quando
            múltiplos filtros são aplicados, apenas os registros que atendem a **TODOS** os critérios selecionados
            são incluídos na análise final.
            
            **Exemplo de Aplicação Sequencial:**
            
            Imagine que você selecionou:
            - Ano: 2024
            - Oficina: "Oficina A" e "Oficina B"
            - Veículo: "Veículo X"
            - Período: "Janeiro"
            
            O sistema primeiro filtra todos os dados de 2024, depois mantém apenas os registros das Oficinas A e B,
            em seguida mantém apenas os registros do Veículo X, e finalmente mantém apenas os registros de Janeiro.
            O resultado final contém apenas os registros que atendem a todos esses critérios simultaneamente.
            """)
        
        with st.expander("⚠️ **REGRA CRÍTICA: Perímetro de Filtros para Volumes**", expanded=False):
            st.markdown("""
            ### ⚠️ REGRA CRÍTICA: Perímetro de Filtros para Volumes
            
            **Princípio Fundamental:** Os volumes devem usar **EXATAMENTE** os mesmos filtros aplicados aos dados
            principais de custo. Esta é uma das regras mais importantes do sistema, pois garante que o cálculo de
            CPU seja preciso e consistente.
            
            **Por que esta regra é crítica?**
            
            O CPU é calculado como a razão entre Custo Total e Volume. Se os filtros aplicados ao custo forem
            diferentes dos filtros aplicados ao volume, o CPU resultante será completamente incorreto.
            
            **Exemplo Ilustrativo:**
            
            Imagine que você filtrou os dados de custo para incluir apenas:
            - Oficina: "Oficina A"
            - Veículo: "Veículo X"
            - Período: "Janeiro"
            
            Se o volume não for filtrado da mesma forma, você poderia estar dividindo:
            - Custo Total (filtrado): R$ 100.000 (apenas Oficina A, Veículo X, Janeiro)
            - Volume Total (não filtrado): 50.000 unidades (todas as oficinas, todos os veículos, todos os períodos)
            - CPU Incorreto: R$ 100.000 / 50.000 = R$ 2,00/unidade [INCORRETO]
            
            O correto seria:
            - Custo Total (filtrado): R$ 100.000 (Oficina A, Veículo X, Janeiro)
            - Volume Total (filtrado): 10.000 unidades (Oficina A, Veículo X, Janeiro)
            - CPU Correto: R$ 100.000 / 10.000 = R$ 10,00/unidade [CORRETO]
            
            **Mecanismo de Sincronização:**
            
            O sistema garante a sincronização dos filtros extraindo os valores únicos das dimensões filtradas dos
            dados principais e aplicando esses mesmos valores aos dados de volume. Isso garante que o perímetro de
            análise seja idêntico para ambos os conjuntos de dados.
            
            **Dimensões Sincronizadas:**
            
            As seguintes dimensões são sempre sincronizadas entre dados de custo e volume:
            - Veículo
            - Oficina
            - USI
            - Centro de Custo
            - Conta Contábil
            - Type 05
            - Type 06
            - Fornecedor
            - Tipo
            - E todos os filtros avançados (Usuário, Material, Data, etc.)
            """)
        
        with st.expander("📊 **Sincronização de Filtros para Budget**", expanded=False):
            st.markdown("""
            ### 📊 Sincronização de Filtros para Budget
            
            **Regra Fundamental:** O Budget deve usar os mesmos filtros aplicados aos dados reais para garantir
            comparações justas e precisas.
            
            **Por que sincronizar filtros do Budget?**
            
            Quando comparamos dados reais com budget, precisamos garantir que estamos comparando "maçãs com maçãs".
            Se os dados reais estão filtrados para uma oficina específica, o budget também deve estar filtrado para
            a mesma oficina, caso contrário a comparação não terá sentido.
            
            **Exemplo:**
            
            Se você filtrar os dados reais para:
            - Oficina: "Oficina A"
            - Veículo: "Veículo X"
            
            O budget também será automaticamente filtrado para:
            - Oficina: "Oficina A"
            - Veículo: "Veículo X"
            
            Isso garante que a comparação entre Real e Budget seja feita no mesmo contexto operacional.
            
            **Mecanismo de Aplicação:**
            
            O sistema extrai os valores únicos de todas as dimensões filtradas dos dados reais e aplica esses mesmos
            valores como filtros ao budget. Isso garante que o perímetro de análise seja idêntico para ambos os
            conjuntos de dados, permitindo comparações precisas e significativas.
            """)

    # EXPANDER 3: Volumes
    with st.expander("📊 **Cálculo de Volumes**", expanded=False):
        with st.expander("📁 **Fonte de Dados de Volume**", expanded=False):
            st.markdown("""
            ### 📁 Fonte de Dados de Volume
            
            Os dados de volume são armazenados em arquivos Parquet otimizados para garantir performance e eficiência
            no processamento. O sistema utiliza diferentes arquivos de volume dependendo do contexto de análise.
            
            **Arquivos de Volume:**
            
            - **`df_vol_historico.parquet`:** Dados históricos consolidados de volume
            - **`df_vol.parquet`:** Dados de volume por ano específico
            - **`df_vol_historico_BUD.parquet`:** Dados de volume do budget histórico
            
            **Estrutura dos Dados de Volume:**
            
            Os arquivos de volume contêm as seguintes colunas obrigatórias:
            - **`Volume`:** Quantidade de unidades produzidas
            - **`Período`:** Período de referência (mês, trimestre, etc.)
            - **`Oficina`:** Identificação da oficina
            - **`Veículo`:** Identificação do veículo
            
            E a seguinte coluna opcional:
            - **`Ano`:** Ano de referência (quando disponível)
            
            **Sincronização com Dados de Custo:**
            
            Os dados de volume são estruturados de forma a permitir sincronização perfeita com os dados de custo,
            garantindo que o cálculo de CPU seja preciso e consistente.
            """)
        
        with st.expander("⚠️ **REGRA CRÍTICA: Filtragem de Volumes**", expanded=False):
            st.markdown("""
            ### ⚠️ REGRA CRÍTICA: Filtragem de Volumes
            
            **Princípio Fundamental:** Os volumes devem usar **EXATAMENTE** os mesmos filtros aplicados aos dados
            principais de custo. Esta regra é absolutamente crítica para garantir a precisão do cálculo de CPU.
            
            **Processo de Filtragem:**
            
            O processo de filtragem de volumes segue os seguintes passos:
            
            1. **Criar conjunto filtrado:** Iniciar com uma cópia dos dados de volume originais
            2. **Aplicar filtros sincronizados:** Usar os valores únicos extraídos dos dados de custo filtrados
            3. **Filtrar por período:** Aplicar filtro específico para o período de análise
            4. **Agrupar e somar:** Agrupar por período e somar os volumes
            
            **Fórmula de Agregação:**
            
            ```
            Volume_Total = Σ(Volume_Individual) agrupado por Período
            ```
            
            Onde `Σ` representa a soma de todos os volumes individuais que atendem aos critérios de filtragem.
            
            **Importância da Consistência:**
            
            A consistência entre os filtros aplicados aos dados de custo e volume é essencial porque:
            - O CPU é calculado como Custo Total / Volume Total
            - Se os filtros forem diferentes, o CPU resultante será completamente incorreto
            - Análises baseadas em CPU incorreto podem levar a decisões de negócio equivocadas
            
            **Exemplo de Impacto:**
            
            Se você filtrar os custos para uma oficina específica mas não filtrar os volumes da mesma forma:
            - Custo Total (filtrado): R$ 50.000 (apenas Oficina A)
            - Volume Total (não filtrado): 100.000 unidades (todas as oficinas)
            - CPU Incorreto: R$ 50.000 / 100.000 = R$ 0,50/unidade [INCORRETO]
            
            O correto seria:
            - Custo Total (filtrado): R$ 50.000 (Oficina A)
            - Volume Total (filtrado): 10.000 unidades (Oficina A)
            - CPU Correto: R$ 50.000 / 10.000 = R$ 5,00/unidade [CORRETO]
            """)
    
    st.markdown("---")
    
    st.markdown("## 📋 Regras Fundamentais: Fixo vs Não‑Fixo")
    
    st.markdown("""
    ### Regra Geral para Custos Fixos
    
    **Princípio:** Custos fixos NÃO variam com o volume de produção.
    
    **Fórmula Geral:**
    ```
    Flex_Fixo = Valor_Original_Fixo
    ```
    
    **Explicação:**
    - Independente da variação de volume, o custo fixo permanece constante
    - Exemplos: Aluguel, salários fixos, depreciação
    - Sensibilidade ao volume: **0%** (zero por cento)
    
    **Implementação:**
    ```python
    # Sempre manter o valor original
    flex_fixo = custo_fixo_original
    ```
    """)
    
    st.markdown("---")
    
    st.markdown("""
    ### Regra Geral para Custos Variáveis
    
    **Princípio:** Custos variáveis variam PROPORCIONALMENTE ao volume de produção.
    
    **Fórmula Geral:**
    ```
    Flex_NãoFixo = Valor_Original_NãoFixo * (Volume_Novo / Volume_Original)
    ```
    
    **Explicação:**
    - Se o volume dobra, o custo variável dobra
    - Se o volume reduz pela metade, o custo variável reduz pela metade
    - Exemplos: Matéria-prima, energia variável, comissões
    - Sensibilidade ao volume: **100%** (cem por cento)
    
    **Implementação:**
    ```python
    # Calcular proporção de volume
    proporcao = volume_novo / volume_original
    
    # Aplicar proporção ao custo variável
    flex_variavel = custo_variavel_original * proporcao
    ```
    """)
    
    st.markdown("---")
    
    st.markdown("""
    ### Identificação de Fixo vs Variável
    
    **Coluna 'Custo' no DataFrame:**
    - Deve conter os valores: `'Fixo'` ou `'Variável'`
    - Cada linha de dados deve ter esta classificação
    
    **Implementação:**
    ```python
    # Separar Fixo e Variável
    if 'Custo' in df.columns:
        custo_fixo = df[df['Custo'] == 'Fixo']['Total'].sum()
        custo_variavel = df[df['Custo'] == 'Variável']['Total'].sum()
    else:
        # Se não tiver coluna Custo, assumir tudo como variável
        custo_fixo = 0
        custo_variavel = df['Total'].sum()
    ```
    """)
    
    st.markdown("---")
    
    # Sub-seções para separar os dois casos
    st.markdown("## 📊 CASO 1: Flex para Comparação Real x Real (Waterfall)")
    
    st.markdown("""
    ### Contexto
    
    Usado na página **1 - Waterfall** para comparar dois períodos reais:
    - **Mês 1** (período inicial real)
    - **Mês 2** (período final real)
    
    **Objetivo:** Calcular o que seria o custo do Mês 1 ajustado pelo volume do Mês 2.
    """)
    
    st.markdown("---")
    
    st.markdown("""
    ### Regras de Cálculo - Real x Real
        
        **Passo 1: Identificar Custos do Mês 1**
        ```python
        # Separar Fixo e Variável do Mês 1
        C1_Fixo = df_m1[df_m1['Custo'] == 'Fixo']['Total'].sum()
        C1_Variavel = df_m1[df_m1['Custo'] == 'Variável']['Total'].sum()
        C1_Total = C1_Fixo + C1_Variavel
        ```
        
        **Passo 2: Obter Volumes**
        ```python
        V1 = volume_real_mes1  # Volume do Mês 1
        V2 = volume_real_mes2  # Volume do Mês 2
        ```
        
        **Passo 3: Calcular Proporção de Volume**
        ```python
        rho = V2 / V1  # Proporção de volume
        ```
        
        **Passo 4: Aplicar Regras de Fixo e Variável**
        
        **Para Custo Fixo:**
        ```python
        # REGRA: Fixo não varia com volume
        Flex_Mes1_Fixo = C1_Fixo
        # Explicação: Mantém o valor original, independente da variação de volume
        ```
        **Fórmula Matemática:**
        ```
        Flex_Mês1_Fixo = C_1_Fixo
        ```
        **Por que não multiplica pela proporção?**
        - Custos fixos são independentes do volume de produção
        - Exemplos: Aluguel, salários fixos, depreciação
        - Mesmo que o volume dobre, o custo fixo permanece igual
        
        **Para Custo Variável:**
        ```python
        # REGRA: Variável varia proporcionalmente ao volume
        Flex_Mes1_Variavel = C1_Variavel * rho
                             = C1_Variavel * (V2 / V1)
        # Explicação: Multiplica pelo fator de proporção de volume
        ```
        **Fórmula Matemática:**
        ```
        Flex_Mês1_Variável = C_1_Variável * rho
                              = C_1_Variável * (V_2 / V_1)
        ```
        **Por que multiplica pela proporção?**
        - Custos variáveis variam proporcionalmente ao volume
        - Se o volume dobra, o custo variável dobra
        - Se o volume reduz pela metade, o custo variável reduz pela metade
        - Exemplos: Matéria-prima, energia variável, comissões
        
        **Passo 5: Calcular Flex Mês 1 Total**
        ```python
        Flex_Mes1_Total = Flex_Mes1_Fixo + Flex_Mes1_Variavel
                         = C1_Fixo + (C1_Variavel * rho)
        ```
        **Fórmula Matemática:**
        ```
        Flex_Mês1_Total = Flex_Mês1_Fixo + Flex_Mês1_Variável
                        = C_1_Fixo + (C_1_Variável * rho)
                        = C_1_Fixo + C_1_Variável * (V_2 / V_1)
        ```
    """)
    
    st.markdown("---")
    
    st.markdown("""
    ### Fórmulas Matemáticas Completas - Real x Real
    
    **Definições:**
    - `V_1` = Volume Real do Mês 1
    - `V_2` = Volume Real do Mês 2
    - `C_1_Fixo` = Custo Total Fixo do Mês 1
    - `C_1_Variável` = Custo Total Variável do Mês 1
    - `C_1_Total` = Custo Total do Mês 1 = `C_1_Fixo + C_1_Variável`
    
    **Proporção de Volume:**
    ```
    rho = V_2 / V_1
    ```
    Onde:
    - `rho > 1` significa que o volume aumentou
    - `rho < 1` significa que o volume diminuiu
    - `rho = 1` significa que o volume permaneceu igual
    
    **Cálculo de Flex Mês 1 (em Custo Total):**
    
    Para **Custo Fixo:**
    ```
    Flex_Mês1_Fixo = C_1_Fixo
    ```
    **Regra Aplicada:** Fixo não varia com volume
    - Valor original mantido: `C_1_Fixo`
    - Não multiplica pela proporção de volume
    - Motivo: Custos fixos são independentes do volume de produção
    
    Para **Custo Variável:**
    ```
    Flex_Mês1_Variável = C_1_Variável * rho
                          = C_1_Variável * (V_2 / V_1)
    ```
    **Regra Aplicada:** Variável varia proporcionalmente ao volume
    - Valor original: `C_1_Variável`
    - Multiplica pela proporção: `rho = V_2 / V_1`
    - Motivo: Custos variáveis aumentam/diminuem na mesma proporção do volume
    
    **Flex Mês 1 Total (em Custo Total):**
    ```
    Flex_Mês1_Total = Flex_Mês1_Fixo + Flex_Mês1_Variável
                    = C_1_Fixo + (C_1_Variável * rho)
                    = C_1_Fixo + C_1_Variável * (V_2 / V_1)
    ```
    **Regra Aplicada:** Soma do Fixo (inalterado) + Variável (ajustado)
    """)
    
    st.markdown("---")
    
    st.markdown("""
    ### Cálculo em CPU (Custo por Unidade) - Real x Real
    
    **IMPORTANTE:** No modo CPU, calcular em Custo Total primeiro, depois converter.
        
        **BUD (Mês 1) em CPU:**
        ```
        BUD_CPU = C_1_Total / V_1
                 = (C_1_Fixo + C_1_Variável) / V_1
        ```
        
        **Flex Mês 1 em CPU:**
        ```
        Flex_Mês1_CPU = Flex_Mês1_Total / V_2
                       = [C_1_Fixo + C_1_Variável * (V_2 / V_1)] / V_2
                       = (C_1_Fixo / V_2) + (C_1_Variável / V_1)
        ```
        
        **Diferença (Flex Mês 1 - Mês 1):**
        ```
        Delta_Flex = Flex_Mês1_CPU - BUD_CPU
               = [(C_1_Fixo / V_2) + (C_1_Variável / V_1)] - [(C_1_Fixo + C_1_Variável) / V_1]
               = (C_1_Fixo / V_2) - (C_1_Fixo / V_1)
               = C_1_Fixo * (1/V_2 - 1/V_1)
               = C_1_Fixo * (V_1 - V_2) / (V_1 * V_2)
        ```
        
        **Interpretação:**
        - Se `V_2 > V_1`: `Delta_Flex < 0` (CPU diminui porque custo fixo é diluído em mais volume)
        - Se `V_2 < V_1`: `Delta_Flex > 0` (CPU aumenta porque custo fixo é concentrado em menos volume)
        - Se `V_2 = V_1`: `Delta_Flex = 0` (sem variação)
    """)
    
    st.markdown("---")
    
    st.markdown("""
    ### Implementação - Real x Real
        
        ```python
        # 1. Obter dados do Mês 1
        df_m1 = df_filtrado[df_filtrado['Período'] == mes_inicial]
        
        # 2. Separar Fixo e Variável
        if 'Custo' in df_m1.columns:
            C1_Fixo = df_m1[df_m1['Custo'] == 'Fixo']['Total'].sum()
            C1_Variavel = df_m1[df_m1['Custo'] == 'Variável']['Total'].sum()
        else:
            C1_Fixo = 0
            C1_Variavel = df_m1['Total'].sum()  # Tudo é variável
        
        C1_Total = C1_Fixo + C1_Variavel
        
        # 3. Obter volumes
        volume_m1 = df_vol_m1['Volume'].sum()
        volume_m2 = df_vol_m2['Volume'].sum()
        
        # 4. Calcular proporção
        rho = volume_m2 / volume_m1 if volume_m1 != 0 else 1.0
        
        # 5. Calcular Flex Mês 1 (em Custo Total)
        Flex_Mes1_Fixo = C1_Fixo  # Não varia
        Flex_Mes1_Variavel = C1_Variavel * rho  # Varia proporcionalmente
        Flex_Mes1_Total = Flex_Mes1_Fixo + Flex_Mes1_Variavel
        
        # 6. Converter para CPU (se necessário)
        if tipo_visualizacao == "CPU (Custo por Unidade)":
            BUD_CPU = C1_Total / volume_m1 if volume_m1 != 0 else 0
            Flex_Mes1_CPU = Flex_Mes1_Total / volume_m2 if volume_m2 != 0 else 0
            Delta_Flex = Flex_Mes1_CPU - BUD_CPU
        ```
    """)
    
    st.markdown("---")
    
    st.markdown("""
    ### Exemplo Prático - Real x Real
        
        **Dados:**
        - Volume Real Mês 1 (`V_1`): 40,848 unidades
        - Volume Real Mês 2 (`V_2`): 60,333 unidades
        - Custo Total Fixo Mês 1 (`C_1_Fixo`): R$ 126.91
        - Custo Total Variável Mês 1 (`C_1_Variável`): R$ 755.36
        - Custo Total Mês 1 (`C_1_Total`): R$ 882.27
        
        **Cálculo:**
        ```
        rho = V_2 / V_1 = 60,333 / 40,848 = 1.482373
        
        Flex_Mês1_Fixo = R$ 126.91
        Flex_Mês1_Variável = R$ 755.36 * 1.482373 = R$ 1,119.72
        Flex_Mês1_Total = R$ 126.91 + R$ 1,119.72 = R$ 1,246.63
        ```
        
        **Em CPU:**
        ```
        BUD_CPU = R$ 882.27 / 40,848 = R$ 0.0216 por unidade
        Flex_Mês1_CPU = R$ 1,246.63 / 60,333 = R$ 0.0207 por unidade
        Delta_Flex = R$ 0.0207 - R$ 0.0216 = -R$ 0.0009 por unidade
        ```
        
        **Interpretação:**
        - O volume aumentou 48.24% (`rho = 1.482373`)
        - O custo variável aumentou proporcionalmente: R$ 755.36 -> R$ 1,119.72
        - O custo fixo permaneceu igual: R$ 126.91
        - Em CPU, o custo por unidade diminuiu porque o custo fixo foi diluído em mais volume
        """)
        
    st.markdown("---")
        
    st.markdown("""
        ### Modos de Comparação - Real x Real
        
        **Mês a Mês:**
        - `V_1` = Volume do mês inicial
        - `V_2` = Volume do mês final
        
        **Ano a Ano:**
        - `V_1` = Volume total do ano inicial
        - `V_2` = Volume total do ano final
        
        **Semestre:**
        - `V_1` = Volume total do semestre inicial
        - `V_2` = Volume total do semestre final
        
        **Quarter:**
        - `V_1` = Volume total do trimestre inicial
        - `V_2` = Volume total do trimestre final
    """)
    
    st.markdown("---")
    
    st.markdown("## 💰 CASO 2: Flex para Comparação Real x Budget (TC Ext)")
    
    st.markdown("""
        ### Contexto
        
        Usado na página **TC Ext** para comparar período real vs budget planejado:
        - **Real** = Dados reais do período
        - **Budget** = Dados planejados do período
        
        **Objetivo:** Calcular o que seria o budget ajustado pelo volume real.
        """)
        
    st.markdown("---")
        
    st.markdown("""
        ### Regras de Cálculo - Real x Budget
        
        **Passo 1: Identificar Custos do Budget**
        ```python
        # Separar Fixo e Variável do Budget
        B_Fixo = df_budget[df_budget['Custo'] == 'Fixo']['Total'].sum()
        B_Variavel = df_budget[df_budget['Custo'] == 'Variável']['Total'].sum()
        B_Total = B_Fixo + B_Variavel
        ```
        
        **Passo 2: Obter Volumes**
        ```python
        V_Budget = volume_budget  # Volume planejado no Budget
        V_Real = volume_real      # Volume real do período
        ```
        
        **Passo 3: Calcular Proporção de Volume**
        ```python
        rho = V_Real / V_Budget  # Proporção de volume real vs planejado
        ```
        
        **Passo 4: Aplicar Regras de Fixo e Variável**
        
        **Para Custo Fixo:**
        ```python
        # REGRA: Fixo não varia com volume
        Flex_Bud_Fixo = B_Fixo
        # Explicação: Mantém o valor do budget, independente da variação de volume
        ```
        
        **Para Custo Variável:**
        ```python
        # REGRA: Variável varia proporcionalmente ao volume
        Flex_Bud_Variavel = B_Variavel * rho
                            = B_Variavel * (V_Real / V_Budget)
        # Explicação: Ajusta o budget variável pela proporção de volume real vs planejado
        ```
        
        **Passo 5: Calcular Flex Bud Total**
        ```python
        Flex_Bud_Total = Flex_Bud_Fixo + Flex_Bud_Variavel
                        = B_Fixo + (B_Variavel * rho)
        ```
        """)
        
    st.markdown("---")
        
    st.markdown("""
        ### Fórmulas Matemáticas Completas - Real x Budget
        
        **Definições:**
        ```python
        # Separar Fixo e Variável do Budget
        B_Fixo = df_budget[df_budget['Custo'] == 'Fixo']['Total'].sum()
        B_Variavel = df_budget[df_budget['Custo'] == 'Variável']['Total'].sum()
        B_Total = B_Fixo + B_Variavel
        ```
        
        **Passo 2: Obter Volumes**
        ```python
        V_Budget = volume_budget  # Volume planejado no Budget
        V_Real = volume_real      # Volume real do período
        ```
        
        **Passo 3: Calcular Proporção de Volume**
        ```python
        rho = V_Real / V_Budget  # Proporção de volume real vs planejado
        ```
        
        **Passo 4: Aplicar Regras de Fixo e Variável**
        
        **Para Custo Fixo:**
        ```python
        # REGRA: Fixo não varia com volume
        Flex_Bud_Fixo = B_Fixo
        # Explicação: Mantém o valor do budget, independente da variação de volume
        ```
        **Fórmula Matemática:**
        ```
        Flex_Bud_Fixo = B_Fixo
        ```
        **Por que não multiplica pela proporção?**
        - Custos fixos são independentes do volume de produção
        - O budget fixo foi planejado e não deve ser ajustado
        - Exemplos: Aluguel, salários fixos, depreciação
        - Mesmo que o volume real seja diferente do planejado, o custo fixo permanece igual
        
        **Para Custo Variável:**
        ```python
        # REGRA: Variável varia proporcionalmente ao volume
        Flex_Bud_Variavel = B_Variavel * rho
                            = B_Variavel * (V_Real / V_Budget)
        # Explicação: Ajusta o budget variável pela proporção de volume real vs planejado
        ```
        **Fórmula Matemática:**
        ```
        Flex_Bud_Variável = B_Variável * rho
                           = B_Variável * (V_Real / V_Budget)
        ```
        **Por que multiplica pela proporção?**
        - Custos variáveis variam proporcionalmente ao volume
        - Se o volume real for maior que o planejado, o custo variável deve aumentar
        - Se o volume real for menor que o planejado, o custo variável deve diminuir
        - Exemplos: Matéria-prima, energia variável, comissões
        - O budget variável precisa ser ajustado para refletir o volume real
        
        **Passo 5: Calcular Flex Bud Total**
        ```python
        Flex_Bud_Total = Flex_Bud_Fixo + Flex_Bud_Variavel
                        = B_Fixo + (B_Variavel * rho)
        ```
        **Fórmula Matemática:**
        ```
        Flex_Bud_Total = Flex_Bud_Fixo + Flex_Bud_Variável
                       = B_Fixo + (B_Variável * rho)
                       = B_Fixo + B_Variável * (V_Real / V_Budget)
        ```
        """)
        
    st.markdown("---")
        
    st.markdown("""
        ### Fórmulas Matemáticas Completas - Real x Budget
        
        **Definições:**
        - `V_Real` = Volume Real do período
        - `V_Budget` = Volume Budget planejado do período
        - `B_Fixo` = Custo Total Fixo do Budget
        - `B_Variável` = Custo Total Variável do Budget
        - `B_Total` = Custo Total do Budget = `B_Fixo + B_Variável`
        - `R_Total` = Custo Total Real do período
        
        **Proporção de Volume:**
        ```
        rho = V_Real / V_Budget
        ```
        Onde:
        - `rho > 1` significa que o volume real foi maior que o planejado
        - `rho < 1` significa que o volume real foi menor que o planejado
        - `rho = 1` significa que o volume real foi exatamente o planejado
        
        **Cálculo de Flex Bud (em Custo Total):**
        
        Para **Custo Fixo:**
        ```
        Flex_Bud_Fixo = B_Fixo
        ```
        **Regra Aplicada:** Fixo não varia com volume
        - Valor do budget mantido: `B_Fixo`
        - Não multiplica pela proporção de volume
        - Motivo: Custos fixos são independentes do volume, então mantém o valor planejado
        
        Para **Custo Variável:**
        ```
        Flex_Bud_Variável = B_Variável * rho
                           = B_Variável * (V_Real / V_Budget)
        ```
        **Regra Aplicada:** Variável varia proporcionalmente ao volume
        - Valor do budget: `B_Variável`
        - Multiplica pela proporção: `rho = V_Real / V_Budget`
        - Motivo: Se o volume real for maior que o planejado, o custo variável deve aumentar proporcionalmente
        
        **Flex Bud Total (em Custo Total):**
        ```
        Flex_Bud_Total = Flex_Bud_Fixo + Flex_Bud_Variável
                       = B_Fixo + (B_Variável * rho)
                       = B_Fixo + B_Variável * (V_Real / V_Budget)
        ```
        **Regra Aplicada:** Soma do Fixo (inalterado) + Variável (ajustado)
        """)
        
    st.markdown("---")
        
    st.markdown("""
        ### Cálculo em CPU (Custo por Unidade) - Real x Budget
        
        **IMPORTANTE:** No modo CPU, calcular em Custo Total primeiro, depois converter.
        
        **BUD em CPU:**
        ```
        BUD_CPU = B_Total / V_Budget
                 = (B_Fixo + B_Variável) / V_Budget
        ```
        
        **Flex Bud em CPU:**
        ```
        Flex_Bud_CPU = Flex_Bud_Total / V_Real
                     = [B_Fixo + B_Variável * (V_Real / V_Budget)] / V_Real
                     = (B_Fixo / V_Real) + (B_Variável / V_Budget)
        ```
        
        **Total Real em CPU:**
        ```
        Total_Real_CPU = R_Total / V_Real
        ```
        
        **Diferenças:**
        
        **Flex Bud - BUD:**
        ```
        Delta_Flex_Bud = Flex_Bud_CPU - BUD_CPU
                   = [(B_Fixo / V_Real) + (B_Variável / V_Budget)] - [(B_Fixo + B_Variável) / V_Budget]
                   = (B_Fixo / V_Real) - (B_Fixo / V_Budget)
                   = B_Fixo * (1/V_Real - 1/V_Budget)
                   = B_Fixo * (V_Budget - V_Real) / (V_Real * V_Budget)
        ```
        
        **Total - Flex Bud:**
        ```
        Delta_Total_Flex = Total_Real_CPU - Flex_Bud_CPU
                     = (R_Total / V_Real) - [(B_Fixo / V_Real) + (B_Variável / V_Budget)]
        ```
        """)
        
    st.markdown("---")
        
    st.markdown("""
        ### Implementação - Real x Budget
        
        ```python
        # 1. Obter dados de Budget
        df_budget = load_budget_data(ano_selecionado)
        
        # 2. Separar Fixo e Variável do Budget
        if 'Custo' in df_budget.columns:
            B_Fixo = df_budget[df_budget['Custo'] == 'Fixo']['Total'].sum()
            B_Variavel = df_budget[df_budget['Custo'] == 'Variável']['Total'].sum()
        else:
            B_Fixo = 0
            B_Variavel = df_budget['Total'].sum()  # Tudo é variável
        
        B_Total = B_Fixo + B_Variavel
        
        # 3. Obter volumes
        volume_real = df_vol_real['Volume'].sum()
        volume_budget = df_vol_budget['Volume'].sum()
        
        # 4. Calcular proporção
        rho = volume_real / volume_budget if volume_budget != 0 else 1.0
        
        # 5. Calcular Flex Bud (em Custo Total)
        Flex_Bud_Fixo = B_Fixo  # Não varia
        Flex_Bud_Variavel = B_Variavel * rho  # Varia proporcionalmente
        Flex_Bud_Total = Flex_Bud_Fixo + Flex_Bud_Variavel
        
        # 6. Converter para CPU (se necessário)
        if tipo_visualizacao == "CPU (Custo por Unidade)":
            BUD_CPU = B_Total / volume_budget if volume_budget != 0 else 0
            Flex_Bud_CPU = Flex_Bud_Total / volume_real if volume_real != 0 else 0
            Total_Real_CPU = df_real['Total'].sum() / volume_real if volume_real != 0 else 0
            
            Delta_Flex_Bud = Flex_Bud_CPU - BUD_CPU
            Delta_Total_Flex = Total_Real_CPU - Flex_Bud_CPU
        ```
        """)
        
    st.markdown("---")
        
    st.markdown("""
        ### Exemplo Prático - Real x Budget
        
        **Dados:**
        - Volume Real (`V_Real`): 50,000 unidades
        - Volume Budget (`V_Budget`): 60,000 unidades
        - Custo Total Fixo Budget (`B_Fixo`): R$ 200,000
        - Custo Total Variável Budget (`B_Variável`): R$ 400,000
        - Custo Total Budget (`B_Total`): R$ 600,000
        - Custo Total Real (`R_Total`): R$ 550,000
        
        **Cálculo Passo a Passo:**
        
        **1. Calcular Proporção de Volume:**
        ```
        rho = V_Real / V_Budget = 50,000 / 60,000 = 0.833333
        ```
        *Interpretação: Volume real foi 16.67% menor que o planejado*
        
        **2. Aplicar Regra para Custo Fixo:**
        ```
        Flex_Bud_Fixo = B_Fixo
                       = R$ 200,000
        ```
        *Regra: Fixo não varia -> mantém valor do budget*
        
        **3. Aplicar Regra para Custo Variável:**
        ```
        Flex_Bud_Variável = B_Variável * rho
                           = R$ 400,000 * 0.833333
                           = R$ 333,333.33
        ```
        *Regra: Variável varia proporcionalmente -> ajusta pelo volume real*
        
        **4. Calcular Total:**
        ```
        Flex_Bud_Total = Flex_Bud_Fixo + Flex_Bud_Variável
                        = R$ 200,000 + R$ 333,333.33
                        = R$ 533,333.33
        ```
        
        **Em CPU:**
        ```
        BUD_CPU = R$ 600,000 / 60,000 = R$ 10.00 por unidade
        Flex_Bud_CPU = R$ 533,333.33 / 50,000 = R$ 10.67 por unidade
        Total_Real_CPU = R$ 550,000 / 50,000 = R$ 11.00 por unidade
        
        Delta_Flex_Bud = R$ 10.67 - R$ 10.00 = R$ 0.67 por unidade
        Delta_Total_Flex = R$ 11.00 - R$ 10.67 = R$ 0.33 por unidade
        ```
        
        **Interpretação:**
        - O volume real foi 16.67% menor que o planejado (`rho = 0.833333`)
        - O budget variável foi ajustado proporcionalmente: R$ 400,000 -> R$ 333,333.33
        - O budget fixo permaneceu igual: R$ 200,000
        - Em CPU, o Flex Bud aumentou porque o custo fixo foi concentrado em menos volume
        - O Total Real está R$ 0.33 acima do Flex Bud, indicando ineficiência operacional
        """)
        
    st.markdown("---")
        
    st.markdown("""
        ### Exemplo Prático 2 - Real x Budget (Volume Real > Volume Budget)
        
        **Dados:**
        - Volume Real (`V_Real`): 62,208 unidades
        - Volume Budget (`V_Budget`): 60,120 unidades
        - Custo Total Fixo Budget (`B_Fixo`): R$ 200,000
        - Custo Total Variável Budget (`B_Variável`): R$ 400,000
        - Custo Total Budget (`B_Total`): R$ 600,000
        
        **Cálculo Passo a Passo:**
        
        **1. Calcular Proporção de Volume:**
        ```
        rho = V_Real / V_Budget = 62,208 / 60,120 = 1.0347
        ```
        *Interpretação: Volume real foi 3.47% maior que o planejado*
        
        **2. Aplicar Regra para Custo Fixo:**
        ```
        Flex_Bud_Fixo = B_Fixo
                       = R$ 200,000
        ```
        *Regra: Fixo não varia -> mantém valor do budget*
        
        **3. Aplicar Regra para Custo Variável:**
        ```
        Flex_Bud_Variável = B_Variável * rho
                           = R$ 400,000 * 1.0347
                           = R$ 413,880
        ```
        *Regra: Variável varia proporcionalmente -> ajusta pelo volume real*
        
        **4. Calcular Total:**
        ```
        Flex_Bud_Total = Flex_Bud_Fixo + Flex_Bud_Variável
                        = R$ 200,000 + R$ 413,880
                        = R$ 613,880
        ```
        *Resultado: Flex_Bud_Total (R$ 613,880) > BUD_Total (R$ 600,000) [CORRETO]*
        
        **Em CPU:**
        ```
        BUD_CPU = R$ 600,000 / 60,120 = R$ 9.98 por unidade
        Flex_Bud_CPU = R$ 613,880 / 62,208 = R$ 9.87 por unidade
        ```
        
        **Diferenças:**
        ```
        Delta_Flex_Bud (Custo Total) = R$ 613,880 - R$ 600,000 = R$ 13,880 (positivo) [CORRETO]
        Delta_Flex_Bud (CPU) = R$ 9.87 - R$ 9.98 = -R$ 0.11 (negativo)
        ```
        
        **Interpretação:**
        - O volume real foi 3.47% maior que o planejado (`rho = 1.0347`)
        - O budget variável foi ajustado proporcionalmente: R$ 400,000 -> R$ 413,880
        - O budget fixo permaneceu igual: R$ 200,000
        - **Em Custo Total:** Flex_Bud_Total > BUD_Total (porque o custo variável aumentou)
        - **Em CPU:** Flex_Bud_CPU < BUD_CPU (porque o custo fixo foi diluído em mais volume)
        """)
        
    st.markdown("---")
        
    st.markdown("""
        ### Comparação: Real x Real vs Real x Budget
        
        | Aspecto | Real x Real (Waterfall) | Real x Budget (TC Ext) |
        |---------|------------------------|------------------------|
        | **Base** | Custo Real Mês 1 | Custo Budget |
        | **Volume Referência** | Volume Real Mês 1 | Volume Budget |
        | **Volume Ajuste** | Volume Real Mês 2 | Volume Real |
        | **Proporção** | `V_2 / V_1` | `V_Real / V_Budget` |
        | **Objetivo** | Ajustar Mês 1 pelo volume do Mês 2 | Ajustar Budget pelo volume Real |
        | **Uso** | Comparar dois períodos reais | Comparar Real vs Planejado |
        """)
        
    st.markdown("---")
        
    st.markdown("""
        ### Regras Gerais Aplicáveis a Ambos os Casos
        
        **1. Custo Fixo:**
        - Sempre mantém o valor original (não varia com volume)
        - `Flex_Fixo = Valor_Original`
        
        **2. Custo Variável:**
        - Varia proporcionalmente ao volume
        - `Flex_Variável = Valor_Original * (Volume_Novo / Volume_Original)`
        
        **3. Ordem de Cálculo:**
        1. Calcular em **Custo Total** primeiro
        2. Separar Fixo e Variável
        3. Aplicar proporção de volume apenas ao Variável
        4. Somar Fixo + Variável ajustado
        5. Se necessário, converter para **CPU** dividindo pelo volume final
        
        **4. Tratamento de Divisão por Zero:**
        - Se `Volume_Original = 0`: usar `rho = 1.0` (sem ajuste)
        - Se `Volume_Final = 0`: usar `Flex_CPU = 0`
        """)

# ==========================================
# TC VEÍCULOS: CÁLCULO POR TABELAS/GRÁFICOS
# ==========================================
elif indice_selecionado == "🧮 Cálculo por Tabelas/Gráficos (Normal vs CPU)" and modulo_doc.startswith("📌 Ambos"):
    st.header("🧮 Cálculo por Tabelas/Gráficos (Normal vs CPU) — TC Ext + TC Veículos")

    st.subheader("📊 TC Estendido")
    _render_doc_tabelas_graficos_tc_ext()

    st.markdown("---")

    st.subheader("🚗 TC Veículos")
    _render_doc_tabelas_graficos_tc_veiculos()

    st.stop()

elif indice_selecionado == "🧮 Cálculo por Tabelas/Gráficos (Normal vs CPU)" and modulo_doc == "🚗 TC Veículos":
    st.header("🧮 Cálculo por Tabelas/Gráficos — TC Veículos")

    _render_doc_tabelas_graficos_tc_veiculos()
    st.stop()

    st.info(
        "📌 **Módulo TC Veículos** — Tabelas e gráficos específicos do TC Veículos."
    )

    with st.expander("📊 **Análise Flex por Categoria**", expanded=True):
        st.markdown("""
        ### 🔍 Modos de Visualização

        - **Fixo/Variável**: Expanders `💰 Fixo` e `💰 Variável`, cada um com sub-expanders por `Type 05` → tabela por `Account`
        - **Total**: Expanders direto por `Type 05` → tabela por `Account`

        **Expander TOTAL:**
        - Re-agrega **todas** as linhas das oficinas por `(Type 05, Type 06, Account, Custo)`
        - Mostra tabela detalhada com todas as contas (não apenas 1 linha sintética)
        - Mesmo layout dos expanders por oficina

        ### 📋 Tabela Flex por Account

        | Coluna | Cálculo |
        |--------|---------|
        | Account | Nome da conta |
        | BUD | Σ Custo FP Budget |
        | Flex Bud − BUD | Flex − BUD |
        | Flex BUD | Fixo: BUD / Variável: BUD × Proporção |
        | Total − Flex Bud | Real − Flex |
        | Total | Σ Custo FP Real |
        | Total / Flex Bud | Real/Flex (com barrinha de progresso) |

        ### 🎨 Barrinha de Progresso
        - 🟢 Verde: ≤ 90%
        - 🟡 Gradiente verde→vermelho: 90%–100%
        - 🔴 Vermelho: ≥ 100%
        """)

    with st.expander("📈 **Gráficos do TC Veículos**", expanded=False):
        st.markdown("""
        ### 📊 Custo FP por Período
        - **Barras**: Real por período com degradê roxo (`scheme='purples'`)
        - **Linha pontilhada**: Flex BUD (laranja, `strokeDash=[10,5]`)
        - **Delta**: Gráfico inferior com `Real − Flex BUD` (verde/vermelho)
        - Biblioteca: **Altair** com `data_transformers.disable_max_rows()`

        ### 🎨 Cores do Best Estimate
        Na página de **Análise BE**, os gráficos por período usam codificação por cor
        na coluna `Tipo` para diferenciar meses:
        - 🟣 **Roxo escuro** (`#4C1D95`): meses **Históricos** (realizados)
        - 🟣 **Roxo claro** (`#C4B5FD`): meses de **Best Estimate** (projetados)

        ### 📊 Volume
        - **Barras**: Volume Realizado (degradê verde)
        - **Linha tracejada**: Volume Budget (laranja)
        - **Por Veículo**: Barras agrupadas por modelo

        ### 📊 Custos por Oficina
        - Barras Custo FP por Oficina
        - Barras Rateio FA por Oficina (verde/vermelho)
        - Tabela BUD vs Flex pivotada Oficina × Período
        """)

    with st.expander("📋 **Tabs Disponíveis**", expanded=False):
        st.markdown("""
        ### 🗂️ Organização em Tabs

        O TC Veículos organiza os dados em **6 tabs**:

        | Tab | Conteúdo |
        |-----|----------|
        | 🚗 TC Veículos | KPIs resumo + Gráfico Custo FP × Flex BUD por período |
        | 📊 Análise Flex | Fixo/Variável com hierarquia Type 05 → Account |
        | 📈 Volume | Budget vs Realizado (por período e por veículo) |
        | 🏢 Custos por Oficina | Custo FP e Rateio FA por oficina |
        | ⏱️ Tempo de Produção | Tempo Veículo vs Tempo FA por oficina |
        | 📋 Dados Detalhados | Tabelas exportáveis de Real e Budget |
        """)

# ==========================================
# SEÇÃO 2: CÁLCULO POR TABELAS/GRÁFICOS
# ==========================================
elif indice_selecionado == "🧮 Cálculo por Tabelas/Gráficos (Normal vs CPU)":
    st.header("🧮 Cálculo por Tabelas/Gráficos — TC Estendido")

    st.markdown(
        "Esta seção explica os pontos que mais geram divergência entre **tabela** e **gráfico** "
        "no TC Ext (Normal vs CPU)."
    )

    _render_doc_tabelas_graficos_tc_ext()
    st.stop()

# ==========================================
# TC VEÍCULOS: ARQUITETURA E ESTRUTURA
# ==========================================
elif indice_selecionado == "🏗️ Arquitetura e Estrutura" and modulo_doc.startswith("📌 Ambos"):
    st.header("🏗️ Arquitetura e Estrutura — TC Ext + TC Veículos")

    st.subheader("📊 TC Estendido")
    _render_doc_arquitetura_tc_ext()

    st.markdown("---")

    st.subheader("🚗 TC Veículos")
    _render_doc_arquitetura_tc_veiculos()

    st.stop()

elif indice_selecionado == "🏗️ Arquitetura e Estrutura" and modulo_doc == "🚗 TC Veículos":
    st.header("🏗️ Arquitetura e Estrutura — TC Veículos")

    _render_doc_arquitetura_tc_veiculos()
    st.stop()

    st.info(
        "📌 **Módulo TC Veículos** — Estrutura de pastas, contratos de dados e pipeline de processamento."
    )

    with st.expander("📁 **Contratos de Dados (Parquets)**", expanded=True):
        st.markdown("""
        ### 📂 Estrutura de Pastas

        ```
        dados/TC_Principal/
        ├── {ano}/
        │   ├── BUD/
        │   │   ├── df_principal_BUD.parquet         # Custo consolidado BUD
        │   │   ├── df_vol_veiculos_BUD.parquet      # Volume por veículo BUD
        │   │   ├── df_veiculos_custo_fp_BUD.parquet  # Custo FP rateado BUD
        │   │   ├── df_veiculos_cpu_BUD.parquet      # CPU por veículo BUD
        │   │   ├── df_tempo_veiculos_BUD.parquet    # Tempo de produção BUD
        │   │   ├── df_dea_dedicado_BUD.parquet      # D&A Dedicado BUD
        │   │   └── df_volume_fa_BUD.parquet         # Volume Fluxo Anexo BUD
        │   ├── df_principal.parquet                 # Custo Real consolidado
        │   ├── df_vol_veiculos_actual.parquet       # Volume Realizado
        │   ├── df_veiculos_custo_fp.parquet         # Custo FP Real rateado
        │   └── df_veiculos_cpu.parquet              # CPU Real
        ├── Forecast/
        │   ├── forecast_completo.parquet            # Consolidado final com histórico + BE + BE Manual
        │   ├── forecast_historico.parquet           # Histórico sem os meses previstos
        │   ├── forecast_previsao.parquet            # Apenas períodos futuros
        │   ├── forecast_veiculos_custo_fp.parquet   # Forecast rateado por veículo
        │   └── config_forecast.json                 # Configurações persistidas do simulador
        └── historico_consolidado/
            ├── df_principal_historico.parquet        # Multi-ano consolidado
            └── BUD/
                └── df_principal_historico_BUD.parquet
        ```

        ### 📋 Schema — Principal BUD

        | Coluna | Tipo | Descrição |
        |--------|------|-----------|
        | `Oficina` | str | Centro de custo (oficina) |
        | `Veículo` | str | Modelo do veículo |
        | `Type 05` | str | Classificação nível 1 |
        | `Type 06` | str | Classificação nível 2 |
        | `Custo` | str | Fixo ou Variável |
        | `Account` | str | Conta contábil (inclui "Redis") |
        | `Período` | str | Mês por extenso |
        | `Despesa Primaria` | float | Despesa primária (R$) |
        | `Custo FA` | float | Custo do Fluxo Anexo |
        | `Custo FP` | float | Custo FP consolidado |
        | `D&A dedicado` | float | D&A dedicada |
        | `FP sem Dedicada` | float | Custo FP sem D&A |

        ### 📋 Schema — Veículos Rateado (BUD)

        | Coluna | Tipo | Descrição |
        |--------|------|-----------|
        | `Oficina` | str | Centro de custo |
        | `Veículo` | str | Modelo do veículo |
        | `Custo Rateado` | float | Custo × percentual do veículo |
        | `D&A dedicado` | float | D&A dedicada direta |
        | `Custo FP Veiculo` | float | Rateado + D&A |
        | `Ano` | int | Ano de referência |

        > O parquet BUD veículos tem `Custo FP Veiculo` (não `Custo FP`). O sistema faz mapeamento automático.
        """)

    with st.expander("🔧 **Módulos e Arquivos**", expanded=False):
        st.markdown("""
        ### 📂 Estrutura do Código

        ```
        tc_principal/
        ├── __init__.py
        ├── shared.py              # Constantes, loaders, helpers, ratear_be_por_veiculo()
        ├── ui_components.py       # Sidebar filters, CSS, KPIs
        └── pages/
            ├── __init__.py
            ├── home_tc.py                      # Página principal (6 tabs) + consumo/análise do Forecast (Real vs BE)
            ├── best_estimate_simulador_tc.py   # Simulador de premissas BE (gera Forecast)
            └── waterfall_tc.py                 # Análise Waterfall (Real + Budget)
        ```

        ### ⚙️ Filtros — Arquitetura Unificada

        ```
        Sidebar filters
             │
             ├── Veículo = "Todos" ──► usar_rateado = False
             │         ├── df_principal_BUD  → df_bud
             │         └── df_principal_Real → df
             │
             └── Veículo = "CC21 biton" ──► usar_rateado = True
                       ├── df_veiculos_custo_fp_BUD → df_bud (filtrado)
                       └── df_veiculos_custo_fp_Real → df (filtrado)
             │
        aplicar_fator_df() + converter_moeda_df()
             │
        calcular_flex_budget()
             │
        ┌─────────────────────────────┐
        │  Todos os tabs usam         │
        │  df_bud, df, df_vol_bud,    │
        │  df_vol_actual, df_flex     │
        └─────────────────────────────┘
        ```
        """)

    with st.expander("⚙️ **ETL e Processamento**", expanded=False):
        st.markdown("""
        ### 📋 Arquivos de Processamento

        | Arquivo | Função |
        |---------|--------|
        | `tc_principal/pages/extracao_dados_tc.py` | Orquestra upload, pré-validação e execução (Real/Budget) |
        | `processamento_dados_veiculos_BUD.py` | Processa Budget (BUD) e grava parquets BUD |
        | `processamento_dados_veiculos.py` | Processa Real (Sapiens/Redis) e grava parquets Real |

        ### 🔄 Pipeline

        1. Extração dos dados brutos (Excel/SAP)
        2. Normalização de colunas e períodos
        3. Cálculo de composição de custos (Desp. Primária → FA → FP)
        4. Rateio por veículo (tempo de produção)
        5. Cálculo de CPU por veículo
        6. Gravação em Parquet na pasta `dados/TC_Principal/{ano}/`

        ### 💾 Cache
        - `@st.cache_data(ttl=3600)` em todos os loaders
        - Botão "🔄 Limpar Cache" na sidebar para forçar recarga
        """)

    with st.expander("🌐 **Configurações Globais**", expanded=False):
        st.markdown("""
        ### 💱 Moeda

        | Código | Símbolo | Conversão |
        |--------|---------|-----------|
        | BRL | R$ | 1.0 (base) |
        | USD | $ | 1/Taxa USD→BRL |
        | EUR | € | 1/Taxa EUR→BRL |

        ### 📊 Fator

        | Opção | Divisor |
        |-------|---------|
        | Nenhum | 1 |
        | K (milhares) | 1.000 |
        | M (milhões) | 1.000.000 |

        ### 👁️ Tipo de Visualização

        | Tipo | Comportamento |
        |------|---------------|
        | Custo Total | Valores absolutos em R$/USD/EUR |
        | CPU | Custo ÷ Volume (fator = Nenhum) |
        """)

# ==========================================
# SEÇÃO 2: ARQUITETURA E ESTRUTURA
# ==========================================
elif indice_selecionado == "🏗️ Arquitetura e Estrutura":
    st.header("🏗️ Arquitetura e Estrutura — TC Estendido")

    _render_doc_arquitetura_tc_ext()
    st.stop()
    
    st.markdown("""
    Esta seção documenta a arquitetura, estrutura de arquivos, tecnologias utilizadas
    e informações sobre a equipe responsável pelo desenvolvimento do projeto.
    """)
    
    st.markdown("---")
    
    # Métricas principais
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.metric("💻 Linhas de Código", "20.000+", "Sistema completo")
    
    with col2:
        st.metric("📊 Páginas", "6", "Funcionalidades completas")
    
    with col3:
        st.metric("⚡ Otimização", "70%+", "Memória reduzida")
    
    with col4:
        st.metric("📁 Arquivos", "Parquet", "Formato otimizado")
        
    # EXPANDER 1: Estrutura de Arquivos
    with st.expander("📁 **Estrutura de Arquivos e Organização do Projeto**", expanded=False):
        st.subheader("📁 Estrutura de Arquivos")
        
        st.markdown("""
        ### Estrutura do Projeto (visão de alto nível)
        
        ```
        TC/
        ├── app.py                     # Portal / Router (menu via st.navigation)
        ├── pages/                     # Páginas legadas (Waterfall/BE Simulador/Extração/Documentação)
        ├── tc_ext/                    # TC Ext (Linhas Secundárias)
        ├── tc_principal/              # TC Veículos (TC Principal)
        ├── tc_core/                   # Shared (paths/portabilidade/períodos/schema/moedas/UI)
        ├── tc_copilot/                # IA (chat + relatório PDF)
        └── dados/                     # Dados organizados por módulo
            ├── TC_Ext/                # dados/TC_Ext/{ANO}/, historico_consolidado/, Forecast/
            └── TC_Principal/          # dados/TC_Principal/{ANO}/, historico_consolidado/, Forecast/
        ```
        
        **Observações:**
        - A estrutura de dados é **por módulo** (o histórico fica em `dados/TC_Ext/historico_consolidado/` e `dados/TC_Principal/historico_consolidado/`).
        - Caminhos canônicos (Dev ↔ EXE) ficam em `tc_core/data/paths.py` + `tc_core/utils/portabilidade.py`.
        """)

        st.markdown("---")
        
        st.subheader("📄 Arquivos Principais")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("""
            **app.py**
            - Portal/roteador do SCI (menu lateral via `st.navigation`)
            - Agrupa páginas de **TC Veículos**, **TC Ext**, **Documentação** e **TC Copilot**
            - Não contém a lógica de cálculo — ela está nos módulos `tc_ext/` e `tc_principal/`

            **tc_ext/pages/home_ext.py**
            - Home do TC Ext (Real/Budget/Flex/CPU)
            - Filtros + gráficos + exportação
            
            **pages/1 - Waterfall.py** (~4.000 linhas)
            - Análise waterfall entre períodos
            - Cálculo Flex Mês 1
            - Gráficos waterfall interativos
            - Tabelas com hierarquia
            
            **pages/2 - Best Estimate - Simulador.py** (~4.300 linhas)
            - Simulação interativa de Best Estimate
            - Ajuste de sensibilidade em tempo real
            - Configuração de inflação
            - Gráficos de premissas
            """)
        
        with col2:
            st.markdown("""
            **tc_ext/pages/be_analise_ext.py**
            - Best Estimate (Análise) no TC Ext (substitui a análise legacy)
            - Mesma base visual e de cálculo da Home (TC Ext)
            - Lê os outputs do simulador em `dados/TC_Ext/Forecast/`
            - Regra de CPU aplicada de forma consistente (Total/Volume)
            
            **(removido) pages/4 - Waterfall_Analysis.py** (página duplicada removida)
            - Análise waterfall entre períodos (legado)
            - Cálculo Flex Mês 1
            - Gráficos waterfall interativos
            
            **pages/5 - Extração de Dados.py** (~600 linhas)
            - Interface para extração e processamento de dados
            - Upload de arquivos
            - Validação de arquivos
            - Execução de notebooks de processamento
            
            **pages/6 - Documentacao.py** (~3.900 linhas)
            - Documentação completa do sistema
            - Regras e cálculos
            - Arquitetura e estrutura
            - Guia de extração de dados
            """)
    
    # Sub-expander: Estrutura da Pasta dados
    with st.expander("📂 **Estrutura e Funcionamento da Pasta `dados/`**", expanded=False):
        st.markdown("""
            ### 📂 Organização da Pasta `dados/`
            
            A pasta `dados/` é o coração do sistema, onde todos os arquivos processados são armazenados.
            Ela é organizada de forma hierárquica para facilitar o gerenciamento e acesso aos dados.
            
            **Estrutura Completa (padronizada por módulo):**
            ```
            dados/
            ├── TC_Ext/                         # 📊 TC Ext (Linhas Secundárias)
            │   ├── {ANO}/
            │   │   ├── df_final.parquet
            │   │   ├── df_vol.parquet
            │   │   ├── df_ke5z_group.parquet
            │   │   ├── Dados SAPIENS.xlsx
            │   │   ├── Reporting fluxo anexo.xlsx
            │   │   └── BUD/
            │   │       ├── df_final_BUD.parquet
            │   │       ├── df_vol_BUD.parquet
            │   │       └── df_ke5z_group_BUD.parquet
            │   ├── historico_consolidado/
            │   │   ├── df_final_historico.parquet
            │   │   ├── df_ke5z_historico.parquet
            │   │   ├── df_vol_historico.parquet
            │   │   └── BUD/
            │   │       ├── df_final_historico_BUD.parquet
            │   │       ├── df_ke5z_historico_BUD.parquet
            │   │       └── df_vol_historico_BUD.parquet
            │   └── Forecast/                   # 🔮 Outputs do Best Estimate / Forecast (TC Ext)
            │
            └── TC_Principal/                   # 🚗 TC Veículos (TC Principal)
                ├── {ANO}/
                │   ├── df_principal.parquet
                │   ├── df_tc_sapiens.parquet
                │   ├── df_veiculos_custo_fp.parquet
                │   ├── df_vol_veiculos_actual.parquet
                │   └── BUD/
                │       ├── df_principal_BUD.parquet
                │       ├── df_veiculos_custo_fp_BUD.parquet
                │       └── df_vol_veiculos_BUD.parquet
                ├── historico_consolidado/
                └── Forecast/                   # 🔮 Outputs do Best Estimate (TC Veículos)
                    ├── forecast_completo.parquet
                    ├── forecast_historico.parquet
                    ├── forecast_previsao.parquet
                    ├── forecast_veiculos_custo_fp.parquet
                    └── config_forecast.json
            ```
            """)
            
        st.markdown("---")
            
        st.markdown("""
            ### 🔄 Como as Pastas São Criadas e Atualizadas
            
            **1. Criação Inicial da Estrutura:**
            
            Quando o sistema é executado pela primeira vez ou quando novos dados são processados,
            o sistema verifica e cria automaticamente as pastas necessárias:
            
            ```python
            # Caminhos canônicos (dev ↔ EXE) — tc_core/data/paths.py
            from tc_core.data.paths import PASTA_TC_EXT, PASTA_TC_PRINCIPAL

            # TC Ext
            pasta_ano_tc_ext = f"{PASTA_TC_EXT}/{ANO_ATUAL}"            # dados/TC_Ext/{ANO}
            pasta_bud_tc_ext = f"{pasta_ano_tc_ext}/BUD"               # dados/TC_Ext/{ANO}/BUD
            pasta_hist_tc_ext = f"{PASTA_TC_EXT}/historico_consolidado" # dados/TC_Ext/historico_consolidado

            # TC Veículos
            pasta_ano_tc_principal = f"{PASTA_TC_PRINCIPAL}/{ANO_ATUAL}" # dados/TC_Principal/{ANO}
            pasta_bud_tc_principal = f"{pasta_ano_tc_principal}/BUD"     # dados/TC_Principal/{ANO}/BUD
            ```
            
            **2. Processo de Atualização:**
            
            **a) Processamento de Dados do Ano:**
            - Os arquivos Excel (`Dados SAPIENS.xlsx`, `Reporting fluxo anexo.xlsx`) do **TC Ext** ficam em `dados/TC_Ext/{ANO}/`
            - O notebook `tc_ext/notebooks/dados.ipynb` processa esses arquivos e gera os arquivos Parquet
            - Os arquivos Parquet são salvos na mesma pasta do ano (`dados/TC_Ext/{ANO}/`)
            - **Simultaneamente**, os dados são consolidados no histórico
            
            **b) Consolidação no Histórico:**
            - Após processar os dados do ano, o sistema **concatena** os novos dados com o histórico existente
            - Os arquivos em `historico_consolidado/` são **atualizados** (não substituídos)
            - Isso permite que o sistema tenha acesso a **todos os dados históricos** em um único lugar
            
            **c) Processamento de Budget:**
            - Similar ao processo de dados do ano, mas os arquivos são processados pelo `tc_ext/notebooks/dados_BUD.ipynb`
            - Os **outputs** de Budget do **TC Ext** são salvos em `dados/TC_Ext/{ANO}/BUD/`
            - O histórico de Budget do **TC Ext** é consolidado em `dados/TC_Ext/historico_consolidado/BUD/`
            
            **d) Processamento de Forecast:**
            - Forecast do **TC Ext**: outputs em `dados/TC_Ext/Forecast/`
            - Forecast do **TC Veículos**: outputs em `dados/TC_Principal/Forecast/`
            """)
        
        st.markdown("---")
        
        st.markdown("""
        ### 🔗 Como as Pastas Funcionam Entre Si
            
            **1. Relação entre Pastas por Ano e Histórico:**
            
            ```
            dados/TC_Ext/2026/df_final.parquet  ──┐
                                                  ├──> Concatena ──> dados/TC_Ext/historico_consolidado/df_final_historico.parquet
            dados/TC_Ext/2025/df_final.parquet  ──┘
            ```
            
            - **Dados do Ano:** Contêm apenas os dados do ano específico (útil para filtros rápidos)
            - **Histórico Consolidado:** Contém **TODOS** os anos concatenados (usado pelo sistema principal)
            - O sistema **prioriza** o histórico consolidado para análises que precisam de múltiplos anos
            
            **2. Fluxo de Dados:**
            
            ```
            Arquivos Excel (entrada)
                │
                ├──> Processamento (tc_ext/notebooks/dados.ipynb) — TC Ext
                │       │
                │       ├──> Salva em dados/TC_Ext/{ANO}/ (dados do ano)
                │       │
                │       └──> Concatena em dados/TC_Ext/historico_consolidado/ (histórico completo)
                │
                └──> Sistema Streamlit lê de dados/TC_Ext/historico_consolidado/ (fonte principal do TC Ext)
            ```
            
            **3. Separação de Budget:**
            
            - **TC Ext (Real):** `dados/TC_Ext/{ANO}/` e `dados/TC_Ext/historico_consolidado/`
            - **TC Ext (Budget):** `dados/TC_Ext/{ANO}/BUD/` e `dados/TC_Ext/historico_consolidado/BUD/`
            - Esta separação evita misturar outputs de Budget com Real
            
            **4. Forecast como Dados Derivados:**
            
            - As pastas `Forecast/` contêm dados **processados e calculados** pelo sistema
            - Não são dados de entrada, mas sim **resultados** de cálculos de forecast
            - São gerados dinamicamente quando o usuário executa o Forecast
            """)
        
        st.markdown("---")
        
        st.markdown("""
        ### ⚙️ Regras de Criação e Atualização
            
            **Regra 1: Criação Automática**
            - Todas as pastas são criadas automaticamente quando necessário
            - O parâmetro `exist_ok=True` garante que não há erro se a pasta já existir
            - Não é necessário criar manualmente nenhuma pasta
            
            **Regra 2: Consolidação Incremental**
            - O histórico é **atualizado** (não substituído) a cada processamento
            - Novos dados são **adicionados** ao histórico existente
            - Isso mantém a integridade dos dados históricos
            
            **Regra 3: Separação por Tipo**
            - Dados Reais e Budget são mantidos **separados** em pastas diferentes
            - Isso evita confusão e permite comparações precisas
            - O sistema sabe qual pasta usar baseado no modo de comparação selecionado
            
            **Regra 4: Formato Parquet**
            - Todos os arquivos processados são salvos em formato **Parquet**
            - Parquet oferece compressão e leitura rápida
            - Formato otimizado para grandes volumes de dados
            """)
    
    # EXPANDER 2: Tecnologias
    with st.expander("💻 **Tecnologias e Bibliotecas**", expanded=False):
        st.subheader("💻 Tecnologias e Bibliotecas")
        
        st.markdown(f"""
        ### Stack Tecnológico
        
        **Framework Principal:**
        - **Streamlit** {st.__version__} - Framework web para aplicações de dados
        
        **Linguagem:**
        - **Python** 3.8+ - Linguagem de programação
        
        **Processamento de Dados:**
        - **Pandas** 2.0.0+ - Manipulação e análise de dados
        - **NumPy** 1.24.0+ - Operações numéricas
        
        **Visualizações:**
        - **Altair** 5.0.0+ - Gráficos interativos
        - **Plotly** - Gráficos waterfall avançados
        
        **Formato de Dados:**
        - **PyArrow** 12.0.0+ - Suporte a Parquet
        - **Parquet** - Formato de dados otimizado
        
        **Exportação:**
        - **OpenPyXL** 3.1.0+ - Geração de arquivos Excel
        """)
        
        st.markdown("---")
        
        st.subheader("🔧 Dependências Principais")
        
        st.code("""
# requirements.txt
streamlit>=1.28.0
pandas>=2.0.0
altair>=5.0.0
numpy>=1.24.0
openpyxl>=3.1.0
pyarrow>=12.0.0
plotly>=5.0.0
        """, language="text")
        
        st.markdown("---")
        
        st.subheader("⚡ Otimizações Implementadas")
        
        st.markdown("""
        **Gestão de Memória:**
        - Cache inteligente com TTL configurável
        - Otimização de tipos: Category para strings repetidas
        - Downcast: Float64 -> Float32, Int64 -> Int32
        - Redução de cópias: Apenas quando necessário
        
        **Operações Vetorizadas:**
        - Substituição de `iterrows()` por merge e `np.where()`
        - Substituição de `apply()` por operações vetorizadas
        - Filtros booleanos ao invés de loops
        - Agrupamento otimizado com `agg()` direto
        
        **Cálculos Otimizados:**
        - CPU calculado após agrupamento (nunca antes)
        - Flex Bud com merge ao invés de loops
        - Volume sincronizado entre tabelas e gráficos
        - Cache de filtros para opções repetidas
        """)
    
    # EXPANDER 3: Desafios e Soluções
    with st.expander("⚠️ **Desafios Principais & Soluções Implementadas**", expanded=False):
        st.markdown("""
        ### 📊 Desafios Identificados
        
        - **📁 Dados grandes:** Milhões de registros causando lentidão
        - **💾 Uso de memória:** Excedia limites de processamento
        - **Instabilidade:** Sistema lento com muitos filtros
        - **🐌 Cálculos complexos:** Flex Bud e Forecast demorados
        - **🔄 Sincronização:** Dados de tabela vs gráficos diferentes
        - **📊 Visualizações:** Gráficos sem gradientes e pouco informativos
        """)
        
        st.markdown("---")
        
        st.markdown("""
        ### ✅ Soluções Implementadas
        
        - **📊 Otimização de dados:** Parquet com tipos categóricos
        - **⚡ Cache estratégico:** TTL configurável por tipo de dado
        - **🔄 Operações vetorizadas:** Substituição de iterrows() e apply()
        - **📈 Cálculos otimizados:** Flex Bud e CPU após agrupamento
        - **🎯 Sincronização:** Mesma fonte de dados para tabelas e gráficos
        - **🎨 Visualizações melhoradas:** Gradientes, delta charts, barras HTML
        """)
        
        st.info("🎆 **Resultado Final:** Sistema 100% estável com performance otimizada e visualizações profissionais!")
    
    # EXPANDER 4: Estatísticas do Sistema
    with st.expander("📊 **Estatísticas e Métricas do Sistema**", expanded=False):
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.markdown("""
            ### 💾 Dados e Performance
            
            **📁 Arquivos Principais:**
            - `df_final_historico.parquet` (dados históricos)
            - `df_vol_historico.parquet` (volumes)
            - `df_final_historico_BUD.parquet` (budget)
            
            **⚡ Otimizações:**
            - Tipos categóricos para strings
            - Downcast de numéricos
            - Compressão Parquet
            - Cache com TTL
            """)
        
        with col2:
            st.markdown("""
            ### 📊 Páginas do Sistema
            
            **📄 Páginas Disponíveis:**
            - `app.py` - Portal / Router (menu via st.navigation)
            - `1 - Waterfall.py` - Análise waterfall (~4.000 linhas)
            - `2 - Best Estimate - Simulador.py` - Simulação (~4.300 linhas)
            - `tc_ext/pages/be_analise_ext.py` - Best Estimate (Análise) (base Home)
            - `4 - Waterfall_Analysis.py` - (removido) página duplicada
            - `5 - Extração de Dados.py` - Extração e processamento (~600 linhas)
            - `6 - Documentacao.py` - Documentação (~3.900 linhas)
            
            **📊 Total:** ~33.000+ linhas de código
            """)
        
        with col3:
            st.markdown(f"""
            ### 🔧 Tecnologias
            
            **Stack Principal:**
            - Streamlit {st.__version__}
            - Pandas {pd.__version__}
            - NumPy {np.__version__}
            - Altair (versão instalada)
            - Plotly (versão instalada)
            - OpenPyXL (versão instalada)
            """)

# ==========================================
# TC VEÍCULOS: ESPECIFICAÇÃO TÉCNICA
# ==========================================
elif indice_selecionado == "🧾 Especificação Técnica" and modulo_doc.startswith("📌 Ambos"):
    st.header("🧾 Especificação Técnica — TC Ext + TC Veículos")

    st.subheader("📊 TC Estendido")
    _render_doc_especificacao_tc_ext()

    st.markdown("---")
    st.subheader("🚗 TC Veículos")
    _render_doc_especificacao_tc_veiculos()

    st.stop()

elif indice_selecionado == "🧾 Especificação Técnica" and modulo_doc == "🚗 TC Veículos":
    st.header("🧾 Especificação Técnica — TC Veículos")

    _render_doc_especificacao_tc_veiculos()

# ==========================================
# SEÇÃO 3: ESPECIFICAÇÃO TÉCNICA (REESCRITA)
# ==========================================
elif indice_selecionado == "🧾 Especificação Técnica":
    st.header("🧾 Especificação Técnica — TC Estendido")

    _render_doc_especificacao_tc_ext()

# ==========================================
# SEÇÃO 4: GUIA DE EXTRAÇÃO DE DADOS
# ==========================================
elif indice_selecionado == "📥 Guia de Extração de Dados":
    st.header("📥 Guia de Extração de Dados — TC Estendido")
    
    st.markdown("""
    <div style="padding: 1.5rem; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); border-radius: 10px; margin-bottom: 2rem; color: white;">
    <h2 style="color: white; margin: 0;">📚 Documentação Completa para IA</h2>
    <p style="color: #f0f0f0; margin: 0.5rem 0 0 0;">
            Todos os Relacionamentos, Processos e Estruturas de Dados
    </p>
    </div>
    """, unsafe_allow_html=True)

    with st.expander("🧭 **Ordem exata da extração e pontos de auditoria**", expanded=True):
        st.markdown(r"""
### Sequência obrigatória do processo

1. Definir ano e tipo de processamento.
2. Resolver caminho dos arquivos de entrada na ordem de prioridade.
3. Validar existência das abas e colunas mínimas.
4. Ler a base principal de custo.
5. Ler a base de classificação e identificar Fixo/Variável.
6. Ler rateio e transformar meses em linhas.
7. Aplicar merges pelas chaves corretas.
8. Calcular valores por veículo ou por linha consolidada.
9. Ler volume e anexar ao mesmo perímetro lógico.
10. Persistir arquivos do ano.
11. Consolidar histórico, sempre por concatenação controlada.

### O que quebra se inverter
- Se o volume entrar antes da normalização de período, o merge falha silenciosamente.
- Se o histórico for salvo antes da validação do ano corrente, a base consolidada fica contaminada.
- Se o rateio for aplicado sem chaves Oficina + Período coerentes, o valor por veículo vira arbitrário.

### Exemplo auditável de merge

Suponha uma linha de custo:
- Oficina = Montagem
- Período = Janeiro
- Valor = R$ 100.000

E rateio do mesmo par Oficina + Período:
- Veículo A = 60%
- Veículo B = 40%

Resultado correto:
- Veículo A = R$ 60.000
- Veículo B = R$ 40.000

Se o merge usasse só Oficina e ignorasse Período, Janeiro poderia capturar o rateio de Fevereiro e o valor final deixaria de ser auditável.
        """)

    with st.expander("🚗 **TC Veículos — Pipeline completo (Real e Budget)**", expanded=False):
        st.markdown("""
        ### 🔄 Fluxo de Processamento (TC Veículos)

        ```
        Arquivos Excel (Entrada)
            │
            ├──> processamento_dados_veiculos_BUD.py (Budget)
            │       ├──> Lê Budget + D&A dedicado + volumes/tempos
            │       ├──> Normaliza colunas e períodos
            │       ├──> Calcula composição (FA/FP) e rateios por veículo
            │       ├──> Grava df_principal_BUD.parquet
            │       ├──> Grava df_veiculos_custo_fp_BUD.parquet
            │       └──> Grava df_veiculos_cpu_BUD.parquet
            │
            └──> processamento_dados_veiculos.py (Real)
                    ├──> Lê Sapiens + Redis + volumes
                    ├──> Normaliza e processa
                    ├──> Grava df_principal.parquet
                    ├──> Grava df_tc_sapiens.parquet (detalhado)
                    └──> Grava df_veiculos_custo_fp.parquet / df_veiculos_cpu.parquet
        ```

        **Página Streamlit que executa o fluxo:** `tc_principal/pages/extracao_dados_tc.py`

        ### ✅ Arquivo de entrada (fonte única)

        - `Reporting veículos.xlsx` em `dados/TC_Principal/{ano}/`
        - A página `extracao_dados_tc.py` permite **upload** com proteção contra sobrescrita (checkbox de confirmação)

        ### 🧾 Abas obrigatórias — Budget (no Excel)

        - `massa primária - BDG`
        - `massa - REDIS`
        - `Volume e EST PdR - BDG`
        - `Volume BDG`
        - `Volume Actual`
        - `EST veículos - BDG`
        - `massa - D&A dedicado`

        ### 🧾 Abas obrigatórias — Real (no Excel)

        - `Sapiens`
        - `Volume e EST PdR - Actual`
        - `Volume Actual`
        - `EST veículos - Actual`

        ### 🔎 Pré-validação (o que o app checa antes de processar)

        - Se as abas obrigatórias existem
        - Budget: colunas mínimas em `massa primária - BDG` (ex.: `Oficina`, `Account`) e `massa - REDIS` (ex.: `Oficina`)
        - Budget: detecção de meses em `Volume BDG` (tentando múltiplos headers)
        - Real: em `Sapiens`, valida colunas mínimas (ex.: `Oficina`, `Account`, `Valor`)
        - Aviso operacional: para o fluxo completo, o Real depende do Budget ter gerado `df_dea_dedicado_BUD.parquet`
        - Rateios manuais (QY/GS/SM): persistidos em `rateios_manuais.json` (usados no cálculo da taxa PdR)

        ### 🧱 Consolidação histórica (multi-ano)

        A página também consolida parquets multi-ano em `dados/TC_Principal/historico_consolidado/`.

        ### 📂 Scripts e Funções

        | Arquivo | Função |
        |---------|--------|
        | `tc_principal/pages/extracao_dados_tc.py` | Orquestra execução e gravação dos parquets (Real/Budget) |
        | `processamento_dados_veiculos_BUD.py` | Processa Budget + gera parquets BUD (principal + por veículo + CPU) |
        | `processamento_dados_veiculos.py` | Processa Real (Sapiens/Redis) + gera parquets Real (principal + por veículo + CPU) |

        ### 🗃️ Principais parquets gerados

        **Budget** (`dados/TC_Principal/{ano}/BUD/`):
        - `df_principal_BUD.parquet`
        - `df_vol_veiculos_BUD.parquet` / `df_vol_veiculos_actual.parquet`
        - `df_tempo_veiculos_BUD.parquet`
        - `df_dea_dedicado_BUD.parquet`
        - `df_veiculos_percentual_rateio_BUD.parquet` / `df_veiculos_custo_rateado_BUD.parquet`
        - `df_veiculos_custo_fp_BUD.parquet` / `df_veiculos_cpu_BUD.parquet`

        **Real** (`dados/TC_Principal/{ano}/`):
        - `df_principal.parquet`
        - `df_tc_sapiens.parquet` (detalhado)
        - `df_vol_veiculos.parquet` / `df_tempo_veiculos.parquet` / `df_dea_dedicado.parquet`
        - `df_veiculos_percentual_rateio.parquet` / `df_veiculos_custo_rateado.parquet`
        - `df_veiculos_custo_fp.parquet` / `df_veiculos_cpu.parquet`
        - `df_comparativo_real_budget.parquet`

        ### 📁 Pastas (entrada e saída)
        - Entrada: `dados/TC_Principal/{ano}/` (Excel/insumos)
        - Saída Real: `dados/TC_Principal/{ano}/` (parquets Real)
        - Saída Budget: `dados/TC_Principal/{ano}/BUD/` (parquets BUD)

        ### 📊 Dados de Volume (usos)
        Os volumes são usados para:
        - CPU
        - Flex Budget (proporção Real/BUD)
        - Gráficos comparativos (BUD vs Real)
        """)
    
    # Índice interno
    st.markdown("## 📋 Índice do Guia")
    st.markdown("""
    ### 📖 Capítulo 1: Estrutura e Processamento dos Notebooks
    1. [Visão Geral](#visao-geral)
    2. [Notebook tc_ext/notebooks/dados.ipynb - Dados REAIS](#dados-reais)
    3. [Notebook tc_ext/notebooks/dados_BUD.ipynb - Dados BUDGET](#dados-budget)
    4. [Estrutura de Arquivos de Entrada](#estrutura-entrada)
    5. [Relacionamentos e Merges](#relacionamentos)
    6. [Colunas e Estrutura Final](#colunas-finais)
    7. [Consolidação do Histórico](#consolidacao)
    8. [Arquivos de Saída](#arquivos-saida)
    9. [Fluxo Completo](#fluxo-completo)
    10. [Tratamento de Erros](#tratamento-erros)
    11. [Checklist para Manutenção](#checklist)
    
    ### 🔄 Capítulo 2: Funcionamento da Atualização e Extração
    1. [Visão Geral do Processo de Atualização](#visao-atualizacao)
    2. [Ordem Cronológica dos Eventos](#ordem-cronologica)
    3. [Sistema de Busca de Arquivos](#busca-arquivos)
    4. [Criação de Pastas e Estrutura](#criacao-pastas)
    5. [Sistema de Upload de Arquivos](#sistema-upload)
    6. [Processamento e Execução](#processamento-execucao)
    7. [Cenários de Uso](#cenarios-uso)
    """)
    
    st.markdown("---")
    
    # ==========================================
    # CAPÍTULO 1: ESTRUTURA E PROCESSAMENTO DOS NOTEBOOKS
    # ==========================================
    
    with st.expander("📖 **Capítulo 1: Estrutura e Processamento dos Notebooks**", expanded=False):
        st.markdown("""
        <div style="padding: 1.5rem; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); border-radius: 10px; margin: 2rem 0; color: white;">
            <h2 style="color: white; margin: 0;">📖 Capítulo 1: Estrutura e Processamento dos Notebooks</h2>
            <p style="color: #f0f0f0; margin: 0.5rem 0 0 0;">
                Documentação Completa dos Notebooks de Extração - Estrutura, Processamento e Relacionamentos
            </p>
        </div>
        """, unsafe_allow_html=True)
        
        # Seção 1: Visão Geral
        st.markdown("## 🎯 VISÃO GERAL {#visao-geral}")
        
        st.markdown("### Objetivo dos Notebooks")
        st.markdown("""
        Os notebooks `tc_ext/notebooks/dados.ipynb` e `tc_ext/notebooks/dados_BUD.ipynb` são responsáveis por:
        - **Carregar** dados de múltiplas fontes (Excel: SAPIENS, Reporting fluxo anexo)
        - **Processar** e **normalizar** dados de diferentes formatos e guias
        - **Unificar** informações através de merges por chaves comuns
        - **Calcular** rateios por veículo (CC21, CC22, CC24, CC24 5L, CC24 7L, J516)
        - **Gerar** arquivos Parquet e Excel otimizados para uso no dashboard
        - **Consolidar** dados históricos para análises multi-anos
        """)
        
        st.markdown("### Diferença entre tc_ext/notebooks/dados.ipynb e tc_ext/notebooks/dados_BUD.ipynb")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("""
            **📊 tc_ext/notebooks/dados.ipynb - Dados REAIS**
        - Processa dados de custos **reais** (executados)
        - Lê guia **"Sapiens"** do Reporting fluxo anexo.xlsx
        - Lê guia **"Rateio"** para rateio por veículo
        - Lê guia **"Volume"** para volumes
        - Salva em: `dados/TC_Ext/{ANO}/`
        - Histórico: `dados/TC_Ext/historico_consolidado/`
        """)
        
        with col2:
            st.markdown("""
            **📈 tc_ext/notebooks/dados_BUD.ipynb - Dados BUDGET**
        - Processa dados de **orçamento/planejamento** (Budget)
        - Lê guia **"Voz de custo BDG"** do Reporting fluxo anexo.xlsx
        - Lê guia **"Rateio BDG"** para rateio por veículo
        - Lê guia **"Volume BDG"** para volumes
        - Salva em: `dados/TC_Ext/{ANO}/BUD/`
        - Histórico: `dados/TC_Ext/historico_consolidado/BUD/`
        """)
        
        st.markdown("### Fluxo Principal")
        st.code("""
        Arquivos Excel (Entrada)
            │
            ├──> tc_ext/notebooks/dados.ipynb (REAL)
            │       ├──> Processamento
            │       ├──> Merges (Account, Nº conta, Centro cst, Oficina+Período)
            │       ├──> Cálculo Rateio por Veículo
            │       ├──> Merge com Volume
            │       └──> Salvar Parquet + Consolidar Histórico
            │
            └──> tc_ext/notebooks/dados_BUD.ipynb (BUDGET)
                    ├──> Processamento (mesma lógica)
                    ├──> Merges (mesmas chaves)
                    ├──> Cálculo Rateio por Veículo
                    ├──> Merge com Volume
                    └──> Salvar Parquet (BUD) + Consolidar Histórico (BUD)
        """, language="text")
        
        st.markdown("---")
        
        # Seção 2: tc_ext/notebooks/dados.ipynb - Dados REAIS
        st.markdown("## 📊 NOTEBOOK tc_ext/notebooks/dados.ipynb - DADOS REAIS {#dados-reais}")
        
        st.markdown("### Estrutura do Processamento")
        
        with st.expander("🔧 **Célula 0: Configuração Inicial**", expanded=False):
            st.markdown("""
            **Objetivo**: Configurar ano, pastas e caminhos
            
            **Processo**:
            1. Solicita ano para processar (padrão: ano atual)
            2. Cria estrutura de pastas:
                    - `dados/TC_Ext/{ANO_ATUAL}/` - Dados do ano específico (TC Ext)
                    - `dados/TC_Ext/historico_consolidado/` - Histórico consolidado (TC Ext)
            3. Verifica arquivos de entrada:
               - `Dados SAPIENS.xlsx`
               - `Reporting fluxo anexo.xlsx`
            4. Define caminhos de entrada e saída
            
            **Variáveis Criadas**:
            - `ANO_ATUAL`: Ano selecionado para processamento
            - `PASTA_ANO`: `dados/TC_Ext/{ANO_ATUAL}/`
            - `PASTA_HISTORICO`: `dados/TC_Ext/historico_consolidado/`
            - `CAMINHO_SAPIENS`: Caminho para Dados SAPIENS.xlsx
            - `CAMINHO_RATEIO`: Caminho para Reporting fluxo anexo.xlsx
            - `CAMINHO_DF_FINAL`: `dados/TC_Ext/{ANO}/df_final.parquet`
            - `CAMINHO_DF_VOL`: `dados/TC_Ext/{ANO}/df_vol.parquet`
            - `CAMINHO_DF_KE5Z_GROUP`: `dados/TC_Ext/{ANO}/df_ke5z_group.parquet`
            """)
        
        with st.expander("📥 **Célula 1: Leitura dos Dados SAPIENS (KE5Z)**", expanded=False):
            st.markdown("""
            **Arquivo**: `Reporting fluxo anexo.xlsx`
            **Guia**: `"Sapiens"`
            **Cabeçalho**: Linha 1 (`header=1`)
            **Colunas**: A até T (20 colunas, `usecols=range(20)`)
            
            **Colunas Lidas**:
            - `Mes`, `Período`, `Nºconta`, `Centrocst`, `Nºdoc.ref.`, `Dt.lçto.`
            - `Valor`, `QTD`, `Type 05`, `Type 06`, `Account` (Type 07)
            - `USI`, `Oficina`, `Doc.compra`, `Texto breve`
            - `Fornecedor`, `Material`, `Usuário`, `Fornec.`, `Tipo`
            
            **DataFrame Criado**: `df_KE5Z`
            
            **Validação**: Soma da coluna `Valor` para verificar leitura
            """)
        
        with st.expander("🔗 **Célula 2: Merge com Base Conso (Custo)**", expanded=False):
            st.markdown("""
            **Arquivo**: `Dados SAPIENS.xlsx`
            **Guia**: `"Base conso"`
            
            **Processo**:
            1. Lê guia "Base conso"
            2. Renomeia `Type 04` → `Custo` (se existir)
            3. Mantém apenas colunas: `Custo`, `Type 07`
            4. Renomeia `Type 07` → `Account`
            5. Faz merge com `df_KE5Z` usando `Account` como chave
            
            **Chave de Merge**: `Account` (Type 07)
            **Tipo**: `left` (mantém todos os registros de KE5Z)
            
            **Resultado**: Adiciona coluna `Custo` ao `df_KE5Z`
            - Valores possíveis: `"Variável"` ou `"Fixo"`
            """)
        
        with st.expander("📊 **Célula 3: Processamento de Rateio**", expanded=False):
            st.markdown("""
            **Arquivo**: `Reporting fluxo anexo.xlsx`
            **Guia**: `"Rateio"`
            
            **Processo**:
            1. Lê guia sem header (`header=None`)
            2. Remove primeira linha (linha de referência)
            3. Usa segunda linha como cabeçalho (meses)
            4. Remove linha usada como cabeçalho
            5. Identifica colunas de meses (janeiro a dezembro)
            6. Usa `melt()` para transformar colunas de meses em linhas
            7. Cria colunas: `Período` (mês) e `Rateio` (valor)
            8. Normaliza `Período` para capitalizado (Janeiro, Fevereiro, etc.)
            9. Filtra: Remove `Oficina == 'Veículos'` e linhas com `Oficina` NaN
            
            **Colunas de Identificação (id_vars)**:
            - `Oficina`, `Veículo` (e outras colunas não-mês)
            
            **Colunas Transformadas (value_vars)**:
            - Meses: Janeiro, Fevereiro, Março, ..., Dezembro
            
            **DataFrame Criado**: `df` (com colunas: `Oficina`, `Veículo`, `Período`, `Rateio`)
            """)
        
        with st.expander("🔄 **Célula 4: Merge KE5Z ↔ Rateio e Cálculo por Veículo**", expanded=False):
            st.markdown("""
            **Chave de Merge**: `['Oficina', 'Período']` (COMPOSTA)
            **Tipo**: `left` (mantém todos os registros de KE5Z)
            
            **Processo**:
            1. Merge `df_KE5Z` com `df` (rateio) usando `['Oficina', 'Período']`
            2. Pivot: Transforma `Veículo` em colunas de `Rateio`
               - Index: `['Oficina', 'Período']`
               - Columns: `Veículo` (CC21, CC22, CC24, CC24 5L, CC24 7L, J516)
               - Values: `Rateio` (percentuais)
               - Aggfunc: `mean` (média para agregar duplicatas)
            3. Renomeia colunas de veículos: adiciona `%` (ex: `CC21%`, `CC22%`)
            4. Merge reverso: `df_KE5Z` com `df_pivot` usando `['Oficina', 'Período']`
            5. Calcula colunas de valores por veículo:
               - `CC21 = CC21% * Valor`
               - `CC22 = CC22% * Valor`
               - `CC24 = CC24% * Valor`
               - `CC24 5L = CC24 5L% * Valor`
               - `CC24 7L = CC24 7L% * Valor`
               - `J516 = J516% * Valor`
            6. Calcula `Soma_Percentuais = CC21% + CC22% + ... + J516%`
            7. Remove colunas de percentual (`CC21%`, `CC22%`, etc.)
            
            **Resultado**: `df_final` com colunas de valores por veículo calculadas
            """)
        
        with st.expander("📈 **Célula 5: Processamento de Volume**", expanded=False):
            st.markdown("""
            **Arquivo**: `Reporting fluxo anexo.xlsx`
            **Guia**: `"Volume"`
            **Cabeçalho**: Linha 51 (`header=50`, 0-indexed)
            
            **Processo**:
            1. Lê guia "Volume" com cabeçalho na linha 51
            2. Identifica colunas de meses (janeiro a dezembro)
            3. Usa `melt()` para transformar colunas de meses em linhas
            4. Cria colunas: `Período` (mês) e `Volume` (valor)
            5. Normaliza `Período` para capitalizado
            6. Converte `Volume` para numérico
            7. Remove linhas onde `Oficina` ou `Período` são NaN
            8. Preenche NaN em `Volume` com 0
            9. Remove duplicatas
            
            **Colunas Finais**: `Oficina`, `Veículo`, `Período`, `Volume`
            
            **DataFrame Criado**: `df_vol`
            """)
        
        with st.expander("🔗 **Célula 6: Merge df_final ↔ df_vol (Volume)**", expanded=False):
            st.markdown("""
            **Chave de Merge**: `['Oficina', 'Período', 'Veículo']` (COMPOSTA)
            **Tipo**: `left` (mantém todos os registros de df_final)
            
            **Processo**:
            1. Verifica se colunas de chave existem em ambos DataFrames
            2. Faz merge adicionando coluna `Volume` ao `df_final`
            3. Preenche NaN em `Volume` com 0 (se não houver match)
            
            **Resultado**: `df_final` com coluna `Volume` adicionada
            """)
        
        with st.expander("💾 **Célula 7: Salvamento e Consolidação**", expanded=False):
            st.markdown("""
            **Arquivos Salvos (Pasta do Ano)**:
            1. `df_final.parquet` - Dados completos com rateio por veículo e volume
            2. `df_vol.parquet` - Dados de volume
            3. `df_ke5z_group.parquet` - Dados agrupados (se aplicável)
            
            **Consolidação do Histórico**:
            1. Carrega histórico existente (se existir):
                    - `dados/TC_Ext/historico_consolidado/df_final_historico.parquet`
                    - `dados/TC_Ext/historico_consolidado/df_vol_historico.parquet`
            2. Adiciona coluna `Ano` aos dados do ano atual
            3. Concatena dados do ano atual com histórico existente
            4. Remove duplicatas (se houver)
            5. Salva histórico atualizado:
                    - `dados/TC_Ext/historico_consolidado/df_final_historico.parquet`
                    - `dados/TC_Ext/historico_consolidado/df_vol_historico.parquet`
            
            **IMPORTANTE**: O histórico é sempre **concatenado**, nunca substituído
            """)
        
        st.markdown("---")
        
        # Seção 3: tc_ext/notebooks/dados_BUD.ipynb - Dados BUDGET
        st.markdown("## 📈 NOTEBOOK tc_ext/notebooks/dados_BUD.ipynb - DADOS BUDGET {#dados-budget}")
        
        st.markdown("### Diferenças Principais em Relação a tc_ext/notebooks/dados.ipynb")
        
        diferencas_bud = {
            "Aspecto": [
                "Guia de Dados Principais",
                "Guia de Rateio",
                "Guia de Volume",
                "Pasta de Saída",
                "Sufixo dos Arquivos",
                "Pasta de Histórico"
            ],
            "tc_ext/notebooks/dados.ipynb (REAL)": [
                '"Sapiens"',
                '"Rateio"',
                '"Volume"',
                "dados/TC_Ext/{ANO}/",
                "Sem sufixo",
                "dados/TC_Ext/historico_consolidado/"
            ],
            "tc_ext/notebooks/dados_BUD.ipynb (BUDGET)": [
                '"Voz de custo BDG"',
                '"Rateio BDG"',
                '"Volume BDG"',
                "dados/TC_Ext/{ANO}/BUD/",
                "_BUD (ex: df_final_BUD.parquet)",
                "dados/TC_Ext/historico_consolidado/BUD/"
            ]
        }
        
        st.dataframe(pd.DataFrame(diferencas_bud), width="stretch", hide_index=True)
        
        st.markdown("### Processo Idêntico")
        st.info("""
        **IMPORTANTE**: O processo de processamento, merges, cálculos e consolidação
        é **IDÊNTICO** ao `tc_ext/notebooks/dados.ipynb`. A única diferença são as guias lidas e os
        caminhos de saída. Todas as transformações, relacionamentos e cálculos seguem
        a mesma lógica.
        """)
        
        st.markdown("---")
        
        # Seção 4: Estrutura de Arquivos de Entrada
        st.markdown("## 📁 ESTRUTURA DE ARQUIVOS DE ENTRADA {#estrutura-entrada}")
        
        st.markdown("### Arquivos Necessários")
        
        with st.expander("📊 **Reporting fluxo anexo.xlsx**", expanded=False):
            st.markdown("""
            **Localização**: `dados/TC_Ext/{ANO}/Reporting fluxo anexo.xlsx` ou raiz do projeto
            
            **Guias Utilizadas (tc_ext/notebooks/dados.ipynb - REAL)**:
            1. **"Sapiens"** (Célula 1)
               - Cabeçalho: Linha 1
               - Colunas: A até T (20 colunas)
               - Dados: Custos reais executados
            
            2. **"Rateio"** (Célula 3)
               - Cabeçalho: Segunda linha (após linha de referência)
               - Colunas de meses: Janeiro a Dezembro
               - Dados: Percentuais de rateio por Oficina, Veículo e Período
            
            3. **"Volume"** (Célula 5)
               - Cabeçalho: Linha 51
               - Colunas de meses: Janeiro a Dezembro
               - Dados: Volumes por Oficina, Veículo e Período
            
            **Guias Utilizadas (tc_ext/notebooks/dados_BUD.ipynb - BUDGET)**:
            1. **"Voz de custo BDG"** (equivalente a "Sapiens")
            2. **"Rateio BDG"** (equivalente a "Rateio")
            3. **"Volume BDG"** (equivalente a "Volume")
            """)
        
        with st.expander("📋 **Dados SAPIENS.xlsx**", expanded=False):
            st.markdown("""
            **Localização**: `dados/TC_Ext/{ANO}/Dados SAPIENS.xlsx` ou raiz do projeto
            
            **Guias Utilizadas**:
            1. **"Base conso"**
               - Colunas: `Type 04` (renomeado para `Custo`), `Type 07` (renomeado para `Account`)
               - Propósito: Mapear Account para tipo de custo (Variável/Fixo)
               - Chave de merge: `Account` (Type 07)
            
            **Observação**: Este arquivo é usado tanto em `tc_ext/notebooks/dados.ipynb` quanto em `tc_ext/notebooks/dados_BUD.ipynb`
            """)
        
        st.markdown("---")
        
        # Seção 5: Relacionamentos e Merges
        st.markdown("## 🔗 RELACIONAMENTOS E MERGES {#relacionamentos}")
        
        st.markdown("### Resumo de Todos os Merges")
        
        resumo_merges = {
            "Merge": [
                "KE5Z ↔ Base Conso",
                "KE5Z ↔ Rateio",
                "KE5Z ↔ Volume",
                "Histórico ↔ Ano Atual"
            ],
            "Chave KE5Z": [
                "Account (Type 07)",
                "['Oficina', 'Período']",
                "['Oficina', 'Período', 'Veículo']",
                "N/A (concatenação)"
            ],
            "Chave Externa": [
                "Account (Type 07)",
                "['Oficina', 'Período']",
                "['Oficina', 'Período', 'Veículo']",
                "N/A (concatenação)"
            ],
            "Tipo": [
                "left",
                "left",
                "left",
                "concat"
            ],
            "Resultado": [
                "Coluna Custo (Variável/Fixo)",
                "Colunas de rateio por veículo (CC21%, CC22%, etc.)",
                "Coluna Volume",
                "Histórico consolidado com todos os anos"
            ]
        }
        
        st.dataframe(pd.DataFrame(resumo_merges), width="stretch", hide_index=True)
        
        st.markdown("### Detalhamento dos Merges")
        
        with st.expander("1. Merge KE5Z ↔ Base Conso (Custo)", expanded=False):
            st.code("""
# Leitura
df_base_conso = pd.read_excel('Dados SAPIENS.xlsx', sheet_name='Base conso')
df_base_conso = df_base_conso.rename(columns={'Type 04': 'Custo', 'Type 07': 'Account'})
df_base_conso = df_base_conso[['Custo', 'Account']]

# Merge
df_KE5Z = pd.merge(df_KE5Z, df_base_conso, on='Account', how='left')
            """, language="python")
            
            st.markdown("""
            **Resultado**: Adiciona coluna `Custo` ao `df_KE5Z`
            - Valores: `"Variável"` ou `"Fixo"`
            - Usado para cálculos de Flex Bud e análises de custos fixos vs variáveis
            """)
        
        with st.expander("2. Merge KE5Z ↔ Rateio (Percentuais por Veículo)", expanded=False):
            st.code("""
# Processamento do Rateio
df_rateio = pd.read_excel('Reporting fluxo anexo.xlsx', sheet_name='Rateio', header=None)
# ... processamento com melt() ...
df_pivot = df_rateio.pivot_table(
        index=['Oficina', 'Período'],
        columns='Veículo',
        values='Rateio',
        aggfunc='mean'
).reset_index()

# Renomear colunas de veículos para adicionar %
veiculos_cols = ['CC21', 'CC22', 'CC24', 'CC24 5L', 'CC24 7L', 'J516']
rename_dict = {col: f"{col}%" for col in veiculos_cols}
df_pivot = df_pivot.rename(columns=rename_dict)

# Merge
df_final = pd.merge(df_KE5Z, df_pivot, on=['Oficina', 'Período'], how='left')
            """, language="python")
            
            st.markdown("""
            **Resultado**: Adiciona colunas de percentuais por veículo
            - `CC21%`, `CC22%`, `CC24%`, `CC24 5L%`, `CC24 7L%`, `J516%`
            - Valores: Percentuais (0.0 a 1.0 ou 0% a 100%)
            - Usado para calcular valores por veículo: `CC21 = CC21% * Valor`
            """)
        
        with st.expander("3. Merge df_final ↔ Volume", expanded=False):
            st.code("""
# Processamento do Volume
df_vol = pd.read_excel('Reporting fluxo anexo.xlsx', sheet_name='Volume', header=50)
# ... processamento com melt() ...
# Colunas finais: Oficina, Veículo, Período, Volume

# Merge
df_final = pd.merge(df_final, df_vol, on=['Oficina', 'Período', 'Veículo'], how='left')
df_final['Volume'] = df_final['Volume'].fillna(0)
            """, language="python")
            
            st.markdown("""
            **Resultado**: Adiciona coluna `Volume` ao `df_final`
            - Valores: Volumes numéricos por veículo
            - Usado para cálculos de CPU (Custo por Unidade)
            """)
        
        st.markdown("---")
        
        # Seção 6: Colunas e Estrutura Final
        st.markdown("## 📊 COLUNAS E ESTRUTURA FINAL {#colunas-finais}")
        
        st.markdown("### Colunas do DataFrame Final (df_final.parquet)")
        
        colunas_finais = {
            "Coluna": [
                "Mes", "Período", "Ano",
                "Nºconta", "Centrocst", "Nºdoc.ref.", "Dt.lçto.",
                "Valor", "QTD", "Volume",
                "Type 05", "Type 06", "Account", "Custo",
                "USI", "Oficina",
                "Doc.compra", "Texto breve",
                "Fornecedor", "Material", "Usuário", "Fornec.", "Tipo",
                "CC21", "CC22", "CC24", "CC24 5L", "CC24 7L", "J516",
                "Soma_Percentuais"
            ],
            "Tipo": [
                "float64", "object", "int64",
                "object", "object", "float64", "object",
                "float64", "float64", "float64",
                "object", "object", "object", "object",
                "object", "object",
                "object", "object",
                "object", "object", "object", "object", "object",
                "float64", "float64", "float64", "float64", "float64", "float64",
                "float64"
            ],
            "Origem": [
                "Sapiens", "Sapiens", "Adicionado na consolidação",
                "Sapiens", "Sapiens", "Sapiens", "Sapiens",
                "Sapiens", "Sapiens", "Volume (merge)",
                "Sapiens", "Sapiens", "Sapiens", "Base conso (merge)",
                "Sapiens", "Sapiens",
                "Sapiens", "Sapiens",
                "Sapiens", "Sapiens", "Sapiens", "Sapiens", "Sapiens",
                "Calculado (CC21% * Valor)", "Calculado", "Calculado", "Calculado", "Calculado", "Calculado",
                "Calculado (soma dos %)"
            ],
            "Descrição": [
                "Mês numérico (1-12)", "Mês por extenso (Janeiro, etc.)", "Ano do registro",
                "Código da conta contábil", "Centro de custo", "Número documento referência", "Data de lançamento",
                "Valor monetário do custo", "Quantidade", "Volume do veículo",
                "Classificação Type 05", "Classificação Type 06", "Account (Type 07)", "Tipo de custo (Variável/Fixo)",
                "Unidade de negócio", "Nome da oficina",
                "Documento de compra", "Descrição breve do material",
                "Nome do fornecedor", "Código do material", "Usuário", "Código fornecedor", "Tipo de lançamento",
                "Valor rateado para CC21", "Valor rateado para CC22", "Valor rateado para CC24", "Valor rateado para CC24 5L", "Valor rateado para CC24 7L", "Valor rateado para J516",
                "Soma de todos os percentuais (validação)"
            ]
        }
        
        st.dataframe(pd.DataFrame(colunas_finais), width="stretch", hide_index=True)
        
        st.markdown("### Colunas do DataFrame de Volume (df_vol.parquet)")
        
        colunas_volume = {
            "Coluna": ["Oficina", "Veículo", "Período", "Volume"],
            "Tipo": ["object", "object", "object", "float64"],
            "Descrição": [
                "Nome da oficina",
                "Código do veículo (CC21, CC22, CC24, CC24 5L, CC24 7L, J516)",
                "Mês por extenso (Janeiro, Fevereiro, etc.)",
                "Volume numérico do veículo no período"
            ]
        }
        
        st.dataframe(pd.DataFrame(colunas_volume), width="stretch", hide_index=True)
        
        st.markdown("### Relacionamento entre Colunas")
        
        st.markdown("""
        **Chaves Primárias para Merges**:
        - `Account` (Type 07) → Merge com Base Conso
        - `['Oficina', 'Período']` → Merge com Rateio
        - `['Oficina', 'Período', 'Veículo']` → Merge com Volume
        
        **Colunas Calculadas**:
        - `CC21 = CC21% * Valor` (e similares para outros veículos)
        - `Soma_Percentuais = CC21% + CC22% + CC24% + CC24 5L% + CC24 7L% + J516%`
        - `CPU = Valor / Volume` (calculado no app.py, não no notebook)
        
        **Normalizações Críticas**:
        - `Período`: Sempre capitalizado (Janeiro, Fevereiro, etc.)
        - `Account`: Mantido como string/object
        - `Volume`: Sempre numérico (float64), NaN preenchido com 0
        """)
        
        st.markdown("---")
        
        # Seção 7: Consolidação do Histórico
        st.markdown("## 📚 CONSOLIDAÇÃO DO HISTÓRICO {#consolidacao}")
        
        st.markdown("### Processo de Consolidação")
        
        st.markdown("""
        **Objetivo**: Manter um histórico completo de todos os anos processados
        
        **Processo**:
        1. **Verificar histórico existente**:
              - Tenta carregar `dados/TC_Ext/historico_consolidado/df_final_historico.parquet`
           - Se não existir, cria DataFrame vazio
        
        2. **Adicionar coluna Ano**:
           - Adiciona `Ano = ANO_ATUAL` aos dados do ano atual
           - Garante que cada registro tenha identificação do ano
        
        3. **Concatenação**:
           - Concatena dados do ano atual com histórico existente
           - Usa `pd.concat([df_historico, df_ano_atual], ignore_index=True)`
        
        4. **Validação**:
           - Verifica se `Volume` é sempre numérico
           - Garante tipos de dados consistentes
        
        5. **Salvamento**:
           - Salva histórico atualizado
           - Mantém histórico sempre completo
        
        **IMPORTANTE**: 
        - O histórico é **sempre concatenado**, nunca substituído
        - Permite análises multi-anos no dashboard
        - O sistema prioriza o histórico consolidado para carregar dados
        """)
        
        st.markdown("### Estrutura do Histórico")
        
        st.code("""
        dados/TC_Ext/historico_consolidado/
        ├── df_final_historico.parquet      # Todos os anos de custos (REAL)
        ├── df_vol_historico.parquet        # Todos os anos de volumes
        ├── df_ke5z_historico.parquet       # Dados KE5Z agrupados
        └── BUD/
            ├── df_final_historico_BUD.parquet  # Todos os anos de custos (BUDGET)
            ├── df_vol_historico_BUD.parquet    # Todos os anos de volumes (BUDGET)
            └── df_ke5z_historico_BUD.parquet   # Dados KE5Z agrupados (BUDGET)
        """, language="text")
        
        st.markdown("---")
        
        # Seção 8: Arquivos de Saída
        st.markdown("## 💾 ARQUIVOS DE SAÍDA {#arquivos-saida}")
        
        st.markdown("### Arquivos Gerados por tc_ext/notebooks/dados.ipynb (REAL)")
        
        arquivos_saida_real = {
            "Arquivo": [
                "df_final.parquet",
                "df_vol.parquet",
                "df_ke5z_group.parquet",
                "df_final_historico.parquet",
                "df_vol_historico.parquet",
                "df_ke5z_historico.parquet"
            ],
            "Localização": [
                "dados/TC_Ext/{ANO}/",
                "dados/TC_Ext/{ANO}/",
                "dados/TC_Ext/{ANO}/",
                "dados/TC_Ext/historico_consolidado/",
                "dados/TC_Ext/historico_consolidado/",
                "dados/TC_Ext/historico_consolidado/"
            ],
            "Conteúdo": [
                "Dados completos com rateio por veículo e volume",
                "Dados de volume por Oficina, Veículo e Período",
                "Dados agrupados KE5Z",
                "Histórico consolidado de todos os anos (REAL)",
                "Histórico consolidado de volumes",
                "Histórico consolidado KE5Z"
            ],
            "Uso": [
                "Dashboard principal (app.py)",
                "Cálculos de CPU e análises de volume",
                "Análises específicas",
                "Análises multi-anos",
                "Análises multi-anos de volume",
                "Análises históricas KE5Z"
            ]
        }
        
        st.dataframe(pd.DataFrame(arquivos_saida_real), width="stretch", hide_index=True)
        
        st.markdown("### Arquivos Gerados por tc_ext/notebooks/dados_BUD.ipynb (BUDGET)")
        
        arquivos_saida_bud = {
            "Arquivo": [
                "df_final_BUD.parquet",
                "df_vol_BUD.parquet",
                "df_ke5z_group_BUD.parquet",
                "df_final_historico_BUD.parquet",
                "df_vol_historico_BUD.parquet",
                "df_ke5z_historico_BUD.parquet"
            ],
            "Localização": [
                "dados/TC_Ext/{ANO}/BUD/",
                "dados/TC_Ext/{ANO}/BUD/",
                "dados/TC_Ext/{ANO}/BUD/",
                "dados/TC_Ext/historico_consolidado/BUD/",
                "dados/TC_Ext/historico_consolidado/BUD/",
                "dados/TC_Ext/historico_consolidado/BUD/"
            ],
            "Conteúdo": [
                "Dados de Budget com rateio por veículo e volume",
                "Dados de volume de Budget",
                "Dados agrupados KE5Z (Budget)",
                "Histórico consolidado de todos os anos (BUDGET)",
                "Histórico consolidado de volumes (Budget)",
                "Histórico consolidado KE5Z (Budget)"
            ],
            "Uso": [
                "Comparações Real vs Budget",
                "Análises de volume Budget",
                "Análises específicas Budget",
                "Análises multi-anos Budget",
                "Análises multi-anos de volume Budget",
                "Análises históricas KE5Z Budget"
            ]
        }
        
        st.dataframe(pd.DataFrame(arquivos_saida_bud), width="stretch", hide_index=True)
        
        st.markdown("---")
        
        # Seção 9: Fluxo Completo
        st.markdown("## 🔄 FLUXO COMPLETO {#fluxo-completo}")
        
        st.markdown("### Diagrama de Fluxo - tc_ext/notebooks/dados.ipynb")
        
        st.code("""
        ┌─────────────────────────────────────┐
        │  Configuração (Célula 0)           │
        │  - Define ANO_ATUAL                 │
        │  - Cria pastas                      │
        │  - Verifica arquivos                │
        └──────────────┬──────────────────────┘
                       │
                       ▼
        ┌─────────────────────────────────────┐
        │  Leitura SAPIENS (Célula 1)        │
        │  - Reporting fluxo anexo.xlsx       │
        │  - Guia "Sapiens"                   │
        │  - Cria df_KE5Z                     │
        └──────────────┬──────────────────────┘
                       │
                       ▼
        ┌─────────────────────────────────────┐
        │  Merge Base Conso (Célula 2)       │
        │  - Dados SAPIENS.xlsx                │
        │  - Guia "Base conso"                 │
        │  - Adiciona coluna Custo             │
        └──────────────┬──────────────────────┘
                       │
                       ▼
        ┌─────────────────────────────────────┐
        │  Processamento Rateio (Célula 3)   │
        │  - Reporting fluxo anexo.xlsx         │
        │  - Guia "Rateio"                     │
        │  - Transforma meses em linhas         │
        │  - Cria df (Oficina, Veículo,        │
        │    Período, Rateio)                  │
        └──────────────┬──────────────────────┘
                       │
                       ▼
        ┌─────────────────────────────────────┐
        │  Merge + Cálculo Veículos (Célula 4)│
        │  - Merge KE5Z ↔ Rateio               │
        │  - Pivot: Veículo → Colunas          │
        │  - Calcula: CC21 = CC21% * Valor     │
        │  - Cria df_final                     │
        └──────────────┬──────────────────────┘
                       │
                       ▼
        ┌─────────────────────────────────────┐
        │  Processamento Volume (Célula 5)   │
        │  - Reporting fluxo anexo.xlsx       │
        │  - Guia "Volume"                     │
        │  - Transforma meses em linhas         │
        │  - Cria df_vol                       │
        └──────────────┬──────────────────────┘
                       │
                       ▼
        ┌─────────────────────────────────────┐
        │  Merge Volume (Célula 6)           │
        │  - Merge df_final ↔ df_vol           │
        │  - Adiciona coluna Volume             │
        └──────────────┬──────────────────────┘
                       │
                       ▼
        ┌─────────────────────────────────────┐
        │  Salvamento + Consolidação (Célula 7)│
        │  - Salva df_final.parquet            │
        │  - Salva df_vol.parquet               │
        │  - Carrega histórico                  │
        │  - Concatena com ano atual            │
        │  - Salva histórico atualizado         │
        └─────────────────────────────────────┘
        """, language="text")
        
        st.markdown("### Sequência de Operações Detalhada")
        
        operacoes_detalhadas = [
            "**Célula 0**: Configuração - Define ano, cria pastas, verifica arquivos de entrada",
            "**Célula 1**: Leitura SAPIENS - Lê guia 'Sapiens' (20 colunas), cria df_KE5Z com dados de custos",
            "**Célula 2**: Merge Base Conso - Adiciona coluna 'Custo' (Variável/Fixo) usando Account como chave",
            "**Célula 3**: Processamento Rateio - Lê guia 'Rateio', transforma meses em linhas (melt), cria df com Oficina, Veículo, Período, Rateio",
            "**Célula 4**: Merge Rateio + Cálculo - Merge KE5Z ↔ Rateio, pivot de Veículo para colunas, calcula valores por veículo (CC21, CC22, etc.), cria df_final",
            "**Célula 5**: Processamento Volume - Lê guia 'Volume' (header=50), transforma meses em linhas, cria df_vol com Oficina, Veículo, Período, Volume",
            "**Célula 6**: Merge Volume - Merge df_final ↔ df_vol usando ['Oficina', 'Período', 'Veículo'], adiciona coluna Volume",
            "**Célula 7**: Salvamento - Salva df_final.parquet, df_vol.parquet na pasta do ano, carrega histórico, concatena, salva histórico consolidado"
        ]
        
        for op in operacoes_detalhadas:
            st.markdown(f"- {op}")
        
        st.markdown("---")
        
        # Seção 10: Tratamento de Erros
        st.markdown("## ⚠️ TRATAMENTO DE ERROS {#tratamento-erros}")
        
        st.markdown("### Erros Comuns e Soluções")
        
        with st.expander("1. Arquivo Não Encontrado", expanded=False):
            st.markdown("""
            **Sintoma**: `FileNotFoundError` ao tentar ler arquivo Excel
            
            **Soluções**:
            - Verificar se arquivo está em `dados/TC_Ext/{ANO}/` ou na raiz do projeto
            - Verificar nomes exatos: `Dados SAPIENS.xlsx` e `Reporting fluxo anexo.xlsx`
            - O notebook tenta copiar da raiz para pasta do ano automaticamente
            """)
        
        with st.expander("2. Guia Não Encontrada", expanded=False):
            st.markdown("""
            **Sintoma**: `ValueError: Worksheet named 'X' not found`
            
            **Soluções**:
            - Verificar nomes exatos das guias (case-sensitive):
              - `tc_ext/notebooks/dados.ipynb`: "Sapiens", "Rateio", "Volume"
              - `tc_ext/notebooks/dados_BUD.ipynb`: "Voz de custo BDG", "Rateio BDG", "Volume BDG"
            - Verificar se guias existem no arquivo Excel
            """)
        
        with st.expander("3. Coluna Não Encontrada Após Merge", expanded=False):
            st.markdown("""
            **Sintoma**: `KeyError: 'Coluna X'` após merge
            
            **Soluções**:
            - Verificar se chaves de merge existem em ambos DataFrames
            - Verificar normalização de `Período` (deve estar capitalizado)
            - Verificar tipos de dados das chaves (devem ser compatíveis)
            - Verificar se merge foi feito com chaves corretas
            """)
        
        with st.expander("4. Volume NaN ou Zerado", expanded=False):
            st.markdown("""
            **Sintoma**: Coluna Volume com muitos NaN ou zeros
            
            **Soluções**:
            - Verificar se merge foi feito com chave composta correta: `['Oficina', 'Período', 'Veículo']`
            - Verificar se dados de volume existem para a combinação Oficina+Período+Veículo
            - Verificar normalização de `Período` (deve estar capitalizado em ambos DataFrames)
            - O notebook preenche NaN com 0 automaticamente
            """)
        
        with st.expander("5. Percentuais de Rateio Não Somam 100%", expanded=False):
            st.markdown("""
            **Sintoma**: `Soma_Percentuais` diferente de 1.0 (ou 100%)
            
            **Soluções**:
            - Verificar se todos os veículos estão incluídos no rateio
            - Verificar se há veículos não mapeados
            - Verificar se pivot foi feito corretamente (aggfunc='mean')
            - Validação: `Soma_Percentuais` deve estar próximo de 1.0
            """)
        
        with st.expander("6. Histórico Não Atualizado", expanded=False):
            st.markdown("""
            **Sintoma**: Histórico não inclui dados do ano atual após processamento
            
            **Soluções**:
            - Verificar se coluna `Ano` foi adicionada aos dados do ano atual
            - Verificar se concatenação foi executada corretamente
            - Verificar se arquivo de histórico foi salvo após concatenação
            - Verificar permissões de escrita na pasta `dados/TC_Ext/historico_consolidado/`
            """)
        
        st.markdown("### Validações Implementadas")
        
        st.markdown("""
        **Validações Automáticas**:
        1. **Validação de Arquivos**: Verifica existência antes de processar
        2. **Validação de Colunas**: Verifica se colunas essenciais existem antes de merge
        3. **Validação de Volume**: Garante que Volume seja sempre numérico
        4. **Validação de Período**: Normaliza para formato capitalizado
        5. **Validação de Histórico**: Verifica tipos de dados ao carregar histórico
        6. **Validação de Soma**: Calcula `Soma_Percentuais` para validar rateios
        """)
        
        st.markdown("---")
        
        # Seção 11: Checklist para Manutenção
        st.markdown("## ✅ CHECKLIST PARA MANUTENÇÃO {#checklist}")
        
        st.markdown("### Antes de Modificar os Notebooks")
        
        checklist_antes = [
            "Verificar se estrutura de pastas está correta",
            "Verificar se nomes de guias estão corretos",
            "Verificar se chaves de merge estão corretas",
            "Verificar se tipos de dados estão consistentes",
            "Verificar se normalização de Período está funcionando",
            "Verificar se cálculo de veículos está correto",
            "Verificar se consolidação de histórico está funcionando"
        ]
        
        for item in checklist_antes:
            st.markdown(f"- [ ] {item}")
        
        st.markdown("### Ao Modificar")
        
        checklist_modificar = [
            "Manter mesma estrutura de chaves de merge",
            "Manter normalização de Período (capitalizado)",
            "Manter tipos de dados consistentes (Volume sempre numérico)",
            "Manter lógica de cálculo de veículos (CC21 = CC21% * Valor)",
            "Manter processo de consolidação de histórico (concatenação, não substituição)",
            "Testar com dados de um ano antes de processar todos",
            "Validar que Volume não está sendo zerado incorretamente",
            "Validar que Soma_Percentuais está próximo de 1.0"
        ]
        
        for item in checklist_modificar:
            st.markdown(f"- [ ] {item}")
        
        st.markdown("### Após Modificar")
        
        checklist_depois = [
            "Verificar se arquivos Parquet foram gerados corretamente",
            "Verificar se histórico foi atualizado",
            "Verificar se Volume está presente e numérico",
            "Verificar se colunas de veículos foram calculadas",
            "Verificar se não há erros de tipo de dados",
            "Testar carregamento no app.py",
            "Validar que dados aparecem corretamente no dashboard"
        ]
        
        for item in checklist_depois:
            st.markdown(f"- [ ] {item}")
        
        st.markdown("### Regras Críticas que NUNCA Devem Ser Alteradas")
        
        st.warning("""
        **⚠️ ATENÇÃO**: As seguintes regras são CRÍTICAS e não devem ser alteradas sem
        análise profunda, pois podem quebrar todo o sistema:
        
        1. **Chaves de Merge**: `['Oficina', 'Período']` para Rateio e `['Oficina', 'Período', 'Veículo']` para Volume
        2. **Normalização de Período**: Sempre capitalizado (Janeiro, Fevereiro, etc.)
        3. **Cálculo de Veículos**: `CC21 = CC21% * Valor` (e similares)
        4. **Consolidação de Histórico**: Sempre concatenar, nunca substituir
        5. **Tipo de Volume**: Sempre numérico (float64), nunca object
        6. **Estrutura de Pastas (TC Ext)**: `dados/TC_Ext/{ANO}/` para ano específico, `dados/TC_Ext/historico_consolidado/` para histórico
        7. **Sufixo BUD**: Arquivos de Budget sempre com sufixo `_BUD` e em pasta `BUD/`
        """)
        
        st.markdown("---")
        
        # Seção Final: Notas Importantes
        st.markdown("## 📝 NOTAS IMPORTANTES PARA IA")
        
        st.markdown("### Quando Fazer Manutenção")
        
        st.markdown("""
        **Faça manutenção quando**:
        - Estrutura dos arquivos Excel de entrada mudar
        - Novas colunas forem adicionadas aos dados
        - Novos veículos forem adicionados ao sistema
        - Lógica de rateio mudar
        - Estrutura de pastas precisar ser alterada
        
        **NÃO faça manutenção quando**:
        - Apenas dados novos forem adicionados (processe normalmente)
        - Apenas valores mudarem (processe normalmente)
        - Apenas anos novos forem processados (processe normalmente)
        """)
        
        st.markdown("### Como Fazer Manutenção Segura")
        
        st.markdown("""
        1. **Sempre teste primeiro**: Processe um ano de teste antes de processar todos
        2. **Mantenha backups**: Faça backup dos arquivos Parquet antes de modificar
        3. **Valide resultados**: Verifique se Volume, valores por veículo e histórico estão corretos
        4. **Documente mudanças**: Adicione comentários explicando alterações
        5. **Mantenha consistência**: Se alterar `tc_ext/notebooks/dados.ipynb`, altere `tc_ext/notebooks/dados_BUD.ipynb` da mesma forma
        6. **Valide merges**: Sempre verifique se chaves de merge existem antes de fazer merge
        7. **Valide tipos**: Sempre verifique tipos de dados após transformações
        """)
        
        st.markdown("### Estrutura de Dependências")
        
        st.code("""
        tc_ext/notebooks/dados.ipynb depende de:
        ├── Reporting fluxo anexo.xlsx
        │   ├── Guia "Sapiens" (dados principais)
        │   ├── Guia "Rateio" (percentuais por veículo)
        │   └── Guia "Volume" (volumes por veículo)
        └── Dados SAPIENS.xlsx
            └── Guia "Base conso" (mapeamento Custo)
        
        tc_ext/notebooks/dados_BUD.ipynb depende de:
        ├── Reporting fluxo anexo.xlsx
        │   ├── Guia "Voz de custo BDG" (dados principais)
        │   ├── Guia "Rateio BDG" (percentuais por veículo)
        │   └── Guia "Volume BDG" (volumes por veículo)
        └── Dados SAPIENS.xlsx
            └── Guia "Base conso" (mapeamento Custo)
        """, language="text")
        
        st.markdown("---")
        
        st.success("""
        **✅ Este guia contém todas as informações necessárias para fazer manutenção**
        nos notebooks de extração sem quebrar o sistema. Sempre consulte este guia
        antes de fazer alterações e siga o checklist de validação.
        """)
    
    # ==========================================
    # CAPÍTULO 2: FUNCIONAMENTO DA ATUALIZAÇÃO E EXTRAÇÃO
    # ==========================================
    
    with st.expander("🔄 **Capítulo 2: Funcionamento da Atualização e Extração**", expanded=False):
        st.markdown("""
        <div style="padding: 1.5rem; background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%); border-radius: 10px; margin: 2rem 0; color: white;">
            <h2 style="color: white; margin: 0;">🔄 Capítulo 2: Funcionamento da Atualização e Extração</h2>
            <p style="color: #f0f0f0; margin: 0.5rem 0 0 0;">
                Processo Completo de Atualização de Dados - Passo a Passo Detalhado
            </p>
        </div>
        """, unsafe_allow_html=True)
        
        # Seção 1: Visão Geral do Processo de Atualização
        st.markdown("## 🎯 VISÃO GERAL DO PROCESSO DE ATUALIZAÇÃO {#visao-atualizacao}")
        
        st.markdown("""
        Este capítulo descreve **como funciona o processo completo de atualização de dados**,
        desde a preparação dos arquivos até a execução do processamento. Entender este fluxo
        é essencial para realizar atualizações corretamente, especialmente quando se trabalha
        com novos anos ou quando se precisa atualizar arquivos existentes.
        """)
        
        st.info("""
        **💡 Importante**: O sistema foi projetado para ser flexível e permitir atualizações
        de diferentes formas: através de upload direto na interface, colocando arquivos na
        raiz do projeto, ou organizando-os nas pastas do ano. O sistema busca automaticamente
        os arquivos na ordem de prioridade definida.
        """)
        
        st.markdown("---")
        
        # Seção 2: Ordem Cronológica dos Eventos
        st.markdown("## ⏱️ ORDEM CRONOLÓGICA DOS EVENTOS {#ordem-cronologica}")
        
        st.markdown("### Sequência Completa do Processo")
        
        with st.expander("**1️⃣ Seleção do Ano e Tipo de Extração**", expanded=False):
            st.markdown("""
            **Onde**: Página "5 - Extração de Dados" (Streamlit)
            
            **Processo**:
            1. Usuário seleciona o **ano** que deseja processar (ex: 2024, 2025, 2026)
            2. Usuário seleciona o **tipo de extração**:
               - 📊 **Dados REAIS** (tc_ext/notebooks/dados.ipynb) - Processa custos reais executados
               - 💰 **Dados BUDGET** (tc_ext/notebooks/dados_BUD.ipynb) - Processa dados de orçamento
               - 🔄 **Ambos** - Processa REAIS e BUDGET sequencialmente
            
            **Resultado**: Sistema sabe qual ano processar e quais notebooks executar
            """)
        
        with st.expander("**2️⃣ Verificação e Preparação de Arquivos**", expanded=False):
            st.markdown("""
            **Onde**: Antes do processamento, na aba "Validação de Arquivos"
            
            **Processo**:
            1. Sistema verifica se os arquivos necessários já existem
            2. Sistema mostra avisos se arquivos já existem (para evitar sobrescrita acidental)
            3. Usuário pode fazer upload de arquivos diretamente na interface
            
            **Arquivos Necessários para Dados REAIS**:
            - `Dados SAPIENS.xlsx` - Base de dados SAPIENS com classificação de custos
            - `Reporting fluxo anexo.xlsx` - Dados de custos, rateio e volumes
            
            **Arquivos Necessários para Dados BUDGET**:
            - `Dados SAPIENS.xlsx` - Base de dados SAPIENS (mesmo arquivo ou versão Budget)
            - `Reporting fluxo anexo.xlsx` - Dados de Budget (guias "Voz de custo BDG", "Rateio BDG", "Volume BDG")
            """)
        
        with st.expander("**3️⃣ Sistema de Upload de Arquivos (Opcional)**", expanded=False):
            st.markdown("""
            **Onde**: Aba "Validação de Arquivos" → Seção "📤 Upload de Arquivos"
            
            **Processo**:
            1. Usuário clica em "Browse Files" para selecionar arquivo
            2. **ANTES do upload**: Sistema verifica se arquivo já existe na pasta de destino
               - Se existe: Mostra aviso ⚠️ informando que será sobrescrito
               - Se não existe: Permite upload direto
            3. Usuário seleciona arquivo do computador
            4. **APÓS seleção**: Sistema verifica novamente se arquivo existe
               - Se existe: Mostra aviso e botão "🔄 Confirmar Sobrescrita"
               - Se não existe: Salva automaticamente
            5. Arquivo é salvo em: `dados/TC_Ext/{ANO_SELECIONADO}/Nome_do_Arquivo.xlsx`
            6. Página recarrega automaticamente (`st.rerun()`) para atualizar status
            
            **Vantagens do Upload**:
            - Não precisa colocar arquivos na raiz do projeto
            - Arquivos são organizados automaticamente na pasta do ano
            - Sistema cria a pasta do ano automaticamente se não existir
            - Avisos preventivos evitam sobrescrita acidental
            """)
        
        with st.expander("**4️⃣ Criação da Estrutura de Pastas**", expanded=False):
            st.markdown("""
            **Onde**: Função `configurar_ano()` ou `configurar_ano_bud()` nos módulos Python
            
            **Processo** (executado automaticamente ao iniciar processamento):
                1. **Cria pasta do ano**: `dados/TC_Ext/{ANO}/`
                    - Exemplo: `dados/TC_Ext/2024/` para ano 2024
                    - Exemplo: `dados/TC_Ext/2026/` para ano 2026 (novo ano)
            
                2. **Para dados REAIS**: Cria apenas `dados/TC_Ext/{ANO}/`
            
                3. **Para dados BUDGET**: Cria também `dados/TC_Ext/{ANO}/BUD/`
                          - Estrutura: `dados/TC_Ext/2024/BUD/` para **outputs** de Budget
            
            4. **Cria pastas de histórico** (se não existirem):
                    - `dados/TC_Ext/historico_consolidado/` - Para dados REAIS
                    - `dados/TC_Ext/historico_consolidado/BUD/` - Para dados BUDGET
            
            **IMPORTANTE**: 
            - Pastas são criadas automaticamente, mesmo que não existam
            - Se a pasta já existe, não há problema (não sobrescreve)
            - Sistema usa `os.makedirs(pasta, exist_ok=True)` para criar com segurança
            """)
        
        with st.expander("**5️⃣ Sistema de Busca de Arquivos**", expanded=False):
            st.markdown("""
            **Onde**: Função `encontrar_arquivo()` nos módulos de processamento
            
            **Ordem de Prioridade de Busca** (do mais prioritário ao menos prioritário):
            
            **Para Dados REAIS**:
                1. **Primeira opção**: `dados/TC_Ext/{ANO}/Nome_do_Arquivo.xlsx`
                    - Exemplo: `dados/TC_Ext/2024/Dados SAPIENS.xlsx`
               - **Esta é a pasta preferencial!** Arquivos aqui têm prioridade máxima
            
            2. **Segunda opção**: `./Nome_do_Arquivo.xlsx` (raiz do projeto)
               - Exemplo: `./Dados SAPIENS.xlsx`
               - Usado quando arquivo não está na pasta do ano
            
            **Para Dados BUDGET**:
                1. **Primeira opção**: `dados/TC_Ext/{ANO}/Nome_do_Arquivo.xlsx`
                    - Exemplo: `dados/TC_Ext/2024/Dados SAPIENS.xlsx`

                2. **Segunda opção**: `./Nome_do_Arquivo.xlsx` (raiz do projeto)

                *(Compatibilidade/legado)*: se existir arquivo em `dados/TC_Ext/{ANO}/BUD/`, ele pode ser **copiado** para `dados/TC_Ext/{ANO}/`.
            
            **Comportamento**:
            - Sistema busca na ordem acima e usa o **primeiro arquivo encontrado**
            - Se arquivo não for encontrado em nenhum local, sistema retorna erro
            - Se arquivo for encontrado na raiz, pode ser copiado para pasta do ano (dependendo da configuração)
            
            **Exemplo Prático - Processando 2026 pela primeira vez**:
            ```
            1. Sistema cria: dados/TC_Ext/2026/
            2. Sistema busca: dados/TC_Ext/2026/Dados SAPIENS.xlsx → ❌ Não encontrado
            3. Sistema busca: ./Dados SAPIENS.xlsx → ✅ Encontrado na raiz
            4. Sistema usa: ./Dados SAPIENS.xlsx (da raiz)
            5. Arquivos de saída são salvos em: dados/TC_Ext/2026/
            ```
            """)
        
        with st.expander("**6️⃣ Execução do Processamento**", expanded=False):
            st.markdown("""
            **Onde**: Aba "Executar Processamento" → Botões de execução
            
            **Processo**:
            1. Usuário clica em botão de execução:
               - "🚀 Executar tc_ext/notebooks/dados.ipynb" (para REAIS)
               - "🚀 Executar tc_ext/notebooks/dados_BUD.ipynb" (para BUDGET)
               - "🚀 Executar Ambos" (para REAIS e BUDGET)
            
            2. Sistema chama função de processamento correspondente:
               - `processar_completo()` para dados REAIS
               - `processar_completo_bud()` para dados BUDGET
            
            3. **Configuração inicial**:
               - Chama `configurar_ano()` ou `configurar_ano_bud()`
               - Cria estrutura de pastas
               - Busca arquivos usando `encontrar_arquivo()`
               - Valida se arquivos existem
            
            4. **Processamento dos dados**:
               - Lê arquivos Excel das guias corretas
               - Faz merges e transformações
               - Calcula valores por veículo
               - Processa volumes
            
            5. **Salvamento**:
               - Salva arquivos Parquet na pasta do ano (ou BUD/)
               - Salva arquivos Excel intermediários (diagnósticos)
               - Consolida histórico (concatena, não substitui)
            
            6. **Feedback ao usuário**:
               - Barra de progresso mostra status
               - Mensagens de log aparecem em tempo real
               - Mensagem de sucesso ao finalizar
            """)
        
        with st.expander("**7️⃣ Consolidação do Histórico**", expanded=False):
            st.markdown("""
            **Onde**: Função `salvar_e_consolidar()` ou `salvar_e_consolidar_bud()`
            
            **Processo**:
            1. **Carrega histórico existente** (se existir):
                    - Tenta carregar: `dados/TC_Ext/historico_consolidado/df_final_historico.parquet`
               - Se não existir, cria DataFrame vazio
            
            2. **Adiciona coluna Ano aos dados atuais**:
               - Adiciona coluna `Ano` com valor do ano processado
               - Exemplo: Se processando 2026, todos os registros recebem `Ano = 2026`
            
            3. **Concatena dados**:
               - Concatena dados do ano atual com histórico existente
               - Usa `pd.concat([historico, dados_atuais])`
            
            4. **Remove duplicatas** (se houver):
               - Verifica e remove registros duplicados
            
            5. **Valida tipos de dados**:
               - Garante que Volume é numérico (float64)
               - Converte tipos se necessário
            
            6. **Salva histórico atualizado**:
                    - Salva em: `dados/TC_Ext/historico_consolidado/df_final_historico.parquet`
               - **IMPORTANTE**: Histórico é sempre **concatenado**, nunca substituído
            
            **Resultado**: Histórico contém dados de todos os anos processados
            """)
        
        st.markdown("---")
        
        # Seção 3: Sistema de Busca de Arquivos (Detalhado)
        st.markdown("## 🔍 SISTEMA DE BUSCA DE ARQUIVOS {#busca-arquivos}")
        
        st.markdown("### Lógica de Busca Detalhada")
        
        st.markdown("""
        O sistema implementa uma **busca hierárquica** que prioriza arquivos organizados
        nas pastas do ano, mas permite flexibilidade ao buscar na raiz do projeto quando
        necessário. Isso facilita o trabalho com novos anos sem precisar mover arquivos manualmente.
        """)
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("""
            **📊 Dados REAIS - Ordem de Busca:**
            
            1. `dados/TC_Ext/{ANO}/Dados SAPIENS.xlsx`
            2. `./Dados SAPIENS.xlsx` (raiz)
            
            1. `dados/TC_Ext/{ANO}/Reporting fluxo anexo.xlsx`
            2. `./Reporting fluxo anexo.xlsx` (raiz)
            """)
        
        with col2:
            st.markdown("""
            **💰 Dados BUDGET - Ordem de Busca:**
            
            1. `dados/TC_Ext/{ANO}/Dados SAPIENS.xlsx`
            2. `./Dados SAPIENS.xlsx` (raiz)
            
            1. `dados/TC_Ext/{ANO}/Reporting fluxo anexo.xlsx`
            2. `./Reporting fluxo anexo.xlsx` (raiz)
            """)
        
        st.markdown("### Exemplos Práticos de Busca")
        
        with st.expander("**Exemplo 1: Processando 2024 (ano existente)**", expanded=False):
            st.markdown("""
            **Cenário**: Pasta `dados/TC_Ext/2024/` já existe com arquivos
            
            **Busca de Dados SAPIENS.xlsx**:
            1. ✅ Encontrado em: `dados/TC_Ext/2024/Dados SAPIENS.xlsx`
            2. Sistema usa este arquivo (para na primeira opção)
            
            **Resultado**: Arquivo da pasta do ano é usado (prioridade máxima)
            """)
        
        with st.expander("**Exemplo 2: Processando 2026 (ano novo)**", expanded=False):
            st.markdown("""
            **Cenário**: Pasta `dados/TC_Ext/2026/` não existe ainda, arquivo está na raiz
            
            **Busca de Dados SAPIENS.xlsx**:
            1. ❌ Não encontrado em: `dados/TC_Ext/2026/Dados SAPIENS.xlsx` (pasta não existe)
            2. ✅ Encontrado em: `./Dados SAPIENS.xlsx` (raiz do projeto)
            3. Sistema usa arquivo da raiz
            
            **Resultado**: 
            - Sistema cria `dados/TC_Ext/2026/` automaticamente
            - Arquivo da raiz é usado para processamento
            - Arquivos de saída são salvos em `dados/TC_Ext/2026/`
            - **Arquivo da raiz permanece na raiz** (não é movido automaticamente)
            """)
        
        with st.expander("**Exemplo 3: Upload de Arquivo para 2026**", expanded=False):
            st.markdown("""
            **Cenário**: Usuário faz upload de arquivo para ano 2026
            
            **Processo**:
            1. Sistema cria `dados/TC_Ext/2026/` (se não existir)
            2. Usuário faz upload de `Dados SAPIENS.xlsx`
            3. Arquivo é salvo em: `dados/TC_Ext/2026/Dados SAPIENS.xlsx`
            
            **Próxima busca**:
            1. ✅ Encontrado em: `dados/TC_Ext/2026/Dados SAPIENS.xlsx`
            2. Sistema usa este arquivo (prioridade máxima)
            
            **Resultado**: Arquivo uploadado tem prioridade sobre arquivo da raiz
            """)
        
        st.markdown("---")
        
        # Seção 4: Criação de Pastas e Estrutura
        st.markdown("## 📁 CRIAÇÃO DE PASTAS E ESTRUTURA {#criacao-pastas}")
        
        st.markdown("### Estrutura Completa de Pastas")
        
        st.code("""
        dados/TC_Ext/
        ├── 2024/                    # Ano 2024 (dados REAIS)
        │   ├── Dados SAPIENS.xlsx
        │   ├── Reporting fluxo anexo.xlsx
        │   ├── df_final.parquet
        │   ├── df_vol.parquet
        │   ├── df_ke5z_group.parquet
        │   └── BUD/                 # Dados BUDGET do ano 2024
        │       ├── Dados SAPIENS.xlsx (opcional)
        │       ├── Reporting fluxo anexo.xlsx (opcional)
        │       ├── df_final_BUD.parquet
        │       ├── df_vol_BUD.parquet
        │       └── df_ke5z_group_BUD.parquet
        │
        ├── 2025/                    # Ano 2025
        │   └── ...
        │
        ├── 2026/                    # Ano 2026 (novo ano)
        │   └── ...                  # Criado automaticamente
        │
        └── historico_consolidado/   # Histórico de todos os anos
            ├── df_final_historico.parquet
            ├── df_vol_historico.parquet
            └── BUD/
                ├── df_final_historico_BUD.parquet
                └── df_vol_historico_BUD.parquet
        """, language="text")
        
        st.markdown("### Quando as Pastas São Criadas")
        
        with st.expander("**Criação Automática**", expanded=False):
            st.markdown("""
            **Momento**: Ao iniciar o processamento (função `configurar_ano()`)
            
            **Pastas criadas automaticamente**:
            - `dados/TC_Ext/{ANO}/` - Sempre criada, mesmo que vazia
            - `dados/TC_Ext/{ANO}/BUD/` - Criada apenas para **outputs** do processamento BUDGET
            - `dados/TC_Ext/historico_consolidado/` - Criada se não existir
            - `dados/TC_Ext/historico_consolidado/BUD/` - Criada se não existir (para BUDGET)
            
            **Comando usado**: `os.makedirs(pasta, exist_ok=True)`
            - `exist_ok=True` significa que não dá erro se pasta já existe
            - Cria todas as pastas intermediárias automaticamente
            """)
        
        with st.expander("**Criação via Upload**", expanded=False):
            st.markdown("""
            **Momento**: Quando usuário faz upload de arquivo
            
            **Processo**:
            1. Usuário seleciona arquivo para upload
            2. Sistema verifica se pasta `dados/TC_Ext/{ANO}/` existe
            3. Se não existe: Cria automaticamente com `os.makedirs(pasta_ano, exist_ok=True)`
            4. Salva arquivo em: `dados/TC_Ext/{ANO}/Nome_do_Arquivo.xlsx`
            
            **Resultado**: Pasta do ano é criada antes mesmo do processamento
            """)
        
        st.markdown("---")
        
        # Seção 5: Sistema de Upload de Arquivos
        st.markdown("## 📤 SISTEMA DE UPLOAD DE ARQUIVOS {#sistema-upload}")
        
        st.markdown("### Funcionalidades do Upload")
        
        st.markdown("""
        O sistema de upload permite que arquivos sejam enviados diretamente pela interface
        web, sem necessidade de colocá-los manualmente na raiz do projeto ou nas pastas.
        Isso facilita especialmente o trabalho com novos anos ou atualizações de arquivos.
        """)
        
        with st.expander("**Interface de Upload**", expanded=False):
            st.markdown("""
            **Localização**: Página "5 - Extração de Dados" → Aba "Validação de Arquivos" → Seção "📤 Upload de Arquivos"
            
            **Componentes**:
            - Uploaders separados por tipo de processamento (REAIS ou BUDGET)
            - Uploaders separados por arquivo (Dados SAPIENS.xlsx e Reporting fluxo anexo.xlsx)
            - Avisos proativos mostrando se arquivo já existe
            - Mensagens de confirmação após upload bem-sucedido
            
            **Layout**: Dois uploaders lado a lado (colunas) para cada tipo de processamento
            """)
        
        with st.expander("**Fluxo Completo de Upload**", expanded=False):
            st.markdown("""
            **Passo 1: Verificação Proativa**
            - Ao carregar a página, sistema verifica se arquivos já existem
            - Se existem: Mostra aviso ⚠️ acima do botão "Browse Files"
            - Aviso informa: "O arquivo já existe e será sobrescrito se você fizer upload"
            
            **Passo 2: Seleção do Arquivo**
            - Usuário clica em "Browse Files"
            - Seleciona arquivo do computador
            - Sistema detecta que arquivo foi selecionado
            
            **Passo 3: Verificação Pós-Seleção**
            - Sistema verifica novamente se arquivo existe na pasta de destino
            - Se existe: Mostra aviso adicional e botão "🔄 Confirmar Sobrescrita"
            - Se não existe: Prossegue para salvamento automático
            
            **Passo 4: Confirmação (se necessário)**
            - Se arquivo existe, usuário deve clicar em "🔄 Confirmar Sobrescrita"
            - Botão só aparece se arquivo realmente existe
            - Confirmação evita sobrescrita acidental
            
            **Passo 5: Salvamento**
            - Arquivo é salvo em: `dados/TC_Ext/{ANO_SELECIONADO}/Nome_do_Arquivo.xlsx`
            - Pasta do ano é criada automaticamente se não existir
            - Mensagem de sucesso é exibida
            
            **Passo 6: Atualização Automática**
            - Página recarrega automaticamente (`st.rerun()`)
            - Status dos arquivos é atualizado
            - Avisos são atualizados (se arquivo agora existe)
            """)
        
        with st.expander("**Vantagens do Sistema de Upload**", expanded=False):
            st.markdown("""
            ✅ **Organização Automática**: Arquivos são salvos na pasta correta automaticamente
            
            ✅ **Flexibilidade**: Não precisa colocar arquivos na raiz do projeto
            
            ✅ **Segurança**: Avisos preventivos evitam sobrescrita acidental
            
            ✅ **Facilidade**: Especialmente útil para novos anos (ex: 2026)
            
            ✅ **Rastreabilidade**: Mensagens claras mostram onde arquivo foi salvo
            
            ✅ **Validação**: Sistema verifica existência antes e depois do upload
            """)
        
        st.markdown("---")
        
        # Seção 6: Processamento e Execução
        st.markdown("## ⚙️ PROCESSAMENTO E EXECUÇÃO {#processamento-execucao}")
        
        st.markdown("### Fluxo de Execução Completo")
        
        st.markdown("""
        O processamento segue uma sequência bem definida, garantindo que todos os passos
        sejam executados na ordem correta e que os dados sejam processados e salvos adequadamente.
        """)
        
        with st.expander("**Fase 1: Preparação**", expanded=False):
            st.markdown("""
            1. **Configuração do Ano**:
               - Chama `configurar_ano()` ou `configurar_ano_bud()`
               - Cria estrutura de pastas
               - Define caminhos de entrada e saída
            
            2. **Busca de Arquivos**:
               - Busca `Dados SAPIENS.xlsx` na ordem de prioridade
               - Busca `Reporting fluxo anexo.xlsx` na ordem de prioridade
               - Valida se arquivos foram encontrados
            
            3. **Validação**:
               - Verifica se todos os arquivos necessários existem
               - Se faltar arquivo: Retorna erro ou aviso (dependendo da configuração)
            """)
        
        with st.expander("**Fase 2: Leitura e Transformação**", expanded=False):
            st.markdown("""
            1. **Leitura dos Dados Principais**:
               - Lê guia "Sapiens" ou "Voz de custo BDG" do Reporting fluxo anexo.xlsx
               - Cria DataFrame inicial (`df_KE5Z`)
            
            2. **Merge com Base Conso**:
               - Lê guia "Base conso" do Dados SAPIENS.xlsx
               - Faz merge adicionando coluna `Custo` (Variável/Fixo)
            
            3. **Processamento de Rateio**:
               - Lê guia "Rateio" ou "Rateio BDG"
               - Transforma colunas de meses em linhas
               - Cria DataFrame com percentuais de rateio por veículo
            
            4. **Merge e Cálculo por Veículo**:
               - Merge com dados principais
               - Calcula valores por veículo (CC21, CC22, etc.)
            
            5. **Processamento de Volume**:
               - Lê guia "Volume" ou "Volume BDG"
               - Transforma colunas de meses em linhas
               - Cria DataFrame com volumes
            
            6. **Merge Final com Volume**:
               - Adiciona coluna Volume ao DataFrame principal
            """)
        
        with st.expander("**Fase 3: Salvamento e Consolidação**", expanded=False):
            st.markdown("""
            1. **Salvamento na Pasta do Ano**:
               - Salva `df_final.parquet` em `dados/TC_Ext/{ANO}/` (ou `BUD/`)
               - Salva `df_vol.parquet`
               - Salva `df_ke5z_group.parquet`
               - Salva arquivos Excel intermediários (diagnósticos)
            
            2. **Consolidação do Histórico**:
               - Carrega histórico existente (se houver)
               - Adiciona coluna `Ano` aos dados atuais
               - Concatena dados atuais com histórico
               - Remove duplicatas
               - Salva histórico atualizado
            
            3. **Validação Final**:
               - Verifica tipos de dados
               - Valida integridade dos arquivos salvos
            """)
        
        st.markdown("---")
        
        # Seção 7: Cenários de Uso
        st.markdown("## 📋 CENÁRIOS DE USO {#cenarios-uso}")
        
        st.markdown("### Casos Práticos Completos")
        
        with st.expander("**Cenário 1: Primeira Vez Processando um Novo Ano (ex: 2026)**", expanded=False):
            st.markdown("""
            **Situação**: Nunca processou dados de 2026, arquivos estão na raiz do projeto
            
            **Passo a Passo**:
            
            1. **Acessar página de extração**:
               - Selecionar ano: 2026
               - Selecionar tipo: "📊 Dados REAIS" ou "🔄 Ambos"
            
            2. **Opção A - Usar Upload** (Recomendado):
               - Ir para aba "Validação de Arquivos"
                    - Fazer upload de `Dados SAPIENS.xlsx` → Salvo em `dados/TC_Ext/2026/`
                    - Fazer upload de `Reporting fluxo anexo.xlsx` → Salvo em `dados/TC_Ext/2026/`
            
            3. **Opção B - Usar Arquivos da Raiz**:
               - Colocar arquivos na raiz do projeto
               - Sistema buscará automaticamente na raiz se não encontrar na pasta do ano
            
            4. **Executar processamento**:
               - Clicar em "🚀 Executar tc_ext/notebooks/dados.ipynb"
                    - Sistema cria `dados/TC_Ext/2026/` automaticamente
               - Sistema busca arquivos (encontra na raiz ou na pasta do ano)
                    - Processa e salva em `dados/TC_Ext/2026/`
               - Consolida histórico
            
            **Resultado**: 
                - Pasta `dados/TC_Ext/2026/` criada com arquivos processados
            - Histórico atualizado com dados de 2026
            """)
        
        with st.expander("**Cenário 2: Atualizar Arquivos de um Ano Existente**", expanded=False):
            st.markdown("""
            **Situação**: Já processou 2024 antes, mas recebeu arquivos atualizados
            
            **Passo a Passo**:
            
            1. **Acessar página de extração**:
               - Selecionar ano: 2024
               - Selecionar tipo: "📊 Dados REAIS"
            
            2. **Verificar arquivos existentes**:
               - Sistema mostra aviso: "⚠️ O arquivo já existe"
               - Aviso aparece antes mesmo de fazer upload
            
            3. **Fazer upload do arquivo atualizado**:
               - Selecionar arquivo atualizado
               - Sistema mostra aviso: "Arquivo será sobrescrito"
               - Clicar em "🔄 Confirmar Sobrescrita"
               - Arquivo é salvo substituindo o anterior
            
            4. **Executar processamento**:
               - Clicar em "🚀 Executar tc_ext/notebooks/dados.ipynb"
                    - Sistema usa arquivo atualizado de `dados/TC_Ext/2024/`
               - Processa e atualiza arquivos Parquet
               - Atualiza histórico (concatena, não substitui)
            
            **Resultado**: 
            - Arquivos de 2024 atualizados
            - Histórico contém versão mais recente
            """)
        
        with st.expander("**Cenário 3: Processar Ambos (REAIS e BUDGET) para Novo Ano**", expanded=False):
            st.markdown("""
            **Situação**: Processar dados REAIS e BUDGET de 2026 pela primeira vez
            
            **Passo a Passo**:
            
            1. **Preparar arquivos REAIS**:
                    - Upload de `Dados SAPIENS.xlsx` (REAIS) → `dados/TC_Ext/2026/`
                    - Upload de `Reporting fluxo anexo.xlsx` (REAIS) → `dados/TC_Ext/2026/`
            
            2. **Preparar arquivos BUDGET** (se diferentes):
                    - Upload de `Dados SAPIENS.xlsx` (BUD) → `dados/TC_Ext/2026/` (mesmo arquivo ou versão BUD)
                    - Upload de `Reporting fluxo anexo.xlsx` (BUD) → `dados/TC_Ext/2026/` (com guias BDG)
            
            3. **Executar processamento**:
               - Selecionar tipo: "🔄 Ambos"
               - Clicar em "🚀 Executar Ambos"
                    - Sistema processa REAIS primeiro → Salva em `dados/TC_Ext/2026/`
                    - Sistema processa BUDGET depois → Salva em `dados/TC_Ext/2026/BUD/`
               - Consolida ambos os históricos
            
            **Resultado**: 
                - Estrutura completa criada: `dados/TC_Ext/2026/` e `dados/TC_Ext/2026/BUD/`
            - Históricos REAIS e BUDGET atualizados
            """)
        
        with st.expander("**Cenário 4: Processar Apenas BUDGET para Ano Existente**", expanded=False):
            st.markdown("""
            **Situação**: Já processou REAIS de 2024, agora quer processar BUDGET
            
            **Passo a Passo**:
            
            1. **Preparar arquivos BUDGET**:
                    - Upload de `Dados SAPIENS.xlsx` (BUD) → `dados/TC_Ext/2024/`
                    - Upload de `Reporting fluxo anexo.xlsx` (BUD) → `dados/TC_Ext/2024/`
            
            2. **Executar processamento BUDGET**:
               - Selecionar tipo: "💰 Dados BUDGET"
               - Clicar em "🚀 Executar tc_ext/notebooks/dados_BUD.ipynb"
                    - Sistema cria `dados/TC_Ext/2024/BUD/` automaticamente
                    - Processa e salva em `dados/TC_Ext/2024/BUD/`
               - Consolida histórico BUDGET
            
            **Resultado**: 
                - Pasta `dados/TC_Ext/2024/BUD/` criada com dados de Budget
            - Histórico BUDGET atualizado
            - Dados REAIS permanecem inalterados
            """)
        
        st.markdown("---")
        
        st.success("""
        **✅ Este capítulo descreve completamente o funcionamento do sistema de atualização**
        e extração de dados. Use estas informações para realizar atualizações de forma
        segura e eficiente, especialmente ao trabalhar com novos anos ou atualizar
        arquivos existentes.
        """)


# ==========================================
# SEÇÃO 5: GUIA DE BEST ESTIMATE
# ==========================================
elif indice_selecionado == "🔮 Guia de Best Estimate":
    st.header("🔮 Guia de Best Estimate — TC Ext + TC Veículos")

    with st.expander("🧮 **Ordem exata do cálculo do Best Estimate**", expanded=True):
        st.markdown(r"""
### Sequência obrigatória

1. Carregar a base histórica consolidada e definir quais períodos entram na média.
2. Remover da base de cálculo qualquer linha já marcada como BE, BE Manual ou Forecast para evitar circularidade.
3. Excluir meses anômalos, se aplicável, mantendo apenas os períodos aprovados para média.
4. Classificar cada linha como Fixo ou Variável a partir do conteúdo textual da coluna Custo.
5. Calcular a média histórica de custo no mesmo nível analítico exibido pelo simulador.
6. Calcular o volume médio histórico no mesmo perímetro usado na média de custo.
7. Definir o volume futuro de cada período a projetar.
8. Calcular a proporção de volume futuro versus volume histórico e transformá-la em variação percentual.
9. Aplicar a sensibilidade da linha sobre a variação de volume.
10. Aplicar o fator monetário com inflação e produtividade.
11. Gravar o valor calculado como BE da linha.
12. Somar BE Manual apenas como linhas adicionais, sem alterar a média histórica nem o forecast calculado pelo modelo.

### Regra econômica real aplicada pelo sistema

- Base do cálculo por linha:

$$
BE_{linha} = Média\ Histórica \times Fator\ de\ Variação \times Fator\ Monetário
$$

- Fator de variação:

$$
Fator\ de\ Variação = 1 + \left(\left(\frac{Volume\ Futuro}{Volume\ Médio\ Histórico} - 1\right) \times Sensibilidade\right)
$$

- Fator monetário:

$$
Fator\ Monetário = (1 + Inflação) \times (1 - Produtividade)
$$

- Quando o volume médio histórico é zero, o sistema força a proporção de volume para 1,0 para não distorcer o cálculo nem dividir por zero.
- A inflação pode ser aplicada por Type 06; se não houver regra específica, o sistema usa fallback global ou zero, conforme a configuração disponível.
- A produtividade é aplicada depois do ajuste de volume, também podendo ser específica por Type 06 ou global.

### Como o sistema decide Fixo vs Variável

- Se a coluna Custo contiver FIXO, FIX ou FIXED, a linha é tratada como Fixo.
- Caso contrário, a linha é tratada como Variável, salvo classificações já normalizadas anteriormente no pipeline.
- Essa distinção é crítica porque altera a sensibilidade padrão:
    - Fixo: sensibilidade padrão 0%
    - Variável: sensibilidade padrão 100%

### Fórmula geral

$$
BE = Média\ Histórica \times (1 + ((\frac{Volume\ Futuro}{Volume\ Histórico} - 1) \times Sensibilidade)) \times (1 + Inflação) \times (1 - Produtividade)
$$

### Exemplo numérico completo

- Média histórica = R$ 100.000
- Volume histórico = 1.000
- Volume futuro = 1.150
- Sensibilidade = 80%
- Inflação = 5%
- Produtividade = 2%

Passos:
1. Proporção de volume = 1.150 / 1.000 = 1,15
2. Variação de volume = 1,15 - 1 = 0,15
3. Variação ajustada = 0,15 × 0,80 = 0,12
4. Fator de volume = 1 + 0,12 = 1,12
5. Fator monetário = (1 + 0,05) × (1 - 0,02) = 1,029
6. BE = 100.000 × 1,12 × 1,029 = R$ 115.248

### BE Manual

Se houver um custo manual de R$ 8.000 para o mesmo período, ele entra como linha separada identificada como BE Manual. O consolidado final do período passa a ser R$ 123.248, mas o valor projetado pelo modelo continua R$ 115.248.

### Leitura correta da ordem do cálculo

1. O sistema nunca parte do consolidado final para calcular média; ele parte apenas das linhas históricas válidas.
2. O efeito de volume entra antes do bloco monetário.
3. Inflação e produtividade são multiplicativas entre si, não aditivas.
4. O BE Manual não recalcula o forecast do modelo; ele só complementa o consolidado final.
        """)

    with st.expander("🚗 **TC Veículos — Resumo operacional (Simulador + consumo na Home)**", expanded=False):
        st.markdown("""
        ### 🔮 O que é o Best Estimate (TC Veículos)

        O Best Estimate (BE) projeta custos futuros a partir da média histórica já realizada,
        ajustada por premissas de **sensibilidade**, **inflação**, **produtividade** e **volume**.

        **Onde configurar e gerar o Forecast:**
        - Página Streamlit: `pages/2 - Best Estimate - Simulador.py`
        - Lógica principal: `tc_principal/pages/best_estimate_simulador_tc.py`

        **Arquivos gerados:**
        - `dados/TC_Principal/Forecast/forecast_completo.parquet` — consolidado final consumido na análise
        - `dados/TC_Principal/Forecast/forecast_historico.parquet` — histórico sem os meses previstos
        - `dados/TC_Principal/Forecast/forecast_previsao.parquet` — apenas períodos futuros de BE e BE Manual
        - `dados/TC_Principal/Forecast/forecast_veiculos_custo_fp.parquet` — base rateada por veículo
        - `dados/TC_Principal/Forecast/config_forecast.json` — parâmetros persistidos do simulador

        **Onde o Forecast é consumido/analisado:**
        - `tc_principal/pages/home_tc.py` (tabs) — compara Real vs BE no layout da Home

        **Pontos de atenção (operacional):**
        - Se o Forecast parecer "não atualizar", confirme que o `forecast_completo.parquet` foi regravado.
        - `forecast_historico.parquet` exclui os meses previstos para evitar duplicação quando é combinado com `forecast_previsao.parquet`.
        - Se a granularidade por veículo depender de rateio, a função `ratear_be_por_veiculo()` (em `tc_principal/shared.py`) gera `forecast_veiculos_custo_fp.parquet` a partir do consolidado usando percentuais de rateio e D&A dedicado do Real.
        """)
    
    st.markdown("""
    <div style="padding: 1.5rem; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); border-radius: 10px; margin-bottom: 2rem; color: white;">
    <h2 style="color: white; margin: 0;">🔮 Documentação Completa do Best Estimate</h2>
    <p style="color: #f0f0f0; margin: 0.5rem 0 0 0;">
            Teoria, Cálculos, Estrutura e Funcionamento do Sistema de Previsão
    </p>
    </div>
    """, unsafe_allow_html=True)
    
    # Índice interno
    st.markdown("## 📋 Índice do Guia")
    st.markdown("""
    ### 📖 Capítulo 1: Teoria e Funcionamento do Best Estimate
    1. [O que é Best Estimate?](#o-que-e-best-estimate)
    2. [Teoria e Conceitos Fundamentais](#teoria-conceitos)
    3. [Cálculo de Médias Históricas](#calculo-medias)
    4. [Sensibilidade, Inflação e Produtividade](#sensibilidade-inflacao)
    5. [Fórmulas e Lógica de Cálculo](#formulas-logica)
    6. [Tipos de Custos: Fixo vs Variável](#tipos-custos)
    7. [Volume e Proporções](#volume-proporcoes)
    
    ### 🔄 Capítulo 2: Estrutura, Atualização e Páginas
    1. [Estrutura de Pastas do Forecast](#estrutura-forecast)
    2. [Ordem Cronológica de Atualização](#ordem-cronologica-forecast)
    3. [Página 2 - Best Estimate Simulador](#pagina-simulador)
    4. [Página - Best Estimate (Análise)](#pagina-analise)
    5. [Fluxo de Dados e Processamento](#fluxo-dados-forecast)
    6. [Arquivos Gerados](#arquivos-gerados-forecast)
    7. [Cenários de Uso](#cenarios-uso-forecast)
    """)
    
    st.markdown("---")
    
    # ==========================================
    # CAPÍTULO 1: TEORIA E FUNCIONAMENTO DO BEST ESTIMATE
    # ==========================================
    
    with st.expander("📖 **Capítulo 1: Teoria e Funcionamento do Best Estimate**", expanded=False):
        st.markdown("""
        <div style="padding: 1.5rem; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); border-radius: 10px; margin: 2rem 0; color: white;">
            <h2 style="color: white; margin: 0;">📖 Capítulo 1: Teoria e Funcionamento do Best Estimate</h2>
            <p style="color: #f0f0f0; margin: 0.5rem 0 0 0;">
                Conceitos, Teoria e Cálculos do Sistema de Previsão de Custos
            </p>
        </div>
        """, unsafe_allow_html=True)
        
        # Seção 1: O que é Best Estimate?
        st.markdown("## 🎯 O QUE É BEST ESTIMATE? {#o-que-e-best-estimate}")
        
        st.markdown("""
        ### Definição e Conceito
        
        **Best Estimate** (Melhor Estimativa) é uma metodologia de previsão de custos que combina:
        - **Dados históricos** (médias de períodos anteriores)
        - **Ajustes por sensibilidade** (resposta a variações de volume)
        - **Ajustes por inflação** (correção monetária)
        - **Ajustes por produtividade** (ganhos de eficiência que reduzem o custo projetado)
        - **Classificação de custos** (Fixo vs Variável)
        
        **Objetivo Principal:**
        Prever os custos futuros com base em padrões históricos, ajustados para refletir mudanças esperadas
        em volume de produção e inflação, permitindo planejamento financeiro mais preciso.
        
        **Aplicação no SCI:**
        O Best Estimate é usado para gerar previsões de custos para períodos futuros, permitindo comparações
        entre o que foi planejado (Budget), o que realmente aconteceu (Real) e o que se espera que aconteça
        (Best Estimate/Forecast).
        """)
        
        st.info("""
        **💡 Importante**: Best Estimate não é uma simples projeção linear. Ele considera a natureza dos custos
        (fixos ou variáveis) e aplica sensibilidades diferentes para cada tipo, resultando em previsões mais
        realistas e acuradas.
        """)
        
        st.markdown("---")
        
        # Seção 2: Teoria e Conceitos Fundamentais
        st.markdown("## 📚 TEORIA E CONCEITOS FUNDAMENTAIS {#teoria-conceitos}")
        
        st.markdown("""
        ### Fundamentos Teóricos
        
        **1. Princípio da Média Histórica:**
        - O Best Estimate parte do pressuposto de que o comportamento histórico é um bom indicador do futuro
        - Médias calculadas sobre períodos selecionados fornecem uma base sólida para previsões
        - Períodos anômalos podem ser excluídos para melhorar a acurácia
        
        **2. Princípio da Sensibilidade:**
        - Custos **fixos** não variam com volume (sensibilidade = 0%)
        - Custos **variáveis** variam proporcionalmente ao volume (sensibilidade = 100%)
        - Sensibilidades intermediárias (0% < sensibilidade < 100%) representam custos semi-variáveis
        
        **3. Princípio da Inflação:**
        - Inflação afeta todos os custos de forma uniforme
        - É aplicada como um fator multiplicador sobre o custo ajustado por sensibilidade
        - Permite correção monetária para períodos futuros
        
        **4. Princípio da Produtividade:**
        - Produtividade representa ganho operacional e atua como redutor do custo previsto
        - Pode ser aplicada globalmente ou por Type 06
        - É combinada multiplicativamente com a inflação, compondo um único fator monetário

        **5. Princípio da Proporcionalidade de Volume:**
        - A variação de volume impacta diferentemente custos fixos e variáveis
        - Custos fixos são "diluídos" quando o volume aumenta (CPU diminui)
        - Custos variáveis aumentam proporcionalmente ao volume
        """)
        
        st.markdown("---")
        
        # Seção 3: Cálculo de Médias Históricas
        st.markdown("## 📊 CÁLCULO DE MÉDIAS HISTÓRICAS {#calculo-medias}")
        
        st.markdown("""
        ### Processo de Cálculo de Médias
        
        **Passo 1: Seleção de Períodos**
        - O usuário seleciona quais períodos históricos serão usados para calcular a média
        - Exemplo: Janeiro 2024, Fevereiro 2024, Março 2024
        - Períodos podem ser excluídos se forem considerados anômalos
        
        **Passo 2: Filtragem de Dados**
        - Aplicam-se os mesmos filtros usados na análise (Oficina, Veículo, Type 05, Type 06, etc.)
        - Garante que a média seja calculada sobre o mesmo contexto operacional
        
        **Passo 3: Agrupamento e Agregação**
        - Dados são agrupados por chaves únicas: `Oficina`, `Veículo`, `Tipo_Custo`, `Type 06`, etc.
        - Para cada grupo, calcula-se a média dos valores históricos
        - Fórmula: `Média_Histórica = Σ(Valores_Históricos) / Número_de_Períodos`
        
        **Passo 4: Volume Médio Histórico**
        - Calcula-se também o volume médio histórico para os mesmos períodos
        - Usado para calcular proporções de volume futuro vs histórico
        - Fórmula: `Volume_Médio_Histórico = Σ(Volumes_Históricos) / Número_de_Períodos`
        
        **Exemplo Prático:**
        ```
        Períodos selecionados: Janeiro 2024, Fevereiro 2024, Março 2024
        
        Para Oficina A, Veículo CC21, Type 06 "Material":
        - Janeiro 2024: R$ 10.000
        - Fevereiro 2024: R$ 12.000
        - Março 2024: R$ 11.000
        
        Média Histórica = (10.000 + 12.000 + 11.000) / 3 = R$ 11.000
        ```
        """)
        
        with st.expander("**🔍 Detalhes Técnicos do Cálculo de Médias**", expanded=False):
            st.markdown("""
            **Agrupamento por Chaves:**
            - O sistema agrupa dados por múltiplas dimensões simultaneamente
            - Chaves padrão: `['Oficina', 'Veículo', 'Tipo_Custo', 'Type 06', ...]`
            - Cada combinação única de chaves gera uma linha no forecast
            
            **Tratamento de Dados Faltantes:**
            - Se um período não tiver dados para uma combinação de chaves, ele é excluído do cálculo
            - A média é calculada apenas sobre períodos com dados disponíveis
            - Isso evita distorções por períodos incompletos
            
            **Normalização de Períodos:**
            - Períodos são normalizados para comparação (ex: "Janeiro 2024" → "janeiro 2024")
            - Permite comparação case-insensitive e tolerante a espaços
            - Anos são extraídos dos períodos para filtragem adicional
            """)
        
        st.markdown("---")
        
        # Seção 4: Sensibilidade, Inflação e Produtividade
        st.markdown("## ⚙️ SENSIBILIDADE, INFLAÇÃO E PRODUTIVIDADE {#sensibilidade-inflacao}")
        
        st.markdown("""
        ### Sensibilidade ao Volume
        
        **Conceito:**
        Sensibilidade mede o quanto um custo responde a variações no volume de produção.
        
        **Tipos de Sensibilidade:**
        
        **1. Sensibilidade Fixa (0%):**
        - Aplicada a custos **fixos**
        - Independente da variação de volume, o custo permanece constante
        - Exemplos: Aluguel, salários fixos, depreciação
        - Fórmula: `Custo_Ajustado = Custo_Original` (sem alteração)
        
        **2. Sensibilidade Variável (100%):**
        - Aplicada a custos **variáveis**
        - Varia proporcionalmente ao volume
        - Se volume aumenta 10%, custo aumenta 10%
        - Exemplos: Matéria-prima, energia variável, comissões
        - Fórmula: `Custo_Ajustado = Custo_Original * (Volume_Novo / Volume_Histórico)`
        
        **3. Sensibilidades Intermediárias (0% < sensibilidade < 100%):**
        - Aplicadas a custos **semi-variáveis**
        - Resposta parcial a variações de volume
        - Exemplo: Se sensibilidade = 50% e volume aumenta 10%, custo aumenta 5%
        - Fórmula: `Custo_Ajustado = Custo_Original * (1 + (Variação_Volume * Sensibilidade))`
        
        **4. Sensibilidade por Type 06:**
        - Cada Type 06 pode ter sua própria sensibilidade específica
        - Permite ajustes finos por categoria de custo
        - Sobrescreve a sensibilidade geral (Fixo/Variável) quando configurada
        """)
        
        st.markdown("""
        ### Inflação
        
        **Conceito:**
        Inflação é aplicada como um ajuste monetário uniforme sobre todos os custos, independente
        de serem fixos ou variáveis.
        
        **Aplicação:**
        - Inflação é configurada como percentual (ex: 5% ao ano)
        - É aplicada após o ajuste por sensibilidade
        - Fórmula: `Custo_Final = Custo_Ajustado_Sensibilidade * (1 + Inflação/100)`
        
        **Exemplo:**
        ```
        Custo médio histórico: R$ 10.000
        Variação de volume: +10%
        Sensibilidade: 50%
        Inflação: 5%
        
        Passo 1: Ajuste por sensibilidade
        Variação_ajustada = 10% * 50% = 5%
        Custo_ajustado = 10.000 * (1 + 0.05) = R$ 10.500
        
        Passo 2: Aplicar inflação
        Custo_final = 10.500 * (1 + 0.05) = R$ 11.025
        ```
        """)

        st.markdown("""
        ### Produtividade

        **Conceito:**
        Produtividade representa ganho de eficiência operacional e reduz o custo previsto.

        **Aplicação:**
        - Pode ser configurada globalmente ou por Type 06
        - É aplicada após o ajuste de volume, no mesmo bloco monetário da inflação
        - Fórmula operacional do sistema: `Fator_Monetário = (1 + Inflação/100) × (1 - Produtividade/100)`

        **Exemplo:**
        ```
        Custo ajustado por volume: R$ 10.500
        Inflação: 5%
        Produtividade: 3%

        Fator_monetário = (1 + 0.05) * (1 - 0.03) = 1.0185
        Custo_final = 10.500 * 1.0185 = R$ 10.694,25
        ```

        **Leitura correta:**
        - Inflação pressiona o custo para cima
        - Produtividade compensa parte dessa pressão
        - O efeito final depende da combinação multiplicativa dos dois fatores
        """)
        
        st.markdown("---")
        
        # Seção 5: Fórmulas e Lógica de Cálculo
        st.markdown("## 🧮 FÓRMULAS E LÓGICA DE CÁLCULO {#formulas-logica}")
        
        st.markdown("""
        ### Fórmula Completa do Best Estimate
        
        **Fórmula Geral (linha a linha):**
        ```
        Best_Estimate = Média_Histórica * Fator_Variação * Fator_Monetário
        ```
        
        **Onde, na implementação real:**
        - `Média_Histórica` = média calculada apenas sobre linhas históricas válidas, excluindo `BE`, `BE Manual` e `Forecast`
        - `Fator_Variação` = `1 + (Variação_Percentual_Volume * Sensibilidade)`
        - `Fator_Monetário` = `(1 + Inflação / 100) × (1 - Produtividade / 100)`
        - `Sensibilidade` pode ser global ou específica por `Type 06`
        - `Inflação` pode ser global ou específica por `Type 06`
        - `Produtividade` pode ser global ou específica por `Type 06`
        
        **Cálculo Detalhado Passo a Passo:**

        **0. Higienizar a base antes da média:**
        ```
        Base_para_média = Base_Histórica sem tipos ['BE', 'BE Manual', 'Forecast']
        ```
        - Este passo evita que um forecast antigo contamine a nova média histórica
        - Meses removidos manualmente pelo usuário também ficam fora da média
        
        **1. Calcular Proporção de Volume:**
        ```
        proporção_volume = Volume_do_Mês_Futuro / Volume_Médio_Histórico
        ```
        - Se `Volume_Médio_Histórico <= 0`, o sistema usa `proporção_volume = 1.0`
        
        **2. Calcular Variação Percentual:**
        ```
        variação_percentual = proporção_volume - 1.0
        ```
        - Se `variação_percentual > 0`: Volume aumentou
        - Se `variação_percentual < 0`: Volume diminuiu
        - Se `variação_percentual = 0`: Volume permaneceu igual
        
        **3. Aplicar Sensibilidade:**
        ```
        variação_ajustada = variação_percentual * sensibilidade
        ```
        - Para custos fixos: sensibilidade padrão = `0.0` → volume não altera o valor, salvo override manual
        - Para custos variáveis: sensibilidade padrão = `1.0` → volume entra integralmente, salvo override manual
        - Para custos semi-variáveis ou ajustes finos: o usuário pode configurar sensibilidades intermediárias
        
        **4. Calcular Fator de Variação:**
        ```
        fator_variação = 1.0 + variação_ajustada
        ```
        
        **5. Calcular Fator Monetário:**
        ```
        fator_monetário = (1.0 + (inflação / 100.0)) * (1.0 - (produtividade / 100.0))
        ```
        - O sistema aplica inflação e produtividade de forma multiplicativa
        - A produtividade reduz o efeito monetário total e entra depois do ajuste de volume
        
        **6. Calcular Best Estimate Final:**
        ```
        Best_Estimate = Média_Histórica * fator_variação * fator_monetário
        ```

        **7. Persistir e consolidar:**
        ```
        Consolidado_Final = Histórico + BE + BE_Manual
        ```
        - O valor `BE` calculado pelo modelo é salvo separadamente do `BE Manual`
        - O consolidado final combina as camadas, mas o cálculo do modelo continua rastreável
        """)
        
        with st.expander("**📐 Exemplo Completo de Cálculo**", expanded=False):
            st.markdown("""
            **Cenário:**
            - Média histórica: R$ 10.000
            - Volume médio histórico: 1.000 unidades
            - Volume do mês futuro: 1.100 unidades
            - Tipo de custo: Variável (sensibilidade = 100%)
            - Inflação: 5%
            - Produtividade: 2%
            
            **Cálculo:**
            
            **Passo 1:** Proporção de volume
            ```
            proporção = 1.100 / 1.000 = 1.1
            ```
            
            **Passo 2:** Variação percentual
            ```
            variação = 1.1 - 1.0 = 0.1 (10% de aumento)
            ```
            
            **Passo 3:** Aplicar sensibilidade
            ```
            variação_ajustada = 0.1 * 1.0 = 0.1 (10%)
            ```
            
            **Passo 4:** Fator de variação
            ```
            fator_variação = 1.0 + 0.1 = 1.1
            ```
            
            **Passo 5:** Fator monetário
            ```
            fator_monetário = (1.0 + 0.05) * (1.0 - 0.02) = 1.029
            ```
            
            **Passo 6:** Best Estimate
            ```
            Best_Estimate = 10.000 * 1.1 * 1.029 = R$ 11.319
            ```
            
            **Interpretação:**
            O custo previsto é R$ 11.319, representando:
            - Aumento de 10% devido ao aumento de volume (de 1.000 para 1.100 unidades)
            - Aumento monetário de 5% por inflação parcialmente compensado por 2% de produtividade
            - Total: 13,19% de aumento sobre a média histórica
            """)
        
        st.markdown("---")
        
        # Seção 6: Tipos de Custos
        st.markdown("## 💰 TIPOS DE CUSTOS: FIXO VS VARIÁVEL {#tipos-custos}")
        
        st.markdown("""
                ### Classificação de Custos

                **Custos Fixos:**
                - **Características:** Não variam com o volume de produção
                - **Sensibilidade:** 0% (zero por cento)
                - **Exemplos:** Aluguel, salários fixos, depreciação, seguros
                - **Comportamento no Best Estimate:**
                    - Média histórica é mantida (sem ajuste por volume)
                    - Aplicam-se apenas os fatores monetários de inflação e produtividade
                    - Fórmula: `Best_Estimate_Fixo = Média_Histórica_Fixo * (1 + Inflação/100) * (1 - Produtividade/100)`

                **Custos Variáveis:**
                - **Características:** Variam proporcionalmente ao volume de produção
                - **Sensibilidade:** 100% (cem por cento)
                - **Exemplos:** Matéria-prima, energia variável, comissões, peças de reposição
                - **Comportamento no Best Estimate:**
                    - Média histórica é ajustada pela proporção de volume
                    - Inflação e produtividade são aplicadas sobre o valor já ajustado por volume
                    - Fórmula: `Best_Estimate_Variável = Média_Histórica_Variável * (Volume_Futuro/Volume_Histórico) * (1 + Inflação/100) * (1 - Produtividade/100)`

                **Identificação no Sistema:**
                                - O sistema procura na coluna `Custo` textos contendo `FIXO`, `FIX` ou `FIXED`
                                - Se encontrar esses termos, a linha é classificada como `Fixo`
                                - Se não encontrar, a linha segue como `Variável` ou mantém a classificação já normalizada anteriormente
                                - Em alguns fluxos a classificação já chega tratada após merges e normalizações, mas a regra textual continua sendo o fallback operacional
                                - Cada linha precisa ter essa classificação correta porque a sensibilidade aplicada depende dela

                                **Como calcular na prática:**
                                - **Fixo:**
                                    `BE_Fixo = Média_Histórica_Fixo × Fator_Monetário`
                                - **Variável:**
                                    `BE_Variável = Média_Histórica_Variável × Fator_Variação × Fator_Monetário`
                                - **Semi-variável / ajuste manual:**
                                    usa a mesma fórmula geral, mudando apenas a sensibilidade para um valor intermediário
                """)
        
        st.markdown("---")
        
        # Seção 7: Volume e Proporções
        st.markdown("## 📈 VOLUME E PROPORÇÕES {#volume-proporcoes}")
        
        st.markdown("""
        ### Importância do Volume no Best Estimate
        
        **Volume como Base de Cálculo:**
        - O volume futuro é usado para calcular a proporção em relação ao volume histórico
        - Esta proporção determina o ajuste aplicado aos custos variáveis
        - Volume médio histórico é calculado sobre os mesmos períodos usados para a média de custos
        
        **Cálculo de Proporção:**
        ```
        proporção = Volume_Mês_Futuro / Volume_Médio_Histórico
        ```
        
        **Interpretação da Proporção:**
        - `proporção > 1.0`: Volume futuro é maior que o histórico → Custos variáveis aumentam
        - `proporção < 1.0`: Volume futuro é menor que o histórico → Custos variáveis diminuem
        - `proporção = 1.0`: Volume futuro igual ao histórico → Sem ajuste por volume, restando apenas os fatores monetários
        
        **Impacto nos Custos:**
        - **Custos Fixos:** Não são afetados pela proporção (sensibilidade = 0%), mas ainda sofrem inflação e produtividade
        - **Custos Variáveis:** São multiplicados pela proporção (sensibilidade = 100%) e depois ajustados por inflação e produtividade
        - **Custos Semi-Variáveis:** São multiplicados por `1 + (proporção - 1) * sensibilidade` e depois passam pelo fator monetário

        **Regras de fallback importantes:**
        - Se o volume médio histórico for zero, o sistema usa proporção = 1,0
        - Isso evita divisão por zero e impede explosões artificiais no forecast
        - Nessa situação, a linha fica dependente apenas da média histórica e do fator monetário
        """)
        
        st.success("""
        **✅ Este capítulo descreve completamente a teoria e funcionamento do Best Estimate.**
        Use estas informações para entender como as previsões são calculadas e como os parâmetros
        (sensibilidade, inflação, produtividade e períodos históricos) impactam os resultados.
        """)
    
    # ==========================================
    # CAPÍTULO 2: ESTRUTURA, ATUALIZAÇÃO E PÁGINAS
    # ==========================================
    
    with st.expander("🔄 **Capítulo 2: Estrutura, Atualização e Páginas**", expanded=False):
        st.markdown("""
        <div style="padding: 1.5rem; background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%); border-radius: 10px; margin: 2rem 0; color: white;">
            <h2 style="color: white; margin: 0;">🔄 Capítulo 2: Estrutura, Atualização e Páginas</h2>
            <p style="color: #f0f0f0; margin: 0.5rem 0 0 0;">
                Estrutura de Pastas, Ordem de Atualização e Funcionalidades das Páginas
            </p>
        </div>
        """, unsafe_allow_html=True)
        
        # Seção 1: Estrutura de Pastas do Forecast
        st.markdown("## 📁 ESTRUTURA DE PASTAS DO FORECAST {#estrutura-forecast}")
        
        st.markdown("""
        ### Organização da Pasta `dados/TC_Ext/Forecast/`
        
        A pasta `Forecast/` é criada automaticamente quando o Best Estimate é gerado e contém
        os arquivos de previsão calculados pelo sistema.
        
        **Estrutura Completa:**
        ```
        dados/TC_Ext/
        └── Forecast/                       # 🔮 Dados de Best Estimate/Forecast
            ├── forecast_completo.parquet   # Forecast completo com todas as linhas
            ├── forecast_historico.parquet  # Histórico de forecasts gerados
            ├── forecast_previsao.parquet   # Previsões futuras
            ├── df_final_historico_forecast.parquet  # Dados históricos filtrados para forecast
            ├── df_vol_historico.parquet    # Volumes históricos para cálculo
            ├── custos_especificos.parquet  # Linhas BE Manual persistidas
            └── config_forecast.json        # Configurações persistidas do simulador
        ```
        
        **Características:**
        - **Criação Automática:** A pasta é criada automaticamente se não existir
        - **Substituição:** Arquivos são substituídos a cada geração (não concatenados)
        - **Prioridade:** Sistema busca arquivos nesta pasta primeiro antes de usar histórico consolidado
        - **Formato:** Todos os arquivos são Parquet para performance otimizada
        """)
        
        st.markdown("---")
        
        # Seção 2: Ordem Cronológica de Atualização
        st.markdown("## ⏱️ ORDEM CRONOLÓGICA DE ATUALIZAÇÃO {#ordem-cronologica-forecast}")
        
        st.markdown("### Sequência Completa do Processo")
        
        with st.expander("**1️⃣ Configuração de Parâmetros**", expanded=False):
                st.markdown("""
                **Onde**: Página 2 (Simulador) ou Página 3 (Análise)

                **Processo**:
                1. Usuário seleciona **períodos históricos** para calcular a média
                    - Exemplo: Janeiro 2024, Fevereiro 2024, Março 2024
                    - Períodos podem ser excluídos se anômalos

                2. Usuário configura **sensibilidades**:
                    - Sensibilidade para custos fixos (geralmente 0%)
                    - Sensibilidade para custos variáveis (geralmente 100%)
                    - Sensibilidades específicas por Type 06 (opcional)

                3. Usuário configura **inflação**:
                    - Percentual de inflação anual (ex: 5%)
                    - Pode ser aplicada globalmente ou por Type 06

                4. Usuário configura **produtividade**:
                    - Ganho de eficiência que reduz o custo projetado
                    - Pode ser aplicada globalmente ou por Type 06
                    - É combinada com a inflação no fator monetário final

                5. Usuário seleciona **períodos futuros** para forecast:
                    - Exemplo: Abril 2024, Maio 2024, Junho 2024
                    - Volumes futuros são informados ou calculados

                **Resultado**: Sistema tem todos os parâmetros necessários para calcular o forecast e persistir as configurações em `config_forecast.json`
                """)
        
        with st.expander("**2️⃣ Carregamento de Dados Históricos**", expanded=False):
            st.markdown("""
            **Onde**: Função `load_data()` nas páginas 2 e 3
            
            **Ordem de Prioridade de Busca**:
                1. **Primeira opção**: `dados/TC_Ext/Forecast/forecast_completo.parquet`
               - Se existir, pode ser usado como base (mas forecast é recalculado)
            
                2. **Segunda opção**: `dados/TC_Ext/historico_consolidado/df_final_historico.parquet`
               - **Fonte principal** de dados históricos
               - Contém todos os anos consolidados
            
                3. **Terceira opção**: `dados/TC_Ext/{ANO}/df_final.parquet`
               - Dados específicos do ano (se filtro de ano aplicado)
            
            **Processo**:
            - Sistema carrega dados históricos completos
            - Aplica filtros selecionados (Oficina, Veículo, Type 05, Type 06, etc.)
            - Filtra pelos períodos selecionados para cálculo de média
            - Remove períodos excluídos (meses_excluir_media)
            """)
        
        with st.expander("**3️⃣ Cálculo de Médias Históricas**", expanded=False):
            st.markdown("""
            **Onde**: Função `calcular_medias_forecast()` nas páginas 2 e 3
            
            **Processo**:
            1. **Filtrar dados pelos períodos selecionados**:
               - Apenas períodos marcados para média são considerados
               - Períodos excluídos são removidos
            
            2. **Agrupar por chaves únicas**:
               - Chaves: `['Oficina', 'Veículo', 'Tipo_Custo', 'Type 06', ...]`
               - Cada combinação única gera uma linha no forecast
            
            3. **Calcular média por grupo**:
               - Soma dos valores históricos / número de períodos
               - Usa coluna `Total` (nunca `Valor`)
            
            4. **Calcular volume médio histórico**:
               - Mesma lógica: agrupa e calcula média de volumes
               - Usado para calcular proporções futuras
            
            **Resultado**: DataFrame com médias históricas por combinação de chaves
            """)
        
        with st.expander("**4️⃣ Cálculo do Forecast**", expanded=False):
                st.markdown("""
                **Onde**: Função `calcular_forecast_completo()` nas páginas 2 e 3

                **Processo (linha a linha)**:
                1. **Para cada linha do forecast**:
                    - Obtém média histórica da combinação de chaves
                    - Obtém volume do mês futuro
                    - Obtém volume médio histórico

                2. **Calcula proporção de volume**:
                    ```
                    proporção = Volume_Mês_Futuro / Volume_Médio_Histórico
                    ```

                3. **Calcula variação percentual**:
                    ```
                    variação = proporção - 1.0
                    ```

                4. **Aplica sensibilidade**:
                    - Se `Tipo_Custo == 'Fixo'`: usa `sensibilidade_fixo`
                    - Se `Tipo_Custo == 'Variável'`: usa `sensibilidade_variavel`
                    - Se modo Type 06: usa sensibilidade específica do Type 06
                    ```
                    variação_ajustada = variação * sensibilidade
                    ```

                5. **Calcula fator monetário e forecast**:
                    ```
                    fator_variação = 1.0 + variação_ajustada
                    fator_monetário = (1.0 + (inflação / 100.0)) * (1.0 - (produtividade / 100.0))
                    forecast = Média_Histórica * fator_variação * fator_monetário
                    ```

                6. **Prioridade das premissas monetárias**:
                    - Se houver configuração detalhada por Type 06, ela prevalece
                    - Na ausência dela, o sistema usa inflação e produtividade globais

                **Resultado**: DataFrame completo com forecast linha a linha
                """)
        
        with st.expander("**5️⃣ Salvamento dos Arquivos**", expanded=False):
                st.markdown("""
                **Onde**: Função de salvamento nas páginas 2 e 3

                **Processo**:
                1. **Verificar/Criar pasta Forecast**:
                    - Verifica se `dados/TC_Ext/Forecast/` existe
                    - Se não existe, cria automaticamente: `os.makedirs(pasta_forecast, exist_ok=True)`

                2. **Salvar forecast_completo.parquet**:
                    - Arquivo principal com todas as linhas do forecast
                    - Substitui arquivo anterior (não concatena)
                    - Localização: `dados/TC_Ext/Forecast/forecast_completo.parquet`

                3. **Salvar forecast_historico.parquet** (se aplicável):
                    - Histórico sem os períodos que estão sendo previstos
                    - Evita duplicação ao ser combinado com `forecast_previsao.parquet`

                4. **Salvar forecast_previsao.parquet** (se aplicável):
                    - Apenas previsões futuras (sem dados históricos)

                5. **Salvar config_forecast.json**:
                    - Persiste o modo global/detalhado
                    - Guarda sensibilidade, inflação e produtividade aplicadas

                **IMPORTANTE**:
                - Arquivos são **substituídos** a cada geração (não concatenados como histórico)
                - Cada geração cria um forecast novo baseado nas configurações atuais
                - O resultado é sobrescrito, mas a configuração aplicada permanece disponível para a próxima abertura
                """)
        
        st.markdown("---")
        
        # Seção 3: Página 2 - Best Estimate Simulador
        st.markdown("## 🔮 PÁGINA 2 - BEST ESTIMATE SIMULADOR {#pagina-simulador}")
        
        st.markdown("""
        ### Funcionalidades Principais
        
        **Objetivo:**
        A página 2 (Best Estimate - Simulador) permite **simular e ajustar** parâmetros do forecast
        em tempo real, visualizando o impacto das mudanças antes de salvar.
        
        **Funcionalidades:**
        
        **1. Configuração Interativa de Parâmetros:**
        - Seleção de períodos históricos para média (multiselect)
        - Exclusão de meses específicos (multiselect)
        - Configuração de sensibilidades (fixo, variável, Type 06)
        - Configuração de inflação (global e por Type 06)
        - Configuração de produtividade (global e por Type 06)
        - Seleção de períodos futuros para forecast
        
        **2. Visualização em Tempo Real:**
        - Gráficos atualizados automaticamente ao alterar parâmetros
        - Tabelas interativas mostrando valores linha a linha
        - Comparação entre diferentes cenários
        
        **3. Ajustes de Volume:**
        - Permite ajustar volumes futuros manualmente
        - Visualiza impacto imediato nos custos previstos
        - Suporta diferentes volumes por período
        
        **4. Salvamento de Forecast:**
        - Botão para salvar forecast calculado
        - Salva em `dados/TC_Ext/Forecast/forecast_completo.parquet`
        - Persiste parâmetros em `dados/TC_Ext/Forecast/config_forecast.json`
        - Substitui forecast anterior
        
        **5. Análise de Sensibilidade:**
        - Permite testar diferentes valores de sensibilidade
        - Visualiza impacto de mudanças de sensibilidade, inflação e produtividade
        - Útil para cenários "what-if"
        
        **6. Custos Específicos (BE Manual):**
        - Permite adicionar custos específicos com valores manuais
        - Suporta dois tipos de aplicação:
          - **Pontual**: Aplicado em meses específicos selecionados
          - **Constante**: Aplicado a partir de um mês inicial em diante
        - Rateio automático por veículo baseado em percentuais do arquivo de rateio
        - Integração automática com Account (Type 07) para buscar Type 06, Type 05, Custo e USI
        - Visualização e exclusão de custos específicos cadastrados
        - Formatação numérica com separador de milhares (formato brasileiro)
        - Tabela interativa com seleção múltipla para exclusão em lote
        - Os custos específicos são marcados como "BE Manual" na coluna Tipo
        - Integrados automaticamente ao forecast final como linhas separadas
        
        **7. Nomenclatura Atualizada:**
        - Coluna "Tipo" agora usa "BE" para forecast normal
        - Coluna "Tipo" usa "BE Manual" para custos específicos/manuais
        - Título atualizado: "Best Estimate - Previsão de Custo Total"
        - Compatibilidade automática com arquivos antigos (conversão de "Forecast" para "BE")
        """)
        
        st.markdown("---")
        
        # Seção 3.1: Custos Específicos - Detalhamento
        st.markdown("### 💰 Custos Específicos (BE Manual) - Detalhamento")
        
        st.markdown("""
        **Funcionalidade:** Permite adicionar custos específicos com valores manuais que são integrados ao forecast.
        
        **Como Funciona:**
        
        **1. Adicionar Custo Específico:**
        - Acesse a aba "➕ Adicionar Custo" na página 2
        - Preencha os campos obrigatórios:
          - **Account (Type 07)**: Seleciona o Account e busca automaticamente Type 06, Type 05, Custo e USI
          - **Oficina**: Seleciona a oficina (sem opção "Todos")
          - **Veículo**: Seleciona veículo específico ou "Todos" para rateio automático
          - **Período**: Seleciona o período (mês e ano)
          - **Tipo de Aplicação**: 
            - **Pontual**: Aplicado apenas nos meses selecionados
            - **Constante**: Aplicado a partir do mês inicial em diante
          - **Valor Total**: Valor total do custo
          - **Descrição**: Descrição opcional do custo
        
        **2. Rateio Automático:**
        - Se "Todos" for selecionado para Veículo, o sistema busca automaticamente os percentuais de rateio do arquivo `Reporting fluxo anexo.xlsx` (aba "Rateio")
        - O rateio é aplicado mês a mês conforme os percentuais do arquivo
        - Se um veículo específico for selecionado, o rateio é 100% para aquele veículo
        - O valor total é distribuído proporcionalmente entre os veículos
        
        **3. Visualizar Custos:**
        - Acesse a aba "📋 Visualizar Custos"
        - Tabela interativa com todas as colunas do formato `df_final_historico_forecast.xlsx`
        - Formatação numérica com 2 casas decimais e separador de milhares (formato brasileiro)
        - Seleção múltipla com checkboxes para exclusão em lote
        - Botão "🗑️ Deletar Selecionadas" para remover custos
        
        **4. Integração com Forecast:**
        - Os custos específicos são automaticamente incluídos no forecast final
        - Aparecem como linhas separadas com Tipo = "BE Manual"
        - Não são somados ao forecast calculado, mas adicionados como linhas independentes
        - Mantém o mesmo formato e estrutura do forecast normal
        
        **5. Persistência:**
        - Os custos específicos são salvos em `dados/TC_Ext/Forecast/custos_especificos.parquet`
        - São carregados automaticamente ao gerar o forecast
        - Permanecem salvos até serem explicitamente excluídos
        
        **6. Formato de Dados:**
        - Os custos específicos seguem exatamente o formato de `df_final_historico_forecast.xlsx`
        - Colunas na ordem: Account, Ano, Centrocst, Custo, Fornec., Fornecedor, Mes, Oficina, Período, Soma_Percentuais, Tipo, Total, Type 05, Type 06, USI, Valor, Veículo
        - Tipo sempre preenchido como "BE Manual" para identificação
        """)
        
        st.markdown("---")
        
        # Seção 4: Página - Best Estimate (Análise)
        st.markdown("## 📊 PÁGINA - BEST ESTIMATE (ANÁLISE) {#pagina-analise}")
        
        st.markdown("""
        ### Funcionalidades Principais
        
        **Objetivo:**
        A página **Best Estimate (Análise)** no menu **TC Ext** substitui a análise legacy e entrega:
        - as **mesmas tabelas/visuais** da Home (TC Ext),
        - porém alimentadas pelos **arquivos de Forecast** gerados pelo simulador.
        
        **Funcionalidades:**
        
        **1. Fonte de dados (Forecast):**
        - Lê `dados/TC_Ext/Forecast/forecast_completo.parquet` (custos) e `dados/TC_Ext/Forecast/df_vol_historico.parquet` (volume)
        - Permite analisar previsões (BE) e histórico no mesmo layout
        - Expander de diagnóstico mostra paths, mtimes e contagens
        
        **2. Visualizações (mesma base da Home):**
        - Gráficos e tabelas por período, oficina, veículo
        - Mesmo padrão de filtros e formatação
        - Sem “corte” de meses futuros quando houver Forecast
        
        **3. Tabelas detalhadas (com TOTAL coerente):**
        - No modo CPU, totais são sempre `CPU = sum(Total) / sum(Volume)` (ponderado)
        - Expander opcional “Volume por período” ajuda a explicar variações do TOTAL mês a mês
        
        **4. Comparações:**
        - Permite comparar BE vs histórico dentro do mesmo layout de análise
        - Facilita validar premissas (sensibilidade/inflação/produtividade) pela variação temporal
        
        **5. Integração com o simulador:**
        - O simulador gera/salva os arquivos em `dados/TC_Ext/Forecast/`
        - A análise lê esses arquivos e atualiza as visualizações
        
        **6. Modos de visualização:**
        - **Custo Total:** Valores absolutos em R$
        - **CPU (Custo por Unidade):** Valores por unidade produzida
        - Permite alternar entre modos para diferentes análises
        """)
        
        st.markdown("---")
        
        # Seção 5: Fluxo de Dados e Processamento
        st.markdown("## 🔄 FLUXO DE DADOS E PROCESSAMENTO {#fluxo-dados-forecast}")
        
        st.markdown("""
        ### Fluxo Completo de Dados
        
        **Diagrama de Fluxo:**
        ```
        Dados Históricos (dados/TC_Ext/historico_consolidado/)
                │
                ├──> Carregamento (load_data)
                │       │
                │       ├──> Aplicar Filtros (Oficina, Veículo, etc.)
                │       │
                │       └──> Filtrar Períodos Selecionados
                │
                ├──> Cálculo de Médias (calcular_medias_forecast)
                │       │
                │       ├──> Agrupar por Chaves Únicas
                │       │
                │       ├──> Calcular Média de Custos
                │       │
                │       └──> Calcular Volume Médio Histórico
                │
                ├──> Cálculo de Forecast (calcular_forecast_completo)
                │       │
                │       ├──> Para cada linha:
                │       │       ├──> Obter Média Histórica
                │       │       ├──> Obter Volume Futuro
                │       │       ├──> Calcular Proporção
                │       │       ├──> Aplicar Sensibilidade
                │       │       ├──> Aplicar Inflação
                │       │       └──> Aplicar Produtividade
                │       │
                │       └──> DataFrame Completo com Forecast
                │
                └──> Salvamento (dados/TC_Ext/Forecast/)
                        │
                        ├──> forecast_completo.parquet
                        ├──> forecast_historico.parquet
                        ├──> forecast_previsao.parquet
                        └──> config_forecast.json
        ```
        
        **Características do Fluxo:**
        - **Tempo Real:** Forecast é calculado em tempo real com configurações atuais
        - **Persistência de Parâmetros:** Sensibilidade, inflação e produtividade aplicadas ficam salvas em `config_forecast.json`
        - **Substituição:** Cada geração substitui o forecast anterior
        - **Independência:** Cada página pode gerar seu próprio forecast
        """)
        
        st.markdown("---")
        
        # Seção 6: Arquivos Gerados
        st.markdown("## 📄 ARQUIVOS GERADOS {#arquivos-gerados-forecast}")
        
        st.markdown("""
        ### Arquivos na Pasta `dados/TC_Ext/Forecast/`
        
        **1. forecast_completo.parquet**
        - **Conteúdo**: Forecast completo com todas as linhas calculadas
        - **Estrutura**: Mesmas colunas dos dados históricos + colunas de forecast
        - **Uso**: Fonte principal para análises e visualizações
        - **Atualização**: Substituído a cada geração de forecast
        
        **2. forecast_historico.parquet**
        - **Conteúdo**: Histórico de forecasts gerados anteriormente
        - **Estrutura**: Base histórica usada no consolidado, sem os períodos que estão sendo previstos
        - **Uso**: Evita duplicação ao juntar histórico com `forecast_previsao.parquet`
        - **Atualização**: Substituído a cada geração
        
        **3. forecast_previsao.parquet**
        - **Conteúdo**: Apenas previsões futuras (sem dados históricos)
        - **Estrutura**: Apenas períodos futuros do forecast
        - **Uso**: Análise focada apenas em previsões
        - **Atualização**: Substituído a cada geração
        
        **4. df_final_historico_forecast.parquet**
        - **Conteúdo**: Dados históricos filtrados usados para calcular o forecast
        - **Estrutura**: Dados históricos após aplicação de filtros e seleção de períodos
        - **Uso**: Referência dos dados que foram usados para calcular a média
        - **Atualização**: Gerado junto com o forecast
        
        **5. df_vol_historico.parquet**
        - **Conteúdo**: Volumes históricos usados para cálculo de proporções
        - **Estrutura**: Volumes por período, oficina, veículo
        - **Uso**: Cálculo de volume médio histórico e proporções
        - **Atualização**: Pode ser copiado do histórico consolidado ou gerado
        
        **6. custos_especificos.parquet**
        - **Conteúdo**: Custos específicos cadastrados manualmente (BE Manual)
        - **Estrutura**: Mesmo formato de `df_final_historico_forecast.xlsx` com coluna Tipo = "BE Manual"
        - **Uso**: Armazena custos específicos que são integrados ao forecast final
        - **Atualização**: Criado/modificado ao adicionar ou excluir custos específicos
        - **Localização**: `dados/TC_Ext/Forecast/custos_especificos.parquet`

        **7. config_forecast.json**
        - **Conteúdo**: Modo de configuração e últimos parâmetros aplicados
        - **Estrutura**: JSON com sensibilidade, inflação e produtividade globais e/ou por Type 06
        - **Uso**: Recarrega automaticamente as premissas na próxima abertura do simulador
        - **Atualização**: Sobrescrito quando o usuário aplica nova configuração
        """)
        
        st.markdown("---")
        
        # Seção 6.1: Nomenclatura e Tipos
        st.markdown("### 🏷️ Nomenclatura e Tipos de Dados")
        
        st.markdown("""
        **Coluna "Tipo" no Forecast:**
        
        O sistema utiliza a coluna "Tipo" para identificar diferentes tipos de dados no forecast:
        
        - **"Histórico"**: Dados históricos reais (não previstos)
        - **"BE"**: Best Estimate - Forecast calculado automaticamente pelo sistema
        - **"BE Manual"**: Best Estimate Manual - Custos específicos adicionados manualmente
        
        **Compatibilidade:**
        - Arquivos antigos com "Forecast" são automaticamente convertidos para "BE" ao carregar
        - Isso garante compatibilidade com versões anteriores do sistema
        
        **Filtros e Separação:**
        - O sistema separa automaticamente histórico, BE e BE Manual ao gerar arquivos
        - `forecast_historico.parquet`: Apenas dados históricos
        - `forecast_previsao.parquet`: Apenas BE e BE Manual (previsões)
        - `df_final_historico_forecast.parquet`: Consolidado com todos os tipos
        """)
        
        st.markdown("---")
        
        # Seção 7: Cenários de Uso
        st.markdown("## 📋 CENÁRIOS DE USO {#cenarios-uso-forecast}")
        
        st.markdown("### Casos Práticos Completos")
        
        with st.expander("**Cenário 1: Gerar Forecast pela Primeira Vez**", expanded=False):
            st.markdown("""
            **Situação**: Nunca gerou forecast, precisa criar previsões para próximos meses
            
            **Passo a Passo**:
            
            1. **Acessar Página 2 (Simulador)**:
               - Selecionar períodos históricos (ex: últimos 3 meses)
               - Configurar sensibilidades (Fixo: 0%, Variável: 100%)
               - Configurar inflação (ex: 5%)
                    - Configurar produtividade (ex: 2%)
               - Selecionar períodos futuros (ex: próximos 6 meses)
            
            2. **Informar Volumes Futuros**:
               - Inserir volumes esperados para cada período futuro
               - Ou usar volumes projetados automaticamente
            
            3. **Visualizar Resultados**:
               - Verificar gráficos e tabelas
               - Ajustar parâmetros se necessário
            
            4. **Salvar Forecast**:
               - Clicar em "Salvar Forecast"
                    - Sistema cria `dados/TC_Ext/Forecast/` automaticamente
               - Salva `forecast_completo.parquet`
            
            5. **Analisar na Página 3**:
               - Acessar Página 3 (Análise)
               - Carregar forecast gerado
               - Visualizar análises detalhadas
            
            **Resultado**: 
            - Pasta `dados/TC_Ext/Forecast/` criada com forecast completo
            - Forecast disponível para análises e comparações
            """)
        
        with st.expander("**Cenário 2: Atualizar Forecast com Novos Dados**", expanded=False):
            st.markdown("""
            **Situação**: Já existe forecast, mas novos dados históricos foram adicionados
            
            **Passo a Passo**:
            
            1. **Atualizar Dados Históricos** (se necessário):
               - Executar extração de dados (Página 5) para incluir novos períodos
               - Histórico consolidado é atualizado automaticamente
            
            2. **Acessar Página 2 ou 3**:
               - Selecionar novos períodos históricos (incluindo os mais recentes)
                    - Manter ou ajustar sensibilidades, inflação e produtividade
            
            3. **Gerar Novo Forecast**:
               - Clicar em "Gerar Forecast" ou "Salvar Forecast"
               - Sistema recalcula com dados atualizados
            
            4. **Forecast Anterior é Substituído**:
               - `forecast_completo.parquet` é sobrescrito
               - Novo forecast reflete dados mais recentes
            
            **Resultado**: 
            - Forecast atualizado com dados mais recentes
            - Previsões mais acuradas baseadas em histórico expandido
            """)
        
        with st.expander("**Cenário 3: Testar Diferentes Cenários (What-If)**", expanded=False):
            st.markdown("""
            **Situação**: Quer testar impacto de diferentes volumes, inflações ou ganhos de produtividade
            
            **Passo a Passo**:
            
            1. **Acessar Página 2 (Simulador)**:
               - Configurar parâmetros base (sensibilidades, períodos históricos)
            
            2. **Testar Cenário 1**:
               - Ajustar volumes futuros (ex: +10%)
               - Visualizar impacto nos custos
               - **NÃO salvar** (apenas visualizar)
            
            3. **Testar Cenário 2**:
               - Ajustar volumes futuros (ex: -5%)
               - Visualizar impacto
               - Comparar com Cenário 1
            
            4. **Testar Diferentes Inflações**:
               - Alterar percentual de inflação
               - Ver impacto em todos os custos
               - Comparar cenários

                5. **Testar Ganhos de Produtividade**:
                    - Alterar produtividade global ou por Type 06
                    - Medir quanto do aumento monetário é compensado por eficiência
                    - Comparar com os cenários anteriores
            
                6. **Salvar Cenário Escolhido**:
               - Após decidir qual cenário usar
               - Configurar parâmetros finais
               - Salvar forecast
            
            **Resultado**: 
            - Múltiplos cenários testados sem salvar
            - Forecast final salvo com cenário escolhido
            """)
        
        with st.expander("**Cenário 4: Análise Detalhada de Forecast Gerado**", expanded=False):
            st.markdown("""
            **Situação**: Forecast já foi gerado, precisa de análises detalhadas
            
            **Passo a Passo**:
            
            1. **Acessar Página 3 (Análise)**:
               - Sistema carrega `forecast_completo.parquet` automaticamente
               - Mostra data de última atualização
            
            2. **Aplicar Filtros**:
               - Filtrar por Oficina, Veículo, Type 05, Type 06, etc.
               - Selecionar períodos específicos
            
            3. **Visualizar Gráficos**:
               - Gráficos de linha mostrando evolução
               - Gráficos de barras comparando períodos
               - Identificar tendências e padrões
            
            4. **Analisar Tabelas**:
               - Tabelas hierárquicas com drill-down
               - Detalhamento linha a linha
               - Identificar maiores custos previstos
            
            5. **Exportar para Excel**:
               - Exportar tabelas para análise externa
               - Compartilhar resultados com equipe
            
            **Resultado**: 
            - Análises detalhadas do forecast
            - Insights para tomada de decisão
            - Documentação dos resultados
            """)
        
        st.markdown("---")
        
        st.success("""
        **✅ Este capítulo descreve completamente a estrutura, atualização e funcionamento**
        das páginas de Best Estimate. Use estas informações para entender como o sistema
        organiza os dados, processa os forecasts e como cada página contribui para o processo.
        """)

# ==========================================
# SEÇÃO 6: TC CLOUD
# ==========================================
elif indice_selecionado == "☁️ TC Cloud":
    st.header("☁️ TC Cloud")

    st.markdown("""
    <div style="padding: 1.5rem; background: linear-gradient(135deg, #0f172a 0%, #1d4ed8 55%, #06b6d4 100%); border-radius: 10px; margin-bottom: 2rem; color: white;">
        <h2 style="color: white; margin: 0;">☁️ SCI no Databricks — Fonte de Verdade Operacional</h2>
        <p style="color: #e5eefc; margin: 0.5rem 0 0 0;">
            Consolidação da arquitetura cloud, tecnologias, fluxo do TC Veículos e regras para sincronização sem regressão.
        </p>
    </div>
    """, unsafe_allow_html=True)

    with st.expander("🌐 Ecossistema tecnológico do TC Cloud", expanded=True):
        st.markdown(
            """
            **Python** executa o app Streamlit, o processamento e as integrações do SCI.

            **GitHub** é a base de versionamento do código e o ponto de controle das mudanças feitas localmente.

            **Databricks Apps** hospeda a interface web em cloud para uso operacional.

            **Databricks Workspace Files** guarda o código publicado do app e os artefatos usados no runtime.

            **Databricks Jobs e notebooks** executam pipelines, validações e rotinas de suporte ao ambiente cloud.

            **Azure / infraestrutura corporativa** sustenta autenticação, rede e serviços de base do workspace Databricks.

            **TC-Cloud e espelhos locais** funcionam como camada de publicação e segurança operacional para sincronizar o que está estável no cloud sem perder o desenvolvimento local.
            """
        )

    with st.expander("🏗️ Arquitetura atual", expanded=False):
        st.markdown(
            """
            - `sci` separado de `sci_app`
            - `dados/`, `jobs/` e `workspace_publish/` fora do app
            - notebooks e pipeline isolados da interface
            - app Streamlit focado em leitura e navegação
            """
        )

    with st.expander("🚗 TC Veículos no cloud", expanded=False):
        st.markdown(
            """
            - pipeline Real e Budget alinhado
            - leitura por Parquet em Workspace Files
            - startup do app definindo `SCI_SHARED_DATA_ROOT`
            - validação operacional pelos notebooks principais
            """
        )

    with st.expander("🔁 Operação segura", expanded=False):
        st.markdown(
            """
            - pull do Databricks para `TC-Cloud`
            - propagação para raiz e espelhos locais
            - upload via SDK quando o CLI não estiver disponível
            - checklist anti-regressão no próprio produto
            """
        )

    _render_doc_tc_cloud()

# ==========================================
# SEÇÃO 7: APRESENTAÇÃO VISUAL
# ==========================================
elif indice_selecionado == "📊 Apresentação Visual":
    st.header("📊 Apresentação Visual")
    _render_doc_apresentacao_visual()

# ==========================================
# SEÇÃO 8: CHATBOT DE DOCUMENTAÇÃO
# ==========================================
elif indice_selecionado == "💬 Chatbot de Documentação":
    st.header("💬 Chatbot de Documentação")
    
    st.markdown("""
    <div style="padding: 1.5rem; background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%); border-radius: 10px; margin-bottom: 2rem; color: white;">
        <h2 style="color: white; margin: 0;">💬 Assistente Virtual de Documentação</h2>
        <p style="color: #f0f0f0; margin: 0.5rem 0 0 0;">
            Faça perguntas sobre o sistema e receba respostas baseadas na documentação completa
        </p>
    </div>
    """, unsafe_allow_html=True)
    
    # Importar chatbot
    try:
        # Adicionar diretório raiz ao path para importar chatbot
        import sys
        base_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        if base_path not in sys.path:
            sys.path.insert(0, base_path)
        from chatbot_documentacao import responder_pergunta
        
        # Inicializar histórico de conversa
        if 'historico_chat' not in st.session_state:
            st.session_state.historico_chat = []
        
        # Exibir histórico
        st.subheader("💬 Conversa")
        
        if st.session_state.historico_chat:
            for i, (pergunta, resposta, score) in enumerate(st.session_state.historico_chat):
                with st.expander(f"❓ {pergunta[:50]}...", expanded=False):
                    st.markdown(f"**Pergunta:** {pergunta}")
                    st.markdown(f"**Resposta:**")
                    st.markdown(resposta)
                    if score > 0:
                        st.caption(f"Relevância: {score:.0%}")
        else:
            st.info("💡 Faça sua primeira pergunta abaixo para começar!")
        
        st.markdown("---")
        
        # Campo de entrada
        st.subheader("📝 Faça uma Pergunta")
        
        pergunta = st.text_input(
            "Digite sua pergunta sobre o sistema:",
            placeholder="Ex: Como funciona o Best Estimate? O que é Flex Bud? Como processar dados?",
            key="input_pergunta"
        )
        
        col1, col2 = st.columns([1, 4])
        
        with col1:
            botao_perguntar = st.button("🔍 Buscar Resposta", type="primary", use_container_width=True)
        
        with col2:
            botao_limpar = st.button("🗑️ Limpar Histórico", use_container_width=True)
        
        if botao_limpar:
            st.session_state.historico_chat = []
            st.rerun()
        
        if botao_perguntar and pergunta:
            with st.spinner("🔍 Buscando na documentação..."):
                resultado = responder_pergunta(pergunta)
                
                if resultado['resposta']:
                    # Adicionar ao histórico
                    st.session_state.historico_chat.append((
                        pergunta,
                        resultado['resposta'],
                        resultado['score']
                    ))
                    
                    # Exibir resposta
                    st.success("✅ Resposta encontrada!")
                    st.markdown("**Resposta:**")
                    st.markdown(resultado['resposta'])
                    
                    if resultado['score'] > 0:
                        st.caption(f"📊 Relevância da resposta: {resultado['score']:.0%}")
                    
                    # Exibir segmentos adicionais se houver
                    if resultado['segmentos_encontrados']:
                        st.markdown("---")
                        st.subheader("📚 Informações Adicionais")
                        for i, segmento in enumerate(resultado['segmentos_encontrados'], 1):
                            with st.expander(f"Informação adicional {i}", expanded=False):
                                st.markdown(segmento)
                    
                    st.rerun()
        
        # Sugestões de perguntas
        st.markdown("---")
        st.subheader("💡 Perguntas Sugeridas")
        
        perguntas_sugeridas = [
            "O que é o Stellantis Cost Intelligence (SCI)?",
            "Como funciona o Best Estimate?",
            "O que é Flex Bud?",
            "Como funciona o rateio por veículo?",
            "Qual a diferença entre TC Ext e TC Veículos?",
            "Como funciona a sensibilidade no simulador?",
            "O que é CPU (Custo por Unidade)?",
            "Como funciona o Waterfall?",
        ]
        
        cols = st.columns(2)
        for i, pergunta_sugerida in enumerate(perguntas_sugeridas):
            with cols[i % 2]:
                if st.button(f"❓ {pergunta_sugerida}", key=f"sug_{i}", use_container_width=True):
                    # Processar pergunta sugerida diretamente
                    with st.spinner("🔍 Buscando na documentação..."):
                        resultado = responder_pergunta(pergunta_sugerida)
                        
                        if resultado['resposta']:
                            # Adicionar ao histórico
                            st.session_state.historico_chat.append((
                                pergunta_sugerida,
                                resultado['resposta'],
                                resultado['score']
                            ))
                            st.rerun()
        
    except ImportError as e:
        st.error(f"❌ Erro ao importar módulo de chatbot: {str(e)}")
        st.info("💡 Certifique-se de que o arquivo chatbot_documentacao.py existe na raiz do projeto.")
    except Exception as e:
        st.error(f"❌ Erro no chatbot: {str(e)}")
        import traceback
        st.code(traceback.format_exc())

# ==========================================
# SEÇÃO 8: SISTEMA DE ALERTAS
# ==========================================
elif indice_selecionado == "🔔 Sistema de Alertas":
    st.header("🔔 Sistema de Alertas")

    st.markdown(
        """
        <div style="padding: 1.5rem; background: linear-gradient(135deg, #ff6b6b 0%, #feca57 100%); border-radius: 10px; margin-bottom: 1.25rem; color: white;">
            <h2 style="color: white; margin: 0;">🔔 Central de Alertas (TC Veículos)</h2>
            <p style="color: #fff; opacity: 0.92; margin: 0.5rem 0 0 0;">
                Monitoramento automático de desvios relevantes no TC Veículos — com ranking consolidado e notificações.
            </p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    st.markdown(
        """
### ✅ Objetivo
Detectar rapidamente **anomalias / perdas** no custo do TC Veículos, priorizando o que mais impacta o resultado.

O sistema gera um **ranking hierárquico**:
**Type 05 → Type 06 → Account → Oficinas** (com texto breve quando disponível).

### 📍 Onde fica no app
No menu lateral do SCI existem duas páginas:
- **Central de Alertas → Monitoramento** (`alertas/alert_ui.py`)
- **Central de Alertas → Configuração de Alertas** (`alertas/alert_config_ui.py`)

### 🔎 O que o alerta compara
O motor suporta dois modos:
1. **Budget Flex × Real** *(principal)*
2. **Mês × Mês Anterior** *(secundário)*

> Observação: no modo **Budget Flex × Real**, o “esperado” vem do cálculo de Flex BUD detalhado (reuso da base do TC Veículos).
        """
    )

    with st.expander("🧠 Como o motor funciona (visão geral)", expanded=False):
        st.markdown(
            """
**Fonte de dados (TC Veículos):** parquets consolidados em `dados/TC_Principal/historico_consolidado/`.

**Etapas (alto nível):**
1. Carrega Real, Volume Real, Budget e Volume Budget
2. Calcula Flex BUD detalhado com dimensões (Oficina / Type 05 / Type 06 / Account)
3. Aplica filtros da regra (Oficina, Type 05, Type 06, Account)
4. Calcula **Real vs Esperado** e ranqueia os **Top N Type 06** com maior perda
5. Consolida em um **card único** (drill-down hierárquico)

**Severidade (padrão):** classificada por desvio percentual absoluto:
- **Crítico:** ≥ 15%
- **Moderado:** ≥ 5%
- **Informativo:** < 5%

**Base técnica principal:** `alertas/alert_engine.py`

**Funções mais importantes:**
- `calcular_ranking_consolidado()` — monta o card hierárquico único usado no monitoramento e nas notificações.
- `gerar_tabela_validacao()` — gera a conferência `Type 05 / Type 06 / Account / Flex BUD / Real / Delta / % Delta`.
- `evaluate_rule()` e `evaluate_all_rules()` — avaliam regras ativas e retornam a estrutura final do alerta.
            """
        )

    with st.expander("⚙️ Configuração de regras", expanded=False):
        st.markdown(
            """
Em **Configuração de Alertas**, é possível criar regras com:
- **Ano** e **modo de comparação**
- **Top N** (quantos Type 06 destacar)
- **Moeda** (BRL / EUR / USD)
- Filtros opcionais em cascata: **Type 05**, **Type 06**, **Account**
- Filtro opcional de **Oficinas** (vazio = todas)

Cada regra pode ser **ativada/desativada** e removida.

**Observação importante:** hoje o SCI trabalha principalmente com **disparo manual ou pós-processamento**. A regra continua armazenando metadados de agenda por compatibilidade, mas o fluxo operacional atual privilegia o acionamento quando a base já foi processada e está pronta para leitura.
            """
        )

    with st.expander("📊 Tabela de validação", expanded=False):
        st.markdown(
            """
Essa tabela existe para permitir **auditoria rápida do cálculo** antes ou depois do envio do alerta.

**Colunas principais:**
- `Type 05`
- `Type 06`
- `Account`
- `Flex BUD`
- `Flex BUD P`
- `Real`
- `Real - Flex BUD P`
- `% Delta`

**Como interpretar:**
- `Flex BUD` = valor esperado integral
- `Flex BUD P` = valor esperado proporcional ao período corrente
- `Real - Flex BUD P` = desvio monetário principal
- `% Delta` = desvio percentual relativo ao esperado proporcional

**Melhoria recente do sistema:**
- o preenchimento de `Type 05` passou a ser preservado a partir do `flex_detalhado` e, quando necessário, complementado por mapeamento controlado com base em `Type 06 + Account`, evitando linhas vazias na validação.

**Uso prático:**
- conferir se a perda identificada no card bate com a linha detalhada
- validar se a regra/filtro aplicado está trazendo o universo correto
- apoiar explicação de desvio antes de acionar stakeholders
            """
        )

    with st.expander("📨 Notificações (E-mail / Teams)", expanded=False):
        st.markdown(
            """
O sistema pode enviar o ranking consolidado para:
- **E-mail (Microsoft Graph / OAuth2)** — com autenticação MSAL
- **Microsoft Teams (Webhook)** — com card hierárquico formatado

Também existe a opção de manter apenas o uso **interno no app** (sem envio).

Na aba **🧪 Testar Envio**, dá para validar rapidamente se o Graph/Webhook estão corretos.

**Teams:**
- envia card consolidado com árvore visual por `Type 05 -> Type 06 -> Account -> Oficina`
- inclui barra textual de representatividade do desvio total
- mantém legibilidade próxima da visualização interna do SCI

**E-mail:**
- usa HTML estruturado com ranking consolidado
- pode incluir a tabela de validação junto do alerta
- autenticação moderna via Microsoft Graph quando configurada
            """
        )

    with st.expander("▶️ Fluxo operacional atual", expanded=False):
        st.markdown(
            """
O fluxo hoje foi simplificado para ficar mais aderente ao uso real do SCI:

- o usuário pode clicar em **Verificar agora** para gerar o ranking consolidado
- o usuário pode clicar em **Disparar alertas ativos** para forçar o envio manual
- o processamento de dados pode acionar a avaliação dos alertas ao final da atualização da base

**Por que isso é melhor:**
- reduz dependência de app aberto em um horário fixo
- garante que o alerta roda sobre a base mais recente
- deixa a operação mais previsível para fechamento e acompanhamento mensal
            """
        )

    with st.expander("🗂️ Persistência e auditoria", expanded=False):
        st.markdown(
            """
As configurações e o histórico ficam salvos em JSON no pacote `alertas/`:
- `alertas/alert_rules.json` — regras + canais de notificação + agenda
- `alertas/alert_log.json` — histórico de execuções (quando e o que foi enviado)

O histórico pode ser consultado na aba **📜 Histórico** da página de configuração.
            """
        )

    with st.expander("🧩 O que o usuário enxerga no monitoramento", expanded=False):
        st.markdown(
            """
Na página de monitoramento, o SCI mostra um **card consolidado** em vez de vários alertas soltos.

**Estrutura do card:**
- agrupamento por `Type 05`
- detalhamento por `Type 06`
- drill-down até `Account` e `Oficina`
- `Texto breve` em lowercase para leitura mais limpa
- barra com percentual de participação no **desvio total** do card

**Resultado para o usuário:**
- mais fácil priorizar o que mais pesa no problema
- mais fácil explicar o desvio para a operação
- mais fácil validar antes de enviar Teams/e-mail
            """
        )

# ==========================================
# ==========================================
# SEÇÃO: GUIA DE BUILD (EXE)
# ==========================================
elif indice_selecionado == "📦 Guia de Build (EXE)":
    st.header("📦 Guia de Build (EXE)")

    st.header("📦 Guia de Build — Empacotamento como Executável Windows")

    st.info(
        "Atualizado em 20/02/2026: este projeto gera EXE Windows usando `streamlit-desktop-app` "
        "(que internamente usa PyInstaller) e empacota o runtime em `_internal/`. "
        "Este passo a passo foi escrito para ser reusado em outro projeto e para uma LLM conseguir "
        "reconstruir o mesmo método com alta fidelidade."
    )

    st.markdown("""
    <div style="background: linear-gradient(135deg, #1a1a2e 0%, #16213e 50%, #0f3460 100%); padding: 20px; border-radius: 10px; margin-bottom: 16px; color: white;">
        <h2 style="color: white; margin: 0;">📦 SCI — Guia de Empacotamento (EXE)</h2>
        <p style="color: #a0c4ff; margin: 0.5rem 0 0 0;">Passo a passo oficial embutido nesta página, sem dependência externa.</p>
    </div>
    """, unsafe_allow_html=True)

    with st.expander("0) Visão geral (método)", expanded=False):
        st.markdown("""
        **Objetivo:** gerar um executável Windows do Streamlit que abre como *desktop app*, sem depender do repositório.

        **Método adotado (o mesmo padrão do projeto referência):**
        - `streamlit-desktop-app build app.py --name <NOME>`
        - Pós-build: copiar `dados/`, módulos/páginas e scripts `.py` avulsos para `dist/<NOME>/_internal/`

        **Observação importante (AgGrid / st_aggrid):**
        - As páginas do Streamlit (multipage) são carregadas em *runtime*.
        - Isso pode fazer o empacotador **não incluir automaticamente** dependências importadas apenas nessas páginas.
        - Solução robusta adotada no SCI: pós-build, copiar o pacote `st_aggrid` do `.venv` para dentro do `_internal/`.

        **Por que isso evita bugs no EXE:**
        - No executável, o caminho “real” do código empacotado é `sys._MEIPASS` (pasta `_internal/`).
        - Qualquer lógica de `sys.path` baseada em `dirname(__file__)` precisa considerar `sys._MEIPASS`.
        """)

    with st.expander("1) Pré-requisitos (ambiente)", expanded=False):
        st.markdown("""
        - Windows 10/11
        - Python (mesma versão usada no projeto, preferencialmente) + `venv`
        - Dependências do projeto instaladas (`pip install -r requirements.txt`)
        """)

    with st.expander("2) Bibliotecas e ferramentas usadas", expanded=False):
        st.markdown("""
        **Ferramenta de build (principal):**
        - `streamlit-desktop-app`

        **Empacotador (indireto):**
        - PyInstaller (chamado pela ferramenta)

        **Desktop container:**
        - `pywebview` (pasta `webview/` aparece no `_internal/`)

        **Framework:**
        - `streamlit`

        > Observação: no nosso caso, o build falha se existir BOM (U+FEFF) no começo do `app.py`.
        """)

    with st.expander("3) Passo crítico: remover BOM (U+FEFF) do app.py", expanded=False):
        st.markdown("""
        Se o build acusar:
        `SyntaxError: invalid non-printable character U+FEFF`

        Remova o BOM regravando em UTF-8 sem BOM.
        """)
        st.code(
            "$c = [System.IO.File]::ReadAllBytes('app.py'); "
            "if ($c.Length -ge 3 -and $c[0] -eq 0xEF -and $c[1] -eq 0xBB -and $c[2] -eq 0xBF) { "
            "  [System.IO.File]::WriteAllBytes('app.py', $c[3..($c.Length-1)]) ; "
            "  'OK: BOM removido' "
            "} else { 'OK: sem BOM' }",
            language="powershell",
        )

    with st.expander("4) Construção do build (comando oficial)", expanded=False):
        st.markdown("""
        Na raiz do projeto (mesma pasta do `app.py`), execute:
        """)
        st.code("streamlit-desktop-app build app.py --name Stellantis-Cost-Intelligence", language="powershell")
        st.markdown("""
        Depois disso, o diretório `dist/Stellantis-Cost-Intelligence/` deve existir.

        > Nota: o `streamlit-desktop-app` **não aceita** `--hidden-import` no CLI.
        > Para garantir dependências de páginas carregadas em runtime (ex.: `st_aggrid`), use o `build_exe.bat`.
        """)

    with st.expander("5) Pós-build obrigatório: copiar recursos para _internal/", expanded=False):
        st.markdown("""
        O runtime do EXE lê tudo de dentro de `dist/<NOME>/_internal/`.

        No SCI, nós copiamos para `_internal/`:
        - `dados/` (parquets, históricos)
        - `pages/`, `tc_core/`, `tc_principal/`, `tc_ext/`, `.streamlit/`
        - scripts `.py` avulsos que são importados em runtime (extração, exports, versionamento)
        - JSONs e imagens necessárias

        **AgGrid (streamlit-aggrid):**
        - Sintoma quando não incluído: `módulo 'st_aggrid' não encontrado` e o sistema entra em fallback.
        - Correção aplicada no SCI: copiar `st_aggrid/` e `streamlit_aggrid-*.dist-info/` do `.venv` para dentro do `_internal/`.

        Exemplo (PowerShell):
        """)
        st.code(
            "$dest = 'dist\\Stellantis-Cost-Intelligence\\_internal'\n"
            "Copy-Item '.venv\\Lib\\site-packages\\st_aggrid' -Destination ($dest + '\\st_aggrid') -Recurse -Force\n"
            "Copy-Item '.venv\\Lib\\site-packages\\streamlit_aggrid-*.dist-info' -Destination $dest -Recurse -Force\n",
            language="powershell",
        )
        st.markdown("""

        **Script oficial:** `build_exe.bat` (na raiz) automatiza isso.
        """)

    with st.expander("9) O que NÃO fazer (armadilha do .spec)", expanded=False):
        st.markdown("""
        Evite tentar rodar `pyinstaller` manualmente a partir do `.spec` gerado automaticamente pelo `streamlit-desktop-app`.

        **Por quê?** Esse `.spec` costuma referenciar um script temporário em `%TEMP%` (ex.: `tmp_xxx.py`).
        Depois do build, esse arquivo pode ser apagado, e o rebuild falha com:
        - `ERROR: script 'C:\\Users\\...\\AppData\\Local\\Temp\\tmp_XXXX.py' not found`

        **Solução adotada:** não rebuildar via `.spec`; em vez disso, fazer pós-build (cópias) para `_internal/`.
        """)

    with st.expander("6) Armadilha comum no EXE: sys.path e _MEIPASS", expanded=False):
        st.markdown("""
        **Sintoma:** no EXE, algumas telas funcionam, mas módulos avulsos (ex.: `processamento_dados_veiculos.py`) “somem”.

        **Causa:** páginas faziam `sys.path.insert(0, dirname(dirname(dirname(__file__))))`.
        No EXE isso aponta para a pasta do `.exe`, não para `_internal/`.

        **Correção padrão (reutilizável):**
        """)
        st.code(
            "import sys\n"
            "import os\n"
            "if hasattr(sys, '_MEIPASS'):\n"
            "    project_root = sys._MEIPASS\n"
            "else:\n"
            "    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))\n"
            "if project_root not in sys.path:\n"
            "    sys.path.insert(0, project_root)\n",
            language="python",
        )

    with st.expander("7) Validação (checklist)", expanded=False):
        st.markdown("""
        - Abrir: `dist\\Stellantis-Cost-Intelligence\\Stellantis-Cost-Intelligence.exe`
        - Confirmar que o app abre (janela desktop) e/ou responde em `http://localhost:8501`
        - Testar uma extração (Budget e Real) e confirmar geração dos parquets por veículo:
          - `df_veiculos_custo_fp.parquet`
          - `df_veiculos_cpu.parquet`
        """
        )

    with st.expander("8) Guia completo embutido", expanded=False):
        _render_doc_build_exe_completo()

# ==========================================
# ==========================================
# SEÇÃO: PRÓXIMOS PASSOS
# ==========================================
elif indice_selecionado == "🚀 Próximos Passos":
    st.header("🚀 Próximos Passos — SCI LATAM (Multiplantas)")

    st.markdown("""
    <div style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); padding: 20px; border-radius: 10px; margin-bottom: 20px;">
        <h2 style="color: white; margin: 0;">🚀 Próximos Passos — SCI LATAM (Multiplantas)</h2>
        <p style="color: rgba(255,255,255,0.9); margin-top: 8px;">Plano multiplantas: 5 fábricas, 2 países, evolução não-regressiva</p>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("""
    Este documento apresenta o **plano completo** para a evolução do **Stellantis Cost Intelligence (SCI)**
    rumo ao modelo **SCI LATAM (Multiplantas)** — cobrindo **5 plantas** em **2 países**:

    🇧🇷 **Brasil:** Betim · Porto Real · Goiana
    🇦🇷 **Argentina:** Córdoba · Palomar

    **Princípio fundamental:** toda evolução é **NÃO REGRESSIVA** — nenhum cálculo, regra ou output
    existente pode ser quebrado. Qualquer mudança é incremental e validada por gates de aprovação.

    > 📌 **Momento atual:** apenas PLANEJAMENTO + DOCUMENTAÇÃO. Nenhuma implementação será executada agora.

    ---

    #### 🔑 Por que um novo workspace (e não continuar no atual)?

    O repositório atual (`TC`) cresceu organicamente e acumulou **dívida estrutural significativa**:

    | Problema | Evidência |
    |----------|-----------|
    | **61 scripts .py soltos na raiz** | `_test_*`, `_verify_*`, `_diag_*`, `corrigir_*`, `debug_*`, `comparar_*` |
    | **19 arquivos temporários `_*.py`** | Testes ad-hoc, scripts de diagnóstico, backups manuais |
    | **16 markdowns + 10 HTMLs + 15 txts na raiz** | Documentação solta, logs de deploy, resultados de sync |
    | **7 planilhas Excel na raiz** | Dados brutos misturados com código |
    | **Notebooks triplicados** | Mesmos 7 notebooks em `notebooks/`, `Databricks/sci/notebooks/` e `Databricks/sci_app/notebooks/` |
    | **Código quadruplicado** | `processamento_dados*.py` repetido em 4 lugares (raiz, `Databricks/sci/`, `Databricks/sci_app/`, `Databricks/pulled_from_workspace/`) |
    | **3 cópias do app inteiro** | `Databricks/sci/`, `Databricks/sci_app/`, `Databricks/pulled_from_workspace/` — todas com `pages/`, `tc_core/`, `tc_ext/`, etc. |
    | **Módulos legacy + novos misturados** | `jobs/` vs `notebooks/` vs `Databricks/sci/notebooks/` fazem coisas parecidas |

    O SCI LATAM é a oportunidade de **começar limpo**, com estrutura planejada, sem arrastar essa dívida.
    Vamos **aproveitar 100% dos cálculos e regras de negócio** do projeto atual — mas organizados corretamente.
    """)

    # ------------------------------------------------------------------
    # 0) Visão SCI LATAM e princípio não-regressivo
    # ------------------------------------------------------------------
    with st.expander("🟦 0) Visão SCI LATAM e Princípio Não-Regressivo", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Contextualizar a evolução do SCI para o modelo multiplantas (SCI LATAM), definindo
        a regra de ouro que norteia todo o plano: **nada do que já funciona pode quebrar**.

        ---

        #### 📋 O que fazer
        1. Documentar formalmente o escopo SCI LATAM: 5 plantas (Betim, Porto Real, Goiana, Córdoba, Palomar)
        2. Registrar a regra de ouro: **evolução não-regressiva** — cálculos, regras e outputs atuais permanecem intactos
        3. Definir a arquitetura conceitual em camadas:
           - **Plant Layer** — segregação lógica por planta (dados, configs, processamento)
           - **Country Layer** — consolidação por país (BR, AR)
           - **Region Layer** — consolidação LATAM (baseada apenas em outputs oficiais)
        4. Mapear o que já existe (SCI atual = planta Betim) e o que será estendido
        5. Estabelecer vocabulário padronizado: plant_id, country_code, region_code

        ---

        #### 📦 Artefatos / Outputs esperados
        - Documento de visão SCI LATAM (1–2 páginas)
        - Diagrama de camadas (Plant → Country → Region)
        - Glossário de termos multiplantas
        - Mapa de funcionalidades existentes vs. novas

        ---

        #### ✅ Critérios de aceite
        - Documento aprovado por pelo menos 1 stakeholder de cada país
        - Diagrama de camadas validado pela equipe técnica
        - Glossário sem ambiguidades entre termos BR e AR

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Escopo cresce além das 5 plantas iniciais | Travar MVP em 5 plantas; novas plantas só após validação completa |
        | Termos diferentes entre BR e AR geram confusão | Glossário oficial + revisão cruzada |
        | Resistência organizacional | Envolver controllers locais desde o início |
        """)

    # ------------------------------------------------------------------
    # 1) O que faremos inicialmente — Novo Workspace Databricks
    # ------------------------------------------------------------------
    with st.expander("🟦 1) O que Faremos Inicialmente — Novo Workspace no Databricks", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Criar um **novo workspace/pasta** no Databricks para o SCI LATAM, **completamente separado**
        do projeto atual (`sci_app` / `sci`), aproveitando a oportunidade para corrigir todos os
        problemas estruturais acumulados.

        ---

        #### 📋 Por que separar e não continuar no workspace atual?

        O workspace atual tem **dívida estrutural grave** que não queremos arrastar:

        **1. Código quadruplicado no repositório local:**
        ```
        processamento_dados.py  → existe em 4 lugares:
          ├── raiz/
          ├── Databricks/sci/
          ├── Databricks/sci_app/
          └── Databricks/pulled_from_workspace/
        ```
        Qual é o "certo"? Ninguém sabe com certeza. No novo workspace, haverá **uma única fonte de verdade**.

        **2. Notebooks triplicados:**
        ```
        00_validar_ambiente_databricks.py → 3 cópias:
          ├── notebooks/
          ├── Databricks/sci/notebooks/
          └── Databricks/sci_app/notebooks/
        ```
        No SCI LATAM, cada notebook existirá em **um único lugar**.

        **3. Três cópias inteiras do app:**
        ```
        Databricks/
          ├── sci/          ← cópia do app + notebooks
          ├── sci_app/      ← outra cópia do app
          └── pulled_from_workspace/  ← mais uma cópia
        ```
        Todas com `pages/`, `tc_core/`, `tc_ext/`, `tc_principal/` duplicados.

        **4. Raiz do repositório poluída:**
        - 61 scripts Python soltos (testes ad-hoc, diagnósticos, correções pontuais)
        - 19 arquivos `_*.py` temporários
        - 16 markdowns, 10 HTMLs, 15 TXTs, 7 planilhas Excel — tudo misturado

        ---

        #### 📋 O que fazer
        1. **Criar novo workspace no Databricks:**
           ```
           /Workspace/Users/u235107@inetpsa.com/sci_latam/
           ```
        2. **Estrutura inicial planejada:**
           ```
           sci_latam/
           ├── app.py                     # Entry point Streamlit
           ├── app.yaml                   # Config do Databricks App
           ├── requirements.txt           # Deps únicas e limpas
           ├── config/
           │   ├── plants_master.json     # Cadastro das 5 plantas
           │   ├── settings_dev.json      # Config ambiente DEV
           │   └── settings_prod.json     # Config ambiente PROD
           ├── notebooks/                 # ÚNICO lugar para notebooks
           │   ├── 00_validar_ambiente.py
           │   ├── 01_criar_tabelas_delta.py
           │   ├── 02_carga_dados.py
           │   ├── 03_processar_publicar.py
           │   ├── 04_prevalidar.py
           │   └── 05_validacao_pos_job.py
           ├── src/
           │   ├── core/                  # Motor de cálculo (migrado do tc_core)
           │   │   ├── backend.py
           │   │   ├── io_delta.py
           │   │   ├── io_excel.py
           │   │   ├── transform.py
           │   │   └── constants.py
           │   ├── finance/               # Regras financeiras (FP, FA, CPU, BE, Flex)
           │   ├── plants/                # Plant Layer — config e lógica por planta
           │   ├── consolidation/         # Consolidação BR e LATAM
           │   └── exports/               # Geração de relatórios e exports
           ├── pages/                     # Páginas Streamlit (UI)
           ├── tests/                     # Testes organizados
           │   ├── test_core/
           │   ├── test_finance/
           │   ├── test_plants/
           │   └── test_regression/       # Golden tests de regressão
           ├── scripts/                   # Automação (sync, deploy, CI)
           └── docs/                      # Documentação vive aqui
           ```
        3. **Não copiar lixo:** nenhum `_test_*`, `_verify_*`, `_diag_*`, `corrigir_*`, `debug_*` solto
        4. **Nenhuma duplicação:** cada arquivo existe em exatamente 1 lugar

        ---

        #### 📦 Artefatos / Outputs esperados
        - Workspace `sci_latam/` criado no Databricks com estrutura vazia
        - Repo GitHub `SCI-LATAM` (ou branch `sci-latam` no repo atual) com a mesma estrutura
        - `README.md` documentando a estrutura e as convenções
        - `.gitignore` limpo (sem `dist/`, `build/`, `__pycache__/`, `*.xlsx`, `*.html` temporários)

        ---

        #### ✅ Critérios de aceite
        - Workspace acessível no Databricks
        - Estrutura de pastas criada conforme o diagrama acima
        - Zero arquivos duplicados
        - `README.md` com instruções de setup completas
        - O repositório antigo (`TC`) continua funcionando normalmente (não tocamos nele)

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Esquecer de migrar algum módulo essencial | Checklist de migração vs. inventário do repo atual |
        | Novo workspace vazio causa confusão | README claro + comunicação com a equipe |
        | Repo antigo continua recebendo commits por hábito | Comunicar que novo desenvolvimento vai no SCI LATAM |
        """)

    # ------------------------------------------------------------------
    # 2) Reaproveitamento de notebooks e módulos existentes
    # ------------------------------------------------------------------
    with st.expander("🟦 2) Reaproveitamento de Notebooks e Módulos Existentes", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Migrar **seletivamente** os notebooks, módulos e regras de negócio do projeto atual
        para o novo workspace, sem arrastar código duplicado, scripts temporários ou dívida técnica.

        ---

        #### 📋 Inventário do que será reaproveitado

        **Notebooks (origem: `notebooks/` — fonte de verdade):**

        | Notebook atual | Destino no SCI LATAM | Ação |
        |---------------|----------------------|------|
        | `00_validar_ambiente_databricks.py` | `notebooks/00_validar_ambiente.py` | Migrar + parametrizar por plant_id |
        | `01_criar_tabelas_delta.py` | `notebooks/01_criar_tabelas_delta.py` | Migrar + schemas multi-planta |
        | `02_carga_snowflake.py` | `notebooks/02_carga_dados.py` | Migrar + adapter por fonte de dados |
        | `03_processar_e_publicar_delta.py` | `notebooks/03_processar_publicar.py` | Migrar + plant_id como parâmetro |
        | `04_prevalidar_excel.py` | `notebooks/04_prevalidar.py` | Migrar como está |
        | `05_validacao_pos_job.py` | `notebooks/05_validacao_pos_job.py` | Migrar + validação multi-planta |
        | `06_ui_consulta_workspace.py` | Absorvido pela UI principal | Não migrar como notebook separado |

        **Módulos core (origem: `tc_core/`):**

        | Módulo atual | Destino | Ação |
        |-------------|---------|------|
        | `tc_core/constants.py` | `src/core/constants.py` | Migrar + adicionar constantes multi-planta |
        | `tc_core/data_source.py` | `src/core/data_source.py` | Migrar + parametrizar paths por plant_id |
        | `tc_core/data_router.py` | `src/core/data_router.py` | Migrar + rotas por planta |
        | `tc_core/finance/` | `src/finance/` | Migrar integralmente (regras de cálculo são o ativo principal) |
        | `tc_core/ui/` | `src/ui/` ou `pages/` | Avaliar; pode ser simplificado |
        | `tc_core/utils/` | `src/core/utils/` | Migrar funções realmente usadas |
        | `src/sci_core/` | Merge com `src/core/` | Unificar; eliminar duplicação `tc_core` vs `src/sci_core` |

        **Módulos de domínio:**

        | Módulo atual | Destino | Ação |
        |-------------|---------|------|
        | `tc_principal/` | `pages/` (componentes UI) | Migrar componentes úteis, descartar acoplamento |
        | `tc_ext/` | `src/finance/ext/` | Migrar lógica de cálculo, reorganizar |
        | `tc_copilot/` | Não migrar agora | Funcionalidade futura; manter referência no repo antigo |
        | `processamento_dados*.py` (raiz) | `src/core/processing.py` | Consolidar 4 variantes em 1 módulo parametrizado |

        **O que NÃO migrar:**
        - 19 scripts `_*.py` temporários
        - Scripts de correção pontual (`corrigir_*.py`, `limpar_*.py`, `remover_*.py`)
        - Scripts de debug (`debug_*.py`, `diag_*.py`, `diagnostico_*.py`)
        - Planilhas Excel da raiz (dados vão para Delta tables)
        - HTMLs de teste, TXTs de log
        - `Databricks/pulled_from_workspace/` (cópia morta)
        - `Databricks/sci_app/` (será substituído pelo novo workspace)
        - `dist/`, `build/` (artefatos de build)

        ---

        #### 📋 O que fazer — passo a passo
        1. Gerar inventário completo do repo atual (script automático)
        2. Classificar cada arquivo: **migrar / descartar / referência futura**
        3. Para cada arquivo a migrar:
           - Copiar para o destino no SCI LATAM
           - Adaptar imports e paths
           - Adicionar parâmetro `plant_id` onde necessário
           - Rodar teste básico de importação
        4. Criar testes de smoke: cada módulo migrado importa sem erro
        5. Documentar o mapeamento origem → destino

        ---

        #### 📦 Artefatos / Outputs esperados
        - Planilha de inventário: arquivo → classificação → destino
        - Módulos migrados e funcionais no novo workspace
        - Testes de smoke passando
        - Documento de mapeamento origem → destino

        ---

        #### ✅ Critérios de aceite
        - Todo módulo migrado importa sem erro no Databricks
        - Zero duplicação: cada função/classe existe em exatamente 1 arquivo
        - Nenhum arquivo do tipo `_test_*`, `corrigir_*`, `debug_*` no novo repo
        - Planilha de inventário 100% preenchida

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Esquecer função importante que estava em script solto | Inventário automático + grep por funções chamadas |
        | Imports quebrados após reorganização | Teste de smoke obrigatório para cada módulo migrado |
        | Resistência a "jogar fora" código | Nada é jogado fora — repo antigo continua intacto como referência |
        """)

    # ------------------------------------------------------------------
    # 3) Passo a passo inicial no Databricks
    # ------------------------------------------------------------------
    with st.expander("🟦 3) Passo a Passo Inicial no Databricks", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Detalhar exatamente o que fazer no Databricks para criar o novo workspace SCI LATAM,
        configurar o ambiente e preparar tudo para receber os módulos migrados.

        ---

        #### 📋 O que fazer — sequência exata

        **Etapa 3.1 — Criar pasta no Workspace**
        ```
        databricks workspace mkdirs /Workspace/Users/u235107@inetpsa.com/sci_latam
        ```
        Subpastas:
        ```
        databricks workspace mkdirs /Workspace/Users/u235107@inetpsa.com/sci_latam/notebooks
        databricks workspace mkdirs /Workspace/Users/u235107@inetpsa.com/sci_latam/config
        databricks workspace mkdirs /Workspace/Users/u235107@inetpsa.com/sci_latam/src
        databricks workspace mkdirs /Workspace/Users/u235107@inetpsa.com/sci_latam/pages
        databricks workspace mkdirs /Workspace/Users/u235107@inetpsa.com/sci_latam/tests
        databricks workspace mkdirs /Workspace/Users/u235107@inetpsa.com/sci_latam/scripts
        ```

        **Etapa 3.2 — Configurar o Databricks App (se aplicável)**
        - Criar `app.yaml` com referência ao novo path
        - Configurar variáveis de ambiente: `SCI_ENV=dev`, `SCI_PLANT=BTM`
        - Definir cluster / serverless compute

        **Etapa 3.3 — Criar o primeiro notebook de smoke test**
        ```python
        # notebooks/00_validar_ambiente.py
        # Databricks notebook source
        import sys
        print(f"Python: {sys.version}")
        print(f"Workspace: sci_latam")
        print(f"Status: OK")

        # Verificar acesso ao catalog/schema
        spark.sql("SHOW DATABASES").display()
        ```

        **Etapa 3.4 — Configurar Delta tables (schema)**
        - Criar catalog/schema dedicado:
          ```sql
          CREATE SCHEMA IF NOT EXISTS sci_latam_dev;
          CREATE SCHEMA IF NOT EXISTS sci_latam_prod;
          ```
        - Definir tabelas Delta por planta:
          ```sql
          CREATE TABLE IF NOT EXISTS sci_latam_dev.raw_data_{plant_id} (...)
          CREATE TABLE IF NOT EXISTS sci_latam_dev.processed_{plant_id} (...)
          ```

        **Etapa 3.5 — Configurar secrets**
        - Criar scope: `databricks secrets create-scope sci-latam`
        - Registrar secrets necessários (Snowflake, APIs, etc.)

        **Etapa 3.6 — Primeiro sync local → Databricks**
        ```powershell
        # Do repositório local para o workspace
        databricks workspace import-dir ./sci_latam /Workspace/Users/u235107@inetpsa.com/sci_latam --overwrite
        ```

        ---

        #### 📦 Artefatos / Outputs esperados
        - Workspace `sci_latam/` ativo no Databricks
        - Notebook 00 executando com sucesso
        - Schemas `sci_latam_dev` e `sci_latam_prod` criados
        - Secret scope configurado
        - Primeiro sync local → Databricks funcionando

        ---

        #### ✅ Critérios de aceite
        - `databricks workspace ls /Workspace/Users/u235107@inetpsa.com/sci_latam` mostra a estrutura
        - Notebook 00 roda sem erro
        - `SHOW TABLES IN sci_latam_dev` retorna resultado (mesmo que vazio)
        - Sync local → Databricks funciona sem erro

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Permissões insuficientes para criar schema | Solicitar ao admin do Databricks antes de começar |
        | Nome de workspace conflita com existente | Usar nome único `sci_latam` (não `sci` nem `sci_app`) |
        | Secret scope com nome duplicado | Verificar scopes existentes antes de criar |
        """)

    # ------------------------------------------------------------------
    # 4) Sincronização com GitHub e GitHub Copilot CLI
    # ------------------------------------------------------------------
    with st.expander("🟦 4) Sincronização com GitHub e GitHub Copilot CLI", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Definir como o novo workspace SCI LATAM será sincronizado com GitHub,
        e como aproveitar o **GitHub CLI** e o **GitHub Copilot CLI** para produtividade.

        ---

        #### 📋 Estratégia de repositório GitHub

        **Opção A — Novo repo (recomendada):**
        ```
        github.com/[org]/SCI-LATAM
        ```
        ✅ Histórico limpo, sem 61 scripts soltos, sem pastas duplicadas
        ✅ CI/CD configurado do zero, sem herança de configs antigas
        ✅ Repo antigo (`TC`) continua como referência read-only

        **Opção B — Branch no repo atual:**
        ```
        github.com/[org]/TC  →  branch: sci-latam
        ```
        ⚠️ Mantém histórico, mas arrastra toda a sujeira do `main`

        ---

        #### 📋 Sincronização Databricks ↔ GitHub

        **Fluxo principal:**
        ```
        Desenvolvedor (VS Code)
            ↓  git push
        GitHub (SCI-LATAM)
            ↓  sync script (PowerShell / GitHub Actions)
        Databricks Workspace (sci_latam/)
        ```

        **Script de sync (evolução do atual `sync_databricks_app.ps1`):**
        ```powershell
        # scripts/sync_to_databricks.ps1
        param(
            [string]$Env = "dev",        # dev ou prod
            [switch]$Watch,              # modo watch (auto-sync)
            [switch]$DryRun              # apenas mostrar o que seria feito
        )
        $LocalPath = "./sci_latam"
        $WorkspacePath = "/Workspace/Users/u235107@inetpsa.com/sci_latam"

        databricks workspace import-dir $LocalPath $WorkspacePath --overwrite
        ```

        **Automação via GitHub Actions:**
        ```yaml
        # .github/workflows/sync-databricks.yml
        on:
          push:
            branches: [main, develop]
        jobs:
          sync:
            runs-on: ubuntu-latest
            steps:
              - uses: actions/checkout@v4
              - name: Sync to Databricks
                run: databricks workspace import-dir ...
        ```

        ---

        #### 📋 GitHub CLI — comandos essenciais

        **Instalação e setup:**
        ```powershell
        winget install GitHub.cli
        gh auth login
        gh repo create SCI-LATAM --private --description "SCI LATAM - Multiplantas"
        ```

        **Workflow diário:**
        ```powershell
        # Criar feature branch
        gh issue create --title "Implementar Plant Layer" --body "..."
        git checkout -b feature/plant-layer

        # Ao terminar
        git push -u origin feature/plant-layer
        gh pr create --title "feat: Plant Layer" --body "Closes #1"

        # Review e merge
        gh pr merge --squash
        ```

        ---

        #### 📋 GitHub Copilot CLI — produtividade no terminal

        **Instalação:**
        ```powershell
        gh extension install github/gh-copilot
        ```

        **Uso prático no SCI LATAM:**
        ```powershell
        # Explicar comando complexo
        gh copilot explain "databricks workspace import-dir ./sci_latam /Workspace/... --overwrite"

        # Sugerir comando
        gh copilot suggest "sync local folder to databricks workspace"

        # Ajuda com git
        gh copilot suggest "squash last 3 commits into one"
        ```

        **Integração com VS Code:**
        - GitHub Copilot Chat no VS Code já conhece o contexto do projeto
        - Usar `@workspace` para perguntas sobre a estrutura
        - Usar `@terminal` para sugestões de comandos

        ---

        #### 📦 Artefatos / Outputs esperados
        - Repositório GitHub criado e configurado
        - Script de sync `scripts/sync_to_databricks.ps1`
        - GitHub Actions workflow para sync automático
        - `.github/PULL_REQUEST_TEMPLATE.md` com checklist
        - Documentação de workflow (branch strategy, PRs, releases)

        ---

        #### ✅ Critérios de aceite
        - `gh repo view` mostra o repo acessível
        - Script de sync funciona: local → Databricks em 1 comando
        - GitHub Actions roda em cada push (pelo menos lint + sync)
        - Copilot CLI instalado e funcional no terminal
        - PR template com checklist de promoção

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Token do Databricks expira e sync falha silenciosamente | Health check no workflow + alerta |
        | Conflito de sync: Databricks editado diretamente sem PR | Política: nunca editar direto no workspace; sempre via repo |
        | GitHub Actions sem acesso à rede corporativa | Runner self-hosted ou sync manual via script local |
        """)

    # ------------------------------------------------------------------
    # 5) Cadastro e Configuração de Plantas (Plant Layer)
    # ------------------------------------------------------------------
    with st.expander("🟦 5) Cadastro e Configuração de Plantas (Plant Layer)", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Criar a estrutura de dados que permite ao SCI reconhecer e segregar informações por planta,
        sem alterar o processamento existente de Betim.

        ---

        #### 📋 O que fazer
        1. Criar tabela-mestra de plantas (`plants_master`):
           - `plant_id` (ex.: BTM, PTR, GOI, CBA, PLM)
           - `plant_name` (nome completo)
           - `country_code` (BR, AR)
           - `region_code` (LATAM)
           - `currency` (BRL, ARS)
           - `timezone`
           - `status` (ativo/inativo)
        2. Configurar parâmetros por planta:
           - Paths de dados de entrada (parquets, CSVs)
           - Paths de outputs
           - Regras de câmbio (se aplicável)
           - Calendário fiscal específico
        3. Garantir retrocompatibilidade: Betim = plant_id BTM, comportamento idêntico ao atual
        4. Criar mecanismo de fallback: se plant_id não informado, assume BTM (backward compatible)

        ---

        #### 📦 Artefatos / Outputs esperados
        - Arquivo `plants_master.json` ou tabela equivalente
        - Schema documentado (tipos, obrigatoriedades, exemplos)
        - Config por planta (`plant_config_{id}.json` ou seção no config central)
        - Script de seed/inicialização com as 5 plantas

        ---

        #### ✅ Critérios de aceite
        - Todas as 5 plantas cadastradas e carregáveis
        - Betim funciona exatamente como antes (teste de regressão)
        - Config de cada planta contém todos os campos obrigatórios
        - Teste unitário: carregar planta por ID retorna dados corretos

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Hardcode de Betim espalhado pelo código | Auditoria de grep + refactor incremental |
        | Plantas AR com calendário fiscal diferente | Campo de calendário fiscal no config |
        | Moedas diferentes geram confusão nos totais | Consolidação sempre em moeda base (BRL) com taxa explícita |
        """)

    # ------------------------------------------------------------------
    # 6) Contratos, ingestão e validações por planta
    # ------------------------------------------------------------------
    with st.expander("🟦 6) Contratos, Ingestão e Validações por Planta", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Definir como os dados de cada planta entram no SCI, quais validações são aplicadas,
        e como garantir que dados corrompidos não contaminem o processamento.

        ---

        #### 📋 O que fazer
        1. Definir o **contrato de dados** por planta:
           - Colunas obrigatórias (schema)
           - Tipos esperados (int, float, string, date)
           - Ranges válidos (ex.: CPU > 0, volumes >= 0)
           - Encoding e formato (UTF-8, parquet, CSV)
        2. Criar pipeline de ingestão parametrizado por `plant_id`:
           - Leitura do path configurado
           - Validação contra o schema
           - Quarentena de registros inválidos (log + isolamento)
           - Stamp de metadados: plant_id, data_ingestão, versão_schema
        3. Reaproveitar o pipeline atual (Betim) como template base
        4. Adicionar coluna `plant_id` em todos os DataFrames processados

        ---

        #### 📦 Artefatos / Outputs esperados
        - Schema de contrato por planta (JSON Schema ou Pydantic model)
        - Script de validação de ingestão
        - Log de quarentena (registros rejeitados + motivo)
        - Dashboard ou relatório de qualidade de ingestão

        ---

        #### ✅ Critérios de aceite
        - Dados de Betim passam 100% na validação (regressão)
        - Dados inválidos sintéticos são corretamente rejeitados
        - Coluna `plant_id` presente em todos os DataFrames pós-ingestão
        - Log de quarentena funcional e consultável

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Plantas AR enviam dados em formato diferente | Schema por planta com adapter de normalização |
        | Volume de dados de 5 plantas causa lentidão | Processamento paralelo por planta + cache |
        | Dados históricos sem plant_id | Migration script: dados antigos recebem plant_id=BTM |
        """)

    # ------------------------------------------------------------------
    # 7) Processamento core unificado (sem fork)
    # ------------------------------------------------------------------
    with st.expander("🟦 7) Processamento Core Unificado (Sem Fork)", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Garantir que o motor de cálculo existente (FP, FA, Redis, CPU, BE, Flex, Rateios)
        seja reutilizado por todas as plantas **sem duplicação de código**.

        ---

        #### 📋 O que fazer
        1. Auditar o código de processamento atual e identificar dependências hardcoded de Betim
        2. Parametrizar o processamento:
           - Todas as funções de cálculo recebem `plant_id` como parâmetro
           - Configs carregados dinamicamente a partir do `plant_config`
           - Paths de entrada/saída resolvidos via config (não hardcoded)
        3. Manter **um único motor de cálculo** — sem fork, sem cópia, sem branch por planta
        4. Criar test harness: executar processamento de Betim antes e depois da parametrização,
           outputs devem ser idênticos (diff zero)
        5. Documentar quais funções foram parametrizadas e quais já eram genéricas

        ---

        #### 📦 Artefatos / Outputs esperados
        - Relatório de auditoria: funções hardcoded vs. parametrizadas
        - Código refatorado com `plant_id` como parâmetro
        - Test harness de regressão (comparação de outputs)
        - Documentação técnica atualizada

        ---

        #### ✅ Critérios de aceite
        - Processamento de Betim gera outputs **idênticos** ao baseline (diff zero)
        - Processamento de uma planta nova (ex.: Porto Real) executa sem erro com dados de teste
        - Zero duplicação de lógica de cálculo entre plantas
        - Cobertura de testes: todas as funções parametrizadas possuem teste

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Refactor quebra cálculos existentes | Test harness obrigatório antes de qualquer merge |
        | Plantas AR têm regras fiscais diferentes | Regras de negócio em config por planta, não no código |
        | Tentação de fazer fork "temporário" | Code review obrigatório; proibido aprovar fork |
        """)

    # ------------------------------------------------------------------
    # 8) UI/Exports por planta + consolidado
    # ------------------------------------------------------------------
    with st.expander("🟦 8) UI / Exports por Planta + Consolidado", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Permitir que a interface (Streamlit) e os exports (Excel, parquet) operem tanto
        por planta individual quanto de forma consolidada (BR, LATAM).

        ---

        #### 📋 O que fazer
        1. Adicionar seletor de planta na UI (sidebar ou filtro global):
           - Lista dinâmica a partir de `plants_master`
           - Opção "Todas" para visão consolidada
           - Opção por país (BR, AR) e por região (LATAM)
        2. Filtrar todos os DataFrames exibidos pelo `plant_id` selecionado
        3. Gráficos e tabelas: respeitar o filtro de planta ativo
        4. Exports (Excel, parquet): incluir coluna `plant_id` e permitir export filtrado ou completo
        5. Manter a experiência atual como default: ao abrir, mostra Betim (ou "Todas" conforme decisão)
        6. Labels e títulos dinâmicos: exibir nome da planta nos headers

        ---

        #### 📦 Artefatos / Outputs esperados
        - Componente de seletor de planta (reutilizável)
        - Lógica de filtragem centralizada (não espalhada por cada página)
        - Templates de export com plant_id
        - Screenshots/mockups da UI com filtro ativo

        ---

        #### ✅ Critérios de aceite
        - Seletor de planta funcional em todas as páginas principais
        - Filtro por Betim mostra dados idênticos ao sistema atual
        - Export Excel inclui coluna plant_id
        - Consolidado (Todas) soma corretamente dados de múltiplas plantas

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | UI fica lenta com dados de 5 plantas | Cache por planta + lazy loading |
        | Usuário confunde visão consolidada com individual | Badge/indicador visual claro da planta ativa |
        | Métricas consolidadas somam moedas diferentes | Conversão para moeda base antes de consolidar |
        """)

    # ------------------------------------------------------------------
    # 9) Governança e trilha de auditoria
    # ------------------------------------------------------------------
    with st.expander("🟦 9) Governança e Trilha de Auditoria", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Garantir rastreabilidade completa: quem processou, quando, qual planta, quais dados,
        qual versão do sistema, e quais outputs foram gerados.

        ---

        #### 📋 O que fazer
        1. Implementar log de processamento por execução:
           - Timestamp, plant_id, versão do SCI, usuário (se aplicável)
           - Contagem de registros de entrada e saída
           - Hash dos arquivos de entrada (integridade)
           - Status: sucesso, falha parcial, erro
        2. Implementar log de auditoria de dados:
           - Registro de qualquer alteração manual (rateios_manuais, ajustes)
           - Quem alterou, quando, valor anterior, valor novo
        3. Manter histórico de versões de outputs:
           - Cada processamento gera outputs versionados (não sobrescreve)
           - Possibilidade de comparar versões
        4. Dashboard de governança: visão consolidada dos logs por planta

        ---

        #### 📦 Artefatos / Outputs esperados
        - Tabela de logs de processamento (`processing_log`)
        - Tabela de auditoria de dados (`audit_trail`)
        - Script de consulta de histórico
        - Dashboard de governança (nova aba ou seção na Documentação)

        ---

        #### ✅ Critérios de aceite
        - Todo processamento gera entrada no log (sem exceção)
        - Alterações manuais são rastreáveis (quem, quando, o quê)
        - Outputs versionados: possível recuperar resultado de qualquer execução anterior
        - Dashboard exibe logs das últimas N execuções por planta

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Logs crescem demais e ocupam espaço | Política de retenção (ex.: 12 meses) + compactação |
        | Usuários esquecem de registrar alterações manuais | Tornar registro automático via sistema (não manual) |
        | Logs sem padrão dificultam consulta | Schema fixo e validado para todos os logs |
        """)

    # ------------------------------------------------------------------
    # 10) Consolidação BR e consolidação LATAM
    # ------------------------------------------------------------------
    with st.expander("🟦 10) Consolidação BR e Consolidação LATAM", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Criar as visões consolidadas por país (BR) e por região (LATAM), baseadas
        **exclusivamente** em outputs oficiais das plantas — sem reprocessamento.

        ---

        #### 📋 O que fazer
        1. Definir quais métricas/outputs são consolidáveis:
           - CPU total, FP total, FA total por oficina/modelo
           - Volumes totais
           - Budget × Real × BE agregados
        2. Implementar consolidação BR:
           - Soma/agregação dos outputs de Betim + Porto Real + Goiana
           - Moeda única: BRL (sem conversão necessária)
        3. Implementar consolidação LATAM:
           - Soma/agregação de BR + AR
           - Conversão de ARS → BRL (ou USD) com taxa configurável
           - Taxa de câmbio versionada e auditável
        4. Regra: consolidação é **read-only** sobre outputs — nunca altera dados das plantas
        5. Gerar outputs consolidados separados (BR_consolidado, LATAM_consolidado)

        ---

        #### 📦 Artefatos / Outputs esperados
        - Módulo de consolidação (`consolidation.py` ou equivalente)
        - Tabela de taxas de câmbio (`exchange_rates`)
        - Outputs: `outputs/BR_consolidado/`, `outputs/LATAM_consolidado/`
        - Relatório de consolidação com breakdown por planta

        ---

        #### ✅ Critérios de aceite
        - Consolidação BR de Betim sozinho = output atual (regressão)
        - Consolidação LATAM com dados de teste de todas as 5 plantas gera resultado correto
        - Taxa de câmbio aplicada corretamente e registrada no log
        - Consolidação não altera outputs individuais das plantas

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Taxa de câmbio desatualizada distorce consolidação | Taxa versionada + alerta se idade > 30 dias |
        | Nem todas as plantas têm dados no mesmo período | Validação de completude antes de consolidar |
        | Dupla contagem por erro de agregação | Testes de reconciliação: soma das partes = total |
        """)

    # ------------------------------------------------------------------
    # 11) Comparativos entre plantas
    # ------------------------------------------------------------------
    with st.expander("🟦 11) Comparativos entre Plantas", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Permitir análises comparativas de performance entre as 5 plantas,
        gerando benchmarks, rankings e deltas.

        ---

        #### 📋 O que fazer
        1. Definir métricas comparáveis entre plantas:
           - CPU/unidade (normalizado por volume)
           - % de desvio Budget × Real
           - FP e FA por oficina (normalizados)
           - Mix de modelos
        2. Criar visão de ranking:
           - Melhor/pior planta por métrica
           - Tendência mensal por planta
        3. Criar visão de delta:
           - Planta A vs. Planta B (selecionável)
           - Destaque de maiores diferenças
        4. Normalização: garantir que comparações sejam justas (moeda, volume, mix)
        5. Visualizações: gráficos de radar, barras comparativas, heatmaps

        ---

        #### 📦 Artefatos / Outputs esperados
        - Módulo de comparativos (`plant_comparison.py`)
        - Página/aba de comparativos na UI
        - Templates de gráficos comparativos
        - Export de ranking em Excel

        ---

        #### ✅ Critérios de aceite
        - Comparativo entre Betim e Betim (mesma planta) mostra delta zero
        - Normalização por volume funciona corretamente
        - Rankings ordenam corretamente (asc e desc)
        - Gráficos renderizam sem erro com 2, 3, 4 ou 5 plantas

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Comparação injusta por mix de modelos diferente | Normalização por volume e mix |
        | Plantas com dados incompletos distorcem ranking | Indicador de completude visível |
        | Sensibilidade política de rankings | Acesso configurável por perfil/permissão |
        """)

    # ------------------------------------------------------------------
    # 12) Ambiente DEV — setup e convenções
    # ------------------------------------------------------------------
    with st.expander("🟦 12) Ambiente DEV — Setup e Convenções", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Definir o ambiente de desenvolvimento (DEV) como espaço de evolução contínua
        e experimentação controlada, separado logicamente de PROD.

        ---

        #### 📋 O que fazer
        1. Definir a separação DEV × PROD:
           - **Paths de dados:** `data/dev/` vs `data/prod/`
           - **Configurações:** `config_dev.json` vs `config_prod.json`
           - **Variável de ambiente:** `SCI_ENV=dev` ou `SCI_ENV=prod`
        2. Branch strategy:
           - `main` = PROD (protegido, merge via PR)
           - `develop` = DEV (desenvolvimento ativo)
           - Feature branches: `feature/plant-layer`, `feature/consolidation`, etc.
        3. Definir convenções de código:
           - Naming: snake_case, prefixos por módulo
           - Commits: conventional commits (feat:, fix:, docs:)
           - Testes obrigatórios para novas funcionalidades
        4. Pipeline de CI em DEV:
           - Lint (flake8/ruff)
           - Testes unitários
           - Validação de schemas

        ---

        #### 📦 Artefatos / Outputs esperados
        - Documento de convenções DEV (`CONTRIBUTING.md`)
        - Arquivo de configuração DEV (`config_dev.json`)
        - Script de setup do ambiente DEV
        - Template de feature branch

        ---

        #### ✅ Critérios de aceite
        - Ambiente DEV funciona isoladamente (sem impactar PROD)
        - Mudanças em DEV não aparecem em PROD até promoção explícita
        - Pipeline CI roda em cada push para develop
        - Novo desenvolvedor consegue configurar DEV em até 30 minutos seguindo o guia

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Dev usa dados de PROD acidentalmente | Variável SCI_ENV obrigatória + validação no boot |
        | Branches divergem demais e merge fica complexo | Merges frequentes develop → feature e vice-versa |
        | Falta de testes em DEV contamina PROD | Gate de cobertura mínima para aprovar PR |
        """)

    # ------------------------------------------------------------------
    # 13) Ambiente PROD — estabilidade operacional
    # ------------------------------------------------------------------
    with st.expander("🟦 13) Ambiente PROD — Estabilidade Operacional", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Garantir que o sistema PROD (em operação, "sempre rodando") tenha estabilidade,
        monitoramento e processos de deploy controlados.

        ---

        #### 📋 O que fazer
        1. Definir política de PROD:
           - Apenas código aprovado via PR + review entra em PROD
           - Deploy controlado (nunca direto na main)
           - Rollback documentado e testado
        2. Monitoramento básico:
           - Health check: sistema responde? Dados carregam?
           - Alerta se processamento falha
           - Log de acessos e erros
        3. Versionamento:
           - Tags semânticas (v2.1.0, v2.2.0)
           - Changelog atualizado a cada release
           - Documentação versionada junto com o código
        4. Backup e recovery:
           - Backup dos dados processados (outputs)
           - Procedimento de recovery documentado
           - Tempo máximo de indisponibilidade aceitável (SLA informal)

        ---

        #### 📦 Artefatos / Outputs esperados
        - Documento de política PROD (`PRODUCTION.md`)
        - Script de deploy controlado
        - Script de rollback
        - Checklist de deploy (pré e pós)
        - CHANGELOG.md

        ---

        #### ✅ Critérios de aceite
        - Deploy em PROD só acontece via processo documentado
        - Rollback funciona e restaura versão anterior em menos de 15 minutos
        - Health check retorna status de todas as plantas
        - Changelog reflete todas as mudanças desde a última release

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Deploy quebra PROD | Checklist + smoke test pós-deploy obrigatório |
        | Ninguém atualiza o changelog | Changelog como requisito do PR template |
        | Dados de PROD corrompidos | Backup antes de cada processamento |
        """)

    # ------------------------------------------------------------------
    # 14) Promoção DEV → PROD (checklist, gates, validações)
    # ------------------------------------------------------------------
    with st.expander("🟦 14) Promoção DEV → PROD (Checklist, Gates, Validações)", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Definir o processo seguro e auditável para promover mudanças de DEV para PROD,
        garantindo que apenas código validado e testado chegue ao sistema em operação.

        ---

        #### 📋 O que fazer
        1. Criar **checklist de promoção** (PR template):
           - [ ] Testes unitários passando
           - [ ] Testes de regressão executados (Betim baseline = OK)
           - [ ] Sem hardcodes de plant_id
           - [ ] Changelog atualizado
           - [ ] Documentação atualizada (se aplicável)
           - [ ] Code review aprovado por pelo menos 1 revisor
           - [ ] Smoke test em DEV bem-sucedido
        2. Criar **gates de aprovação**:
           - Gate técnico: CI verde + testes OK
           - Gate funcional: validação de outputs por usuário-chave
           - Gate de governança: log de promoção registrado
        3. Definir **pipeline de release**:
           - Merge develop → main via PR
           - Tag automática com versão semântica
           - Deploy automático ou semi-automático para PROD
        4. Documentar procedimento de rollback emergencial

        ---

        #### 📦 Artefatos / Outputs esperados
        - PR template com checklist de promoção
        - Pipeline de CI/CD (GitHub Actions ou equivalente)
        - Documento de rollback emergencial
        - Log de promoções (histórico)

        ---

        #### ✅ Critérios de aceite
        - Nenhuma mudança chega a PROD sem passar pelo checklist completo
        - Pipeline de CI bloqueia merge se testes falham
        - Histórico de promoções auditável
        - Rollback testado e funcional

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Pressão para pular gates ("é urgente") | Processo documentado; exceções exigem justificativa formal |
        | CI verde mas output incorreto | Gate funcional obrigatório (validação por humano) |
        | Rollback falha sob pressão | Simular rollback periodicamente (drill) |
        """)

    # ------------------------------------------------------------------
    # 15) Mecanismos anti-regressão
    # ------------------------------------------------------------------
    with st.expander("🟦 15) Mecanismos Anti-Regressão", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Implementar salvaguardas automáticas que detectem qualquer regressão nos cálculos,
        regras ou outputs do SCI antes que ela chegue a PROD.

        ---

        #### 📋 O que fazer
        1. **Baseline de outputs:** salvar snapshot dos outputs atuais de Betim como referência
           - Parquets de FP, FA, CPU, BE, Flex, Rateios
           - Totais por oficina e modelo
           - Hash dos arquivos
        2. **Comparativo automático DEV vs baseline:**
           - Após cada processamento em DEV, comparar outputs com baseline
           - Alertar se qualquer diferença for detectada
           - Diferenciar: diferença esperada (nova funcionalidade) vs. regressão
        3. **Testes de sanidade por planta:**
           - Soma de outputs > 0
           - Número de registros dentro do range esperado
           - Métricas-chave dentro de bounds (ex.: CPU entre X e Y)
        4. **Golden test:** processar dados de Betim com código novo e comparar, byte a byte,
           com output anterior
        5. **Dashboard de regressão:** visão rápida do status de regressão por planta

        ---

        #### 📦 Artefatos / Outputs esperados
        - Snapshot de baseline (versionado)
        - Script de comparação automática (`regression_check.py`)
        - Relatório de diferenças (se houver)
        - Dashboard de status de regressão
        - Integração com CI: falha automática se regressão detectada

        ---

        #### ✅ Critérios de aceite
        - Golden test de Betim passa com diff zero no baseline
        - Alteração intencional (nova planta) não dispara falso positivo
        - Regressão real é detectada e bloqueia promoção
        - Dashboard de regressão acessível e atualizado

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Baseline desatualizado após mudança legítima | Processo de atualização de baseline documentado |
        | Falsos positivos por diferenças de arredondamento | Tolerância configurável (epsilon) |
        | Testes de regressão muito lentos | Executar em paralelo; subset rápido para CI, completo para release |
        """)

    # ------------------------------------------------------------------
    # 16) Azure como caminho opcional (condicional)
    # ------------------------------------------------------------------
    with st.expander("🟦 16) Azure como Caminho Opcional (Condicional)", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Documentar o caminho de migração/expansão para Azure como infraestrutura cloud,
        **caso e quando** a Stellantis disponibilize essa opção para o projeto SCI.

        > ⚠️ **Este passo é condicional.** Só se aplica se/quando a infraestrutura Azure
        > corporativa for disponibilizada. O SCI LATAM funciona sem Azure.

        ---

        #### 📋 O que fazer (quando aplicável)
        1. Mapear serviços Azure equivalentes ao stack atual:
           - **Azure Blob Storage** → armazenamento de parquets e outputs
           - **Azure SQL / Cosmos DB** → dados estruturados (se necessário)
           - **Azure App Service / Container Apps** → hospedagem do Streamlit
           - **Azure DevOps / GitHub Actions** → CI/CD
           - **Azure Monitor** → monitoramento e alertas
        2. Definir estratégia de migração:
           - Migração incremental (não big-bang)
           - Dual-run: manter PROD local funcionando enquanto Azure é validado
           - Critérios de cutover: quando mudar para Azure como primário
        3. Estimar custos mensais de Azure por tier
        4. Documentar requisitos de segurança: rede, identidade, certificados

        ---

        #### 📦 Artefatos / Outputs esperados
        - Documento de arquitetura Azure proposta
        - Mapa de serviços: stack atual → serviço Azure
        - Estimativa de custos (3 cenários: mínimo, médio, máximo)
        - Plano de migração incremental
        - Checklist de segurança Azure

        ---

        #### ✅ Critérios de aceite
        - Documento aprovado pela TI/IS local
        - Estimativa de custos validada
        - Plano de migração não exige downtime
        - Todos os requisitos de segurança mapeados

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Azure não é aprovado pela TI | Manter stack local como fallback permanente |
        | Custos maiores que o esperado | Estimativa de 3 cenários + revisão trimestral |
        | Migração causa downtime | Dual-run obrigatório durante transição |
        """)

    # ------------------------------------------------------------------
    # 17) Roadmap e cronograma macro
    # ------------------------------------------------------------------
    with st.expander("🟦 17) Roadmap e Cronograma Macro", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Estabelecer o roadmap de alto nível para a implementação do SCI LATAM,
        com fases, dependências e marcos de entrega.

        ---

        #### 📋 O que fazer
        1. Definir fases do roadmap:
           - **Fase 0 — Fundação** (atual): planejamento, documentação, glossário, Plant Layer conceitual
           - **Fase 1 — Novo Workspace:** criar workspace `sci_latam/` no Databricks, repo GitHub, estrutura de pastas, CI básico
           - **Fase 2 — Migração Seletiva:** migrar módulos core (`tc_core`, `finance`, notebooks), eliminar duplicações, testes de smoke
           - **Fase 3 — Parametrização:** refactor do core para aceitar plant_id, test harness, baseline de Betim
           - **Fase 4 — Segunda planta (piloto):** onboarding de Porto Real como planta piloto
           - **Fase 5 — Brasil completo:** onboarding de Goiana, consolidação BR
           - **Fase 6 — Argentina:** onboarding de Córdoba e Palomar, consolidação LATAM
           - **Fase 7 — Comparativos e otimização:** dashboards comparativos, benchmarks, refinamentos
        2. Definir dependências entre fases (não paralelizável: Fase N depende de N-1)
        3. Definir marcos de entrega (milestones) por fase
        4. Identificar riscos de cronograma e planos de contingência

        ---

        #### 📦 Artefatos / Outputs esperados
        - Roadmap visual (diagrama de Gantt simplificado ou timeline)
        - Lista de milestones por fase
        - Matriz de dependências
        - Plano de contingência por fase

        ---

        #### ✅ Critérios de aceite
        - Roadmap aprovado pelos stakeholders principais
        - Cada fase tem milestone e critério de conclusão claros
        - Dependências explícitas (nenhuma fase começa sem a anterior concluída)
        - Plano de contingência documentado para as 2 fases mais críticas

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Cronograma irrealista | Fases incrementais com revisão ao final de cada uma |
        | Dependência de times externos (TI, AR) | Identificar dependências cedo e engajar stakeholders |
        | Scope creep durante execução | Cada fase tem escopo congelado após aprovação |
        """)

    # ------------------------------------------------------------------
    # 18) Conclusão executiva e decisões pendentes
    # ------------------------------------------------------------------
    with st.expander("🟦 18) Conclusão Executiva e Decisões Pendentes", expanded=False):
        st.markdown("""
        #### 🎯 Objetivo
        Consolidar a visão do plano SCI LATAM, registrar decisões já tomadas e
        listar decisões pendentes que dependem de stakeholders.

        ---

        #### 📋 Decisões já tomadas
        - ✅ SCI LATAM é a próxima evolução do SCI atual
        - ✅ 5 plantas no escopo inicial (3 BR + 2 AR)
        - ✅ Evolução não-regressiva: cálculos, regras e outputs atuais permanecem intactos
        - ✅ Código único (sem fork por planta)
        - ✅ Dois ambientes: DEV e PROD
        - ✅ Consolidação baseada em outputs oficiais
        - ✅ Azure é condicional (não bloqueante)
        - ✅ **Novo workspace no Databricks** (`sci_latam/`) — separado do atual
        - ✅ **Novo repositório GitHub** (ou branch dedicada) — estrutura limpa desde o início
        - ✅ **Aproveitamento seletivo** dos módulos existentes — sem arrastar dívida técnica
        - ✅ **Notebooks migrados e unificados** — uma única cópia de cada
        - ✅ **Sincronização via script + GitHub Actions** — não editar direto no workspace

        #### 📋 Decisões pendentes
        - ❓ Moeda base para consolidação LATAM: BRL ou USD?
        - ❓ Fonte oficial de taxa de câmbio ARS/BRL
        - ❓ Responsáveis por planta (data owners) na Argentina
        - ❓ Calendário fiscal AR vs BR: alinhamento ou execução separada?
        - ❓ Nível de acesso: cada planta vê apenas seus dados ou todas veem tudo?
        - ❓ SLA informal de disponibilidade do sistema PROD
        - ❓ Quem aprova o gate funcional de promoção DEV → PROD?
        - ❓ Novo repo (`SCI-LATAM`) ou branch (`sci-latam`) no repo atual?
        - ❓ Cluster Databricks: compartilhado ou dedicado para SCI LATAM?
        - ❓ Catalog Unity Catalog: schema único ou por planta?

        #### 📋 Próximos gates
        1. **Gate de Planejamento:** aprovação deste documento pelos stakeholders
        2. **Gate de Fundação:** Plant Layer e baseline implementados e validados
        3. **Gate de Piloto:** Porto Real onboarded e processando corretamente
        4. **Gate BR:** 3 plantas BR funcionando + consolidação BR validada
        5. **Gate LATAM:** 5 plantas + consolidação LATAM validada

        ---

        #### 📦 Artefatos / Outputs esperados
        - Ata de decisões (atualizada a cada gate)
        - Lista de pendências com responsáveis e prazo
        - Relatório de status por gate

        ---

        #### ✅ Critérios de aceite
        - Todas as decisões pendentes têm responsável atribuído
        - Cronograma de resolução das pendências definido
        - Stakeholders cientes e comprometidos com o plano
        - Documento assinado / aprovado formalmente

        ---

        #### ⚠️ Riscos comuns + Mitigação
        | Risco | Mitigação |
        |-------|-----------|
        | Decisões pendentes travam o início da Fase 1 | Priorizar decisões bloqueantes; assumir defaults onde possível |
        | Falta de engajamento da Argentina | Envolver AR desde a Fase 0 (planejamento) |
        | Documento fica obsoleto | Revisão obrigatória a cada gate |
        """)

# Rodapé
st.markdown("---")
mes_atual = obter_mes_atual()
ano_atual = datetime.now().year
versao_atual = obter_versao_atual()
st.markdown(f"""
<div style='text-align: center; color: #666; padding: 20px;'>
    📚 Stellantis Cost Intelligence (SCI) | Versão {versao_atual} | {mes_atual} {ano_atual}
    <br>
    <small>Desenvolvido por Hudson Cardin, Lauro Paiva e Frederico Cesar de Jesus</small>
</div>
""", unsafe_allow_html=True)

