from __future__ import annotations

import base64
import html
import json
from dataclasses import dataclass, field
from datetime import datetime
from io import BytesIO
from pathlib import Path

import streamlit as st
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas as pdf_canvas
from reportlab.platypus import Paragraph


ROOT = Path(__file__).resolve().parent.parent
SCI_FAIXA_PATH = ROOT / "SCI_faixa.png"
EQUIPE_JSON_PATH = ROOT / "dados_equipe.json"

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
SLIDE_WIDTH_PT = SLIDE_WIDTH_IN * 72
SLIDE_HEIGHT_PT = SLIDE_HEIGHT_IN * 72
SLIDE_PAGE_SIZE = (SLIDE_WIDTH_PT, SLIDE_HEIGHT_PT)

PRIMARY_HEX = "#13294B"
ACCENT_HEX = "#E95D0F"
TEAL_HEX = "#00A3A3"
INK_HEX = "#172033"
SURFACE_HEX = "#F7F9FC"
MUTED_HEX = "#64748B"
LINE_HEX = "#D8E1EE"
PANEL_HEX = "#EAF0FB"

PRIMARY = colors.HexColor(PRIMARY_HEX)
ACCENT = colors.HexColor(ACCENT_HEX)
TEAL = colors.HexColor(TEAL_HEX)
INK = colors.HexColor(INK_HEX)
SURFACE = colors.HexColor(SURFACE_HEX)
MUTED = colors.HexColor(MUTED_HEX)
LINE = colors.HexColor(LINE_HEX)
PANEL = colors.HexColor(PANEL_HEX)
WHITE = colors.white


@dataclass
class SlideSpec:
    id: str
    title: str
    subtitle: str
    duration: str
    key_message: str
    bullets: list[str]
    highlights: list[str] = field(default_factory=list)
    visual_title: str = ""
    visual_lines: list[str] = field(default_factory=list)
    presenter_note: str = ""
    cover: bool = False
    label: str = ""


@dataclass
class TeamMember:
    key: str
    name: str
    role: str
    description: str
    photo_b64: str | None = None


LANGUAGE_OPTIONS = {
    "pt": "PortuguÃªs",
    "en": "English",
}


def _ui_texts(language: str) -> dict[str, str]:
    if language == "en":
        return {
            "section_badge": "Visual presentation updated",
            "section_title": "Stellantis Cost Intelligence visual pitch",
            "section_subtitle": "The page remains the visual reference; PPT and PDF exports now follow the same widescreen slide design.",
            "language_label": "Presentation language",
            "download_pdf": "ðŸ“„ Export presentation as PDF",
            "download_ppt": "ðŸ“Š Export presentation as PPT",
            "tabs_script": "ðŸŽ¤ Script (5 min)",
            "tabs_slides": "ðŸ§© Visual slides",
            "script_intro": "**Suggested speaking sequence**",
            "slides_info": "Each expander below corresponds to one slide and is the source for the in-app view, the PPT, and the widescreen PDF.",
            "what_will_be_shown": "**What will be shown**",
            "narrative": "**Suggested narrative**",
            "key_points": "Key points",
            "visual_support": "Visual support",
            "why_it_matters": "**Why this matters**",
            "speaker_note": "Speaker note:",
            "cover_badge": "Executive SCI presentation",
            "cover_highlights": "Highlights",
            "header_points": "KEY POINTS",
            "header_why": "WHY THIS MATTERS",
            "pdf_title": "SCI Executive Presentation",
            "pdf_visual_support": "VISUAL SUPPORT",
            "pdf_speaker_note": "Speaker note:",
        }

    return {
        "section_badge": "Apresentacao visual revisada",
        "section_title": "Pitch visual do Stellantis Cost Intelligence",
        "section_subtitle": "A page segue como referencia visual; os exports em PPT e PDF agora nascem do mesmo desenho de slide widescreen.",
        "language_label": "Idioma da apresentacao",
        "download_pdf": "ðŸ“„ Exportar Apresentacao em PDF",
        "download_ppt": "ðŸ“Š Exportar Apresentacao em PPT",
        "tabs_script": "ðŸŽ¤ Roteiro (5 min)",
        "tabs_slides": "ðŸ§© Slides visuais",
        "script_intro": "**Sequencia sugerida para a fala**",
        "slides_info": "Cada expander abaixo corresponde a um slide e serve de base para a visualizacao no app, para o PPT e para o PDF widescreen.",
        "what_will_be_shown": "**O que sera mostrado**",
        "narrative": "**Narrativa sugerida**",
        "key_points": "Pontos-chave",
        "visual_support": "Apoio visual",
        "why_it_matters": "**Por que isso importa**",
        "speaker_note": "Fala sugerida:",
        "cover_badge": "Apresentacao executiva SCI",
        "cover_highlights": "Destaques",
        "header_points": "PONTOS-CHAVE",
        "header_why": "POR QUE ISSO IMPORTA",
        "pdf_title": "SCI Apresentacao Executiva",
        "pdf_visual_support": "APOIO VISUAL",
        "pdf_speaker_note": "Fala sugerida:",
    }


def _translate_team_role(role: str, language: str) -> str:
    if language != "en":
        return role

    role_map = {
        "Full-Stack Developer": "Full-Stack Developer",
        "Tech Advisor": "Tech Advisor",
        "Equipe SCI": "SCI Team",
        "ESPECIALISTA CONTROLE E GESTAO": "Management Control Specialist",
    }
    return role_map.get(role, role)


def _translate_team_description(description: str, language: str) -> str:
    if language != "en":
        return description

    replacements = {
        "Responsavel pelo desenvolvimento end-to-end da plataforma, integrando interface, logica de negocio, calculos e experiencia do usuario.": "Responsible for the platform end-to-end development, integrating interface, business logic, calculations, and user experience.",
        "Atuacao no desenvolvimento full-stack do SCI, conectando camada visual, regras de negocio e confianca dos indicadores para o usuario final.": "Works across the SCI full-stack development, connecting visual layer, business rules, and indicator reliability for the end user.",
        "Conecta a visao de negocio e controladoria ao produto, garantindo aderencia da ferramenta a tomada de decisao industrial.": "Connects business and controlling vision to the product, ensuring the tool stays aligned with industrial decision-making.",
        "Atuacao chave na evolucao do SCI.": "Key contributor to the SCI evolution.",
    }
    return replacements.get(description, description)


_TEAM_ORDER = [
    ("hudson", "Hudson Cardin"),
    ("lauro", "Lauro Paiva Junior"),
    ("frederico", "Frederico Cesar de Jesus"),
]


def _hex_to_rgb(hex_value: str) -> tuple[int, int, int]:
    value = hex_value.lstrip("#")
    return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))


def _safe_html(text: str) -> str:
    return html.escape(text, quote=True)


def _image_file_data_uri(path: Path) -> str | None:
    if not path.exists():
        return None
    encoded = base64.b64encode(path.read_bytes()).decode("ascii")
    suffix = path.suffix.lower()
    mime = "image/png" if suffix == ".png" else "image/jpeg"
    return f"data:{mime};base64,{encoded}"


def _load_team_members(language: str = "pt") -> list[TeamMember]:
    if not EQUIPE_JSON_PATH.exists():
        return []

    try:
        data = json.loads(EQUIPE_JSON_PATH.read_text(encoding="utf-8"))
    except Exception:
        return []

    members: list[TeamMember] = []
    for key, fallback_name in _TEAM_ORDER:
        raw = data.get(key, {})
        members.append(
            TeamMember(
                key=key,
                name=raw.get("nome") or fallback_name,
                role=_translate_team_role(raw.get("papel_projeto") or raw.get("cargo") or "Equipe SCI", language),
                description=_translate_team_description(raw.get("descricao_papel") or "Atuacao chave na evolucao do SCI.", language),
                photo_b64=raw.get("foto"),
            )
        )
    return members


def _photo_bytes(photo_b64: str | None, max_px: int = 320) -> bytes | None:
    if not photo_b64:
        return None

    try:
        raw = base64.b64decode(photo_b64)
    except Exception:
        return None

    try:
        from PIL import Image as PILImage

        image = PILImage.open(BytesIO(raw))
        if image.mode != "RGB":
            image = image.convert("RGB")
        image.thumbnail((max_px, int(max_px * 1.2)), PILImage.LANCZOS)
        out = BytesIO()
        image.save(out, format="JPEG", quality=88)
        out.seek(0)
        return out.read()
    except Exception:
        return raw


def _photo_data_uri(photo_b64: str | None) -> str | None:
    data = _photo_bytes(photo_b64)
    if not data:
        return None
    return "data:image/jpeg;base64," + base64.b64encode(data).decode("ascii")


def _cover_slide(version: str, data_atualizacao: str | None) -> SlideSpec:
    periodo = datetime.now().strftime("%B %Y")
    return SlideSpec(
        id="cover",
        label="Capa",
        title="Stellantis Cost Intelligence",
        subtitle=f"Apresentacao executiva do SCI | Versao {version} | {periodo.title()}",
        duration="0:00",
        key_message="Uma plataforma unica para transformar custo industrial em leitura rapida, explicacao rastreavel e acao orientada ao negocio.",
        bullets=[
            "Visao integrada de Home, TC Ext, TC Veiculos, Best Estimate, Alertas, TC Copilot e Documentacao.",
            "Identidade visual alinhada com a capa do relatorio, agora com a equipe do projeto em evidencia.",
            f"Dados atualizados em: {data_atualizacao or 'informacao nao disponivel'}.",
        ],
        highlights=[
            "Analise executiva",
            "Explicacao de desvios",
            "Monitoramento proativo",
        ],
        visual_title="Narrativa da apresentacao",
        visual_lines=[
            "Extrair e validar dados",
            "Ler o numero certo na Home",
            "Explicar causas com Waterfall e Flex",
            "Projetar cenarios com Best Estimate",
            "Acionar decisao com Alertas e Relatorios",
        ],
        presenter_note="Abrir com proposta de valor, reforcando que o SCI conecta fechamento, analise e acao em uma unica linguagem.",
        cover=True,
    )


def _renumber_slides(slides: list[SlideSpec]) -> list[SlideSpec]:
    counter = 1
    for slide in slides:
        if slide.cover:
            slide.label = "Capa"
            continue
        slide.label = f"Slide {counter}"
        counter += 1
    return slides


def build_presentation_slides(version: str, data_atualizacao: str | None, language: str = "pt") -> list[SlideSpec]:
    if language == "en":
        slides = [
            SlideSpec(
                id="cover",
                label="Cover",
                title="Stellantis Cost Intelligence",
                subtitle=f"SCI executive presentation | Version {version} | {datetime.now().strftime('%B %Y').title()}",
                duration="0:00",
                key_message="A single platform that turns industrial cost data into fast reading, traceable explanations, and action-oriented decisions.",
                bullets=[
                    "Integrated view of Home, TC Ext, TC Vehicles, Best Estimate, Alerts, TC Copilot, and Documentation.",
                    "Visual identity aligned with the report cover, now highlighting the project team as well.",
                    f"Data last updated on: {data_atualizacao or 'information not available'}.",
                ],
                highlights=["Executive analysis", "Deviation explanation", "Proactive monitoring"],
                visual_title="Presentation narrative",
                visual_lines=[
                    "Extract and validate data",
                    "Read the right number on Home",
                    "Explain root causes with Waterfall and Flex",
                    "Project scenarios with Best Estimate",
                    "Trigger decisions with Alerts and Reports",
                ],
                presenter_note="Open with the value proposition, showing that SCI connects closing, analysis, and action through one unified language.",
                cover=True,
            ),
            SlideSpec(
                id="overview",
                title="What SCI is",
                subtitle="One platform to understand industrial cost end to end",
                duration="0:00-0:30",
                key_message="SCI brings executive reading, deviation explanation, projection, and monitoring together in one unified experience.",
                bullets=[
                    "TC Extended provides an aggregated view by shop, period, and cost dimension.",
                    "TC Vehicles details the FA -> FP -> D&A chain by model and by shop.",
                    "Shared capabilities connect extraction, Flex Budget, Waterfall, Best Estimate, alerts, and reports.",
                ],
                highlights=["Fewer parallel spreadsheets", "More context in the same flow", "One narrative for the user"],
                visual_title="Value map",
                visual_lines=["Understand the number", "Explain the cause", "Simulate the future", "Be warned before the issue escalates"],
                presenter_note="Position SCI as a complete journey, not a set of isolated pages.",
            ),
            SlideSpec(
                id="home",
                title="Home - executive reading for the month",
                subtitle="KPIs, Flex Bud, comparisons, and filters that lead straight to the issue",
                duration="0:30-1:00",
                key_message="Home concentrates the information users need to read the result quickly and decide where to drill deeper.",
                bullets=[
                    "KPIs for Budget, Flex Bud, Real or Best Estimate, plus percentage changes, in one fast-reading row.",
                    "Shop, vehicle, currency, and period filters take the same analysis to the exact operational context.",
                    "Period charts and flex analysis blocks help users move from summary to drill-down without changing mental context.",
                ],
                highlights=["Read in seconds", "Flex Bud built in", "Starting point for the full journey"],
                visual_title="What the user does on Home",
                visual_lines=["Read the main KPI", "Validate Real vs Flex Bud", "Filter by shop and vehicle", "Navigate to Waterfall, Flex, or Alerts"],
                presenter_note="Show Home as the user's cockpit, where the number gains context immediately.",
            ),
            SlideSpec(
                id="extraction",
                title="Data extraction and processing",
                subtitle="From Excel to Parquet with validation, logs, and post-processing alert triggering",
                duration="1:00-1:35",
                key_message="Extraction removes manual effort and ensures data reaches the dashboard with structure, traceability, and control.",
                bullets=[
                    "Guided upload of source files plus pre-validation of sheets, required columns, and month structure before processing.",
                    "Separate execution for Actuals, Budget, or both, with status, logs, and notebook synchronization checks.",
                    "At the end, SCI can trigger the Alert Center on the freshly processed base, reducing the risk of reading outdated data.",
                ],
                highlights=["Less input error", "Auditable process", "Direct bridge to monitoring"],
                visual_title="Operational flow",
                visual_lines=["Upload Excel files", "Pre-validation", "Process Actuals and Budget", "Parquets and history", "Post-extraction alerts"],
                presenter_note="Reinforce that ETL complexity is encapsulated for the end user.",
            ),
            SlideSpec(
                id="extended",
                title="TC Extended",
                subtitle="Aggregated view for total cost and CPU with correct Flex Budget logic",
                duration="1:35-2:05",
                key_message="In TC Extended, SCI quickly answers how consolidated cost behaves by period, shop, and spending dimension.",
                bullets=[
                    "Total Cost and CPU modes, with the critical rule that CPU is always recalculated after aggregation.",
                    "Flex Budget separates volume effect from cost effect and avoids distorted deviation readings.",
                    "Advanced filters by USI, Type 05, Type 06, Account, material, order, and origin enable controlled detail.",
                ],
                highlights=["CPU is never summed", "Flex preserves the correct reading", "Consolidated business radar"],
                visual_title="How the user navigates",
                visual_lines=["Choose year, currency, and period", "Switch between Total and CPU", "Drill down from macro view to account"],
                presenter_note="Position this module as the reference for macro reading and fast comparison.",
            ),
            SlideSpec(
                id="vehicles",
                title="TC Vehicles",
                subtitle="Full cost chain with allocation proportional to production time",
                duration="2:05-2:40",
                key_message="In TC Vehicles, SCI brings financial analysis closer to operations by showing how cost reaches each vehicle.",
                bullets=[
                    "Main chain: Primary Expense -> FA -> FP -> Dedicated D&A.",
                    "Vehicle allocation proportional to production time by shop, activated when a specific model is selected.",
                    "Home organized in tabs for KPIs, volume, shop costs, flex analysis, production time, and detailed data.",
                ],
                highlights=["Vehicle selection activates allocation", "All = consolidated without distortion", "Closer reading of the plant"],
                visual_title="Questions this module answers",
                visual_lines=["Which vehicle drove cost up?", "Which shop is pressuring the result?", "Where does allocation explain the difference?"],
                presenter_note="Present the module as a bridge between financial numbers and operational reality.",
            ),
            SlideSpec(
                id="shared",
                title="Waterfall and Best Estimate",
                subtitle="Explain the past and project the future with the same business language",
                duration="2:40-3:15",
                key_message="SCI does not only show the number; it explains what changed and helps project what comes next.",
                bullets=[
                    "Waterfall breaks period-to-period variation down and isolates volume effect versus cost effect through Flex Budget.",
                    "Best Estimate uses sensitivity, inflation, and volume to generate a reproducible and auditable forecast.",
                    "History and projection live in the same reading, supporting closing and planning without context switching.",
                ],
                highlights=["Diagnosis + projection", "Transparent formula", "Same journey for reading and planning"],
                visual_title="Analytical cycle",
                visual_lines=["Today: KPI and deviation", "Yesterday: explanation in Waterfall", "Tomorrow: simulation in Best Estimate"],
                presenter_note="Show that the user does not need to leave SCI to move from diagnosis to planned action.",
            ),
            SlideSpec(
                id="copilot_reports",
                title="TC Copilot and reports",
                subtitle="Smart chat, annual report, and PDF library inside the same ecosystem",
                duration="3:15-3:50",
                key_message="TC Copilot turns data into executive narrative, either through live conversation or ready-to-share reports.",
                bullets=[
                    "The Chatbot tab answers questions about the parquet data in natural language with year, month, and currency context.",
                    "The Vehicle Report tab generates automatic or AI-assisted reports, including annual PDF, monthly PDFs, and a centralized library.",
                    "A dedicated configuration area controls API key, model, language, and feature availability without mixing that with analysis.",
                ],
                highlights=["Executive narrative ready", "Annual and monthly downloads", "AI when useful, automation when scale matters"],
                visual_title="TC Copilot flow",
                visual_lines=["Ask about the data", "Generate automatic or AI-assisted reports", "Download annual or monthly PDF", "Consult the historical library"],
                presenter_note="Give special weight to this slide because reporting is a high-value deliverable.",
            ),
            SlideSpec(
                id="alerts",
                title="Alert Center",
                subtitle="Proactive monitoring with hierarchical ranking, validation, and notification ready for action",
                duration="3:50-4:20",
                key_message="The Alert Center transforms SCI from a consultation tool into an active monitoring system.",
                bullets=[
                    "Consolidated tree ranking: Type 05 -> Type 06 -> Account -> Shop -> Brief text.",
                    "The validation table allows the user to audit Flex Bud, Real, monetary delta, and percentage before sending.",
                    "Manual triggering, Teams and email tests, dependent filters, and history make the flow more predictable for monthly closing.",
                ],
                highlights=["Share of total deviation", "Teams and email aligned with monitoring", "More prioritization, less noise"],
                visual_title="Alert flow",
                visual_lines=["Detect relevant deviation", "Validate in the table", "Notify stakeholders", "Register history and follow-up"],
                presenter_note="Close the bridge between analysis and operational action, highlighting trust in the calculation.",
            ),
            SlideSpec(
                id="architecture",
                title="Architecture and governance",
                subtitle="Modular foundation to grow without losing consistency, performance, and traceability",
                duration="4:20-4:40",
                key_message="SCI's structure sustains user trust because it organizes rules, data, and interface in a modular way.",
                bullets=[
                    "Separate layers for app, business modules, shared core, and common pages.",
                    "Parquet as canonical base, SQLite for exchange rates, automatic versioning, and smart cache for performance.",
                    "Documentation integrated into the product itself accelerates onboarding, auditing, and continuous maintenance.",
                ],
                highlights=["Standardized data", "Growth with governance", "Safer maintenance"],
                visual_title="Functional architecture",
                visual_lines=["Excel input", "Python processing", "Historical and forecast parquet", "Dashboard + Copilot + Alerts"],
                presenter_note="Show that technical robustness sustains business reading.",
            ),
            SlideSpec(
                id="user_journey",
                title="User journey",
                subtitle="From raw data to decision in a few steps, without breaking the mental flow",
                duration="4:40-4:55",
                key_message="SCI's experience was designed to reduce closing time, simplify explanation, and accelerate action.",
                bullets=[
                    "Guided extraction and processing with validation.",
                    "Executive reading on Home, deeper analysis and simulation when needed.",
                    "Alerts and reports close the loop with follow-up and communication ready for the business.",
                ],
                highlights=["Less operational friction", "More speed to act", "Decision-oriented experience"],
                visual_title="Time to insight",
                visual_lines=["Extract", "Process", "Read KPI", "Explain deviation", "Act with confidence"],
                presenter_note="Close the narrative by emphasizing ease of use over a sophisticated foundation.",
            ),
            SlideSpec(
                id="impact",
                title="Impact and closing",
                subtitle="More speed, more traceability, and more autonomy for decision-making",
                duration="4:55-5:00",
                key_message="SCI reduces manual work, shortens the time to understand variation, and strengthens confidence in the number.",
                bullets=[
                    "Closing and analysis happen in minutes, not in hours of manual consolidation.",
                    "The number becomes easier to defend with documented rules, debug support, and validation tables.",
                    "A continuously evolving ecosystem connects dashboards, alerts, reports, and documentation.",
                ],
                highlights=["Hours -> minutes", "From KPI to action plan", "Confidence to scale adoption"],
                visual_title="Final message",
                visual_lines=["SCI = understand + explain + project + monitor + decide"],
                presenter_note="End with the concrete benefit for business and operations.",
            ),
        ]
        return _renumber_slides(slides)

    slides = [
        _cover_slide(version, data_atualizacao),
        SlideSpec(
            id="overview",
            title="O que e o SCI",
            subtitle="Uma plataforma unica para entender custo industrial de ponta a ponta",
            duration="0:00-0:30",
            key_message="O SCI consolida leitura executiva, explicacao de desvios, projecao e monitoramento em uma experiencia unificada.",
            bullets=[
                "TC Estendido entrega visao agregada por oficina, periodo e dimensao de custo.",
                "TC Veiculos detalha a cadeia FA -> FP -> D&A por modelo e por oficina.",
                "Funcionalidades transversais conectam extracao, Flex Budget, Waterfall, Best Estimate, alertas e relatorios.",
            ],
            highlights=[
                "Menos planilhas paralelas",
                "Mais contexto no mesmo fluxo",
                "Uma narrativa unica para o usuario",
            ],
            visual_title="Mapa de valor",
            visual_lines=[
                "Entender o numero",
                "Explicar a causa",
                "Simular o futuro",
                "Ser avisado antes do desvio escalar",
            ],
            presenter_note="Posicionar o SCI como jornada completa, nao como conjunto de paginas isoladas.",
        ),
        SlideSpec(
            id="home",
            title="Home - leitura executiva do mes",
            subtitle="KPIs, Flex Bud, comparativos e filtros que levam direto ao problema",
            duration="0:30-1:00",
            key_message="A Home concentra as informacoes que o usuario precisa para bater o olho no resultado e decidir para onde aprofundar.",
            bullets=[
                "KPIs de BUD, Flex Bud, Real ou Best Estimate e variacoes percentuais em uma linha de leitura rapida.",
                "Filtros de oficina, veiculo, moeda e periodo levam a mesma analise para o contexto exato da operacao.",
                "Graficos por periodo e blocos de analise flex ajudam a sair do resumo para o drill-down sem trocar de raciocinio.",
            ],
            highlights=[
                "Leitura em segundos",
                "Flex Bud integrado",
                "Ponto de partida para toda a jornada",
            ],
            visual_title="O que o usuario faz na Home",
            visual_lines=[
                "Ler o KPI principal",
                "Validar Real vs Flex Bud",
                "Filtrar por oficina e veiculo",
                "Navegar para Waterfall, Flex ou Alertas",
            ],
            presenter_note="Trazer a Home como cockpit do usuario, onde o numero ganha contexto imediatamente.",
        ),
        SlideSpec(
            id="extraction",
            title="Extracao e processamento de dados",
            subtitle="Do Excel ao Parquet com validacao, logs e disparo de alertas pos-processamento",
            duration="1:00-1:35",
            key_message="A extracao tira o trabalho manual do caminho e garante que o dado chegue ao dashboard com estrutura, rastreabilidade e criterio.",
            bullets=[
                "Upload guiado dos arquivos-base e pre-validacao de abas, colunas obrigatorias e estrutura de meses antes do processamento.",
                "Execucao separada para REAIS, BUDGET ou ambos, com status, logs e verificacao de sincronizacao dos notebooks convertidos.",
                "Ao final, o SCI pode acionar a Central de Alertas sobre a base recem-processada, reduzindo risco de leitura em dado desatualizado.",
            ],
            highlights=[
                "Menos erro de entrada",
                "Processo auditavel",
                "Ponte direta para monitoramento",
            ],
            visual_title="Fluxo operacional",
            visual_lines=[
                "Upload dos Excels",
                "Pre-validacao",
                "Processamento REAIS e BUDGET",
                "Parquets e historico",
                "Alertas pos-extracao",
            ],
            presenter_note="Reforcar que a complexidade do ETL fica encapsulada para o usuario final.",
        ),
        SlideSpec(
            id="extended",
            title="TC Estendido",
            subtitle="Visao agregada para custo total e CPU com logica correta de Flex Budget",
            duration="1:35-2:05",
            key_message="No TC Ext, o SCI responde rapidamente como o custo consolidado esta se comportando por periodo, oficina e dimensao de gasto.",
            bullets=[
                "Modo Custo Total e modo CPU, com regra critica de CPU sempre recalculado apos agregacao.",
                "Flex Budget separa efeito volume de efeito custo e evita leituras distorcidas do desvio.",
                "Filtros avancados por USI, Type 05, Type 06, Account, material, pedido e origem permitem detalhamento controlado.",
            ],
            highlights=[
                "CPU nunca e somado",
                "Flex preserva a leitura correta",
                "Radar consolidado do negocio",
            ],
            visual_title="Como o usuario navega",
            visual_lines=[
                "Escolhe ano, moeda e periodo",
                "Alterna entre Total e CPU",
                "Desce da visao macro ate a conta",
            ],
            presenter_note="Posicionar o modulo como referencia para leitura macro e comparacao rapida.",
        ),
        SlideSpec(
            id="vehicles",
            title="TC Veiculos",
            subtitle="Cadeia completa de custos com rateio proporcional ao tempo de producao",
            duration="2:05-2:40",
            key_message="No TC Veiculos, o SCI aproxima a analise financeira da operacao ao mostrar como o custo chega ao veiculo.",
            bullets=[
                "Cadeia principal: Despesa Primaria -> FA -> FP -> D&A dedicado.",
                "Rateio por veiculo proporcional ao tempo de producao por oficina, ativado quando um modelo especifico e selecionado.",
                "Home organizada em tabs para KPIs, volume, custos por oficina, analise flex, tempo de producao e dados detalhados.",
            ],
            highlights=[
                "Selecao de veiculo aciona rateio",
                "Todos = consolidado sem distorcao",
                "Leitura mais proxima da fabrica",
            ],
            visual_title="Perguntas que o modulo responde",
            visual_lines=[
                "Qual veiculo puxou o custo?",
                "Qual oficina pressiona o resultado?",
                "Onde o rateio explica a diferenca?",
            ],
            presenter_note="Trazer o modulo como ponte entre numero financeiro e realidade operacional.",
        ),
        SlideSpec(
            id="shared",
            title="Waterfall e Best Estimate",
            subtitle="Explicar o passado e projetar o futuro com a mesma linguagem de negocio",
            duration="2:40-3:15",
            key_message="O SCI nao apenas mostra o numero; ele explica o que mudou e ajuda a projetar o que vem pela frente.",
            bullets=[
                "Waterfall decompoe a variacao entre periodos e isola efeito volume versus efeito custo via Flex Budget.",
                "Best Estimate usa sensibilidade, inflacao e volume para gerar forecast reproduzivel e auditavel.",
                "Historico e projecao convivem na mesma leitura para apoiar fechamento e planejamento sem ruptura.",
            ],
            highlights=[
                "Diagnostico + projecao",
                "Formula transparente",
                "Mesma jornada para ler e planejar",
            ],
            visual_title="Ciclo analitico",
            visual_lines=[
                "Hoje: KPI e desvio",
                "Ontem: explicacao no Waterfall",
                "Amanha: simulacao no Best Estimate",
            ],
            presenter_note="Mostrar que o usuario nao precisa sair do SCI para passar de diagnostico a acao planejada.",
        ),
        SlideSpec(
            id="copilot_reports",
            title="TC Copilot e relatorios",
            subtitle="Chat inteligente, relatorio anual e biblioteca de PDFs no mesmo ecossistema",
            duration="3:15-3:50",
            key_message="O TC Copilot transforma dados em narrativa executiva, seja por conversa ao vivo, seja por relatorios prontos para compartilhar.",
            bullets=[
                "Aba Chatbot responde perguntas sobre os parquets em linguagem natural com contexto de ano, mes e moeda.",
                "Aba Relatorio Veic. gera relatorios automaticos ou com IA, incluindo PDF anual, PDFs mensais e biblioteca centralizada.",
                "Configuracao dedicada permite controlar API key, modelo, idioma e disponibilidade do recurso sem misturar isso com a analise do usuario.",
            ],
            highlights=[
                "Narrativa executiva pronta",
                "Downloads anual e mensal",
                "IA quando faz sentido, automacao quando precisa escalar",
            ],
            visual_title="Fluxo do TC Copilot",
            visual_lines=[
                "Perguntar sobre o dado",
                "Gerar relatorio automatico ou com IA",
                "Baixar PDF anual ou mensal",
                "Consultar biblioteca historica",
            ],
            presenter_note="Dar peso especial a este slide, porque o relatorio e uma entrega de alto valor percebido.",
        ),
        SlideSpec(
            id="alerts",
            title="Central de Alertas",
            subtitle="Monitoramento proativo com ranking hierarquico, validacao e notificacao pronta para acao",
            duration="3:50-4:20",
            key_message="A Central de Alertas transforma o SCI de ferramenta de consulta em sistema de monitoramento ativo.",
            bullets=[
                "Ranking consolidado em arvore: Type 05 -> Type 06 -> Account -> Oficina -> Texto breve.",
                "Tabela de validacao permite auditar Flex Bud, Real, delta monetario e percentual antes do envio.",
                "Disparo manual, testes de Teams e e-mail, filtros dependentes e historico tornam o fluxo mais previsivel para o fechamento.",
            ],
            highlights=[
                "Percentual do desvio total",
                "Teams e e-mail alinhados ao monitoramento",
                "Mais prioridade, menos ruido",
            ],
            visual_title="Fluxo do alerta",
            visual_lines=[
                "Detectar desvio relevante",
                "Validar na tabela",
                "Notificar stakeholders",
                "Registrar historico e acompanhamento",
            ],
            presenter_note="Fechar a ponte entre analise e acao operacional, destacando a confianca do calculo.",
        ),
        SlideSpec(
            id="architecture",
            title="Arquitetura e governanca",
            subtitle="Base modular para crescer sem perder consistencia, performance e rastreabilidade",
            duration="4:20-4:40",
            key_message="A estrutura do SCI sustenta a confianca do usuario final porque organiza regra, dado e interface de forma modular.",
            bullets=[
                "Camadas separadas entre app, modulos de negocio, core compartilhado e pages comuns.",
                "Parquet como base canonica, SQLite para cambio, versionamento automatico e cache inteligente para performance.",
                "Documentacao integrada ao proprio produto acelera onboarding, auditoria e manutencao continua.",
            ],
            highlights=[
                "Dados padronizados",
                "Crescimento com governanca",
                "Manutencao mais segura",
            ],
            visual_title="Arquitetura funcional",
            visual_lines=[
                "Entrada Excel",
                "Processamento Python",
                "Parquet historico e forecast",
                "Dashboard + Copilot + Alertas",
            ],
            presenter_note="Mostrar que a robustez tecnica sustenta a leitura de negocio.",
        ),
        SlideSpec(
            id="user_journey",
            title="Jornada do usuario",
            subtitle="Do dado bruto a decisao em poucos passos, sem quebrar o fluxo mental",
            duration="4:40-4:55",
            key_message="A experiencia do SCI foi desenhada para reduzir tempo de fechamento, facilitar explicacao e acelerar a acao.",
            bullets=[
                "Extracao guiada e processamento com validacao.",
                "Leitura executiva na Home, aprofundamento em analise e simulacao quando necessario.",
                "Alertas e relatorios fecham o ciclo com acompanhamento e comunicacao prontos para o negocio.",
            ],
            highlights=[
                "Menos atrito operacional",
                "Mais velocidade para agir",
                "Experiencia orientada a decisao",
            ],
            visual_title="Tempo para insight",
            visual_lines=[
                "Extrair",
                "Processar",
                "Ler KPI",
                "Explicar desvio",
                "Agir com seguranca",
            ],
            presenter_note="Fechar a narrativa enfatizando simplicidade de uso sobre uma base sofisticada.",
        ),
        SlideSpec(
            id="impact",
            title="Impacto e encerramento",
            subtitle="Mais velocidade, mais rastreabilidade e mais autonomia para a decisao",
            duration="4:55-5:00",
            key_message="O SCI reduz trabalho manual, encurta o tempo para entender a variacao e fortalece a confianca no numero.",
            bullets=[
                "Fechamento e analise em minutos, nao em horas de consolidacao manual.",
                "Numero mais facil de defender com regra documentada, debug e tabela de validacao.",
                "Ecossistema em evolucao continua com dashboards, alertas, relatorios e documentacao conectados.",
            ],
            highlights=[
                "Horas -> minutos",
                "Do KPI ao plano de acao",
                "Confianca para escalar o uso",
            ],
            visual_title="Mensagem final",
            visual_lines=[
                "SCI = entender + explicar + projetar + monitorar + decidir",
            ],
            presenter_note="Encerrar com beneficio concreto para negocio e operacao.",
        ),
    ]
    return _renumber_slides(slides)


def _presentation_css() -> None:
    st.markdown(
        """
        <style>
        :root {
            --sci-navy: #13294b;
            --sci-orange: #e95d0f;
            --sci-teal: #00a3a3;
            --sci-ink: #172033;
            --sci-surface: #f7f9fc;
            --sci-line: rgba(19, 41, 75, 0.14);
        }
        .sci-pres-shell {
            background:
                radial-gradient(circle at top right, rgba(233,93,15,0.16), transparent 28%),
                radial-gradient(circle at top left, rgba(0,163,163,0.12), transparent 24%),
                linear-gradient(180deg, #ffffff 0%, #f6f8fc 100%);
            border: 1px solid var(--sci-line);
            border-radius: 24px;
            padding: 1.35rem 1.45rem;
            margin-bottom: 1.05rem;
            box-shadow: 0 24px 48px rgba(19,41,75,0.08);
        }
        .sci-cover {
            background:
                radial-gradient(circle at right top, rgba(233,93,15,0.22), transparent 24%),
                linear-gradient(135deg, rgba(19,41,75,0.98), rgba(34,74,124,0.98));
            color: white;
            border-radius: 28px;
            padding: 1.8rem;
            overflow: hidden;
            border: 1px solid rgba(255,255,255,0.08);
        }
        .sci-cover-top {
            display: grid;
            grid-template-columns: 1.7fr 0.9fr;
            gap: 1rem;
            align-items: start;
        }
        .sci-kicker {
            display: inline-block;
            font-size: 0.76rem;
            font-weight: 800;
            letter-spacing: 0.08em;
            text-transform: uppercase;
            color: #dbeafe;
            margin-bottom: 0.7rem;
        }
        .sci-cover h2 {
            color: white;
            margin: 0;
            font-size: 2rem;
            line-height: 1.05;
        }
        .sci-cover-subtitle {
            margin: 0.45rem 0 0 0;
            font-size: 1rem;
            font-weight: 600;
            color: rgba(255,255,255,0.92);
        }
        .sci-cover-message {
            margin-top: 0.95rem;
            max-width: 760px;
            color: rgba(255,255,255,0.94);
            font-size: 1rem;
            line-height: 1.45;
        }
        .sci-cover-panel, .sci-slide-panel {
            background: rgba(255,255,255,0.1);
            border: 1px solid rgba(255,255,255,0.16);
            border-radius: 18px;
            padding: 0.95rem 1rem;
            backdrop-filter: blur(8px);
        }
        .sci-cover-band {
            width: min(240px, 100%);
            margin: 0.85rem 0 0 auto;
            display: block;
            border-radius: 14px;
            box-shadow: 0 16px 32px rgba(0,0,0,0.16);
        }
        .sci-chip {
            display: inline-block;
            background: rgba(255,255,255,0.12);
            color: white;
            padding: 0.34rem 0.74rem;
            border-radius: 999px;
            font-size: 0.76rem;
            margin: 0 0.4rem 0.4rem 0;
            font-weight: 700;
        }
        .sci-team-grid {
            margin-top: 1.25rem;
            display: grid;
            grid-template-columns: repeat(3, minmax(0, 1fr));
            gap: 0.9rem;
        }
        .sci-team-card {
            background: rgba(255,255,255,0.12);
            border: 1px solid rgba(255,255,255,0.16);
            border-radius: 18px;
            padding: 0.85rem;
            min-height: 168px;
        }
        .sci-team-avatar {
            width: 68px;
            height: 68px;
            border-radius: 16px;
            object-fit: cover;
            display: block;
            margin-bottom: 0.7rem;
            border: 2px solid rgba(255,255,255,0.14);
            background: rgba(255,255,255,0.14);
        }
        .sci-team-fallback {
            width: 68px;
            height: 68px;
            border-radius: 16px;
            display: flex;
            align-items: center;
            justify-content: center;
            margin-bottom: 0.7rem;
            background: rgba(255,255,255,0.14);
            color: white;
            font-size: 1.15rem;
            font-weight: 800;
        }
        .sci-team-name {
            color: white;
            font-size: 0.94rem;
            font-weight: 800;
            margin-bottom: 0.2rem;
        }
        .sci-team-role {
            color: #ffd6bf;
            font-size: 0.8rem;
            font-weight: 700;
            margin-bottom: 0.35rem;
        }
        .sci-team-desc {
            color: rgba(255,255,255,0.9);
            font-size: 0.78rem;
            line-height: 1.35;
        }
        .sci-slide-grid {
            display: grid;
            grid-template-columns: 1.6fr 1fr;
            gap: 1rem;
        }
        .sci-slide-box {
            background: white;
            border: 1px solid var(--sci-line);
            border-radius: 18px;
            padding: 1rem 1.05rem;
        }
        .sci-visual-card {
            background: linear-gradient(180deg, rgba(19,41,75,0.04), rgba(0,163,163,0.06));
            border: 1px solid rgba(19,41,75,0.12);
            border-radius: 18px;
            padding: 1rem;
        }
        .sci-label {
            color: var(--sci-orange);
            text-transform: uppercase;
            letter-spacing: 0.08em;
            font-size: 0.72rem;
            font-weight: 800;
            margin-bottom: 0.42rem;
        }
        .sci-slide-title {
            font-size: 1.36rem;
            font-weight: 800;
            color: var(--sci-navy);
            margin-bottom: 0.22rem;
            line-height: 1.1;
        }
        .sci-slide-subtitle {
            color: var(--sci-ink);
            font-size: 0.96rem;
            margin-bottom: 0.74rem;
        }
        .sci-key-message {
            font-size: 1rem;
            font-weight: 700;
            color: var(--sci-ink);
            margin-bottom: 0;
            line-height: 1.45;
        }
        .sci-slide-box ul, .sci-visual-card ul {
            margin: 0;
            padding-left: 1.08rem;
        }
        .sci-highlights {
            margin-top: 0.85rem;
            display: grid;
            grid-template-columns: repeat(3, minmax(0, 1fr));
            gap: 0.75rem;
        }
        .sci-slide-panel {
            background: white;
            border-color: var(--sci-line);
            box-shadow: inset 0 1px 0 rgba(255,255,255,0.4);
        }
        .sci-note {
            border-left: 4px solid var(--sci-teal);
            padding: 0.75rem 0.9rem;
            margin-top: 0.8rem;
            background: rgba(0,163,163,0.08);
            border-radius: 0 14px 14px 0;
            color: var(--sci-ink);
            font-size: 0.9rem;
        }
        @media (max-width: 900px) {
            .sci-cover-top, .sci-slide-grid, .sci-team-grid, .sci-highlights {
                grid-template-columns: 1fr;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def _team_cards_html() -> str:
    cards = []
    for member in _load_team_members():
        photo_uri = _photo_data_uri(member.photo_b64)
        avatar = (
            f'<img class="sci-team-avatar" src="{photo_uri}" alt="{_safe_html(member.name)}" />'
            if photo_uri
            else f'<div class="sci-team-fallback">{_safe_html(member.name[:2].upper())}</div>'
        )
        cards.append(
            "".join(
                [
                    '<div class="sci-team-card">',
                    avatar,
                    f'<div class="sci-team-name">{_safe_html(member.name)}</div>',
                    f'<div class="sci-team-role">{_safe_html(member.role)}</div>',
                    f'<div class="sci-team-desc">{_safe_html(member.description)}</div>',
                    "</div>",
                ]
            )
        )
    return "".join(cards)


def _team_cards_html_for_language(language: str) -> str:
    cards = []
    for member in _load_team_members(language):
        photo_uri = _photo_data_uri(member.photo_b64)
        avatar = (
            f'<img class="sci-team-avatar" src="{photo_uri}" alt="{_safe_html(member.name)}" />'
            if photo_uri
            else f'<div class="sci-team-fallback">{_safe_html(member.name[:2].upper())}</div>'
        )
        cards.append(
            "".join(
                [
                    '<div class="sci-team-card">',
                    avatar,
                    f'<div class="sci-team-name">{_safe_html(member.name)}</div>',
                    f'<div class="sci-team-role">{_safe_html(member.role)}</div>',
                    f'<div class="sci-team-desc">{_safe_html(member.description)}</div>',
                    "</div>",
                ]
            )
        )
    return "".join(cards)


def _render_cover(slide: SlideSpec, language: str = "pt") -> None:
    ui = _ui_texts(language)
    faixa_uri = _image_file_data_uri(SCI_FAIXA_PATH)
    chips_html = "".join(f'<span class="sci-chip">{_safe_html(item)}</span>' for item in slide.highlights)
    faixa_html = (
        f'<img class="sci-cover-band" src="{faixa_uri}" alt="SCI faixa" />'
        if faixa_uri
        else ""
    )

    st.markdown(
        f"""
        <div class="sci-cover">
            <div class="sci-cover-top">
                <div>
                    <div class="sci-kicker">{_safe_html(ui['cover_badge'])}</div>
                    <h2>{_safe_html(slide.title)}</h2>
                    <div class="sci-cover-subtitle">{_safe_html(slide.subtitle)}</div>
                    <div class="sci-cover-message">{_safe_html(slide.key_message)}</div>
                </div>
                <div class="sci-cover-panel">
                    <div class="sci-label" style="color:#ffd6bf;">{_safe_html(ui['cover_highlights'])}</div>
                    {chips_html}
                    {faixa_html}
                </div>
            </div>
            <div class="sci-team-grid">
                {_team_cards_html_for_language(language)}
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    c1, c2 = st.columns([1.35, 1])
    with c1:
        st.markdown(ui["what_will_be_shown"])
        for bullet in slide.bullets:
            st.markdown(f"- {bullet}")
    with c2:
        st.markdown(ui["narrative"])
        for line in slide.visual_lines:
            st.markdown(f"- {line}")


def _render_slide(slide: SlideSpec, language: str = "pt") -> None:
    ui = _ui_texts(language)
    st.markdown(
        f"""
        <div class="sci-pres-shell">
            <div class="sci-label">{_safe_html(slide.label)} | {_safe_html(slide.duration)}</div>
            <div class="sci-slide-title">{_safe_html(slide.title)}</div>
            <div class="sci-slide-subtitle">{_safe_html(slide.subtitle)}</div>
            <div class="sci-key-message">{_safe_html(slide.key_message)}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    c1, c2 = st.columns([1.6, 1])
    with c1:
        left_html = [
            '<div class="sci-slide-box">',
            f'<div class="sci-label">{_safe_html(ui["key_points"])}</div>',
            '<ul>',
        ]
        left_html.extend(f"<li>{_safe_html(bullet)}</li>" for bullet in slide.bullets)
        left_html.extend(["</ul>", "</div>"])
        st.markdown("".join(left_html), unsafe_allow_html=True)

    with c2:
        right_html = [
            '<div class="sci-visual-card">',
            f'<div class="sci-label">{_safe_html(slide.visual_title or ui["visual_support"])}</div>',
            '<ul>',
        ]
        right_html.extend(f"<li>{_safe_html(line)}</li>" for line in slide.visual_lines)
        right_html.extend(["</ul>", "</div>"])
        st.markdown("".join(right_html), unsafe_allow_html=True)

    if slide.highlights:
        st.markdown(ui["why_it_matters"])
        cards_html = ["<div class=\"sci-highlights\">"]
        for item in slide.highlights:
            cards_html.append(
                f'<div class="sci-slide-panel"><strong>{_safe_html(item)}</strong></div>'
            )
        cards_html.append("</div>")
        st.markdown("".join(cards_html), unsafe_allow_html=True)

    if slide.presenter_note:
        st.markdown(
            f'<div class="sci-note"><strong>{_safe_html(ui["speaker_note"])}</strong> {_safe_html(slide.presenter_note)}</div>',
            unsafe_allow_html=True,
        )


def _pdf_styles() -> dict[str, ParagraphStyle]:
    base = getSampleStyleSheet()
    return {
        "cover_kicker": ParagraphStyle(
            "CoverKicker",
            parent=base["Normal"],
            fontSize=9,
            leading=11,
            textColor=colors.HexColor("#DBEAFE"),
            fontName="Helvetica-Bold",
            alignment=TA_LEFT,
        ),
        "cover_title": ParagraphStyle(
            "CoverTitle",
            parent=base["Title"],
            fontSize=24,
            leading=26,
            textColor=WHITE,
            fontName="Helvetica-Bold",
            alignment=TA_LEFT,
        ),
        "cover_subtitle": ParagraphStyle(
            "CoverSubtitle",
            parent=base["Normal"],
            fontSize=12,
            leading=15,
            textColor=colors.HexColor("#E6ECF5"),
            alignment=TA_LEFT,
        ),
        "cover_message": ParagraphStyle(
            "CoverMessage",
            parent=base["Normal"],
            fontSize=14,
            leading=18,
            textColor=colors.HexColor("#FFD6BF"),
            fontName="Helvetica-Bold",
            alignment=TA_LEFT,
        ),
        "cover_body": ParagraphStyle(
            "CoverBody",
            parent=base["Normal"],
            fontSize=10.5,
            leading=13,
            textColor=WHITE,
            alignment=TA_LEFT,
        ),
        "tiny_label": ParagraphStyle(
            "TinyLabel",
            parent=base["Normal"],
            fontSize=8.5,
            leading=10,
            textColor=ACCENT,
            fontName="Helvetica-Bold",
            alignment=TA_LEFT,
        ),
        "slide_title": ParagraphStyle(
            "SlideTitle",
            parent=base["Heading1"],
            fontSize=21,
            leading=23,
            textColor=PRIMARY,
            fontName="Helvetica-Bold",
            alignment=TA_LEFT,
        ),
        "slide_subtitle": ParagraphStyle(
            "SlideSubtitle",
            parent=base["Normal"],
            fontSize=11,
            leading=14,
            textColor=MUTED,
            alignment=TA_LEFT,
        ),
        "key_message": ParagraphStyle(
            "KeyMessage",
            parent=base["Normal"],
            fontSize=13,
            leading=17,
            textColor=INK,
            fontName="Helvetica-Bold",
            alignment=TA_LEFT,
        ),
        "body": ParagraphStyle(
            "Body",
            parent=base["Normal"],
            fontSize=10.5,
            leading=13,
            textColor=INK,
            alignment=TA_LEFT,
        ),
        "highlight": ParagraphStyle(
            "Highlight",
            parent=base["Normal"],
            fontSize=10,
            leading=12,
            textColor=INK,
            alignment=TA_CENTER,
            fontName="Helvetica-Bold",
        ),
        "note": ParagraphStyle(
            "Note",
            parent=base["Normal"],
            fontSize=9.5,
            leading=12,
            textColor=INK,
            alignment=TA_LEFT,
        ),
        "team_name": ParagraphStyle(
            "TeamName",
            parent=base["Normal"],
            fontSize=9.5,
            leading=11,
            textColor=WHITE,
            fontName="Helvetica-Bold",
            alignment=TA_LEFT,
        ),
        "team_role": ParagraphStyle(
            "TeamRole",
            parent=base["Normal"],
            fontSize=8,
            leading=10,
            textColor=colors.HexColor("#FFD6BF"),
            fontName="Helvetica-Bold",
            alignment=TA_LEFT,
        ),
        "team_desc": ParagraphStyle(
            "TeamDesc",
            parent=base["Normal"],
            fontSize=7.2,
            leading=9,
            textColor=WHITE,
            alignment=TA_LEFT,
        ),
    }


def _ppt_add_box(slide, left, top, width, height, fill_rgb, line_rgb, radius=True, transparency=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if transparency is not None:
        shape.fill.transparency = transparency
    shape.line.color.rgb = line_rgb
    return shape


def _ppt_add_textbox(slide, left, top, width, height, text, size, color, bold=False, font_name="Aptos", align=None):
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    frame.margin_top = Pt(4)
    frame.margin_bottom = Pt(4)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align == "center":
        paragraph.alignment = PP_ALIGN.CENTER
    else:
        paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return box


def _ppt_add_list(slide, left, top, width, height, items, size, color, line_spacing=1.05):
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(8)
    frame.margin_right = Pt(8)
    frame.margin_top = Pt(6)
    frame.margin_bottom = Pt(6)

    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"â€¢ {item}"
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(6)
        paragraph.line_spacing = line_spacing
        run = paragraph.runs[0]
        run.font.size = Pt(size)
        run.font.name = "Aptos"
        run.font.color.rgb = color
    return box


def _cover_layout_ppt(slide, slide_spec: SlideSpec, faixa_bytes: bytes | None, language: str = "pt") -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    ui = _ui_texts(language)
    navy = RGBColor(*_hex_to_rgb(PRIMARY_HEX))
    orange = RGBColor(*_hex_to_rgb(ACCENT_HEX))
    white = RGBColor(255, 255, 255)
    pale = RGBColor(230, 236, 245)
    warm = RGBColor(255, 214, 191)
    card_fill = RGBColor(43, 66, 102)
    card_line = RGBColor(173, 195, 224)
    card_band = RGBColor(75, 104, 145)

    _ppt_add_box(slide, 0, 0, SLIDE_WIDTH_IN, SLIDE_HEIGHT_IN, navy, navy, radius=False)
    accent = _ppt_add_box(slide, 9.55, 0, 3.78, SLIDE_HEIGHT_IN, orange, orange, radius=False, transparency=0.78)
    accent.line.fill.background()

    _ppt_add_textbox(slide, 0.62, 0.48, 4.2, 0.35, ui["cover_badge"].upper(), 10.5, RGBColor(219, 234, 254), bold=True)
    _ppt_add_textbox(slide, 0.62, 0.92, 6.7, 0.9, slide_spec.title, 24, white, bold=True)
    _ppt_add_textbox(slide, 0.62, 1.88, 7.6, 0.45, slide_spec.subtitle, 12, pale)
    _ppt_add_textbox(slide, 0.62, 2.45, 7.5, 0.82, slide_spec.key_message, 15, warm, bold=True)
    _ppt_add_list(slide, 0.62, 3.32, 5.45, 1.35, slide_spec.bullets, 11, white)

    if faixa_bytes:
        slide.shapes.add_picture(BytesIO(faixa_bytes), Inches(9.3), Inches(0.75), width=Inches(2.95))

    members = _load_team_members(language)
    left_positions = [0.62, 4.28, 7.94]
    for left, member in zip(left_positions, members):
        card = _ppt_add_box(slide, left, 5.02, 3.1, 1.98, card_fill, card_line, transparency=0.08)
        card.line.width = Inches(0.02)
        top_band = _ppt_add_box(slide, left + 0.1, 5.12, 2.9, 0.16, card_band, card_band, transparency=0.02)
        top_band.line.fill.background()
        photo = _photo_bytes(member.photo_b64)
        if photo:
            slide.shapes.add_picture(BytesIO(photo), Inches(left + 0.18), Inches(5.33), width=Inches(0.86), height=Inches(0.98))
        else:
            _ppt_add_box(slide, left + 0.18, 5.33, 0.86, 0.98, RGBColor(79, 98, 132), RGBColor(79, 98, 132))
            _ppt_add_textbox(slide, left + 0.24, 5.61, 0.72, 0.25, member.name[:2].upper(), 10, white, bold=True, align="center")
        _ppt_add_textbox(slide, left + 1.1, 5.3, 1.8, 0.34, member.name, 10.2, white, bold=True)
        _ppt_add_textbox(slide, left + 1.1, 5.66, 1.78, 0.24, member.role, 8.1, warm, bold=True)
        _ppt_add_textbox(slide, left + 1.1, 5.92, 1.78, 0.72, member.description[:110], 7.2, pale)


def _content_layout_ppt(slide, slide_spec: SlideSpec, language: str = "pt") -> None:
    from pptx.dml.color import RGBColor

    ui = _ui_texts(language)
    navy = RGBColor(*_hex_to_rgb(PRIMARY_HEX))
    orange = RGBColor(*_hex_to_rgb(ACCENT_HEX))
    teal = RGBColor(*_hex_to_rgb(TEAL_HEX))
    ink = RGBColor(*_hex_to_rgb(INK_HEX))
    muted = RGBColor(*_hex_to_rgb(MUTED_HEX))
    line = RGBColor(*_hex_to_rgb(LINE_HEX))
    panel = RGBColor(*_hex_to_rgb(PANEL_HEX))
    white = RGBColor(255, 255, 255)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = navy

    shell = _ppt_add_box(slide, 0.42, 0.38, 12.48, 1.76, white, line)
    shell.shadow.inherit = False
    ribbon = _ppt_add_box(slide, 0.58, 0.55, 2.18, 0.36, orange, orange)
    ribbon.line.fill.background()
    _ppt_add_textbox(slide, 0.64, 0.57, 2.05, 0.24, f"{slide_spec.label} | {slide_spec.duration}", 9.5, white, bold=True)
    _ppt_add_textbox(slide, 0.58, 0.9, 7.6, 0.42, slide_spec.title, 20.5, navy, bold=True)
    _ppt_add_textbox(slide, 0.58, 1.3, 8.4, 0.28, slide_spec.subtitle, 11, muted)
    _ppt_add_textbox(slide, 0.58, 1.58, 8.65, 0.42, slide_spec.key_message, 13.2, ink, bold=True)

    _ppt_add_box(slide, 0.42, 2.36, 7.45, 2.78, white, line)
    _ppt_add_box(slide, 8.02, 2.36, 4.89, 2.78, panel, line)
    _ppt_add_textbox(slide, 0.62, 2.52, 2.4, 0.24, ui["header_points"], 9.5, orange, bold=True)
    _ppt_add_list(slide, 0.58, 2.78, 6.95, 2.06, slide_spec.bullets, 10.8, ink)
    _ppt_add_textbox(slide, 8.2, 2.52, 3.8, 0.24, slide_spec.visual_title or ui["pdf_visual_support"], 9.5, teal, bold=True)
    _ppt_add_list(slide, 8.16, 2.78, 4.35, 1.98, slide_spec.visual_lines, 10.2, ink)

    if slide_spec.highlights:
        _ppt_add_textbox(slide, 0.58, 5.3, 3.6, 0.24, ui["header_why"], 9.5, orange, bold=True)
        gap = 0.18
        count = max(1, len(slide_spec.highlights))
        total_width = 12.48
        card_width = (total_width - (gap * (count - 1))) / count
        left = 0.42
        for item in slide_spec.highlights:
            _ppt_add_box(slide, left, 5.56, card_width, 0.64, white, line)
            _ppt_add_textbox(slide, left + 0.06, 5.68, card_width - 0.12, 0.28, item, 9.5, ink, bold=True, align="center")
            left += card_width + gap

    if slide_spec.presenter_note:
        note = _ppt_add_box(slide, 0.42, 6.4, 12.48, 0.68, RGBColor(228, 247, 247), RGBColor(168, 222, 222))
        note.line.color.rgb = teal
        _ppt_add_textbox(slide, 0.52, 6.52, 12.0, 0.32, f"{ui['speaker_note']} {slide_spec.presenter_note}", 9, ink)


def export_presentation_pptx(slides: list[SlideSpec], language: str = "pt") -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError("python-pptx nao esta instalado no ambiente.") from exc

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    faixa_bytes = SCI_FAIXA_PATH.read_bytes() if SCI_FAIXA_PATH.exists() else None

    for slide_spec in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if slide_spec.cover:
            _cover_layout_ppt(slide, slide_spec, faixa_bytes, language)
        else:
            _content_layout_ppt(slide, slide_spec, language)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def _pdf_top_to_y(top_in: float, height_in: float) -> float:
    return SLIDE_HEIGHT_PT - ((top_in + height_in) * 72)


def _pdf_draw_round_rect(pdf, left_in, top_in, width_in, height_in, fill_color, stroke_color, radius=14, stroke=1):
    x = left_in * 72
    y = _pdf_top_to_y(top_in, height_in)
    pdf.setFillColor(fill_color)
    pdf.setStrokeColor(stroke_color)
    pdf.setLineWidth(stroke)
    pdf.roundRect(x, y, width_in * 72, height_in * 72, radius, fill=1, stroke=1)


def _pdf_draw_paragraph(pdf, text: str, style: ParagraphStyle, left_in: float, top_in: float, width_in: float) -> float:
    paragraph = Paragraph(html.escape(text), style)
    width_pt = width_in * 72
    _, height_pt = paragraph.wrap(width_pt, 1000)
    y = SLIDE_HEIGHT_PT - (top_in * 72) - height_pt
    paragraph.drawOn(pdf, left_in * 72, y)
    return height_pt / 72


def _pdf_draw_list(pdf, items: list[str], style: ParagraphStyle, left_in: float, top_in: float, width_in: float, gap_in: float = 0.08) -> float:
    current_top = top_in
    for item in items:
        used = _pdf_draw_paragraph(pdf, f"â€¢ {item}", style, left_in, current_top, width_in)
        current_top += used + gap_in
    return current_top


def _pdf_draw_image(pdf, image_bytes: bytes, left_in: float, top_in: float, width_in: float, height_in: float | None = None) -> None:
    if not image_bytes:
        return
    reader = ImageReader(BytesIO(image_bytes))
    width_pt = width_in * 72
    if height_in is None:
        img_width, img_height = reader.getSize()
        ratio = img_height / img_width if img_width else 0.4
        height_pt = width_pt * ratio
        height_in = height_pt / 72
    x = left_in * 72
    y = _pdf_top_to_y(top_in, height_in)
    pdf.drawImage(reader, x, y, width=width_pt, height=height_in * 72, preserveAspectRatio=True, mask="auto")


def _draw_cover_pdf(pdf, slide: SlideSpec, styles: dict[str, ParagraphStyle], faixa_bytes: bytes | None, language: str = "pt") -> None:
    ui = _ui_texts(language)
    pdf.setFillColor(PRIMARY)
    pdf.rect(0, 0, SLIDE_WIDTH_PT, SLIDE_HEIGHT_PT, fill=1, stroke=0)
    pdf.setFillColor(colors.HexColor("#30598B"))
    pdf.rect(9.55 * 72, 0, 3.78 * 72, SLIDE_HEIGHT_PT, fill=1, stroke=0)

    _pdf_draw_paragraph(pdf, ui["cover_badge"].upper(), styles["cover_kicker"], 0.62, 0.48, 4.2)
    _pdf_draw_paragraph(pdf, slide.title, styles["cover_title"], 0.62, 0.88, 6.8)
    _pdf_draw_paragraph(pdf, slide.subtitle, styles["cover_subtitle"], 0.62, 1.86, 7.2)
    _pdf_draw_paragraph(pdf, slide.key_message, styles["cover_message"], 0.62, 2.4, 7.2)
    _pdf_draw_list(pdf, slide.bullets, styles["cover_body"], 0.68, 3.36, 5.2, gap_in=0.06)

    if faixa_bytes:
        _pdf_draw_image(pdf, faixa_bytes, 9.3, 0.78, 2.95)

    members = _load_team_members(language)
    left_positions = [0.62, 4.28, 7.94]
    for left, member in zip(left_positions, members):
        _pdf_draw_round_rect(pdf, left, 5.15, 3.08, 1.82, colors.HexColor("#355887"), colors.HexColor("#4A6A95"))
        photo = _photo_bytes(member.photo_b64)
        if photo:
            _pdf_draw_image(pdf, photo, left + 0.18, 5.34, 0.8, 0.88)
        else:
            _pdf_draw_round_rect(pdf, left + 0.18, 5.34, 0.8, 0.88, colors.HexColor("#4A6A95"), colors.HexColor("#4A6A95"))
        _pdf_draw_paragraph(pdf, member.name, styles["team_name"], left + 1.02, 5.28, 1.78)
        _pdf_draw_paragraph(pdf, member.role, styles["team_role"], left + 1.02, 5.58, 1.78)
        _pdf_draw_paragraph(pdf, member.description[:120], styles["team_desc"], left + 1.02, 5.86, 1.75)


def _draw_content_pdf(pdf, slide: SlideSpec, styles: dict[str, ParagraphStyle], language: str = "pt") -> None:
    ui = _ui_texts(language)
    pdf.setFillColor(PRIMARY)
    pdf.rect(0, 0, SLIDE_WIDTH_PT, SLIDE_HEIGHT_PT, fill=1, stroke=0)

    _pdf_draw_round_rect(pdf, 0.42, 0.38, 12.48, 1.76, WHITE, LINE)
    _pdf_draw_round_rect(pdf, 0.58, 0.55, 2.18, 0.36, ACCENT, ACCENT)
    _pdf_draw_paragraph(pdf, f"{slide.label} | {slide.duration}", styles["cover_kicker"], 0.68, 0.56, 1.92)
    _pdf_draw_paragraph(pdf, slide.title, styles["slide_title"], 0.58, 0.88, 7.1)
    _pdf_draw_paragraph(pdf, slide.subtitle, styles["slide_subtitle"], 0.58, 1.3, 8.3)
    _pdf_draw_paragraph(pdf, slide.key_message, styles["key_message"], 0.58, 1.56, 8.5)

    _pdf_draw_round_rect(pdf, 0.42, 2.36, 7.45, 2.78, WHITE, LINE)
    _pdf_draw_round_rect(pdf, 8.02, 2.36, 4.89, 2.78, PANEL, LINE)
    _pdf_draw_paragraph(pdf, ui["header_points"], styles["tiny_label"], 0.66, 2.5, 1.9)
    _pdf_draw_list(pdf, slide.bullets, styles["body"], 0.62, 2.8, 6.85, gap_in=0.05)
    _pdf_draw_paragraph(pdf, slide.visual_title or ui["pdf_visual_support"], styles["tiny_label"], 8.24, 2.5, 2.6)
    _pdf_draw_list(pdf, slide.visual_lines, styles["body"], 8.2, 2.8, 4.2, gap_in=0.05)

    if slide.highlights:
        _pdf_draw_paragraph(pdf, ui["header_why"], styles["tiny_label"], 0.58, 5.28, 3.0)
        gap = 0.18
        count = len(slide.highlights)
        card_width = (12.48 - (gap * (count - 1))) / count
        left = 0.42
        for item in slide.highlights:
            _pdf_draw_round_rect(pdf, left, 5.56, card_width, 0.64, WHITE, LINE)
            _pdf_draw_paragraph(pdf, item, styles["highlight"], left + 0.08, 5.72, card_width - 0.16)
            left += card_width + gap

    if slide.presenter_note:
        _pdf_draw_round_rect(pdf, 0.42, 6.4, 12.48, 0.68, colors.HexColor("#E4F7F7"), colors.HexColor("#A8DEDE"))
        _pdf_draw_paragraph(pdf, f"{ui['pdf_speaker_note']} {slide.presenter_note}", styles["note"], 0.55, 6.58, 11.9)


def export_presentation_pdf(slides: list[SlideSpec], language: str = "pt") -> bytes:
    buffer = BytesIO()
    pdf = pdf_canvas.Canvas(buffer, pagesize=SLIDE_PAGE_SIZE)
    pdf.setTitle(_ui_texts(language)["pdf_title"])
    styles = _pdf_styles()
    faixa_bytes = SCI_FAIXA_PATH.read_bytes() if SCI_FAIXA_PATH.exists() else None

    for index, slide in enumerate(slides):
        if slide.cover:
            _draw_cover_pdf(pdf, slide, styles, faixa_bytes, language)
        else:
            _draw_content_pdf(pdf, slide, styles, language)
        if index < len(slides) - 1:
            pdf.showPage()

    pdf.save()
    return buffer.getvalue()


def render_presentation_section(version: str, data_atualizacao: str | None) -> None:
    language = st.selectbox(
        "Presentation language / Idioma da apresentacao",
        options=list(LANGUAGE_OPTIONS.keys()),
        format_func=lambda value: LANGUAGE_OPTIONS[value],
        index=0,
        key="presentation_docs_language",
    )
    ui = _ui_texts(language)
    slides = build_presentation_slides(version, data_atualizacao, language)
    _presentation_css()

    st.markdown(
        f"""
        <div class="sci-pres-shell" style="padding:1.55rem;">
            <div class="sci-label">{_safe_html(ui['section_badge'])}</div>
            <div class="sci-slide-title">{_safe_html(ui['section_title'])}</div>
            <div class="sci-slide-subtitle">{_safe_html(ui['section_subtitle'])}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    pdf_bytes = export_presentation_pdf(slides, language)
    ppt_bytes = export_presentation_pptx(slides, language)

    c1, c2 = st.columns([1, 1])
    with c1:
        st.download_button(
            ui["download_pdf"],
            data=pdf_bytes,
            file_name="SCI_Apresentacao_Executiva.pdf",
            mime="application/pdf",
            use_container_width=True,
        )
    with c2:
        st.download_button(
            ui["download_ppt"],
            data=ppt_bytes,
            file_name="SCI_Apresentacao_Executiva.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )

    tab_roteiro, tab_slides = st.tabs([ui["tabs_script"], ui["tabs_slides"]])
    with tab_roteiro:
        st.markdown(ui["script_intro"])
        for slide in slides[1:]:
            st.markdown(f"- **{slide.label} ({slide.duration})**: {slide.key_message}")

    with tab_slides:
        st.info(ui["slides_info"])
        for index, slide in enumerate(slides):
            with st.expander(f"{slide.label} â€” {slide.title}", expanded=index == 0):
                if slide.cover:
                    _render_cover(slide, language)
                else:
                    _render_slide(slide, language)

